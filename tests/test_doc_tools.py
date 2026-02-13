# Copyright (c) 2025 Kenneth Stott
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Tests for PDF document loading in doc_tools module."""

import io
import pytest
from pathlib import Path
from unittest.mock import Mock, patch

from constat.core.config import Config
from constat.discovery.doc_tools import DocumentDiscoveryTools
from constat.discovery.doc_tools._file_extractors import (
    _extract_pdf_text,
    _extract_pdf_text_from_bytes,
    _extract_docx_text,
    _extract_docx_text_from_bytes,
    _extract_xlsx_text,
    _extract_xlsx_text_from_bytes,
    _extract_pptx_text,
    _extract_pptx_text_from_bytes,
)


# =============================================================================
# Helper Functions for Creating Test PDFs
# =============================================================================


def create_minimal_pdf(num_pages: int = 1) -> bytes:
    """
    Create a minimal valid PDF file.
    
    Args:
        num_pages: Number of pages to create
        
    Returns:
        PDF file bytes
    """
    from pypdf import PdfWriter, PageObject
    
    writer = PdfWriter()
    for _ in range(num_pages):
        page = PageObject.create_blank_page(width=612, height=792)
        writer.add_page(page)
    
    output = io.BytesIO()
    writer.write(output)
    output.seek(0)
    return output.read()


# =============================================================================
# Fixtures
# =============================================================================


@pytest.fixture
def mock_pdf_reader():
    """Create a mock PdfReader for controlled testing."""
    def _create_mock(page_texts: list[str]):
        """
        Create a mock PdfReader with specified page texts.
        
        Args:
            page_texts: List of text content for each page (empty string for empty pages)
        """
        mock_reader = Mock()
        mock_pages = []
        
        for text in page_texts:
            mock_page = Mock()
            mock_page.extract_text.return_value = text
            mock_pages.append(mock_page)
        
        mock_reader.pages = mock_pages
        return mock_reader
    
    return _create_mock


@pytest.fixture
def simple_pdf_bytes():
    """Create a simple multi-page PDF."""
    return create_minimal_pdf(2)


@pytest.fixture
def pdf_file(tmp_path, simple_pdf_bytes):
    """Create a PDF file on disk."""
    pdf_path = tmp_path / "test_document.pdf"
    pdf_path.write_bytes(simple_pdf_bytes)
    return pdf_path


@pytest.fixture
def config_with_pdf_file(pdf_file):
    """Config with a PDF file document."""
    return Config(
        documents={
            "test_pdf": {
                "type": "file",
                "path": str(pdf_file),
                "description": "Test PDF document",
            }
        }
    )


@pytest.fixture
def config_with_pdf_type(pdf_file):
    """Config with a direct pdf type document."""
    return Config(
        documents={
            "direct_pdf": {
                "type": "pdf",
                "path": str(pdf_file),
                "description": "Direct PDF document",
            }
        }
    )


@pytest.fixture
def config_with_http_pdf():
    """Config with an HTTP PDF document."""
    return Config(
        documents={
            "http_pdf": {
                "type": "http",
                "url": "https://example.com/document.pdf",
                "description": "HTTP PDF document",
            }
        }
    )


@pytest.fixture
def config_with_pdf_url():
    """Config with a PDF type and URL."""
    return Config(
        documents={
            "url_pdf": {
                "type": "pdf",
                "url": "https://example.com/report.pdf",
                "description": "URL-based PDF document",
            }
        }
    )


# =============================================================================
# PDF Text Extraction Tests (using mocks for reliable testing)
# =============================================================================


class TestPdfTextExtraction:
    """Tests for _extract_pdf_text and _extract_pdf_text_from_bytes methods."""

    def test_extract_pdf_text_from_file_with_mock(self, tmp_path, mock_pdf_reader):
        """Test extracting text from a PDF file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        # Create a minimal PDF file
        pdf_path = tmp_path / "test.pdf"
        pdf_bytes = create_minimal_pdf(2)
        pdf_path.write_bytes(pdf_bytes)
        
        mock_reader = mock_pdf_reader(["Page one content.", "Page two content."])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text(pdf_path)
        
        assert "[Page 1]" in result
        assert "[Page 2]" in result
        assert "Page one content" in result
        assert "Page two content" in result

    def test_extract_pdf_text_from_bytes_with_mock(self, mock_pdf_reader):
        """Test extracting text from PDF bytes."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(2)
        mock_reader = mock_pdf_reader(["First page text.", "Second page text."])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        assert "[Page 1]" in result
        assert "[Page 2]" in result
        assert "First page text" in result
        assert "Second page text" in result

    def test_page_markers_format(self, mock_pdf_reader):
        """Test that page markers follow expected format [Page N]."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(3)
        mock_reader = mock_pdf_reader(["A", "B", "C"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Check page marker format
        import re
        page_markers = re.findall(r'\[Page \d+\]', result)
        assert len(page_markers) == 3
        assert page_markers[0] == "[Page 1]"
        assert page_markers[1] == "[Page 2]"
        assert page_markers[2] == "[Page 3]"

    def test_pages_separated_by_double_newline(self, mock_pdf_reader):
        """Test that pages are separated by double newlines."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(2)
        mock_reader = mock_pdf_reader(["Page 1", "Page 2"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Pages should be joined by double newlines
        assert "\n\n" in result

    def test_empty_pages_are_skipped(self, mock_pdf_reader):
        """Test that empty pages are not included in output."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(4)
        # Pages 2 and 4 are empty
        mock_reader = mock_pdf_reader(["Content on page 1", "", "Content on page 3", ""])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Should have page 1 and page 3
        assert "[Page 1]" in result
        assert "[Page 3]" in result
        # Empty pages should be skipped
        assert "[Page 2]" not in result
        assert "[Page 4]" not in result

    def test_whitespace_only_pages_are_skipped(self, mock_pdf_reader):
        """Test that whitespace-only pages are treated as empty."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(3)
        # Page 2 has only whitespace
        mock_reader = mock_pdf_reader(["Content", "   \n\t  ", "More content"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        assert "[Page 1]" in result
        assert "[Page 3]" in result
        assert "[Page 2]" not in result


class TestPdfErrorHandling:
    """Tests for error handling in PDF extraction."""

    def test_corrupt_pdf_raises_error(self):
        """Test that corrupt PDF data raises an error."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        corrupt_data = b"This is not a valid PDF file"
        
        with pytest.raises(Exception):  # pypdf raises various exceptions
            _extract_pdf_text_from_bytes(corrupt_data)

    def test_corrupt_pdf_file_raises_error(self, tmp_path):
        """Test that corrupt PDF file raises an error."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        corrupt_file = tmp_path / "corrupt.pdf"
        corrupt_file.write_bytes(b"Not a PDF")
        
        with pytest.raises(Exception):
            _extract_pdf_text(corrupt_file)

    def test_missing_pdf_file_raises_error(self):
        """Test that missing PDF file raises FileNotFoundError."""
        config = Config(
            documents={
                "missing": {
                    "type": "pdf",
                    "path": "/nonexistent/file.pdf",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("missing")
        
        assert "error" in result
        assert "not found" in result["error"].lower() or "Failed to load" in result["error"]


class TestPdfLoadingViaHttpType:
    """Tests for loading PDFs through type='http' with PDF content-type."""

    def test_load_pdf_via_http_content_type(self, config_with_http_pdf, mock_pdf_reader):
        """Test that PDFs are detected via HTTP content-type header."""
        tools = DocumentDiscoveryTools(config_with_http_pdf)
        
        mock_reader = mock_pdf_reader(["HTTP PDF content"])
        
        mock_response = Mock()
        mock_response.content = b"fake pdf bytes"
        mock_response.headers = {"content-type": "application/pdf"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response), \
             patch('pypdf.PdfReader', return_value=mock_reader):
            result = tools.get_document("http_pdf")
            
            assert "content" in result
            assert "[Page 1]" in result["content"]
            assert result["format"] == "text"

    def test_load_pdf_via_http_url_extension(self, mock_pdf_reader):
        """Test that PDFs are detected via .pdf URL extension."""
        config = Config(
            documents={
                "url_pdf": {
                    "type": "http",
                    "url": "https://example.com/document.pdf",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_reader = mock_pdf_reader(["URL extension PDF content"])
        
        mock_response = Mock()
        mock_response.content = b"fake pdf bytes"
        mock_response.headers = {"content-type": "application/octet-stream"}  # Not PDF content-type
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response), \
             patch('pypdf.PdfReader', return_value=mock_reader):
            result = tools.get_document("url_pdf")
            
            assert "content" in result
            assert "[Page 1]" in result["content"]
            assert result["format"] == "text"

    def test_http_pdf_with_custom_headers(self, mock_pdf_reader):
        """Test that custom headers are passed when fetching HTTP PDFs."""
        config = Config(
            documents={
                "auth_pdf": {
                    "type": "http",
                    "url": "https://example.com/secure.pdf",
                    "headers": {
                        "Authorization": "Bearer token123",
                        "X-Custom": "value",
                    }
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_reader = mock_pdf_reader(["Secure PDF content"])
        
        mock_response = Mock()
        mock_response.content = b"fake pdf bytes"
        mock_response.headers = {"content-type": "application/pdf"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response) as mock_get, \
             patch('pypdf.PdfReader', return_value=mock_reader):
            result = tools.get_document("auth_pdf")
            
            # Verify headers were passed
            call_kwargs = mock_get.call_args[1]
            assert call_kwargs["headers"]["Authorization"] == "Bearer token123"
            assert call_kwargs["headers"]["X-Custom"] == "value"
            assert "content" in result

    def test_http_pdf_fetch_error(self, config_with_http_pdf):
        """Test error handling when HTTP PDF fetch fails."""
        tools = DocumentDiscoveryTools(config_with_http_pdf)
        
        mock_response = Mock()
        mock_response.raise_for_status = Mock(side_effect=Exception("404 Not Found"))
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("http_pdf")
            
            assert "error" in result
            assert "Failed to load" in result["error"]


# =============================================================================
# Multi-page PDF Tests
# =============================================================================


class TestMultiPagePdf:
    """Tests for multi-page PDF handling."""

    def test_multi_page_pdf_all_pages_extracted(self, mock_pdf_reader):
        """Test that all pages are extracted from multi-page PDF."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(5)
        mock_reader = mock_pdf_reader([
            "Content on page 1",
            "Content on page 2",
            "Content on page 3",
            "Content on page 4",
            "Content on page 5",
        ])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        for i in range(1, 6):
            assert f"[Page {i}]" in result
            assert f"Content on page {i}" in result

    def test_page_order_preserved(self, mock_pdf_reader):
        """Test that page order is preserved in output."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(3)
        mock_reader = mock_pdf_reader([
            "MARKER_ALPHA",
            "MARKER_BETA",
            "MARKER_GAMMA",
        ])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Check that markers appear in correct order
        pos_alpha = result.find("MARKER_ALPHA")
        pos_beta = result.find("MARKER_BETA")
        pos_gamma = result.find("MARKER_GAMMA")
        
        assert pos_alpha < pos_beta < pos_gamma


# =============================================================================
# Edge Cases
# =============================================================================


class TestPdfEdgeCases:
    """Edge case tests for PDF handling."""

    def test_single_page_pdf(self, mock_pdf_reader):
        """Test handling of single-page PDF."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(1)
        mock_reader = mock_pdf_reader(["Only page content"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        assert "[Page 1]" in result
        assert "Only page content" in result
        # Should not have page 2
        assert "[Page 2]" not in result

    def test_pdf_with_special_characters(self, mock_pdf_reader):
        """Test PDF with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(1)
        mock_reader = mock_pdf_reader(["Special: < > & \" ' $ % @"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        assert "Special:" in result

    def test_all_empty_pages_pdf(self, mock_pdf_reader):
        """Test PDF where all pages are empty."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(3)
        mock_reader = mock_pdf_reader(["", "", ""])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Should return empty string
        assert result == ""

    def test_pdf_with_none_text_extraction(self):
        """Test PDF where extract_text returns None."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(2)
        
        # Create mock where extract_text returns None
        mock_reader = Mock()
        mock_page1 = Mock()
        mock_page1.extract_text.return_value = None
        mock_page2 = Mock()
        mock_page2.extract_text.return_value = "Valid content"
        mock_reader.pages = [mock_page1, mock_page2]
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Should only have page 2 (page 1 returned None)
        assert "[Page 1]" not in result
        assert "[Page 2]" in result
        assert "Valid content" in result

    def test_pdf_page_numbering_starts_at_one(self, mock_pdf_reader):
        """Test that page numbers start at 1, not 0."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(1)
        mock_reader = mock_pdf_reader(["First page"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        assert "[Page 1]" in result
        assert "[Page 0]" not in result

    def test_large_pdf_many_pages(self, mock_pdf_reader):
        """Test PDF with many pages."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        num_pages = 100
        pdf_bytes = create_minimal_pdf(num_pages)
        mock_reader = mock_pdf_reader([f"Page {i} content" for i in range(1, num_pages + 1)])
        
        with patch('pypdf.PdfReader', return_value=mock_reader):
            result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Verify first, middle, and last pages
        assert "[Page 1]" in result
        assert "[Page 50]" in result
        assert "[Page 100]" in result
        assert "Page 1 content" in result
        assert "Page 50 content" in result
        assert "Page 100 content" in result


# =============================================================================
# Document Config Variations
# =============================================================================


class TestDocumentConfigVariations:
    """Test various DocumentConfig settings for PDFs."""

    def test_pdf_with_description(self, tmp_path, mock_pdf_reader):
        """Test PDF document with description."""
        pdf_path = tmp_path / "described.pdf"
        pdf_path.write_bytes(create_minimal_pdf(1))
        
        config = Config(
            documents={
                "described_pdf": {
                    "type": "pdf",
                    "path": str(pdf_path),
                    "description": "Important business document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        # Check listing
        docs = tools.list_documents()
        assert docs[0]["description"] == "Important business document"

    def test_pdf_with_tags(self, tmp_path, mock_pdf_reader):
        """Test PDF document with tags."""
        pdf_path = tmp_path / "tagged.pdf"
        pdf_path.write_bytes(create_minimal_pdf(1))
        
        config = Config(
            documents={
                "tagged_pdf": {
                    "type": "pdf",
                    "path": str(pdf_path),
                    "tags": ["finance", "quarterly", "2024"],
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        # Check listing
        docs = tools.list_documents()
        assert "finance" in docs[0]["tags"]
        assert "quarterly" in docs[0]["tags"]
        assert "2024" in docs[0]["tags"]


# =============================================================================
# Real PDF Integration Test (without mocking)
# =============================================================================


class TestRealPdfExtraction:
    """Tests using real PDF files without mocking, for full integration coverage."""

    def test_extract_from_blank_pdf(self):
        """Test extracting from a PDF with no text content."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(2)
        
        # Blank pages should result in empty string
        result = _extract_pdf_text_from_bytes(pdf_bytes)
        
        # Should be empty or minimal
        assert isinstance(result, str)
        # Blank pages should be skipped
        assert "[Page" not in result or result.strip() == ""

    def test_pdf_reader_called_with_bytes_io(self, mock_pdf_reader):
        """Test that PdfReader is called with BytesIO wrapper for bytes input."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_bytes = create_minimal_pdf(1)
        mock_reader = mock_pdf_reader(["test content"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader) as mock_class:
            _extract_pdf_text_from_bytes(pdf_bytes)
            
            # Verify PdfReader was called with a BytesIO object
            assert mock_class.call_count == 1
            call_args = mock_class.call_args[0][0]
            assert isinstance(call_args, io.BytesIO)

    def test_pdf_reader_called_with_path_for_file(self, tmp_path, mock_pdf_reader):
        """Test that PdfReader is called with Path for file input."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(create_minimal_pdf(1))
        mock_reader = mock_pdf_reader(["test content"])
        
        with patch('pypdf.PdfReader', return_value=mock_reader) as mock_class:
            _extract_pdf_text(pdf_path)
            
            # Verify PdfReader was called with the path
            assert mock_class.call_count == 1
            call_args = mock_class.call_args[0][0]
            assert call_args == pdf_path


# =============================================================================
# Office Document Tests (Word, Excel, PowerPoint)
# =============================================================================

# =============================================================================
# Helper Functions for Creating Test Office Documents
# =============================================================================


def create_minimal_docx() -> bytes:
    """
    Create a minimal valid Word document with basic content.
    
    Returns:
        DOCX file bytes
    """
    from docx import Document
    from io import BytesIO
    
    doc = Document()
    doc.add_paragraph("Sample content")
    
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.read()


def create_docx_with_content(
    paragraphs: list[str] = None,
    headings: list[tuple[str, int]] = None,  # (text, level)
    tables: list[list[list[str]]] = None,  # List of tables, each table is rows of cells
) -> bytes:
    """
    Create a Word document with specified content.
    
    Args:
        paragraphs: List of paragraph texts
        headings: List of (heading_text, level) tuples
        tables: List of tables, where each table is a list of rows, each row is a list of cell values
        
    Returns:
        DOCX file bytes
    """
    from docx import Document
    from io import BytesIO
    
    doc = Document()
    
    if headings:
        for heading_text, level in headings:
            doc.add_heading(heading_text, level=level)
    
    if paragraphs:
        for para in paragraphs:
            doc.add_paragraph(para)
    
    if tables:
        for table_data in tables:
            if table_data:
                num_rows = len(table_data)
                num_cols = len(table_data[0]) if table_data else 0
                table = doc.add_table(rows=num_rows, cols=num_cols)
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_value in enumerate(row_data):
                        table.rows[row_idx].cells[col_idx].text = cell_value
    
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.read()


def create_minimal_xlsx() -> bytes:
    """
    Create a minimal valid Excel spreadsheet.
    
    Returns:
        XLSX file bytes
    """
    from openpyxl import Workbook
    from io import BytesIO
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Sample data"
    
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output.read()


def create_xlsx_with_content(
    sheets: dict[str, list[list]] = None,  # {sheet_name: [[row1], [row2], ...]}
) -> bytes:
    """
    Create an Excel spreadsheet with specified content.
    
    Args:
        sheets: Dict mapping sheet names to row data (list of lists)
        
    Returns:
        XLSX file bytes
    """
    from openpyxl import Workbook
    from io import BytesIO
    
    wb = Workbook()
    
    # Remove default sheet
    default_sheet = wb.active
    
    if sheets:
        first = True
        for sheet_name, rows in sheets.items():
            if first:
                ws = default_sheet
                ws.title = sheet_name
                first = False
            else:
                ws = wb.create_sheet(title=sheet_name)
            
            for row_idx, row in enumerate(rows, 1):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)
    
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output.read()


def create_minimal_pptx() -> bytes:
    """
    Create a minimal valid PowerPoint presentation.
    
    Returns:
        PPTX file bytes
    """
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add a text box with content
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sample slide content"
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def create_pptx_with_content(
    slides: list[dict] = None,  # [{"texts": [...], "table": [[row1], [row2]]}]
) -> bytes:
    """
    Create a PowerPoint presentation with specified content.
    
    Args:
        slides: List of slide dicts with 'texts' (list of strings) and optional 'table' (list of rows)
        
    Returns:
        PPTX file bytes
    """
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout
    
    if slides:
        for slide_data in slides:
            slide = prs.slides.add_slide(slide_layout)
            
            texts = slide_data.get("texts", [])
            top_offset = 0.5
            for text in texts:
                left = Inches(0.5)
                top = Inches(top_offset)
                width = Inches(8)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = text
                top_offset += 0.7
            
            table_data = slide_data.get("table")
            if table_data and len(table_data) > 0:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0
                left = Inches(0.5)
                top = Inches(top_offset + 0.5)
                width = Inches(8)
                height = Inches(rows * 0.4)
                
                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_value in enumerate(row_data):
                        table.cell(row_idx, col_idx).text = str(cell_value)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# =============================================================================
# Word Document (DOCX) Tests
# =============================================================================


class TestDocxTextExtraction:
    """Tests for _extract_docx_text and _extract_docx_text_from_bytes methods."""

    def test_extract_basic_text(self):
        """Test extracting basic text from a Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(paragraphs=["Hello, World!", "This is a test."])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "Hello, World!" in result
        assert "This is a test." in result

    def test_extract_heading_level_1(self):
        """Test that Heading 1 is converted to markdown # format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(headings=[("Main Title", 1)])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "# Main Title" in result

    def test_extract_heading_level_2(self):
        """Test that Heading 2 is converted to markdown ## format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(headings=[("Section Title", 2)])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "## Section Title" in result

    def test_extract_heading_level_3(self):
        """Test that Heading 3 is converted to markdown ### format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(headings=[("Subsection", 3)])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "### Subsection" in result

    def test_extract_multiple_heading_levels(self):
        """Test extracting multiple heading levels."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(headings=[
            ("Title", 1),
            ("Section", 2),
            ("Subsection", 3),
        ])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "# Title" in result
        assert "## Section" in result
        assert "### Subsection" in result

    def test_extract_table(self):
        """Test extracting tables from Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        table_data = [
            ["Name", "Age", "City"],
            ["Alice", "30", "NYC"],
            ["Bob", "25", "LA"],
        ]
        docx_bytes = create_docx_with_content(tables=[table_data])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        # Tables should use pipe-separated format
        assert "Name | Age | City" in result
        assert "Alice | 30 | NYC" in result
        assert "Bob | 25 | LA" in result

    def test_extract_multiple_tables(self):
        """Test extracting multiple tables from Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        table1 = [["A1", "B1"], ["A2", "B2"]]
        table2 = [["X1", "Y1"], ["X2", "Y2"]]
        docx_bytes = create_docx_with_content(tables=[table1, table2])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "A1 | B1" in result
        assert "X1 | Y1" in result

    def test_extract_mixed_content(self):
        """Test extracting document with headings, paragraphs, and tables."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(
            headings=[("Document Title", 1), ("Data Section", 2)],
            paragraphs=["Introduction text here."],
            tables=[[["Col1", "Col2"], ["Data1", "Data2"]]],
        )
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "# Document Title" in result
        assert "## Data Section" in result
        assert "Introduction text here." in result
        assert "Col1 | Col2" in result

    def test_extract_from_file_path(self, tmp_path):
        """Test extracting text from a Word document file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(paragraphs=["File content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)
        
        result = _extract_docx_text(docx_path)
        
        assert "File content" in result

    def test_paragraphs_separated_by_double_newline(self):
        """Test that paragraphs are separated by double newlines."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(paragraphs=["Para 1", "Para 2"])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "\n\n" in result


class TestDocxLoadingViaFileType:
    """Tests for loading Word documents through type='file' with .docx extension."""

    def test_load_docx_via_file_type(self, tmp_path):
        """Test that DOCX files are loaded correctly via file type auto-detection."""
        docx_bytes = create_docx_with_content(paragraphs=["Auto-detected content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)
        
        config = Config(
            documents={
                "test_doc": {
                    "type": "file",
                    "path": str(docx_path),
                    "description": "Test Word document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("test_doc")
        
        assert "content" in result
        assert "Auto-detected content" in result["content"]
        assert result["format"] == "text"


class TestDocxLoadingViaDocxType:
    """Tests for loading Word documents through type='docx'."""

    def test_load_docx_via_docx_type_with_path(self, tmp_path):
        """Test that DOCX files are loaded correctly via docx type with path."""
        docx_bytes = create_docx_with_content(paragraphs=["Direct type content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)
        
        config = Config(
            documents={
                "direct_doc": {
                    "type": "docx",
                    "path": str(docx_path),
                    "description": "Direct docx type document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("direct_doc")
        
        assert "content" in result
        assert "Direct type content" in result["content"]
        assert result["format"] == "text"

    def test_load_docx_via_docx_type_with_url(self):
        """Test loading Word document via docx type with URL."""
        docx_bytes = create_docx_with_content(paragraphs=["URL content"])
        
        config = Config(
            documents={
                "url_doc": {
                    "type": "docx",
                    "url": "https://example.com/document.docx",
                    "description": "URL-based Word document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_doc")
            
            assert "content" in result
            assert "URL content" in result["content"]
            assert result["format"] == "text"

    def test_missing_docx_file_raises_error(self):
        """Test that missing Word document raises error."""
        config = Config(
            documents={
                "missing": {
                    "type": "docx",
                    "path": "/nonexistent/file.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("missing")
        
        assert "error" in result


class TestDocxLoadingViaHttpType:
    """Tests for loading Word documents through type='http' with DOCX content-type."""

    def test_load_docx_via_http_content_type(self):
        """Test that DOCX is detected via HTTP content-type header."""
        docx_bytes = create_docx_with_content(paragraphs=["HTTP content-type detection"])
        
        config = Config(
            documents={
                "http_doc": {
                    "type": "http",
                    "url": "https://example.com/document.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.headers = {"content-type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("http_doc")
            
            assert "content" in result
            assert "HTTP content-type detection" in result["content"]
            assert result["format"] == "text"

    def test_load_docx_via_http_url_extension(self):
        """Test that DOCX is detected via .docx URL extension."""
        docx_bytes = create_docx_with_content(paragraphs=["URL extension detection"])
        
        config = Config(
            documents={
                "url_doc": {
                    "type": "http",
                    "url": "https://example.com/document.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.headers = {"content-type": "application/octet-stream"}  # Generic content type
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_doc")
            
            assert "content" in result
            assert "URL extension detection" in result["content"]


class TestDocxEdgeCases:
    """Edge case tests for Word document handling."""

    def test_empty_docx_document(self):
        """Test handling of empty Word document."""
        from docx import Document
        from io import BytesIO
        
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        doc = Document()
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        docx_bytes = output.read()
        
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        # Should return empty or whitespace-only string
        assert isinstance(result, str)

    def test_docx_with_special_characters(self):
        """Test Word document with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        docx_bytes = create_docx_with_content(
            paragraphs=["Special chars: < > & \" ' $ % @ \u00e9\u00e8\u00f1"]
        )
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        assert "Special chars" in result

    def test_docx_with_empty_cells_in_table(self):
        """Test Word document with empty table cells."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        table_data = [
            ["Header1", "", "Header3"],
            ["Data1", "Data2", ""],
        ]
        docx_bytes = create_docx_with_content(tables=[table_data])
        result = _extract_docx_text_from_bytes(docx_bytes)
        
        # Empty cells should still be represented
        assert "Header1" in result
        assert "Header3" in result


# =============================================================================
# Excel Spreadsheet (XLSX) Tests
# =============================================================================


class TestXlsxTextExtraction:
    """Tests for _extract_xlsx_text and _extract_xlsx_text_from_bytes methods."""

    def test_extract_basic_content(self):
        """Test extracting basic content from Excel spreadsheet."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Sheet1": [["Hello", "World"], ["Data", "Here"]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "Hello | World" in result
        assert "Data | Here" in result

    def test_extract_with_sheet_markers(self):
        """Test that sheets have [Sheet: name] markers."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"MySheet": [["Content"]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "[Sheet: MySheet]" in result

    def test_extract_multi_sheet(self):
        """Test extracting content from multiple sheets."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={
                "Sales": [["Product", "Revenue"], ["Widget", "1000"]],
                "Costs": [["Item", "Cost"], ["Materials", "500"]],
            }
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "[Sheet: Sales]" in result
        assert "[Sheet: Costs]" in result
        assert "Product | Revenue" in result
        assert "Item | Cost" in result

    def test_skip_empty_rows(self):
        """Test that empty rows are skipped."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Sheet1": [
                ["Data1", "Data2"],
                [None, None],  # Empty row
                ["Data3", "Data4"],
            ]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "Data1 | Data2" in result
        assert "Data3 | Data4" in result
        # Empty row should not produce extra blank lines

    def test_extract_numeric_values(self):
        """Test extracting numeric values from Excel."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Numbers": [[1, 2, 3], [4.5, 5.5, 6.5]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "1 | 2 | 3" in result
        assert "4.5 | 5.5 | 6.5" in result

    def test_extract_from_file_path(self, tmp_path):
        """Test extracting content from Excel file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["File", "Content"]]}
        )
        xlsx_path = tmp_path / "test.xlsx"
        xlsx_path.write_bytes(xlsx_bytes)
        
        result = _extract_xlsx_text(xlsx_path)
        
        assert "File | Content" in result


class TestXlsxLoadingViaFileType:
    """Tests for loading Excel files through type='file' with .xlsx extension."""

    def test_load_xlsx_via_file_type(self, tmp_path):
        """Test that XLSX files are loaded correctly via file type auto-detection."""
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["Auto-detected", "Excel"]]}
        )
        xlsx_path = tmp_path / "test.xlsx"
        xlsx_path.write_bytes(xlsx_bytes)
        
        config = Config(
            documents={
                "test_excel": {
                    "type": "file",
                    "path": str(xlsx_path),
                    "description": "Test Excel spreadsheet",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("test_excel")
        
        assert "content" in result
        assert "Auto-detected | Excel" in result["content"]
        assert result["format"] == "text"


class TestXlsxLoadingViaXlsxType:
    """Tests for loading Excel files through type='xlsx'."""

    def test_load_xlsx_via_xlsx_type_with_path(self, tmp_path):
        """Test that XLSX files are loaded correctly via xlsx type with path."""
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["Direct", "Type"]]}
        )
        xlsx_path = tmp_path / "test.xlsx"
        xlsx_path.write_bytes(xlsx_bytes)
        
        config = Config(
            documents={
                "direct_excel": {
                    "type": "xlsx",
                    "path": str(xlsx_path),
                    "description": "Direct xlsx type document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("direct_excel")
        
        assert "content" in result
        assert "Direct | Type" in result["content"]
        assert result["format"] == "text"

    def test_load_xlsx_via_xlsx_type_with_url(self):
        """Test loading Excel file via xlsx type with URL."""
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["URL", "Content"]]}
        )
        
        config = Config(
            documents={
                "url_excel": {
                    "type": "xlsx",
                    "url": "https://example.com/spreadsheet.xlsx",
                    "description": "URL-based Excel spreadsheet",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = xlsx_bytes
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_excel")
            
            assert "content" in result
            assert "URL | Content" in result["content"]
            assert result["format"] == "text"

    def test_missing_xlsx_file_raises_error(self):
        """Test that missing Excel file raises error."""
        config = Config(
            documents={
                "missing": {
                    "type": "xlsx",
                    "path": "/nonexistent/file.xlsx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("missing")
        
        assert "error" in result


class TestXlsxLoadingViaHttpType:
    """Tests for loading Excel files through type='http' with XLSX content-type."""

    def test_load_xlsx_via_http_content_type(self):
        """Test that XLSX is detected via HTTP content-type header."""
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["HTTP", "Detection"]]}
        )
        
        config = Config(
            documents={
                "http_excel": {
                    "type": "http",
                    "url": "https://example.com/spreadsheet.xlsx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = xlsx_bytes
        mock_response.headers = {"content-type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("http_excel")
            
            assert "content" in result
            assert "HTTP | Detection" in result["content"]

    def test_load_xlsx_via_http_url_extension(self):
        """Test that XLSX is detected via .xlsx URL extension."""
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["URL", "Extension"]]}
        )
        
        config = Config(
            documents={
                "url_excel": {
                    "type": "http",
                    "url": "https://example.com/spreadsheet.xlsx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = xlsx_bytes
        mock_response.headers = {"content-type": "application/octet-stream"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_excel")
            
            assert "content" in result
            assert "URL | Extension" in result["content"]


class TestXlsxEdgeCases:
    """Edge case tests for Excel spreadsheet handling."""

    def test_empty_xlsx_spreadsheet(self):
        """Test handling of empty Excel spreadsheet."""
        from openpyxl import Workbook
        from io import BytesIO
        
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        wb = Workbook()
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        xlsx_bytes = output.read()
        
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        # Empty spreadsheet should return empty string
        assert isinstance(result, str)

    def test_xlsx_with_special_characters(self):
        """Test Excel spreadsheet with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Special": [["< > & \" '", "\u00e9\u00e8\u00f1"]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "< > & \"" in result

    def test_xlsx_with_none_values(self):
        """Test Excel spreadsheet with None values."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Data": [["Value", None, "Another"]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        # None values should be represented as empty strings
        assert "Value" in result
        assert "Another" in result

    def test_xlsx_single_cell(self):
        """Test Excel spreadsheet with single cell."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        xlsx_bytes = create_xlsx_with_content(
            sheets={"Single": [["Only cell"]]}
        )
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        assert "Only cell" in result

    def test_xlsx_many_sheets(self):
        """Test Excel spreadsheet with many sheets."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        sheets = {f"Sheet{i}": [[f"Content {i}"]] for i in range(1, 11)}
        xlsx_bytes = create_xlsx_with_content(sheets=sheets)
        result = _extract_xlsx_text_from_bytes(xlsx_bytes)
        
        for i in range(1, 11):
            assert f"[Sheet: Sheet{i}]" in result
            assert f"Content {i}" in result


# =============================================================================
# PowerPoint Presentation (PPTX) Tests
# =============================================================================


class TestPptxTextExtraction:
    """Tests for _extract_pptx_text and _extract_pptx_text_from_bytes methods."""

    def test_extract_basic_text(self):
        """Test extracting basic text from PowerPoint presentation."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Hello, World!", "Slide content here"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "Hello, World!" in result
        assert "Slide content here" in result

    def test_extract_with_slide_markers(self):
        """Test that slides have [Slide N] markers."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Content"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "[Slide 1]" in result

    def test_extract_multi_slide(self):
        """Test extracting content from multiple slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[
                {"texts": ["Slide 1 Title", "Slide 1 Content"]},
                {"texts": ["Slide 2 Title", "Slide 2 Content"]},
                {"texts": ["Slide 3 Title", "Slide 3 Content"]},
            ]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "[Slide 1]" in result
        assert "[Slide 2]" in result
        assert "[Slide 3]" in result
        assert "Slide 1 Title" in result
        assert "Slide 2 Content" in result
        assert "Slide 3 Title" in result

    def test_extract_table_from_slide(self):
        """Test extracting tables from PowerPoint slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{
                "texts": ["Table Slide"],
                "table": [
                    ["Header1", "Header2", "Header3"],
                    ["Data1", "Data2", "Data3"],
                ]
            }]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "Header1 | Header2 | Header3" in result
        assert "Data1 | Data2 | Data3" in result

    def test_extract_from_file_path(self, tmp_path):
        """Test extracting content from PowerPoint file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["File content"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)
        
        result = _extract_pptx_text(pptx_path)
        
        assert "File content" in result


class TestPptxLoadingViaFileType:
    """Tests for loading PowerPoint files through type='file' with .pptx extension."""

    def test_load_pptx_via_file_type(self, tmp_path):
        """Test that PPTX files are loaded correctly via file type auto-detection."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Auto-detected presentation"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)
        
        config = Config(
            documents={
                "test_pptx": {
                    "type": "file",
                    "path": str(pptx_path),
                    "description": "Test PowerPoint presentation",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("test_pptx")
        
        assert "content" in result
        assert "Auto-detected presentation" in result["content"]
        assert result["format"] == "text"


class TestPptxLoadingViaPptxType:
    """Tests for loading PowerPoint files through type='pptx'."""

    def test_load_pptx_via_pptx_type_with_path(self, tmp_path):
        """Test that PPTX files are loaded correctly via pptx type with path."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Direct type presentation"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)
        
        config = Config(
            documents={
                "direct_pptx": {
                    "type": "pptx",
                    "path": str(pptx_path),
                    "description": "Direct pptx type document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("direct_pptx")
        
        assert "content" in result
        assert "Direct type presentation" in result["content"]
        assert result["format"] == "text"

    def test_load_pptx_via_pptx_type_with_url(self):
        """Test loading PowerPoint file via pptx type with URL."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["URL presentation"]}]
        )
        
        config = Config(
            documents={
                "url_pptx": {
                    "type": "pptx",
                    "url": "https://example.com/presentation.pptx",
                    "description": "URL-based PowerPoint presentation",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_pptx")
            
            assert "content" in result
            assert "URL presentation" in result["content"]
            assert result["format"] == "text"

    def test_missing_pptx_file_raises_error(self):
        """Test that missing PowerPoint file raises error."""
        config = Config(
            documents={
                "missing": {
                    "type": "pptx",
                    "path": "/nonexistent/file.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        result = tools.get_document("missing")
        
        assert "error" in result


class TestPptxLoadingViaHttpType:
    """Tests for loading PowerPoint files through type='http' with PPTX content-type."""

    def test_load_pptx_via_http_content_type(self):
        """Test that PPTX is detected via HTTP content-type header."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["HTTP content-type detection"]}]
        )
        
        config = Config(
            documents={
                "http_pptx": {
                    "type": "http",
                    "url": "https://example.com/presentation.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.headers = {"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("http_pptx")
            
            assert "content" in result
            assert "HTTP content-type detection" in result["content"]

    def test_load_pptx_via_http_url_extension(self):
        """Test that PPTX is detected via .pptx URL extension."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["URL extension detection"]}]
        )
        
        config = Config(
            documents={
                "url_pptx": {
                    "type": "http",
                    "url": "https://example.com/presentation.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.headers = {"content-type": "application/octet-stream"}
        mock_response.raise_for_status = Mock()
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("url_pptx")
            
            assert "content" in result
            assert "URL extension detection" in result["content"]


class TestPptxEdgeCases:
    """Edge case tests for PowerPoint presentation handling."""

    def test_empty_pptx_presentation(self):
        """Test handling of empty PowerPoint presentation."""
        from pptx import Presentation
        from io import BytesIO
        
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        prs = Presentation()
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.read()
        
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        # Empty presentation should return empty string
        assert isinstance(result, str)
        assert result == ""

    def test_pptx_with_special_characters(self):
        """Test PowerPoint presentation with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Special: < > & \" ' \u00e9\u00e8\u00f1"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "Special:" in result

    def test_pptx_slide_with_only_table(self):
        """Test PowerPoint slide with only a table (no text shapes)."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{
                "texts": [],
                "table": [["Table", "Only"], ["Row", "Data"]]
            }]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "[Slide 1]" in result
        assert "Table | Only" in result

    def test_pptx_many_slides(self):
        """Test PowerPoint presentation with many slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        slides = [{"texts": [f"Slide {i} content"]} for i in range(1, 21)]
        pptx_bytes = create_pptx_with_content(slides=slides)
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        for i in range(1, 21):
            assert f"[Slide {i}]" in result
            assert f"Slide {i} content" in result

    def test_pptx_slide_numbering_starts_at_one(self):
        """Test that slide numbers start at 1, not 0."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["First slide"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)
        
        assert "[Slide 1]" in result
        assert "[Slide 0]" not in result


# =============================================================================
# Cross-format Integration Tests
# =============================================================================


@pytest.mark.usefixtures("clear_document_embeddings")
class TestOfficeDocumentIntegration:
    """Integration tests for Office documents in the discovery workflow."""

    def test_list_documents_includes_office_docs(self, tmp_path):
        """Test that Office documents appear in document listing."""
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(create_minimal_docx())
        xlsx_path = tmp_path / "test.xlsx"
        xlsx_path.write_bytes(create_minimal_xlsx())
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(create_minimal_pptx())
        
        config = Config(
            documents={
                "word_doc": {
                    "type": "docx",
                    "path": str(docx_path),
                    "description": "Word document",
                },
                "excel_doc": {
                    "type": "xlsx",
                    "path": str(xlsx_path),
                    "description": "Excel spreadsheet",
                },
                "ppt_doc": {
                    "type": "pptx",
                    "path": str(pptx_path),
                    "description": "PowerPoint presentation",
                },
            }
        )
        tools = DocumentDiscoveryTools(config)
        result = tools.list_documents()
        
        names = [doc["name"] for doc in result]
        assert "word_doc" in names
        assert "excel_doc" in names
        assert "ppt_doc" in names

    def test_office_content_searchable(self, tmp_path):
        """Test that Office document content can be searched."""
        docx_bytes = create_docx_with_content(paragraphs=["unique searchable keyword"])
        docx_path = tmp_path / "searchable.docx"
        docx_path.write_bytes(docx_bytes)
        
        config = Config(
            documents={
                "searchable_doc": {
                    "type": "docx",
                    "path": str(docx_path),
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        # First load the document
        tools.get_document("searchable_doc")
        
        # Then search
        result = tools.search_documents("unique searchable keyword", limit=5)
        
        assert len(result) > 0
        assert "searchable_doc" in result[0]["document"]

    def test_office_doc_loaded_only_once(self, tmp_path):
        """Test that Office documents are loaded once and cached."""
        docx_bytes = create_docx_with_content(paragraphs=["cached content"])
        docx_path = tmp_path / "cached.docx"
        docx_path.write_bytes(docx_bytes)
        
        config = Config(
            documents={
                "cached_doc": {
                    "type": "docx",
                    "path": str(docx_path),
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        # Load twice
        result1 = tools.get_document("cached_doc")
        result2 = tools.get_document("cached_doc")
        
        # Should be the same cached content
        assert result1["content"] == result2["content"]
        assert result1["loaded_at"] == result2["loaded_at"]


# =============================================================================
# Error Handling Tests
# =============================================================================


class TestOfficeDocumentErrorHandling:
    """Error handling tests for Office document loading."""

    def test_corrupt_docx_raises_error(self):
        """Test that corrupt DOCX data raises an error."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        corrupt_data = b"This is not a valid DOCX file"
        
        with pytest.raises(Exception):
            _extract_docx_text_from_bytes(corrupt_data)

    def test_corrupt_xlsx_raises_error(self):
        """Test that corrupt XLSX data raises an error."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        corrupt_data = b"This is not a valid XLSX file"
        
        with pytest.raises(Exception):
            _extract_xlsx_text_from_bytes(corrupt_data)

    def test_corrupt_pptx_raises_error(self):
        """Test that corrupt PPTX data raises an error."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)
        
        corrupt_data = b"This is not a valid PPTX file"
        
        with pytest.raises(Exception):
            _extract_pptx_text_from_bytes(corrupt_data)

    def test_http_office_fetch_error(self):
        """Test error handling when HTTP Office document fetch fails."""
        config = Config(
            documents={
                "http_doc": {
                    "type": "http",
                    "url": "https://example.com/document.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)
        
        mock_response = Mock()
        mock_response.raise_for_status = Mock(side_effect=Exception("404 Not Found"))
        
        with patch('requests.get', return_value=mock_response):
            result = tools.get_document("http_doc")
            
            assert "error" in result
            assert "Failed to load" in result["error"]
