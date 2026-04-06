# Copyright (c) 2025 Kenneth Stott
# Canary: 6a396d74-c4eb-4d41-8169-9268f9f268c6
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Tests for Word/DOCX document loading in doc_tools module."""

from __future__ import annotations

import io
import pytest
from pathlib import Path
from unittest.mock import Mock, patch

from constat.core.config import Config
from constat.discovery.doc_tools import DocumentDiscoveryTools
from constat.discovery.doc_tools import _transport as _doc_transport
from constat.discovery.doc_tools._file_extractors import (
    _extract_docx_text,
    _extract_docx_text_from_bytes,
    _extract_pptx_text,
    _extract_pptx_text_from_bytes,
)


# Reset module-level HTTP session singleton between tests
@pytest.fixture(autouse=True)
def _reset_http_session():
    _doc_transport._http_session = None
    yield
    _doc_transport._http_session = None


# =============================================================================
# Helper Functions for Creating Test Word Documents
# =============================================================================


def create_minimal_docx() -> bytes:
    """
    Create a minimal valid Word document with basic content.

    Returns:
        DOCX file bytes
    """
    from docx import Document
    from io import BytesIO

    doc = Document()
    doc.add_paragraph("Sample content")

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.read()


def create_docx_with_content(
    paragraphs: list[str] = None,
    headings: list[tuple[str, int]] = None,  # (text, level)
    tables: list[list[list[str]]] = None,  # List of tables, each table is rows of cells
) -> bytes:
    """
    Create a Word document with specified content.

    Args:
        paragraphs: List of paragraph texts
        headings: List of (heading_text, level) tuples
        tables: List of tables, where each table is a list of rows, each row is a list of cell values

    Returns:
        DOCX file bytes
    """
    from docx import Document
    from io import BytesIO

    doc = Document()

    if headings:
        for heading_text, level in headings:
            doc.add_heading(heading_text, level=level)

    if paragraphs:
        for para in paragraphs:
            doc.add_paragraph(para)

    if tables:
        for table_data in tables:
            if table_data:
                num_rows = len(table_data)
                num_cols = len(table_data[0]) if table_data else 0
                table = doc.add_table(rows=num_rows, cols=num_cols)
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_value in enumerate(row_data):
                        table.rows[row_idx].cells[col_idx].text = cell_value

    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output.read()


# =============================================================================
# Word Document (DOCX) Tests
# =============================================================================


class TestDocxTextExtraction:
    """Tests for _extract_docx_text and _extract_docx_text_from_bytes methods."""

    def test_extract_basic_text(self):
        """Test extracting basic text from a Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(paragraphs=["Hello, World!", "This is a test."])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "Hello, World!" in result
        assert "This is a test." in result

    def test_extract_heading_level_1(self):
        """Test that Heading 1 is converted to markdown # format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(headings=[("Main Title", 1)])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "# Main Title" in result

    def test_extract_heading_level_2(self):
        """Test that Heading 2 is converted to markdown ## format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(headings=[("Section Title", 2)])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "## Section Title" in result

    def test_extract_heading_level_3(self):
        """Test that Heading 3 is converted to markdown ### format."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(headings=[("Subsection", 3)])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "### Subsection" in result

    def test_extract_multiple_heading_levels(self):
        """Test extracting multiple heading levels."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(headings=[
            ("Title", 1),
            ("Section", 2),
            ("Subsection", 3),
        ])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "# Title" in result
        assert "## Section" in result
        assert "### Subsection" in result

    def test_extract_table(self):
        """Test extracting tables from Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        table_data = [
            ["Name", "Age", "City"],
            ["Alice", "30", "NYC"],
            ["Bob", "25", "LA"],
        ]
        docx_bytes = create_docx_with_content(tables=[table_data])
        result = _extract_docx_text_from_bytes(docx_bytes)

        # Tables should use pipe-separated format
        assert "Name | Age | City" in result
        assert "Alice | 30 | NYC" in result
        assert "Bob | 25 | LA" in result

    def test_extract_multiple_tables(self):
        """Test extracting multiple tables from Word document."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        table1 = [["A1", "B1"], ["A2", "B2"]]
        table2 = [["X1", "Y1"], ["X2", "Y2"]]
        docx_bytes = create_docx_with_content(tables=[table1, table2])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "A1 | B1" in result
        assert "X1 | Y1" in result

    def test_extract_mixed_content(self):
        """Test extracting document with headings, paragraphs, and tables."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(
            headings=[("Document Title", 1), ("Data Section", 2)],
            paragraphs=["Introduction text here."],
            tables=[[["Col1", "Col2"], ["Data1", "Data2"]]],
        )
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "# Document Title" in result
        assert "## Data Section" in result
        assert "Introduction text here." in result
        assert "Col1 | Col2" in result

    def test_extract_from_file_path(self, tmp_path):
        """Test extracting text from a Word document file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(paragraphs=["File content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)

        result = _extract_docx_text(docx_path)

        assert "File content" in result

    def test_paragraphs_separated_by_double_newline(self):
        """Test that paragraphs are separated by double newlines."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(paragraphs=["Para 1", "Para 2"])
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "\n\n" in result


class TestDocxLoadingViaFileType:
    """Tests for loading Word documents through type='file' with .docx extension."""

    def test_load_docx_via_file_type(self, tmp_path):
        """Test that DOCX files are loaded correctly via file type auto-detection."""
        docx_bytes = create_docx_with_content(paragraphs=["Auto-detected content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)

        config = Config(
            documents={
                "test_doc": {
                    "type": "file",
                    "path": str(docx_path),
                    "description": "Test Word document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("test_doc")

        assert "content" in result
        assert "Auto-detected content" in result["content"]
        assert result["format"] == "text"


class TestDocxLoadingViaDocxType:
    """Tests for loading Word documents through type='docx'."""

    def test_load_docx_via_docx_type_with_path(self, tmp_path):
        """Test that DOCX files are loaded correctly via docx type with path."""
        docx_bytes = create_docx_with_content(paragraphs=["Direct type content"])
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)

        config = Config(
            documents={
                "direct_doc": {
                    "type": "docx",
                    "path": str(docx_path),
                    "description": "Direct docx type document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("direct_doc")

        assert "content" in result
        assert "Direct type content" in result["content"]
        assert result["format"] == "text"

    def test_load_docx_via_docx_type_with_url(self):
        """Test loading Word document via docx type with URL."""
        docx_bytes = create_docx_with_content(paragraphs=["URL content"])

        config = Config(
            documents={
                "url_doc": {
                    "type": "docx",
                    "url": "https://example.com/document.docx",
                    "description": "URL-based Word document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("url_doc")

            assert "content" in result
            assert "URL content" in result["content"]
            assert result["format"] == "text"

    def test_missing_docx_file_raises_error(self):
        """Test that missing Word document raises error."""
        config = Config(
            documents={
                "missing": {
                    "type": "docx",
                    "path": "/nonexistent/file.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("missing")

        assert "error" in result


class TestDocxLoadingViaHttpType:
    """Tests for loading Word documents through type='http' with DOCX content-type."""

    def test_load_docx_via_http_content_type(self):
        """Test that DOCX is detected via HTTP content-type header."""
        docx_bytes = create_docx_with_content(paragraphs=["HTTP content-type detection"])

        config = Config(
            documents={
                "http_doc": {
                    "type": "http",
                    "url": "https://example.com/document.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.headers = {"content-type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"}
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("http_doc")

            assert "content" in result
            assert "HTTP content-type detection" in result["content"]
            assert result["format"] == "text"

    def test_load_docx_via_http_url_extension(self):
        """Test that DOCX is detected via .docx URL extension."""
        docx_bytes = create_docx_with_content(paragraphs=["URL extension detection"])

        config = Config(
            documents={
                "url_doc": {
                    "type": "http",
                    "url": "https://example.com/document.docx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = docx_bytes
        mock_response.headers = {"content-type": "application/octet-stream"}  # Generic content type
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("url_doc")

            assert "content" in result
            assert "URL extension detection" in result["content"]


class TestDocxEdgeCases:
    """Edge case tests for Word document handling."""

    def test_empty_docx_document(self):
        """Test handling of empty Word document."""
        from docx import Document
        from io import BytesIO

        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        doc = Document()
        output = BytesIO()
        doc.save(output)
        output.seek(0)
        docx_bytes = output.read()

        result = _extract_docx_text_from_bytes(docx_bytes)

        # Should return empty or whitespace-only string
        assert isinstance(result, str)

    def test_docx_with_special_characters(self):
        """Test Word document with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        docx_bytes = create_docx_with_content(
            paragraphs=["Special chars: < > & \" ' $ % @ \u00e9\u00e8\u00f1"]
        )
        result = _extract_docx_text_from_bytes(docx_bytes)

        assert "Special chars" in result

    def test_docx_with_empty_cells_in_table(self):
        """Test Word document with empty table cells."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        table_data = [
            ["Header1", "", "Header3"],
            ["Data1", "Data2", ""],
        ]
        docx_bytes = create_docx_with_content(tables=[table_data])
        result = _extract_docx_text_from_bytes(docx_bytes)

        # Empty cells should still be represented
        assert "Header1" in result
        assert "Header3" in result


# =============================================================================
# Helper Functions for Creating Test PowerPoint Files
# =============================================================================


def create_minimal_pptx() -> bytes:
    """
    Create a minimal valid PowerPoint presentation.

    Returns:
        PPTX file bytes
    """
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a text box with content
    left = Inches(1)
    top = Inches(1)
    width = Inches(5)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sample slide content"

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def create_pptx_with_content(
    slides: list[dict] = None,  # [{"texts": [...], "table": [[row1], [row2]]}]
) -> bytes:
    """
    Create a PowerPoint presentation with specified content.

    Args:
        slides: List of slide dicts with 'texts' (list of strings) and optional 'table' (list of rows)

    Returns:
        PPTX file bytes
    """
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # Blank layout

    if slides:
        for slide_data in slides:
            slide = prs.slides.add_slide(slide_layout)

            texts = slide_data.get("texts", [])
            top_offset = 0.5
            for text in texts:
                left = Inches(0.5)
                top = Inches(top_offset)
                width = Inches(8)
                height = Inches(0.5)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = text
                top_offset += 0.7

            table_data = slide_data.get("table")
            if table_data and len(table_data) > 0:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0
                left = Inches(0.5)
                top = Inches(top_offset + 0.5)
                width = Inches(8)
                height = Inches(rows * 0.4)

                table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_value in enumerate(row_data):
                        table.cell(row_idx, col_idx).text = str(cell_value)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


# =============================================================================
# PowerPoint Presentation (PPTX) Tests
# =============================================================================


class TestPptxTextExtraction:
    """Tests for _extract_pptx_text and _extract_pptx_text_from_bytes methods."""

    def test_extract_basic_text(self):
        """Test extracting basic text from PowerPoint presentation."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Hello, World!", "Slide content here"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "Hello, World!" in result
        assert "Slide content here" in result

    def test_extract_with_slide_markers(self):
        """Test that slides have [Slide N] markers."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Content"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "[Slide 1]" in result

    def test_extract_multi_slide(self):
        """Test extracting content from multiple slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[
                {"texts": ["Slide 1 Title", "Slide 1 Content"]},
                {"texts": ["Slide 2 Title", "Slide 2 Content"]},
                {"texts": ["Slide 3 Title", "Slide 3 Content"]},
            ]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "[Slide 1]" in result
        assert "[Slide 2]" in result
        assert "[Slide 3]" in result
        assert "Slide 1 Title" in result
        assert "Slide 2 Content" in result
        assert "Slide 3 Title" in result

    def test_extract_table_from_slide(self):
        """Test extracting tables from PowerPoint slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{
                "texts": ["Table Slide"],
                "table": [
                    ["Header1", "Header2", "Header3"],
                    ["Data1", "Data2", "Data3"],
                ]
            }]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "Header1 | Header2 | Header3" in result
        assert "Data1 | Data2 | Data3" in result

    def test_extract_from_file_path(self, tmp_path):
        """Test extracting content from PowerPoint file path."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["File content"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)

        result = _extract_pptx_text(pptx_path)

        assert "File content" in result


class TestPptxLoadingViaFileType:
    """Tests for loading PowerPoint files through type='file' with .pptx extension."""

    def test_load_pptx_via_file_type(self, tmp_path):
        """Test that PPTX files are loaded correctly via file type auto-detection."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Auto-detected presentation"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)

        config = Config(
            documents={
                "test_pptx": {
                    "type": "file",
                    "path": str(pptx_path),
                    "description": "Test PowerPoint presentation",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("test_pptx")

        assert "content" in result
        assert "Auto-detected presentation" in result["content"]
        assert result["format"] == "text"


class TestPptxLoadingViaPptxType:
    """Tests for loading PowerPoint files through type='pptx'."""

    def test_load_pptx_via_pptx_type_with_path(self, tmp_path):
        """Test that PPTX files are loaded correctly via pptx type with path."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Direct type presentation"]}]
        )
        pptx_path = tmp_path / "test.pptx"
        pptx_path.write_bytes(pptx_bytes)

        config = Config(
            documents={
                "direct_pptx": {
                    "type": "pptx",
                    "path": str(pptx_path),
                    "description": "Direct pptx type document",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("direct_pptx")

        assert "content" in result
        assert "Direct type presentation" in result["content"]
        assert result["format"] == "text"

    def test_load_pptx_via_pptx_type_with_url(self):
        """Test loading PowerPoint file via pptx type with URL."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["URL presentation"]}]
        )

        config = Config(
            documents={
                "url_pptx": {
                    "type": "pptx",
                    "url": "https://example.com/presentation.pptx",
                    "description": "URL-based PowerPoint presentation",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("url_pptx")

            assert "content" in result
            assert "URL presentation" in result["content"]
            assert result["format"] == "text"

    def test_missing_pptx_file_raises_error(self):
        """Test that missing PowerPoint file raises error."""
        config = Config(
            documents={
                "missing": {
                    "type": "pptx",
                    "path": "/nonexistent/file.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        result = tools.get_document("missing")

        assert "error" in result


class TestPptxLoadingViaHttpType:
    """Tests for loading PowerPoint files through type='http' with PPTX content-type."""

    def test_load_pptx_via_http_content_type(self):
        """Test that PPTX is detected via HTTP content-type header."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["HTTP content-type detection"]}]
        )

        config = Config(
            documents={
                "http_pptx": {
                    "type": "http",
                    "url": "https://example.com/presentation.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.headers = {"content-type": "application/vnd.openxmlformats-officedocument.presentationml.presentation"}
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("http_pptx")

            assert "content" in result
            assert "HTTP content-type detection" in result["content"]

    def test_load_pptx_via_http_url_extension(self):
        """Test that PPTX is detected via .pptx URL extension."""
        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["URL extension detection"]}]
        )

        config = Config(
            documents={
                "url_pptx": {
                    "type": "http",
                    "url": "https://example.com/presentation.pptx",
                }
            }
        )
        tools = DocumentDiscoveryTools(config)

        mock_response = Mock()
        mock_response.content = pptx_bytes
        mock_response.headers = {"content-type": "application/octet-stream"}
        mock_response.raise_for_status = Mock()

        with patch('constat.discovery.doc_tools._transport._get_http_session', return_value=Mock(get=Mock(return_value=mock_response))):
            result = tools.get_document("url_pptx")

            assert "content" in result
            assert "URL extension detection" in result["content"]


class TestPptxEdgeCases:
    """Edge case tests for PowerPoint presentation handling."""

    def test_empty_pptx_presentation(self):
        """Test handling of empty PowerPoint presentation."""
        from pptx import Presentation
        from io import BytesIO

        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        prs = Presentation()
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        pptx_bytes = output.read()

        result = _extract_pptx_text_from_bytes(pptx_bytes)

        # Empty presentation should return empty string
        assert isinstance(result, str)
        assert result == ""

    def test_pptx_with_special_characters(self):
        """Test PowerPoint presentation with special characters."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["Special: < > & \" ' \u00e9\u00e8\u00f1"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "Special:" in result

    def test_pptx_slide_with_only_table(self):
        """Test PowerPoint slide with only a table (no text shapes)."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{
                "texts": [],
                "table": [["Table", "Only"], ["Row", "Data"]]
            }]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "[Slide 1]" in result
        assert "Table | Only" in result

    def test_pptx_many_slides(self):
        """Test PowerPoint presentation with many slides."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        slides = [{"texts": [f"Slide {i} content"]} for i in range(1, 21)]
        pptx_bytes = create_pptx_with_content(slides=slides)
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        for i in range(1, 21):
            assert f"[Slide {i}]" in result
            assert f"Slide {i} content" in result

    def test_pptx_slide_numbering_starts_at_one(self):
        """Test that slide numbers start at 1, not 0."""
        config = Config(documents={})
        tools = DocumentDiscoveryTools(config)

        pptx_bytes = create_pptx_with_content(
            slides=[{"texts": ["First slide"]}]
        )
        result = _extract_pptx_text_from_bytes(pptx_bytes)

        assert "[Slide 1]" in result
        assert "[Slide 0]" not in result
