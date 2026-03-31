# Copyright (c) 2025 Kenneth Stott
# Canary: 62c1b5ad-61d0-47ef-95d9-23a1191bbe7a
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.

"""Tests for embedded image extraction from PDFs and Office documents."""

import io
import logging
import logging.handlers
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from PIL import Image as PILImage

from constat.discovery.doc_tools._file_extractors import (
    MIN_IMAGE_BYTES,
    MIN_IMAGE_DIMENSION,
    ExtractedImage,
    _extract_images_from_document,
    _extract_pdf_images,
    _extract_docx_images,
    _extract_pptx_images,
    _extract_xlsx_images,
    _fetch_linked_image,
    _filter_and_dedup,
    _guess_mime,
    _is_emf_wmf,
    _passes_size_filter,
)


# ---------------------------------------------------------------------------
# Helpers — programmatic fixture generators
# ---------------------------------------------------------------------------

def _make_png_bytes(width: int = 200, height: int = 200, color: str = "red") -> bytes:
    """Create a real PNG image in memory with noise to avoid tiny compressed size."""
    import random
    img = PILImage.new("RGB", (width, height), color=color)
    # Add noise so PNG doesn't compress to < MIN_IMAGE_BYTES
    pixels = img.load()
    rng = random.Random(42)
    for y in range(height):
        for x in range(width):
            r, g, b = pixels[x, y]
            pixels[x, y] = (
                min(255, r + rng.randint(0, 30)),
                min(255, g + rng.randint(0, 30)),
                min(255, b + rng.randint(0, 30)),
            )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def _make_jpeg_bytes(width: int = 200, height: int = 200, color: str = "blue") -> bytes:
    """Create a real JPEG image in memory."""
    img = PILImage.new("RGB", (width, height), color=color)
    buf = io.BytesIO()
    img.save(buf, format="JPEG")
    return buf.getvalue()


def _make_pdf_with_images(num_pages: int = 2, images_per_page: int = 1) -> bytes:
    """Create a PDF with embedded PNG images using reportlab."""
    import random
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    rng = random.Random(42)
    pdf_buf = io.BytesIO()
    c = canvas.Canvas(pdf_buf, pagesize=letter)
    for page_idx in range(num_pages):
        for img_idx in range(images_per_page):
            img = PILImage.new("RGB", (100, 100))
            pixels = [
                (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
                for _ in range(100 * 100)
            ]
            img.putdata(pixels)
            c.drawImage(ImageReader(img), 72 + img_idx * 110, 600, width=100, height=100)
        c.showPage()
    c.save()
    return pdf_buf.getvalue()


def _make_docx_with_images(count: int = 2) -> bytes:
    """Create a DOCX with embedded images."""
    from docx import Document
    from docx.shared import Inches

    doc = Document()
    doc.add_paragraph("Test document with images")
    for i in range(count):
        color = f"#{(i * 70) % 256:02x}80{(i * 40) % 256:02x}"
        img_bytes = _make_png_bytes(150, 150, color=color)
        doc.add_picture(io.BytesIO(img_bytes), width=Inches(1.5))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_with_images(slides_with_images: list[int] | None = None) -> bytes:
    """Create a PPTX with images on specified slides."""
    from pptx import Presentation
    from pptx.util import Inches

    if slides_with_images is None:
        slides_with_images = [1]

    prs = Presentation()
    for slide_num in range(1, max(slides_with_images) + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        if slide_num in slides_with_images:
            color = f"#{(slide_num * 60) % 256:02x}9090"
            img_bytes = _make_png_bytes(200, 200, color=color)
            slide.shapes.add_picture(
                io.BytesIO(img_bytes), Inches(1), Inches(1), Inches(2), Inches(2),
            )
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx_with_images(count: int = 1) -> bytes:
    """Create an XLSX with embedded images."""
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XlImage

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws["A1"] = "Test"

    for i in range(count):
        color = f"#{(i * 80) % 256:02x}70{(i * 50) % 256:02x}"
        img_bytes = _make_png_bytes(120, 120, color=color)
        img_buf = io.BytesIO(img_bytes)
        img = XlImage(img_buf)
        ws.add_image(img, f"B{i + 2}")

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Unit tests — helper functions
# ---------------------------------------------------------------------------

class TestGuessAndFilter:
    def test_guess_mime_png(self):
        data = _make_png_bytes(10, 10)
        assert _guess_mime(data) == "image/png"

    def test_guess_mime_jpeg(self):
        data = _make_jpeg_bytes(10, 10)
        assert _guess_mime(data) == "image/jpeg"

    def test_guess_mime_unknown(self):
        assert _guess_mime(b"\x00\x00\x00\x00") == "image/png"

    def test_is_emf_wmf(self):
        assert _is_emf_wmf("logo.emf", "image/x-emf")
        assert _is_emf_wmf("chart.wmf", "image/x-wmf")
        assert _is_emf_wmf("logo.EMF", "image/png")
        assert not _is_emf_wmf("photo.png", "image/png")

    def test_passes_size_filter_too_small_bytes(self):
        tiny = _make_png_bytes(5, 5)
        # tiny PNG is likely under MIN_IMAGE_BYTES, but let's be explicit
        if len(tiny) < MIN_IMAGE_BYTES:
            assert not _passes_size_filter(tiny)
        else:
            # Still should fail on dimension
            assert not _passes_size_filter(tiny)

    def test_passes_size_filter_ok(self):
        data = _make_png_bytes(200, 200)
        assert _passes_size_filter(data)

    def test_passes_size_filter_small_dimensions(self):
        data = _make_png_bytes(30, 30)
        # 30x30 is below MIN_IMAGE_DIMENSION=64, but we need to also check bytes
        if len(data) >= MIN_IMAGE_BYTES:
            assert not _passes_size_filter(data)


class TestFilterAndDedup:
    def test_dedup_identical_images(self):
        img_data = _make_png_bytes(100, 100, color="green")
        images = [
            ExtractedImage(name="img_1.png", data=img_data, mime_type="image/png", page=1, index=1),
            ExtractedImage(name="img_2.png", data=img_data, mime_type="image/png", page=2, index=1),
        ]
        result = _filter_and_dedup(images)
        assert len(result) == 1
        assert result[0].name == "img_1.png"

    def test_emf_wmf_skip(self):
        img_data = _make_png_bytes(100, 100)
        images = [
            ExtractedImage(name="chart.emf", data=img_data, mime_type="image/x-emf", page=None, index=1),
            ExtractedImage(name="ok.png", data=img_data, mime_type="image/png", page=None, index=2),
        ]
        _logger = logging.getLogger("constat.discovery.doc_tools._file_extractors")
        handler = logging.handlers.MemoryHandler(capacity=100)
        _logger.addHandler(handler)
        try:
            result = _filter_and_dedup(images)
        finally:
            _logger.removeHandler(handler)
        assert len(result) == 1
        assert result[0].name == "ok.png"
        assert any("Skipping EMF/WMF" in r.getMessage() for r in handler.buffer)

    def test_size_filter(self):
        # Create a very small image (< MIN_IMAGE_BYTES)
        tiny_img = PILImage.new("RGB", (2, 2), color="red")
        buf = io.BytesIO()
        tiny_img.save(buf, format="PNG")
        tiny_data = buf.getvalue()

        images = [
            ExtractedImage(name="tiny.png", data=tiny_data, mime_type="image/png", page=1, index=1),
        ]
        result = _filter_and_dedup(images)
        assert len(result) == 0


# ---------------------------------------------------------------------------
# Unit tests — format-specific extractors
# ---------------------------------------------------------------------------

class TestExtractPdfImages:
    def test_extract_pdf_images(self):
        pdf_bytes = _make_pdf_with_images(num_pages=2, images_per_page=1)
        images = _extract_pdf_images(None, pdf_bytes)
        assert len(images) >= 1
        for img in images:
            assert img.mime_type.startswith("image/")
            assert img.page is not None
            assert img.page >= 1
            assert "page_" in img.name
            assert len(img.data) > 0

    def test_extract_pdf_images_from_path(self, tmp_path):
        pdf_bytes = _make_pdf_with_images(num_pages=1, images_per_page=1)
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(pdf_bytes)
        images = _extract_pdf_images(pdf_path, None)
        assert len(images) >= 1


class TestExtractDocxImages:
    def test_extract_docx_images(self):
        docx_bytes = _make_docx_with_images(count=2)
        images = _extract_docx_images(None, docx_bytes)
        assert len(images) == 2
        for img in images:
            assert img.mime_type.startswith("image/")
            assert len(img.data) > 0

    def test_extract_docx_images_from_path(self, tmp_path):
        docx_bytes = _make_docx_with_images(count=1)
        docx_path = tmp_path / "test.docx"
        docx_path.write_bytes(docx_bytes)
        images = _extract_docx_images(docx_path, None)
        assert len(images) == 1


class TestExtractPptxImages:
    def test_extract_pptx_images(self):
        pptx_bytes = _make_pptx_with_images(slides_with_images=[1, 3])
        images = _extract_pptx_images(None, pptx_bytes)
        assert len(images) == 2
        # Check slide numbering in names
        names = [img.name for img in images]
        assert any("slide_1" in n for n in names)
        assert any("slide_3" in n for n in names)

    def test_extract_pptx_page_numbers(self):
        pptx_bytes = _make_pptx_with_images(slides_with_images=[2])
        images = _extract_pptx_images(None, pptx_bytes)
        assert len(images) == 1
        assert images[0].page == 2


class TestExtractXlsxImages:
    def test_extract_xlsx_images(self):
        xlsx_bytes = _make_xlsx_with_images(count=1)
        images = _extract_xlsx_images(None, xlsx_bytes)
        assert len(images) == 1
        assert "sheet_Data" in images[0].name
        assert images[0].mime_type.startswith("image/")


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

class TestDispatcher:
    def test_dispatch_pdf(self):
        pdf_bytes = _make_pdf_with_images(1, 1)
        images = _extract_images_from_document(None, pdf_bytes, "pdf")
        assert len(images) >= 1

    def test_dispatch_docx(self):
        docx_bytes = _make_docx_with_images(1)
        images = _extract_images_from_document(None, docx_bytes, "docx")
        assert len(images) == 1

    def test_dispatch_pptx(self):
        pptx_bytes = _make_pptx_with_images([1])
        images = _extract_images_from_document(None, pptx_bytes, "pptx")
        assert len(images) == 1

    def test_dispatch_xlsx(self):
        xlsx_bytes = _make_xlsx_with_images(1)
        images = _extract_images_from_document(None, xlsx_bytes, "xlsx")
        assert len(images) == 1

    def test_dispatch_unknown_returns_empty(self):
        assert _extract_images_from_document(None, b"data", "html") == []

    def test_dispatch_with_config_dir(self):
        docx_bytes = _make_docx_with_images(1)
        images = _extract_images_from_document(None, docx_bytes, "docx", config_dir="/tmp")
        assert len(images) == 1


# ---------------------------------------------------------------------------
# Linked image fetching
# ---------------------------------------------------------------------------

class TestFetchLinkedImage:
    def test_fetch_http(self):
        img_bytes = _make_png_bytes(100, 100)
        mock_resp = MagicMock()
        mock_resp.content = img_bytes
        mock_resp.raise_for_status = MagicMock()

        mock_session = MagicMock()
        mock_session.get.return_value = mock_resp

        with patch(
            "constat.discovery.doc_tools._transport._get_http_session",
            return_value=mock_session,
        ):
            result = _fetch_linked_image("https://example.com/logo.png")

        assert result == img_bytes
        mock_session.get.assert_called_once_with("https://example.com/logo.png", timeout=15)

    def test_fetch_file_absolute(self, tmp_path):
        img_bytes = _make_png_bytes(80, 80)
        img_path = tmp_path / "logo.png"
        img_path.write_bytes(img_bytes)
        result = _fetch_linked_image(str(img_path))
        assert result == img_bytes

    def test_fetch_file_relative(self, tmp_path):
        img_bytes = _make_png_bytes(80, 80)
        (tmp_path / "images").mkdir()
        img_path = tmp_path / "images" / "logo.png"
        img_path.write_bytes(img_bytes)
        result = _fetch_linked_image("images/logo.png", config_dir=str(tmp_path))
        assert result == img_bytes

    def test_fetch_file_uri(self, tmp_path):
        img_bytes = _make_png_bytes(80, 80)
        img_path = tmp_path / "logo.png"
        img_path.write_bytes(img_bytes)
        result = _fetch_linked_image(f"file://{img_path}")
        assert result == img_bytes

    def test_fetch_not_found(self):
        _logger = logging.getLogger("constat.discovery.doc_tools._file_extractors")
        handler = logging.handlers.MemoryHandler(capacity=100)
        _logger.addHandler(handler)
        try:
            result = _fetch_linked_image("/nonexistent/path/img.png")
        finally:
            _logger.removeHandler(handler)
        assert result is None
        assert any("Failed to fetch" in r.getMessage() for r in handler.buffer)

    def test_fetch_http_error(self):
        mock_session = MagicMock()
        mock_session.get.side_effect = ConnectionError("fail")
        _logger = logging.getLogger("constat.discovery.doc_tools._file_extractors")
        handler = logging.handlers.MemoryHandler(capacity=100)
        _logger.addHandler(handler)
        try:
            with patch(
                "constat.discovery.doc_tools._transport._get_http_session",
                return_value=mock_session,
            ):
                result = _fetch_linked_image("https://bad.example.com/img.png")
        finally:
            _logger.removeHandler(handler)
        assert result is None
        assert any("Failed to fetch" in r.getMessage() for r in handler.buffer)


# ---------------------------------------------------------------------------
# Integration — _load_document with extract_images=True
# ---------------------------------------------------------------------------

class TestIntegrationLoadDocument:
    def test_load_document_extracts_images(self, tmp_path):
        """Verify _load_document creates child LoadedDocument entries for embedded images."""
        from constat.core.config import DocumentConfig
        from constat.discovery.doc_tools._core import _CoreMixin
        from constat.discovery.doc_tools._transport import FetchResult
        from constat.discovery.doc_tools._image import OcrResult

        pdf_bytes = _make_pdf_with_images(num_pages=1, images_per_page=1)
        pdf_path = tmp_path / "doc.pdf"
        pdf_path.write_bytes(pdf_bytes)

        mixin = object.__new__(_CoreMixin)
        mixin._loaded_documents = {}
        mixin._image_labels = []
        mixin._router = None
        mixin.config = MagicMock()
        mixin.config.config_dir = str(tmp_path)
        mixin.config.documents = {}

        doc_config = DocumentConfig(
            path=str(pdf_path),
            extract_images=True,
        )
        mixin.config.documents["test_doc"] = doc_config

        fetch_result = FetchResult(
            data=pdf_bytes,
            source_path=str(pdf_path),
            detected_mime="application/pdf",
        )

        with (
            patch.object(mixin, "_resolve_doc_config", return_value=doc_config),
            patch(
                "constat.discovery.doc_tools._core.fetch_document",
                return_value=fetch_result,
            ),
            patch(
                "constat.discovery.doc_tools._core.infer_transport",
                return_value="file",
            ),
            patch(
                "constat.discovery.doc_tools._core._infer_structured_schema",
                return_value=None,
            ),
            patch("constat.discovery.doc_tools._image._ocr_extract") as mock_ocr,
        ):
            mock_ocr.return_value = OcrResult(text="", mean_confidence=0.0, word_count=0)
            mixin._load_document("test_doc")

        assert "test_doc" in mixin._loaded_documents
        child_names = [n for n in mixin._loaded_documents if n.startswith("test_doc:")]
        assert len(child_names) >= 1
        for child_name in child_names:
            child = mixin._loaded_documents[child_name]
            assert child.format == "markdown"
            assert "# Image:" in child.content

    def test_add_document_from_file_image_with_vision(self, tmp_path):
        """Verify add_document_from_file generates vision description and chunks it."""
        from constat.core.config import Config, DocumentConfig
        from constat.discovery.doc_tools._core import _CoreMixin
        from constat.discovery.doc_tools._image import OcrResult

        # Create a test image
        img_bytes = _make_png_bytes(200, 200, color="blue")
        img_path = tmp_path / "landscape.png"
        img_path.write_bytes(img_bytes)

        # Set up mixin with a mock router that returns a vision description
        mixin = object.__new__(_CoreMixin)
        mixin._loaded_documents = {}
        mixin._image_labels = []
        mixin.config = MagicMock()
        mixin.config.config_dir = str(tmp_path)

        mock_router = MagicMock()
        mock_router.generate_vision.return_value = (
            '{"description": "A landscape with mountains under a blue sky",'
            ' "subcategory": "photograph",'
            ' "labels": ["mountains", "landscape", "sky"]}'
        )
        mixin._router = mock_router

        # Mock vector store and embedding model
        mixin._vector_store = MagicMock()
        mixin._vector_store.add_chunks = MagicMock()
        mock_model = MagicMock()
        import numpy as np
        mock_model.encode.return_value = np.zeros((1, 1024))
        mixin._model = mock_model
        mixin._model_lock = __import__("threading").Lock()
        mixin.CHUNK_SIZE = 1500
        mixin.TABLE_CHUNK_LIMIT = 6000

        with patch("constat.discovery.doc_tools._image._ocr_extract") as mock_ocr:
            mock_ocr.return_value = OcrResult(text="", mean_confidence=0.0, word_count=0)
            with patch("pathlib.Path.read_text", return_value="Describe this image"):
                success, msg = mixin.add_document_from_file(
                    str(img_path),
                    name="images:landscape",
                    description="Test image",
                    domain_id="__base__",
                    skip_entity_extraction=True,
                )

        assert success, msg
        assert "images:landscape" in mixin._loaded_documents
        doc = mixin._loaded_documents["images:landscape"]
        assert "mountains" in doc.content
        assert "landscape" in doc.content
        assert "## Description" in doc.content
        assert "## Labels" in doc.content

        # Verify it was chunked and embedded
        assert mixin._vector_store.add_chunks.called
        chunks = mixin._vector_store.add_chunks.call_args[0][0]
        chunk_text = " ".join(c.content for c in chunks)
        assert "mountains" in chunk_text

        # Verify image labels were collected for NER
        assert "mountains" in mixin._image_labels
        assert "landscape" in mixin._image_labels
        assert "sky" in mixin._image_labels

    def test_add_document_from_file_image_with_ocr(self, tmp_path):
        """Verify OCR text flows through add_document_from_file into chunks."""
        from constat.discovery.doc_tools._core import _CoreMixin
        from constat.discovery.doc_tools._image import OcrResult
        import numpy as np

        img_bytes = _make_png_bytes(200, 200, color="white")
        img_path = tmp_path / "scanned.png"
        img_path.write_bytes(img_bytes)

        mixin = object.__new__(_CoreMixin)
        mixin._loaded_documents = {}
        mixin._image_labels = []
        mixin._router = None
        mixin.config = MagicMock()
        mixin.config.config_dir = str(tmp_path)
        mixin._vector_store = MagicMock()
        mixin._vector_store.add_chunks = MagicMock()
        mock_model = MagicMock()
        mock_model.encode.return_value = np.zeros((1, 1024))
        mixin._model = mock_model
        mixin._model_lock = __import__("threading").Lock()
        mixin.CHUNK_SIZE = 1500
        mixin.TABLE_CHUNK_LIMIT = 6000

        ocr_text = " ".join(f"word{i}" for i in range(50)) + " EXECUTIVE SUMMARY Regional Sales Report Q4 2025"
        with patch("constat.discovery.doc_tools._image._ocr_extract") as mock_ocr:
            mock_ocr.return_value = OcrResult(
                text=ocr_text, mean_confidence=72.0, word_count=55,
            )
            success, msg = mixin.add_document_from_file(
                str(img_path),
                name="images:scanned",
                description="Scanned report",
                domain_id="__base__",
                skip_entity_extraction=True,
            )

        assert success, msg
        doc = mixin._loaded_documents["images:scanned"]
        # text-primary images render OCR text directly under the name heading
        assert "EXECUTIVE SUMMARY" in doc.content
        assert "Regional Sales Report" in doc.content

        # Verify OCR text made it into the vectorized chunks
        assert mixin._vector_store.add_chunks.called
        chunks = mixin._vector_store.add_chunks.call_args[0][0]
        chunk_text = " ".join(c.content for c in chunks)
        assert "EXECUTIVE SUMMARY" in chunk_text
        assert "Regional" in chunk_text

    def test_add_document_from_file_image_no_router(self, tmp_path):
        """Without router, image docs still get basic metadata chunked."""
        from constat.core.config import Config, DocumentConfig
        from constat.discovery.doc_tools._core import _CoreMixin
        from constat.discovery.doc_tools._image import OcrResult

        img_bytes = _make_png_bytes(200, 200, color="red")
        img_path = tmp_path / "photo.png"
        img_path.write_bytes(img_bytes)

        mixin = object.__new__(_CoreMixin)
        mixin._loaded_documents = {}
        mixin._image_labels = []
        mixin._router = None  # no vision
        mixin.config = MagicMock()
        mixin.config.config_dir = str(tmp_path)
        mixin._vector_store = MagicMock()
        mixin._vector_store.add_chunks = MagicMock()
        mock_model = MagicMock()
        import numpy as np
        mock_model.encode.return_value = np.zeros((1, 1024))
        mixin._model = mock_model
        mixin._model_lock = __import__("threading").Lock()
        mixin.CHUNK_SIZE = 1500
        mixin.TABLE_CHUNK_LIMIT = 6000

        with patch("constat.discovery.doc_tools._image._ocr_extract") as mock_ocr:
            mock_ocr.return_value = OcrResult(text="", mean_confidence=0.0, word_count=0)
            success, msg = mixin.add_document_from_file(
                str(img_path),
                name="images:photo",
                description="Test photo",
                domain_id="__base__",
                skip_entity_extraction=True,
            )

        assert success, msg
        doc = mixin._loaded_documents["images:photo"]
        # No description without router
        assert "## Description" not in doc.content
        assert "image-primary" in doc.content

    def test_load_document_no_images_when_disabled(self, tmp_path):
        """extract_images=False should not create child documents."""
        from constat.core.config import DocumentConfig
        from constat.discovery.doc_tools._core import _CoreMixin
        from constat.discovery.doc_tools._transport import FetchResult

        pdf_bytes = _make_pdf_with_images(num_pages=1, images_per_page=1)
        pdf_path = tmp_path / "doc.pdf"
        pdf_path.write_bytes(pdf_bytes)

        mixin = object.__new__(_CoreMixin)
        mixin._loaded_documents = {}
        mixin._image_labels = []
        mixin._router = None
        mixin.config = MagicMock()
        mixin.config.config_dir = str(tmp_path)
        mixin.config.documents = {}

        doc_config = DocumentConfig(
            path=str(pdf_path),
            extract_images=False,
        )

        fetch_result = FetchResult(
            data=pdf_bytes,
            source_path=str(pdf_path),
            detected_mime="application/pdf",
        )

        with (
            patch.object(mixin, "_resolve_doc_config", return_value=doc_config),
            patch(
                "constat.discovery.doc_tools._core.fetch_document",
                return_value=fetch_result,
            ),
            patch(
                "constat.discovery.doc_tools._core.infer_transport",
                return_value="file",
            ),
            patch(
                "constat.discovery.doc_tools._core._infer_structured_schema",
                return_value=None,
            ),
        ):
            mixin._load_document("test_doc")

        child_names = [n for n in mixin._loaded_documents if n.startswith("test_doc:")]
        assert len(child_names) == 0
