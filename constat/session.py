# Copyright (c) 2025 Kenneth Stott
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Session orchestration for multi-step plan execution."""

from __future__ import annotations

import json
import logging
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Optional, Any

logger = logging.getLogger(__name__)

from constat.core.config import Config
from constat.core.models import Plan, PlannerResponse, PostValidation, Step, StepResult, StepStatus, StepType, TaskType, ValidationOnFail
from constat.core.resources import SessionResources
from constat.storage.datastore import DataStore
from constat.storage.history import SessionHistory
from constat.storage.learnings import LearningStore, LearningCategory, LearningSource
from constat.storage.registry import ConstatRegistry
from constat.storage.registry_datastore import RegistryAwareDataStore
from constat.execution.executor import PythonExecutor, format_error_for_retry
from constat.execution.planner import Planner
from constat.execution.scratchpad import Scratchpad
from constat.execution.fact_resolver import (
    FactResolver,
    FactSource,
    format_source_attribution,
)
from constat.execution.mode import (
    Mode,
    PlanApproval,
    PlanApprovalRequest,
    PlanApprovalResponse,
    Phase,
    PrimaryIntent,
    SubIntent,
    TurnIntent,
    ConversationState,
)
from constat.execution.intent_classifier import IntentClassifier
from constat.execution.parallel_scheduler import ExecutionContext
from constat.execution import RETRY_PROMPT_TEMPLATE
from constat.providers import TaskRouter
import constat.llm
from constat.catalog.schema_manager import SchemaManager
from constat.catalog.api_schema_manager import APISchemaManager
from constat.catalog.preload_cache import MetadataPreloadCache
from constat.discovery.doc_tools import DocumentDiscoveryTools
from constat.discovery.concept_detector import ConceptDetector
from constat.email import create_send_email
from constat.embedding_loader import EmbeddingModelLoader
from constat.context import ContextEstimator, ContextCompactor, ContextStats, CompactionResult
from constat.visualization import create_viz_helper
from constat.commands import get_help_markdown


# Meta-question patterns and prompts loaded from external files
from constat.prompts import load_prompt, load_yaml

META_QUESTION_PATTERNS = load_yaml("meta_question_patterns.yaml")["patterns"]


def is_meta_question(query: str) -> bool:
    """Check if query is a meta-question about capabilities."""
    query_lower = query.lower()
    return any(pattern in query_lower for pattern in META_QUESTION_PATTERNS)


# Question classification types
class QuestionType:
    DATA_ANALYSIS = "data_analysis"  # Requires database queries
    GENERAL_KNOWLEDGE = "general_knowledge"  # LLM can answer directly
    META_QUESTION = "meta_question"  # About system capabilities


@dataclass
class DetectedIntent:
    """A single detected intent with confidence."""
    intent: str  # Intent name (REDO, MODIFY_FACT, etc.)
    confidence: float = 0.8
    extracted_value: Optional[str] = None  # Context extracted (e.g., "threshold=$50k")


@dataclass
class QuestionAnalysis:
    """Combined result of question analysis (facts + classification + intent + mode)."""
    question_type: str  # QuestionType value
    extracted_facts: list = field(default_factory=list)  # List of Fact objects
    cached_fact_answer: Optional[str] = None  # Answer from cached facts if applicable
    intents: list = field(default_factory=list)  # List of DetectedIntent objects
    fact_modifications: list = field(default_factory=list)  # [{fact_name, new_value, action}]
    scope_refinements: list = field(default_factory=list)  # ["California", "Q4"]
    wants_brief: bool = False  # User wants brief/concise output (skip insights)
    recommended_mode: Optional[str] = None  # EXPLORATORY or PROOF
    mode_reasoning: Optional[str] = None  # Why this mode was recommended


STEP_SYSTEM_PROMPT = load_prompt("step_system_prompt.md")
STEP_PROMPT_TEMPLATE = load_prompt("step_prompt_template.md")


# Type for approval callback: (request) -> response
ApprovalCallback = Callable[[PlanApprovalRequest], PlanApprovalResponse]


@dataclass
class ClarificationQuestion:
    """A single clarification question with optional suggested answers."""
    text: str  # The question text
    suggestions: list[str] = field(default_factory=list)  # Suggested answers


@dataclass
class ClarificationRequest:
    """Request for clarification before planning."""
    original_question: str
    ambiguity_reason: str  # Why clarification is needed
    questions: list[ClarificationQuestion]  # Questions with suggestions


@dataclass
class ClarificationResponse:
    """User's response to clarification request."""
    answers: dict[str, str]  # question -> answer mapping
    skip: bool = False  # If True, proceed without clarification


# Type for clarification callback: (request) -> response
ClarificationCallback = Callable[[ClarificationRequest], ClarificationResponse]



@dataclass
class SessionConfig:
    """Configuration for a session."""
    max_retries_per_step: int = 10
    verbose: bool = False

    # Plan approval settings
    require_approval: bool = True  # If True, require approval before execution
    max_replan_attempts: int = 3  # Max attempts to replan with user feedback
    auto_approve: bool = False  # If True, auto-approve plans (for testing/scripts)

    # Clarification settings
    ask_clarifications: bool = True  # If True, ask for clarification on ambiguous requests
    skip_clarification: bool = False  # If True, skip clarification (for testing/scripts)

    # Insight/synthesis settings
    enable_insights: bool = True  # If True, synthesize answer and generate suggestions
    show_raw_output: bool = True  # If True, show raw step output before synthesis

    # Default execution mode (overrides LLM mode selection if set)
    default_mode: Optional[Mode] = None  # If set, always use this mode


@dataclass
class StepEvent:
    """Event emitted during step execution."""
    event_type: str  # step_start, generating, executing, step_complete, step_error
    step_number: int
    data: dict = field(default_factory=dict)


class Session:
    """
    Orchestrates multi-step plan execution with step isolation.

    The session:
    1. Takes a problem and generates a plan
    2. Executes each step sequentially in isolation
    3. State is shared ONLY via DuckDB datastore (no in-memory sharing)
    4. Handles errors and retries
    5. Records history for review/resumption

    Step Isolation:
    - Each step runs independently with only the datastore for state sharing
    - Steps use store.save_dataframe() / store.load_dataframe() for DataFrames
    - Steps use store.set_state() / store.get_state() for simple values
    - No in-memory state dict is shared between steps
    """

    def __init__(
        self,
        config: Config,
        session_id: str,
        session_config: Optional[SessionConfig] = None,
        history: Optional[SessionHistory] = None,
        progress_callback: Optional[Callable[[str, int, int, str], None]] = None,
        user_id: Optional[str] = None,
        data_dir: Optional[Path] = None,
    ):
        self.session_id = session_id  # Always provided by client

        self.config = config
        self.session_config = session_config or SessionConfig()
        self.user_id = user_id or "default"
        self.data_dir = data_dir or Path(".constat")

        # Start loading embedding model in background immediately
        # This allows the model to load while other initialization happens
        EmbeddingModelLoader.get_instance().start_loading()

        # Initialize components with timing
        t0 = time.time()
        self.schema_manager = SchemaManager(config)
        self.schema_manager.initialize(progress_callback=progress_callback)
        logger.debug(f"Session init: SchemaManager took {time.time() - t0:.2f}s")

        # API schema manager (GraphQL introspection, REST metadata)
        t0 = time.time()
        self.api_schema_manager = APISchemaManager(config)
        if config.apis:
            self.api_schema_manager.initialize(progress_callback=progress_callback)
        logger.debug(f"Session init: APISchemaManager took {time.time() - t0:.2f}s")

        # Metadata preload cache for faster context loading
        t0 = time.time()
        self.preload_cache = MetadataPreloadCache(config)
        self._preloaded_context: Optional[str] = None
        self._load_preloaded_context()
        logger.debug(f"Session init: MetadataPreloadCache took {time.time() - t0:.2f}s")

        # Document discovery tools (for reference documents)
        t0 = time.time()
        self.doc_tools = DocumentDiscoveryTools(config)
        logger.debug(f"Session init: DocumentDiscoveryTools took {time.time() - t0:.2f}s")

        # Entity extraction is handled by session_manager.refresh_entities_async()
        # after session creation — not during __init__ to avoid dual extraction race
        self._entities_extracted = False

        # Task router for model routing with escalation
        self.router = TaskRouter(config.llm)

        constat.llm.set_backend(self.router)
        constat.llm.on_call(self._handle_llm_call_event)

        # Default provider (for backward compatibility - e.g., fact resolver)
        self.llm = self.router._get_provider(
            self.router.routing_config.get_models_for_task("general")[0]
        )

        self.planner = Planner(
            config, self.schema_manager, self.router,
            doc_tools=self.doc_tools,
            api_schema_manager=self.api_schema_manager,
        )


        self.executor = PythonExecutor(
            timeout_seconds=config.execution.timeout_seconds,
            allowed_imports=config.execution.allowed_imports or None,
        )

        self.history = history or SessionHistory(user_id=self.user_id)

        # Session state
        # Note: self.session_id is set at the top of __init__ (required parameter)
        # Store server_session_id for history mapping (allows find_session_by_server_id)
        self.server_session_id: Optional[str] = session_id  # Server UUID for reverse lookup
        self.plan: Optional[Plan] = None
        self.scratchpad = Scratchpad()
        self.datastore: Optional[RegistryAwareDataStore] = None  # Persistent storage (only shared state between steps)

        # Central registry for tables and artifacts (shared across sessions)
        self.registry = ConstatRegistry(base_dir=Path(".constat"))

        # Session-scoped data sources (added via /database and /file commands)
        self.session_databases: dict[str, dict] = {}  # name -> {type, uri, description}
        self.session_files: dict[str, dict] = {}  # name -> {uri, auth, description}

        # Project APIs (added when projects are activated)
        self._project_apis: dict[str, Any] = {}  # name -> ApiConfig

        # Consolidated view of all available resources (single source of truth)
        self.resources = SessionResources()
        self._init_resources_from_config()

        # Pass resources to planner (after resources are initialized)
        self.planner.resources = self.resources

        # Fact resolver for auditable mode
        self.fact_resolver = FactResolver(
            llm=self.llm,
            schema_manager=self.schema_manager,
            config=self.config,
            event_callback=self._handle_fact_resolver_event,
            doc_tools=self.doc_tools,  # Enable document-based fact resolution
        )

        # Learning store for corrections and patterns
        self.learning_store = LearningStore(user_id=self.user_id)

        # Pass learning store to planner for injecting learned rules
        self.planner.set_learning_store(self.learning_store)

        # Role manager for user-defined roles ({data_dir}/{user_id}/roles.yaml)
        from constat.core.roles import RoleManager
        self.role_manager = RoleManager(user_id=self.user_id, base_dir=self.data_dir)

        # Role matcher for dynamic role selection based on query
        from constat.core.role_matcher import RoleMatcher
        self.role_matcher = RoleMatcher(self.role_manager)
        # Initialize lazily on first use to avoid blocking startup

        # Skill manager: loads system, project, and user skills in precedence order
        from constat.core.skills import SkillManager
        system_skills_dir = Path(config.config_dir) / "skills" if config.config_dir else None
        self.skill_manager = SkillManager(
            user_id=self.user_id, base_dir=self.data_dir,
            system_skills_dir=system_skills_dir,
        )

        # Skill matcher for dynamic skill selection based on query
        from constat.core.skill_matcher import SkillMatcher
        self.skill_matcher = SkillMatcher(self.skill_manager)
        # Initialize lazily on first use to avoid blocking startup

        # Wire skill manager into planner so active skills appear in planning prompts
        self.planner.set_skill_manager(self.skill_manager)

        # Track current role for this query (None = shared context)
        self._current_role_id: Optional[str] = None

        # Event callbacks for monitoring
        self._event_handlers: list[Callable[[StepEvent], None]] = []

        # Approval callback (set via set_approval_callback)
        self._approval_callback: Optional[ApprovalCallback] = None

        # Clarification callback (set via set_clarification_callback)
        self._clarification_callback: Optional[ClarificationCallback] = None

        # Tool response cache for schema tools (cleared on refresh)
        self._tool_cache: dict[str, any] = {}

        # Cached proof result (set after prove_conversation completes)
        self.last_proof_result: Optional[dict] = None

        # Concept detector for conditional prompt injection
        t0 = time.time()
        self._concept_detector = ConceptDetector()
        self._concept_detector.initialize()
        logger.debug(f"Session init: ConceptDetector took {time.time() - t0:.2f}s")

        # Phase 3: Conversation state and intent classifier
        # Initialize conversation state with idle phase
        self._conversation_state: ConversationState = ConversationState(
            phase=Phase.IDLE,
        )

        # Intent classifier for turn-level intent detection
        # Pass router as LLM provider for fallback classification
        self._intent_classifier: IntentClassifier = IntentClassifier(
            llm_provider=self.router,
        )

        # Phase 4: Execution Control
        # Cancellation flag for stopping execution mid-flight
        self._cancelled: bool = False

        # Intent queue for messages received during execution
        # Queue behavior per intent type:
        # - plan_new: Queue 1, latest wins (new request replaces queued)
        # - control: Queue in order, process after execution
        # - query: No queue, answered in parallel immediately
        # Each entry is a tuple of (TurnIntent, user_input_string)
        self._intent_queue: list[tuple[TurnIntent, str]] = []

        # Execution context for cancellation signaling to scheduler
        self._execution_context: ExecutionContext = ExecutionContext()

    def _init_resources_from_config(self) -> None:
        """Initialize resources from base config."""
        # Add databases from config
        for name, db_config in self.config.databases.items():
            self.resources.add_database(
                name=name,
                description=db_config.description or "",
                db_type=db_config.type or "sql",
                source="config",
            )

        # Add APIs from config
        if self.config.apis:
            for name, api_config in self.config.apis.items():
                self.resources.add_api(
                    name=name,
                    description=api_config.description or "",
                    api_type=api_config.type or "graphql",
                    source="config",
                )

        # Add documents from config
        if self.config.documents:
            for name, doc_config in self.config.documents.items():
                self.resources.add_document(
                    name=name,
                    description=doc_config.description or "",
                    doc_type=doc_config.type or "file",
                    source="config",
                )

    def add_project_resources(
        self,
        project_filename: str,
        databases: dict = None,
        apis: dict = None,
        documents: dict = None,
    ) -> None:
        """Add resources from a project.

        Args:
            project_filename: Project filename for source tracking
            databases: Dict of database configs
            apis: Dict of API configs
            documents: Dict of document configs
        """
        source = f"project:{project_filename}"

        if databases:
            for name, db_config in databases.items():
                self.resources.add_database(
                    name=name,
                    description=getattr(db_config, 'description', '') or "",
                    db_type=getattr(db_config, 'type', 'sql') or "sql",
                    source=source,
                )

        if apis:
            for name, api_config in apis.items():
                self.resources.add_api(
                    name=name,
                    description=getattr(api_config, 'description', '') or "",
                    api_type=getattr(api_config, 'type', 'graphql') or "graphql",
                    source=source,
                )

        if documents:
            for name, doc_config in documents.items():
                self.resources.add_document(
                    name=name,
                    description=getattr(doc_config, 'description', '') or "",
                    doc_type=getattr(doc_config, 'type', 'file') or "file",
                    source=source,
                )

    def remove_project_resources(self, project_filename: str) -> None:
        """Remove all resources from a project.

        Args:
            project_filename: Project filename
        """
        source = f"project:{project_filename}"
        self.resources.remove_by_source(source)

    def sync_resources_to_history(self) -> None:
        """Sync current resources to session history (session.json).

        Call this after loading/unloading projects to keep history in sync.
        """
        if self.session_id and self.history:
            self.history.update_resources(
                session_id=self.session_id,
                databases=self.resources.database_names,
                apis=self.resources.api_names,
                documents=self.resources.document_names,
            )

    def set_approval_callback(self, callback: ApprovalCallback) -> None:
        """
        Set the callback for plan approval.

        The callback receives a PlanApprovalRequest and must return a PlanApprovalResponse.

        Args:
            callback: Function that handles approval requests
        """
        self._approval_callback = callback

    def set_clarification_callback(self, callback: ClarificationCallback) -> None:
        """
        Set the callback for requesting clarification on ambiguous questions.

        The callback receives a ClarificationRequest and must return a ClarificationResponse.

        Args:
            callback: Function that handles clarification requests
        """
        self._clarification_callback = callback

    def on_event(self, handler: Callable[[StepEvent], None]) -> None:
        """Register an event handler for step events."""
        self._event_handlers.append(handler)

    def _emit_event(self, event: StepEvent) -> None:
        """Emit an event to all handlers."""
        for handler in self._event_handlers:
            handler(event)

    def _handle_fact_resolver_event(self, event_type: str, data: dict) -> None:
        """Convert fact resolver events to StepEvents and emit them."""
        self._emit_event(StepEvent(
            event_type=event_type,
            step_number=data.get("step", 0),
            data=data,
        ))

    def _sync_user_facts_to_planner(self) -> None:
        """Sync current user facts to the planner for use in planning prompts."""
        try:
            all_facts = self.fact_resolver.get_all_facts()
            # Convert Fact objects to simple name -> value dict
            facts_dict = {name: fact.value for name, fact in all_facts.items()}
            self.planner.set_user_facts(facts_dict)
        except Exception as e:
            logger.debug(f"Failed to sync user facts to planner: {e}")

    def _sync_available_roles_to_planner(self) -> None:
        """Sync available roles to the planner for role-based step assignment."""
        try:
            role_names = self.role_manager.list_roles()  # Returns list[str]
            # Convert to list of dicts with name and description
            roles_list = []
            for name in role_names:
                role = self.role_manager.get_role(name)
                if role:
                    roles_list.append({"name": name, "description": role.description or ""})
            logger.info(f"[ROLES] Syncing {len(roles_list)} roles to planner: {[r['name'] for r in roles_list]}")
            self.planner.set_available_roles(roles_list)
        except Exception as e:
            logger.warning(f"Failed to sync roles to planner: {e}")

    def _is_unclear_input(self, text: str) -> bool:
        """Check if input appears to be unclear, garbage, or a copy-paste error.

        Detects:
        - Terminal prompts (e.g., "(.venv) user@host % command")
        - File paths without context
        - Very short input with no meaningful words
        - Copy-paste errors with shell syntax
        """
        import re

        text = text.strip()

        # Empty or very short without meaningful content
        if len(text) < 3:
            return True

        # Terminal prompt patterns
        terminal_patterns = [
            r'^\(.+\)\s*\(.+\)\s*\w+@\w+',  # (.venv) (base) user@host
            r'^\w+@[\w\-]+\s*[%$#>]',  # user@hostname %
            r'^[%$#>]\s*\w+',  # % command or $ command
            r'^\(.+\)\s*%',  # (.venv) %
            r'constat\s+repl',  # constat repl command
            r'^pip\s+install',  # pip install
            r'^python\s+-',  # python -m
            r'^cd\s+/',  # cd /path
            r'^\s*\$\s*\w+',  # $ command
        ]

        for pattern in terminal_patterns:
            if re.search(pattern, text, re.IGNORECASE):
                return True

        # Looks like a file path without a question
        if re.match(r'^[/~][\w/\-\.]+$', text) or re.match(r'^[\w]:\\', text):
            return True

        # Contains mostly special characters
        alpha_ratio = sum(1 for c in text if c.isalpha()) / max(len(text), 1)
        if alpha_ratio < 0.3 and len(text) > 10:
            return True

        return False

    def _build_step_prompt(self, step: Step) -> str:
        """Build the prompt for generating step code."""
        # Format datastore tables info with column metadata
        if self.datastore:
            tables = self.datastore.list_tables()
            if tables:
                table_lines = ["Available in `store` (load with `store.load_dataframe('name')` or query with SQL):"]
                for t in tables:
                    schema = self.datastore.get_table_schema(t['name'])
                    if schema:
                        col_names = [c['name'] for c in schema]
                        table_lines.append(f"  - {t['name']}: {t['row_count']} rows, columns: {col_names}")
                    else:
                        table_lines.append(f"  - {t['name']}: {t['row_count']} rows")
                datastore_info = "\n".join(table_lines)
            else:
                datastore_info = "(no tables saved yet)"
        else:
            datastore_info = "(no datastore)"

        # Get scratchpad from datastore (persistent) - source of truth for isolation
        if self.datastore:
            scratchpad_context = self.datastore.get_scratchpad_as_markdown()
        else:
            scratchpad_context = self.scratchpad.get_recent_context(max_steps=5)

        # Build source context with semantic search for step-relevant tables
        ctx = self._build_source_context(query=step.goal)

        # Build codegen learnings section - only for code generation steps
        # Skip for summarization, planning, intent classification etc.
        learnings_text = ""
        code_gen_types = {TaskType.PYTHON_ANALYSIS, TaskType.SQL_GENERATION}
        if step.task_type in code_gen_types:
            try:
                learnings_text = self._get_codegen_learnings(step.goal, step.task_type)
            except Exception as e:
                logger.debug(f"Failed to get codegen learnings: {e}")

        # Detect relevant concepts and inject specialized sections
        injected_sections = self._concept_detector.get_sections_for_prompt(
            query=step.goal,
            target="step",
        )

        # Build published artifacts context for name reuse
        published_artifacts_text = ""
        if self.datastore:
            existing_artifacts = self.datastore.list_artifacts()
            named = [a for a in existing_artifacts if a.get("type") not in ("code", "output", "error")]
            if named:
                artifact_lines = ["Published artifacts (reuse these names with viz.save_* to update):"]
                for a in named:
                    artifact_lines.append(f"  - {a['name']} ({a.get('type', 'unknown')})")
                published_artifacts_text = "\n".join(artifact_lines)

        return STEP_PROMPT_TEMPLATE.format(
            system_prompt=STEP_SYSTEM_PROMPT,
            injected_sections=injected_sections,
            schema_overview=ctx["schema_overview"],
            api_overview=ctx["api_overview"],
            domain_context=self._get_system_prompt() or "No additional context.",
            user_facts=ctx["user_facts"],
            learnings=learnings_text,
            datastore_tables=datastore_info,
            published_artifacts=published_artifacts_text,
            scratchpad=scratchpad_context,
            step_number=step.number,
            total_steps=len(self.plan.steps) if self.plan else 1,
            goal=step.goal,
            inputs=", ".join(step.expected_inputs) if step.expected_inputs else "(none)",
            outputs=", ".join(step.expected_outputs) if step.expected_outputs else "(none)",
        )

    def _get_codegen_learnings(self, step_goal: str, task_type: TaskType = None) -> str:
        """Get relevant codegen learnings showing what didn't work vs what did work.

        Args:
            step_goal: The goal of the current step for context matching
            task_type: The task type to filter learnings (SQL vs Python)

        Returns:
            Formatted learnings text for prompt injection
        """
        if not self.learning_store:
            return ""

        # Determine if this is SQL or Python based on task type
        is_sql = task_type == TaskType.SQL_GENERATION

        lines = []

        # Get rules (compacted learnings) for codegen errors
        rules = self.learning_store.list_rules(
            category=LearningCategory.CODEGEN_ERROR,
            min_confidence=0.6,
        )
        if rules:
            # Filter rules by type - SQL rules mention SQL/query patterns
            sql_keywords = {'sql', 'query', 'select', 'join', 'table', 'column', 'duckdb'}
            filtered_rules = []
            for rule in rules:
                summary_lower = rule.get('summary', '').lower()
                rule_is_sql = any(kw in summary_lower for kw in sql_keywords)
                if is_sql == rule_is_sql:
                    filtered_rules.append(rule)

            if filtered_rules:
                label = "SQL" if is_sql else "Code"
                lines.append(f"\n## {label} Generation Rules (apply these)")
                for rule in filtered_rules[:5]:
                    lines.append(f"- {rule['summary']}")

        # Get recent raw learnings with full context (error vs fix)
        raw_learnings = self.learning_store.list_raw_learnings(
            category=LearningCategory.CODEGEN_ERROR,
            limit=10,
            include_promoted=False,
        )
        if raw_learnings:
            # Filter by type and relevance
            sql_code_patterns = {'SELECT', 'FROM', 'WHERE', 'JOIN', 'INSERT', 'UPDATE', 'store.query'}
            relevant = []
            for l in raw_learnings:
                ctx = l.get("context", {})
                code = ctx.get("original_code", "") + ctx.get("fixed_code", "")
                learning_is_sql = any(pat in code for pat in sql_code_patterns)
                if is_sql == learning_is_sql and self._is_learning_relevant(l, step_goal):
                    relevant.append(l)
            relevant = relevant[:3]  # Limit to 3 detailed examples

            if relevant:
                label = "SQL" if is_sql else "Codegen"
                lines.append(f"\n## Recent {label} Fixes (learn from these)")
                for learning in relevant:
                    ctx = learning.get("context", {})
                    original = ctx.get("original_code", "")
                    fixed = ctx.get("fixed_code", "")
                    error_msg = ctx.get("error_message", "")

                    # Show the contrast
                    lines.append(f"\n### {learning['correction'][:80]}")
                    if error_msg:
                        lines.append(f"**Error:** {error_msg[:100]}")
                    lang = "sql" if is_sql else "python"
                    if original:
                        lines.append(f"**Broken code:**\n```{lang}\n{original[:300]}\n```")
                    if fixed:
                        lines.append(f"**Fixed code:**\n```{lang}\n{fixed[:300]}\n```")

        return "\n".join(lines) if lines else ""

    def _is_learning_relevant(self, learning: dict, step_goal: str) -> bool:
        """Check if a learning is relevant to the current step goal."""
        # Simple keyword overlap check
        goal_words = set(step_goal.lower().split())
        learning_goal = learning.get("context", {}).get("step_goal", "")
        learning_words = set(learning_goal.lower().split())
        correction_words = set(learning.get("correction", "").lower().split())

        # Check for meaningful keyword overlap
        common_words = {"the", "a", "an", "to", "from", "for", "with", "in", "on", "of", "and", "or"}
        goal_keywords = goal_words - common_words
        learning_keywords = (learning_words | correction_words) - common_words

        overlap = goal_keywords & learning_keywords
        return len(overlap) >= 1  # At least one meaningful keyword match

    def _cached_get_table_schema(self, table: str) -> dict:
        """Get table schema with caching."""
        cache_key = f"schema:{table}"
        if cache_key not in self._tool_cache:
            self._tool_cache[cache_key] = self.schema_manager.get_table_schema(table)
        return self._tool_cache[cache_key]

    def _cached_find_relevant_tables(self, query: str, top_k: int = 5) -> list[dict]:
        """Find relevant tables with caching and document enrichment."""
        cache_key = f"relevant:{query}:{top_k}"
        if cache_key not in self._tool_cache:
            self._tool_cache[cache_key] = self.schema_manager.find_relevant_tables(
                query, top_k, doc_tools=self.doc_tools
            )
        return self._tool_cache[cache_key]

    def _cached_find_relevant_apis(self, query: str, limit: int = 5) -> list[dict]:
        """Find relevant APIs with caching."""
        cache_key = f"apis:{query}:{limit}"
        if cache_key not in self._tool_cache:
            self._tool_cache[cache_key] = self.api_schema_manager.find_relevant_apis(
                query, limit=limit
            )
        return self._tool_cache[cache_key]

    def find_relevant_sources(
        self,
        query: str,
        table_limit: int = 5,
        doc_limit: int = 3,
        api_limit: int = 3,
        min_similarity: float = 0.3,
    ) -> dict:
        """Find all relevant sources (tables, documents, APIs) for a query.

        Performs unified semantic search across all configured data sources
        and returns ranked results from each category.

        Args:
            query: Natural language query
            table_limit: Max tables to return
            doc_limit: Max documents to return
            api_limit: Max API endpoints to return
            min_similarity: Minimum similarity threshold

        Returns:
            Dict with 'tables', 'documents', 'apis' keys, each containing
            a list of relevant sources with similarity scores.
        """
        results = {
            "tables": [],
            "documents": [],
            "apis": [],
        }

        # Find relevant tables (with doc enrichment)
        tables = self._cached_find_relevant_tables(query, top_k=table_limit)
        results["tables"] = [
            {
                "source_type": "table",
                "name": t["full_name"],
                "database": t["database"],
                "summary": t["summary"],
                "relevance": t["relevance"],
                "documentation": t.get("documentation", []),
            }
            for t in tables
            if t.get("relevance", 0) >= min_similarity
        ]

        # Find relevant documents
        if self.doc_tools:
            docs = self.doc_tools.search_documents(query, limit=doc_limit)
            results["documents"] = [
                {
                    "source_type": "document",
                    "name": d["document"],
                    "section": d.get("section"),
                    "excerpt": d.get("excerpt", ""),
                    "relevance": d.get("relevance", 0),
                }
                for d in docs
                if d.get("relevance", 0) >= min_similarity
            ]

        # Find relevant APIs
        if not self.config.apis:
            logger.debug(f"[find_relevant_sources] No APIs configured")
        if self.config.apis:
            apis = self._cached_find_relevant_apis(query, limit=api_limit)
            logger.debug(f"[find_relevant_sources] API search for '{query[:50]}': {len(apis)} results, min_sim={min_similarity}")
            for a in apis[:3]:
                logger.debug(f"[find_relevant_sources]   - {a.get('api_name')}.{a.get('endpoint')}: sim={a.get('similarity', 0):.3f}")
            results["apis"] = [
                {
                    "source_type": "api",
                    "name": f"{a['api_name']}.{a['endpoint']}",
                    "api_name": a["api_name"],
                    "endpoint": a["endpoint"],
                    "type": a["type"],
                    "description": a.get("description"),
                    "fields": a.get("fields", []),
                    "relevance": a.get("similarity", 0),
                }
                for a in apis
                if a.get("similarity", 0) >= min_similarity
            ]

        return results

    def _find_entity(self, name: str, limit: int = 3) -> dict:
        """Find all occurrences of an entity across schema and documents."""
        from constat.discovery.schema_tools import SchemaDiscoveryTools
        tools = SchemaDiscoveryTools(self.schema_manager, self.doc_tools)
        return tools.find_entity(name, limit)

    def _get_tool_handlers(self) -> dict:
        """Get schema tool handlers with caching."""
        handlers = {
            "get_table_schema": self._cached_get_table_schema,
            "find_relevant_tables": self._cached_find_relevant_tables,
            "find_entity": self._find_entity,
        }

        # Add API schema tools if APIs are configured
        if self.config.apis:
            handlers["get_api_schema_overview"] = self._get_api_schema_overview
            handlers["get_api_query_schema"] = self._get_api_query_schema
            handlers["find_relevant_apis"] = self._cached_find_relevant_apis

        return handlers

    def _get_api_schema_overview(self, api_name: str) -> dict:
        """Get overview of an API's schema (queries/endpoints)."""
        cache_key = f"api_overview:{api_name}"
        if cache_key not in self._tool_cache:
            from constat.catalog.api_executor import APIExecutor
            executor = APIExecutor(self.config)
            self._tool_cache[cache_key] = executor.get_schema_overview(api_name)
        return self._tool_cache[cache_key]

    def _get_api_query_schema(self, api_name: str, query_name: str) -> dict:
        """Get detailed schema for a specific API query/endpoint."""
        cache_key = f"api_query:{api_name}:{query_name}"
        if cache_key not in self._tool_cache:
            from constat.catalog.api_executor import APIExecutor
            executor = APIExecutor(self.config)
            self._tool_cache[cache_key] = executor.get_query_schema(api_name, query_name)
        return self._tool_cache[cache_key]

    def refresh_metadata(self, force_full: bool = False) -> dict:
        """Refresh all metadata: schema, documents, APIs, and preload cache.

        Args:
            force_full: If True, force full rebuild of all caches

        Returns:
            Dict with refresh statistics
        """
        self._tool_cache.clear()
        self.schema_manager.refresh()

        # Refresh API schema manager
        api_count = 0
        if self.config.apis:
            self.api_schema_manager.initialize()  # Re-introspect APIs
            api_count = len(self.api_schema_manager.metadata_cache)

        # Pass schema entities to doc_tools for entity extraction
        # IMPORTANT: Keep schema and API entities separate to avoid type confusion
        if self.doc_tools:
            schema_entities = set(self.schema_manager.get_entity_names())
            api_entities = list(self._get_api_entity_names())
            self.doc_tools.set_schema_entities(schema_entities)
            # Set API entities separately for proper type assignment
            if api_entities:
                self.doc_tools.set_openapi_entities(api_entities, api_entities)

            # Process schema metadata (names + descriptions) through NER
            schema_metadata = self.schema_manager.get_description_text()
            if schema_metadata:
                self.doc_tools.process_metadata_through_ner(schema_metadata, source_type="schema")

            # Process API metadata through NER
            api_metadata = self.api_schema_manager.get_description_text()
            if api_metadata:
                self.doc_tools.process_metadata_through_ner(api_metadata, source_type="api")

        # Refresh document vector index (incremental by default)
        doc_stats = {}
        if self.doc_tools:
            doc_stats = self.doc_tools.refresh(force_full=force_full)

        # Rebuild preload cache with fresh metadata
        self._rebuild_preload_cache()

        return {
            "preloaded_tables": self.get_preloaded_tables_count(),
            "documents": doc_stats,
            "api_endpoints": api_count,
        }

    def _get_api_entity_names(self) -> set[str]:
        """Get API endpoint names for entity extraction."""
        if not self.config.apis:
            return set()

        entities = set()
        for meta in self.api_schema_manager.metadata_cache.values():
            # Add endpoint name
            entities.add(meta.endpoint_name)
            # Add API.endpoint as full name
            entities.add(meta.full_name)
            # Add field names
            for field in meta.fields:
                entities.add(field.name)
        return entities

    def extract_entities(self, project_ids: list[str] | None = None) -> int:
        """Extract NER entities from all chunks for this session.

        Should be called after session_id is set. Extracts entities from
        base + specified project chunks.

        Args:
            project_ids: Project IDs to include (in addition to base)

        Returns:
            Number of entities extracted
        """
        if not self.session_id:
            logger.warning("extract_entities called without session_id")
            return 0

        if self._entities_extracted:
            logger.debug("Entities already extracted for this session")
            return 0

        t0 = time.time()
        schema_terms = list(self.schema_manager.get_entity_names())
        api_terms = list(self._get_api_entity_names())

        entity_count = self.doc_tools._vector_store.extract_entities_for_session(
            session_id=self.session_id,
            project_ids=project_ids,
            schema_terms=schema_terms,
            api_terms=api_terms,
        )

        self._entities_extracted = True
        logger.debug(f"Entity extraction took {time.time() - t0:.2f}s ({entity_count} entities)")
        return entity_count

    def rebuild_entities(self, project_ids: list[str] | None = None) -> int:
        """Rebuild entity catalog (e.g., when projects change).

        Clears existing entities for this session and re-extracts.

        Args:
            project_ids: Project IDs to include (in addition to base)

        Returns:
            Number of entities extracted
        """
        if not self.session_id:
            logger.warning("rebuild_entities called without session_id")
            return 0

        self._entities_extracted = False
        return self.extract_entities(project_ids)

    def add_project_api(self, name: str, api_config: Any) -> None:
        """Register an API from an active project.

        Args:
            name: API name
            api_config: API configuration (ApiConfig object)
        """
        self._project_apis[name] = api_config
        logger.info(f"Registered project API: {name}")

    def clear_project_apis(self) -> None:
        """Clear all registered project APIs."""
        self._project_apis.clear()

    def get_all_apis(self) -> dict[str, Any]:
        """Get all APIs (config + project).

        Returns:
            Dict of API name to config, combining config.apis and project APIs
        """
        all_apis = dict(self.config.apis) if self.config.apis else {}
        all_apis.update(self._project_apis)
        return all_apis

    def _load_preloaded_context(self) -> None:
        """Load preloaded metadata context from cache if available."""
        if self.config.context_preload.seed_patterns:
            self._preloaded_context = self.preload_cache.get_context_string()

    def _rebuild_preload_cache(self) -> None:
        """Rebuild the preload cache with current metadata."""
        if self.config.context_preload.seed_patterns:
            self.preload_cache.build(self.schema_manager)
            self._preloaded_context = self.preload_cache.get_context_string()

    def get_preloaded_tables_count(self) -> int:
        """Get the number of tables in the preload cache."""
        return len(self.preload_cache.get_cached_tables())

    def _get_brief_schema_summary(self) -> str:
        """Get a brief summary of available databases without listing all tables.

        This is used when no preload config is set. It provides enough context
        for the LLM to know what databases exist without bloating the prompt
        with potentially hundreds of tables.
        """
        lines = ["Available databases:"]

        # Group tables by database to get counts
        by_db: dict[str, int] = {}
        total_rows_by_db: dict[str, int] = {}
        for table_meta in self.schema_manager.metadata_cache.values():
            by_db[table_meta.database] = by_db.get(table_meta.database, 0) + 1
            total_rows_by_db[table_meta.database] = (
                total_rows_by_db.get(table_meta.database, 0) + table_meta.row_count
            )

        # Build database descriptions with descriptions if available
        db_descriptions = {
            db_name: db_config.description
            for db_name, db_config in self.config.databases.items()
            if db_config.description
        }

        for db_name in sorted(by_db.keys()):
            table_count = by_db[db_name]
            row_count = total_rows_by_db[db_name]
            if db_name in db_descriptions:
                lines.append(f"  {db_name} (connection: db_{db_name}): {db_descriptions[db_name]} ({table_count} tables, ~{row_count:,} rows)")
            else:
                lines.append(f"  {db_name} (connection: db_{db_name}): {table_count} tables, ~{row_count:,} rows")

        return "\n".join(lines)

    def _format_relevant_tables(self, tables: list[dict]) -> str:
        """Format semantically matched tables with database connection names."""
        if not tables:
            return self._get_brief_schema_summary()

        lines = ["Relevant tables (use pd.read_sql(query, db_<name>) with the correct connection):"]

        # Group by database for clarity
        by_db: dict[str, list] = {}
        for t in tables:
            by_db.setdefault(t["database"], []).append(t)

        for db_name, db_tables in sorted(by_db.items()):
            lines.append(f"\n  Database '{db_name}' — connection: db_{db_name}")
            for t in db_tables:
                full_name = t["name"]
                table_meta = self.schema_manager.metadata_cache.get(full_name) if self.schema_manager else None
                if table_meta:
                    col_names = ", ".join(c.name for c in table_meta.columns[:15])
                    if len(table_meta.columns) > 15:
                        col_names += f", ... (+{len(table_meta.columns) - 15} more)"
                    lines.append(f"    {table_meta.name}({col_names}) ~{table_meta.row_count} rows")
                else:
                    lines.append(f"    {t.get('name', '')}: {t.get('summary', '')}")

        lines.append("\nUse `find_relevant_tables(query)` or `get_table_schema(table)` for other tables.")
        return "\n".join(lines)

    def _format_relevant_apis(self, apis: list[dict]) -> str:
        """Format semantically matched APIs."""
        if not apis:
            # Fall back to listing all configured APIs
            if self.resources.has_apis():
                api_lines = ["\n## Available APIs"]
                for name, api_info in self.resources.apis.items():
                    api_type = api_info.api_type.upper()
                    desc = api_info.description or f"{api_type} endpoint"
                    api_lines.append(f"- **{name}** ({api_type}): {desc}")
                return "\n".join(api_lines)
            return ""

        lines = ["\n## Relevant APIs"]
        for a in apis:
            name = a.get("name", "")
            desc = a.get("description", "")
            api_type = a.get("type", "").upper()
            lines.append(f"- **{name}** ({api_type}): {desc}")
            if a.get("fields"):
                lines.append(f"  Fields: {', '.join(a['fields'][:10])}")
        return "\n".join(lines)

    def _format_relevant_docs(self, docs: list[dict]) -> str:
        """Format semantically matched documents."""
        if not docs:
            # Fall back to listing all configured documents
            if self.resources.has_documents():
                doc_lines = ["\n## Reference Documents"]
                for name, doc_info in self.resources.documents.items():
                    desc = doc_info.description or doc_info.doc_type
                    doc_lines.append(f"- **{name}**: {desc}")
                return "\n".join(doc_lines)
            return ""

        lines = ["\n## Relevant Documents"]
        for d in docs:
            name = d.get("name", "")
            section = d.get("section", "")
            excerpt = d.get("excerpt", "")
            section_info = f" (section: {section})" if section else ""
            lines.append(f"- **{name}**{section_info}")
            if excerpt:
                lines.append(f"  > {excerpt[:200]}")
        return "\n".join(lines)

    def _get_schema_tools(self) -> list[dict]:
        """Get schema tool definitions loaded from schema_tools.yaml."""
        import copy

        data = load_yaml("schema_tools.yaml")
        tools = copy.deepcopy(data["base_tools"])

        if self.config.apis:
            api_names = list(self.config.apis.keys())
            api_tools = copy.deepcopy(data["api_tools"])
            for tool in api_tools:
                tool["description"] = tool["description"].format(api_names=api_names)
                for prop in tool["input_schema"]["properties"].values():
                    if prop.get("type") == "string" and "api_name" in tool["input_schema"].get("required", []):
                        if "enum" not in prop and prop.get("description", "").startswith("Name of the API"):
                            prop["enum"] = api_names
            tools.extend(api_tools)

        return tools

    def _resolve_llm_knowledge(self, question: str) -> int | float | str:
        """Resolve a fact from LLM general knowledge.

        Args:
            question: The question or fact description to resolve

        Returns:
            The value (number, string, or ISO date string)

        Raises:
            Exception: If the LLM response cannot be parsed
        """
        import json
        from datetime import datetime

        json_key = question.replace(" ", "_").lower()[:30]

        knowledge_prompt = load_prompt("knowledge_prompt.md").format(
            question=question, json_key=json_key,
        )

        result = self.router.execute(
            task_type=TaskType.SYNTHESIS,
            system="You output ONLY valid JSON with a single value (number, string, or ISO date). No explanations.",
            user_message=knowledge_prompt,
            max_tokens=self.router.max_output_tokens,
        )

        response = result.content.strip()
        logger.debug(f"[LLM_KNOWLEDGE] Response for {question}: {response[:200]}")

        # Parse JSON response
        json_str = response
        if "```" in json_str:
            json_str = re.sub(r'```\w*\n?', '', json_str).strip()

        if not json_str.startswith("{"):
            raise Exception(f"Could not parse LLM response: {response[:100]}")

        data = json.loads(json_str)
        if not data:
            raise Exception(f"Empty JSON response: {response[:100]}")

        raw_value = list(data.values())[0]

        # Return typed value
        if isinstance(raw_value, (int, float)):
            return raw_value
        elif isinstance(raw_value, str):
            # Check if it's an ISO date (validate but return as string)
            if re.match(r'^\d{4}-\d{2}-\d{2}(T\d{2}:\d{2}:\d{2})?', raw_value):
                try:
                    datetime.fromisoformat(raw_value.replace('Z', '+00:00'))
                except ValueError:
                    pass  # Not a valid date, return as regular string
            return raw_value
        else:
            return raw_value  # Other types (bool, etc.)

    def _execute_dag_node(
        self,
        node: "FactNode",
        dag: "ExecutionDAG",
        problem: str,
        detailed_schema: str,
        premises: list[dict],
        inferences: list[dict],
        resolved_premises: dict,
        resolved_inferences: dict,
        inference_names: dict,
    ) -> tuple[any, float, str]:
        """Execute a single DAG node (premise or inference).

        Called by DAGExecutor for each node. Handles both:
        - Leaf nodes (premises): resolve from database/document/cache
        - Internal nodes (inferences): execute Python code using dependency values

        Args:
            node: The FactNode to execute
            dag: The full DAG (for accessing dependency values)
            problem: The original problem being solved
            detailed_schema: Schema info for SQL generation
            premises: Original premise list for reference
            inferences: Original inference list for reference
            resolved_premises: Dict to store resolved premise facts
            resolved_inferences: Dict to store resolved inference results
            inference_names: Dict mapping inference ID to table name

        Returns:
            Tuple of (value, confidence, source)
        """
        import re
        import pandas as pd
        from constat.execution.fact_resolver import Fact, FactSource

        if node.is_leaf:
            # === PREMISE RESOLUTION ===
            fact_id = node.fact_id
            fact_name = node.name
            fact_desc = node.description
            source = f"{node.source}:{node.source_db}" if node.source_db else node.source or "database"

            # Check for pre-resolved node (embedded values handled by DAG parser)
            if node.value is not None:
                fact = Fact(
                    name=fact_name,
                    value=node.value,
                    confidence=node.confidence,
                    source=FactSource.LLM_KNOWLEDGE,
                    reasoning="Embedded value from plan",
                )
                resolved_premises[fact_id] = fact
                self.fact_resolver.add_user_fact(
                    fact_name=fact_name,
                    value=node.value,
                    reasoning="Embedded value",
                    source=FactSource.LLM_KNOWLEDGE,
                )
                if self.history:
                    self.history.save_inference_premise(
                        self.session_id, fact_id, fact_name,
                        node.value, "embedded", fact_desc or ""
                    )
                return node.value, node.confidence, "user"

            # Check cache
            cached_fact = self.fact_resolver.get_fact(fact_name)
            if cached_fact and cached_fact.value is not None:
                resolved_premises[fact_id] = cached_fact
                source_str = format_source_attribution(
                    cached_fact.source, cached_fact.source_name, cached_fact.api_endpoint
                )
                return cached_fact.value, cached_fact.confidence, source_str

            fact = None
            sql = None

            # Route based on source type (source is required, validated by DAG parser)
            logger.debug(f"[DAG] Premise {fact_id} '{fact_name}' routing with source='{source}'")

            if source == FactSource.CACHE.value:
                # Cached data resolution - look up in fact cache or datastore
                cached_fact = self.fact_resolver.get_fact(fact_name)
                if cached_fact:
                    logger.debug(f"[DAG] Found cached fact: {fact_name} = {str(cached_fact.value)[:50]}...")
                    resolved_premises[fact_id] = cached_fact
                    source_str = format_source_attribution(
                        cached_fact.source, cached_fact.source_name, cached_fact.api_endpoint
                    )
                    return cached_fact.value, cached_fact.confidence, source_str

                # Try datastore table
                if self.datastore:
                    tables = self.datastore.list_tables()
                    for table in tables:
                        if table["name"].lower() == fact_name.lower():
                            # Found table in datastore
                            row_count = table.get("row_count", 0)
                            value_str = f"({table['name']}) {row_count} rows"
                            fact = Fact(
                                name=fact_name,
                                value=value_str,
                                source=FactSource.CACHE,
                                reasoning=f"Cached table from previous analysis",
                                confidence=0.95,
                                table_name=table["name"],
                                row_count=row_count,
                            )
                            resolved_premises[fact_id] = fact
                            self.fact_resolver.add_user_fact(
                                fact_name=fact_name,
                                value=value_str,
                                reasoning="Retrieved from session cache",
                                source=FactSource.CACHE,
                                table_name=table["name"],
                                row_count=row_count,
                            )
                            return value_str, 0.95, "cache"

                raise ValueError(f"Cached data not found: {fact_name}")

            elif source.startswith(FactSource.DATABASE.value):
                # Database resolution
                db_name = source.split(":", 1)[1].strip() if ":" in source else None
                if not db_name:
                    available_dbs = list(self.schema_manager.connections.keys())
                    if available_dbs:
                        db_name = available_dbs[0]
                    else:
                        raise Exception("No SQL databases configured")

                # Check if fact_name matches a table name - return table reference instead of loading
                fact_name_lower = fact_name.lower().strip()
                cache_keys = list(self.schema_manager.metadata_cache.keys())
                logger.debug(f"[DAG] Checking table match for '{fact_name_lower}' in {len(cache_keys)} tables: {cache_keys[:5]}")
                for full_name, table_meta in self.schema_manager.metadata_cache.items():
                    if table_meta.name.lower() == fact_name_lower:
                        # Table match - return reference without loading data
                        columns = [c.name for c in table_meta.columns]
                        row_info = f"{table_meta.row_count:,} rows" if table_meta.row_count else "table"
                        # Use parentheses instead of brackets (Rich interprets [] as markup)
                        value_str = f"({table_meta.database}.{table_meta.name}) {row_info}"
                        fact = Fact(
                            name=fact_name,
                            value=value_str,
                            source=FactSource.DATABASE,
                            source_name=table_meta.database,
                            reasoning=f"Table '{table_meta.name}' from database '{table_meta.database}'. Columns: {columns}",
                            confidence=0.95,
                            table_name=table_meta.name,
                            row_count=table_meta.row_count,
                        )
                        resolved_premises[fact_id] = fact
                        self.fact_resolver.add_user_fact(
                            fact_name=fact_name,
                            value=value_str,
                            reasoning=f"Table reference: {table_meta.database}.{table_meta.name}",
                            source=FactSource.DATABASE,
                        )
                        return value_str, 0.95, f"database:{db_name}"

                engine = self.schema_manager.get_sql_connection(db_name)
                max_retries = 7
                last_error = None

                for attempt in range(max_retries):
                    # Emit SQL generation event
                    self._emit_event(StepEvent(
                        event_type="sql_generating",
                        step_number=0,
                        data={
                            "fact_name": fact_name,
                            "database": db_name,
                            "attempt": attempt + 1,
                            "max_attempts": max_retries,
                            "is_retry": attempt > 0,
                            "retry_reason": last_error[:50] if last_error else None,
                        }
                    ))

                    error_context = f"\nPREVIOUS ERROR: {last_error}\nFix the query." if last_error else ""
                    sql_learnings = self._get_codegen_learnings(fact_desc, TaskType.SQL_GENERATION)

                    sql_prompt = f"""Generate a SQL query to retrieve: {fact_desc}

Schema:
{detailed_schema}
{sql_learnings}
{error_context}
RULES:
- Always SELECT primary key columns for joins
- Always quote identifiers with double quotes (e.g., "group", "order") to avoid reserved word conflicts"""

                    sql_result = self.router.execute(
                        task_type=TaskType.SQL_GENERATION,
                        system="Output raw SQL only. No markdown.",
                        user_message=sql_prompt,
                        max_tokens=self.router.max_output_tokens,
                    )

                    sql = sql_result.content.strip()
                    code_block = re.search(r'```(?:sql)?\s*\n?(.*?)\n?```', sql, re.DOTALL | re.IGNORECASE)
                    if code_block:
                        sql = code_block.group(1).strip()

                    if "sqlite" in str(engine.url):
                        sql = re.sub(rf'\b{db_name}\.(\w+)', r'\1', sql)

                    # Log generated SQL for debugging
                    logger.debug(f"[SQL] Generated SQL for '{fact_name}' (attempt {attempt + 1}/{max_retries}):\n{sql}")

                    # Emit SQL executing event
                    self._emit_event(StepEvent(
                        event_type="sql_executing",
                        step_number=0,
                        data={
                            "fact_name": fact_name,
                            "database": db_name,
                            "attempt": attempt + 1,
                            "max_attempts": max_retries,
                        }
                    ))

                    try:
                        result_df = pd.read_sql(sql, engine)
                        row_count = len(result_df)
                        node.sql_query = sql

                        if row_count == 1 and len(result_df.columns) == 1:
                            scalar_value = result_df.iloc[0, 0]
                            if hasattr(scalar_value, 'item'):
                                scalar_value = scalar_value.item()
                            fact_value = scalar_value
                        else:
                            fact_value = f"{row_count} rows"

                        table_name = fact_name.lower().replace(' ', '_').replace('-', '_')
                        if row_count > 0 and self.datastore:
                            self.datastore.save_dataframe(table_name, result_df)
                            node.row_count = row_count

                        fact = Fact(
                            name=fact_name,
                            value=fact_value,
                            confidence=0.9,
                            source=FactSource.DATABASE,
                            query=sql,
                            table_name=table_name if row_count > 1 else None,
                            row_count=row_count if row_count > 1 else None,
                        )
                        break
                    except Exception as sql_err:
                        last_error = str(sql_err)
                        will_retry = attempt < max_retries - 1
                        # Log full error for debugging
                        logger.warning(f"[SQL] Error for '{fact_name}' (attempt {attempt + 1}/{max_retries}): {sql_err}")
                        logger.debug(f"[SQL] Failed query:\n{sql}")
                        if will_retry:
                            logger.debug(f"[SQL] Will retry ({max_retries - attempt - 1} attempts remaining)")
                        # Emit SQL error event
                        self._emit_event(StepEvent(
                            event_type="sql_error",
                            step_number=0,
                            data={
                                "fact_name": fact_name,
                                "database": db_name,
                                "error": str(sql_err),  # Full error, not truncated
                                "sql": sql[:500],  # Include SQL (truncated to 500 for events)
                                "attempt": attempt + 1,
                                "max_attempts": max_retries,
                                "will_retry": will_retry,
                            }
                        ))
                        if attempt == max_retries - 1:
                            raise Exception(f"SQL error after {max_retries} attempts: {sql_err}")

            elif source.startswith(FactSource.LLM_KNOWLEDGE.value):
                value = self._resolve_llm_knowledge(fact_desc)
                fact = Fact(
                    name=fact_name,
                    value=value,
                    confidence=0.7,
                    source=FactSource.LLM_KNOWLEDGE,
                    reasoning=f"LLM estimate: {fact_desc}",
                )

            elif source.startswith(FactSource.API.value):
                # API resolution via fact resolver
                api_name = source.split(":", 1)[1].strip() if ":" in source else None
                logger.debug(f"[DAG] Resolving API premise {fact_id} '{fact_name}' from API: {api_name}")
                fact, _ = self.fact_resolver.resolve_tiered(fact_name, fact_description=fact_desc)
                if fact and fact.source != FactSource.API:
                    # Fact resolver may have used a different source - update if API was requested
                    logger.debug(f"[DAG] API resolution fell back to {fact.source.value}")

            else:
                # Generic resolution (document, user, or other sources)
                logger.debug(f"[DAG] Using tiered resolution for {fact_id} '{fact_name}' (source={source})")
                fact, _ = self.fact_resolver.resolve_tiered(fact_name, fact_description=fact_desc)

            if fact and fact.value is not None:
                resolved_premises[fact_id] = fact
                if self.history:
                    self.history.save_inference_premise(
                        self.session_id, fact_id, fact_name,
                        fact.value,
                        fact.source.value if hasattr(fact.source, 'value') else str(fact.source),
                        fact_desc or ""
                    )
                self.fact_resolver.add_user_fact(
                    fact_name=fact_name,
                    value=fact.value,
                    query=sql,
                    reasoning=fact.reasoning,
                    source=fact.source,
                    table_name=getattr(fact, 'table_name', None),
                    row_count=getattr(fact, 'row_count', None),
                    context=f"SQL: {sql}" if sql else None,
                )
                source_str = format_source_attribution(
                    fact.source, fact.source_name, fact.api_endpoint
                )
                return fact.value, fact.confidence, source_str
            else:
                raise ValueError(f"Failed to resolve premise: {fact_name}")

        else:
            # === INFERENCE EXECUTION ===
            inf_id = node.fact_id
            operation = node.operation
            explanation = node.description
            inf_name = node.name
            table_name = inf_name.lower().replace(' ', '_').replace('-', '_')
            inference_names[inf_id] = table_name

            # Handle verify_exists operation
            if operation and operation.startswith("verify_exists("):
                ref_match = re.match(r'verify_exists\((\w+)\)', operation)
                if ref_match:
                    ref_id = ref_match.group(1)
                    ref_table = inference_names.get(ref_id, ref_id.lower())
                    if ref_id.startswith("P"):
                        p_idx = int(ref_id[1:]) - 1
                        if p_idx < len(premises):
                            ref_table = premises[p_idx]['name'].lower().replace(' ', '_')

                    if self.datastore:
                        try:
                            profile = self._profile_table(ref_table)
                            row_count = profile["row_count"]

                            # Build profile summary
                            profile_lines = [f"**{ref_table}**: {row_count} rows, {profile['column_count']} columns"]
                            issues = []
                            for col in profile["columns"]:
                                col_line = f"  - `{col['name']}` ({col['type']}): {col['distinct']} distinct"
                                if col["null_pct"] > 0:
                                    col_line += f", {col['null_pct']:.0f}% null"
                                if col["all_null"]:
                                    col_line += " **ALL NULL**"
                                    issues.append(col["name"])
                                profile_lines.append(col_line)

                            if profile["duplicate_rows"] > 0:
                                profile_lines.append(f"\n  Duplicate rows: {profile['duplicate_rows']}")
                                issues.append(f"{profile['duplicate_rows']} duplicate rows")

                            profile_text = "\n".join(profile_lines)
                            logger.info(f"[VERIFY] {inf_id} profile for '{ref_table}':\n{profile_text}")

                            if issues:
                                logger.warning(f"[VERIFY] {inf_id} data quality issues in '{ref_table}': {issues}")

                            # Save profile as artifact
                            if self.datastore:
                                inf_step = 1000 + int(inf_id[1:])
                                self.datastore.add_artifact(
                                    inf_step, 0, "markdown", profile_text,
                                    name=f"profile_{ref_table}",
                                    title=f"Data Profile: {ref_table}",
                                    content_type="text/markdown",
                                )

                            resolved_inferences[inf_id] = f"Verified: {row_count} records"
                            self.fact_resolver.add_user_fact(
                                fact_name=inf_name,
                                value=f"{row_count} records verified" + (f" (issues: {', '.join(issues)})" if issues else ""),
                                reasoning=profile_text,
                                source=FactSource.DERIVED,
                            )
                            confidence = 0.7 if issues else 0.95
                            return f"Verified: {row_count} records", confidence, "derived"
                        except Exception as ve:
                            raise ValueError(f"Verification failed: {ve}")

            # Build context from dependencies including column names
            scalars = []
            tables = []
            referenced_tables = []  # Tables that need to be queried from original database
            api_sources = []  # Data that needs to be fetched from APIs
            for dep_name in node.dependencies:
                dep_node = dag.get_node(dep_name)
                if not dep_node:
                    raise ValueError(f"Dependency '{dep_name}' not found in DAG")
                if dep_node.value is None:
                    raise ValueError(
                        f"Dependency '{dep_name}' ({dep_node.fact_id}) has no value. "
                        f"Status: {dep_node.status}, Error: {dep_node.error}"
                    )
                val_str = str(dep_node.value)
                val_lower = val_str.lower()
                # Check if this is a loaded table (value contains row count)
                is_loaded_table = "rows" in val_lower or (dep_node.row_count and dep_node.row_count > 1)
                # Check if this is a referenced database table (format: (db.table) X rows)
                is_db_referenced = val_str.startswith("(") and ")" in val_str and "." in val_str.split(")")[0]
                # Check if this is an API-sourced reference (format: (api_name) endpoint_display)
                is_api_referenced = (
                    val_str.startswith("(") and ")" in val_str
                    and "." not in val_str.split(")")[0]
                ) or (dep_node.source and dep_node.source.startswith("api"))

                if is_loaded_table or is_db_referenced or is_api_referenced:
                    dep_table = dep_name.lower().replace(' ', '_').replace('-', '_')
                    columns_info = ""

                    if is_api_referenced and not is_loaded_table:
                        # API-sourced reference: extract api name
                        api_match = re.match(r'\((\w+)\)\s*(.*)', val_str)
                        api_name = None
                        if api_match:
                            api_name = api_match.group(1)
                        elif dep_node.source and ":" in dep_node.source:
                            api_name = dep_node.source.split(":", 1)[1].strip()
                        api_sources.append(
                            f"- {dep_node.fact_id}: fetch from API 'api_{api_name or 'unknown'}' into '{dep_table}'"
                        )
                    elif is_db_referenced:
                        # Referenced table: get metadata from original database
                        # Extract table name from the value format (db.table)
                        match = re.match(r'\(([^.]+)\.([^)]+)\)', val_str)
                        if match:
                            ref_db, ref_table = match.groups()
                            dep_table = ref_table  # Use actual table name

                        # Try to get column info from the original database table
                        if self.schema_manager:
                            try:
                                # Find matching table in schema
                                for db_name in self.schema_manager.connections.keys():
                                    table_meta = self.schema_manager.get_table_metadata(db_name, dep_table)
                                    if table_meta:
                                        cols = [c.name for c in table_meta.columns]
                                        columns_info = f" columns: {cols}"
                                        referenced_tables.append(
                                            f"- {dep_node.fact_id}: query from database table '{dep_table}'{columns_info}"
                                        )
                                        break
                            except Exception as e:
                                logger.debug(f"Failed to get table metadata for {dep_table}: {e}")

                        if not columns_info:
                            # Fallback: still mark as referenced
                            referenced_tables.append(
                                f"- {dep_node.fact_id}: referenced table '{dep_table}' (query from database)"
                            )
                    else:
                        # Regular table in datastore
                        sample_info = ""
                        if self.datastore:
                            try:
                                schema_df = self.datastore.query(f"DESCRIBE {dep_table}")
                                if len(schema_df) > 0:
                                    cols = list(schema_df['column_name']) if 'column_name' in schema_df.columns else list(schema_df.iloc[:, 0])
                                    columns_info = f" columns: {cols}"
                            except Exception as e:
                                logger.debug(f"DESCRIBE failed for {dep_table}, trying SELECT: {e}")
                                try:
                                    sample = self.datastore.query(f"SELECT * FROM {dep_table} LIMIT 1")
                                    columns_info = f" columns: {list(sample.columns)}"
                                except Exception as e2:
                                    logger.debug(f"Failed to get columns for {dep_table}: {e2}")
                            # Fetch sample row to show actual data shape
                            try:
                                sample_df = self.datastore.query(f"SELECT * FROM {dep_table} LIMIT 1")
                                if len(sample_df) > 0:
                                    row_dict = {}
                                    for col in sample_df.columns:
                                        val = sample_df.iloc[0][col]
                                        val_str = str(val)
                                        if len(val_str) > 120:
                                            val_str = val_str[:120] + "..."
                                        row_dict[col] = val_str
                                    sample_info = f"\n    sample row: {row_dict}"
                                    row_count = self.datastore.query(f"SELECT COUNT(*) AS cnt FROM {dep_table}").iloc[0, 0]
                                    sample_info += f"\n    total rows: {row_count}"
                            except Exception as e:
                                logger.debug(f"Failed to get sample for {dep_table}: {e}")
                        tables.append(f"- {dep_node.fact_id}: stored as '{dep_table}'{columns_info}{sample_info}")
                else:
                    if dep_node.source == "document":
                        # Document premises: tell LLM to use doc_read() instead of variable reference
                        # Use source_db (from "document:<name>") if available, else fall back to premise name
                        doc_name = dep_node.source_db or dep_name.lower().replace(' ', '_').replace('-', '_')
                        # Validate the doc name against configured documents
                        if self.doc_tools:
                            configured_docs = list(self.doc_tools._loaded_documents.keys()) + list(getattr(self.doc_tools, '_document_configs', {}).keys())
                            if doc_name not in configured_docs and configured_docs:
                                # Try to find a matching configured document
                                for cfg_name in configured_docs:
                                    if cfg_name in dep_name.lower() or dep_name.lower() in cfg_name:
                                        doc_name = cfg_name
                                        break
                        scalars.append(
                            f"- {dep_node.fact_id} ({dep_name}): [DOCUMENT] "
                            f"Load at runtime with `doc_read('{doc_name}')` — do NOT reference {dep_node.fact_id} as a variable"
                        )
                    else:
                        scalars.append(f"- {dep_node.fact_id} ({dep_name}): {dep_node.value}")

            # Build referenced tables section for prompt
            referenced_section = ""
            if referenced_tables:
                referenced_section = f"""
REFERENCED TABLES (query with pd.read_sql(sql, db_<name>)):
{chr(10).join(referenced_tables)}
"""

            # Build API sources section for prompt (with schema info)
            api_sources_section = ""
            api_schema_cache: dict[str, dict] = {}
            if api_sources:
                # Fetch schema for referenced APIs to include in prompt
                schema_lines = list(api_sources)
                all_apis = self.get_all_apis()
                if all_apis:
                    from constat.catalog.api_executor import APIExecutor
                    api_executor = APIExecutor(self.config, project_apis=self._project_apis)
                    for api_name, api_config in all_apis.items():
                        # Only fetch schema for APIs referenced in this inference
                        if not any(f"api_{api_name}" in s for s in api_sources):
                            continue
                        try:
                            if api_config.type == "graphql":
                                overview = api_executor.get_schema_overview(api_name)
                                api_schema_cache[api_name] = overview
                                schema_lines.append(f"\nSchema for api_{api_name} (GraphQL):")
                                for q in overview.get("queries", []):
                                    args_str = ", ".join(q.get("args", []))
                                    schema_lines.append(f"  query {q['name']}({args_str}) -> {q.get('returns', '?')}")
                                    # Get detailed return fields for this query
                                    try:
                                        detail = api_executor.get_query_schema(api_name, q["name"])
                                        if detail.get("return_fields"):
                                            GRAPHQL_SCALARS = {"String", "Int", "Float", "Boolean", "ID"}
                                            field_strs = []
                                            for f in detail["return_fields"]:
                                                ftype = f['type']
                                                # Extract base type name (strip [], !)
                                                base = ftype.replace("[", "").replace("]", "").replace("!", "").strip()
                                                if base in GRAPHQL_SCALARS:
                                                    field_strs.append(f"{f['name']}: {ftype} (scalar, NO subfields)")
                                                else:
                                                    field_strs.append(f"{f['name']}: {ftype}")
                                            schema_lines.append(f"    fields: {', '.join(field_strs)}")
                                    except Exception:
                                        pass
                            else:
                                # REST API: include endpoint info if available
                                if self.schema_manager and hasattr(self.schema_manager, 'api_schema_manager'):
                                    endpoints = self.schema_manager.api_schema_manager.get_api_schema(api_name)
                                    if endpoints:
                                        schema_lines.append(f"\nEndpoints for api_{api_name} (REST):")
                                        for ep in endpoints[:10]:
                                            schema_lines.append(f"  {ep.method} {ep.path} - {ep.description or ep.endpoint_name}")
                                schema_lines.append(f"\nIMPORTANT: api_{api_name} is REST. Response is often paginated:")
                                schema_lines.append(f"  response = api_{api_name}('GET /endpoint', {{params}})")
                                schema_lines.append(f"  # Extract array from wrapper: check 'data', 'results', 'items' keys")
                        except Exception as e:
                            logger.debug(f"Failed to get API schema for {api_name}: {e}")

                api_sources_section = f"""
API SOURCES (fetch with api_<name>() functions):
{chr(10).join(schema_lines)}

IMPORTANT: api_<name>(query) returns the 'data' dict from the GraphQL response.
Example: result = api_countries('{{ country(code: "GB") {{ name languages {{ name }} currency }} }}')
         result == {{"country": {{"name": "United Kingdom", "languages": [{{"name": "English"}}], "currency": "GBP"}}}}
"""

            # Build dynamic data source descriptions
            data_source_apis = []
            data_source_apis.append("- store.query(sql) -> pd.DataFrame (for datastore tables)")
            data_source_apis.append("- store.save_dataframe(name, df)")

            # SQL databases
            from constat.catalog.sql_transpiler import TranspilingConnection
            if self.schema_manager:
                for db_name in self.schema_manager.connections.keys():
                    conn = self.schema_manager.get_connection(db_name)
                    if isinstance(conn, TranspilingConnection):
                        data_source_apis.append(f"- pd.read_sql(query, db_{db_name}) -> DataFrame")
                        data_source_apis.append(f"- sql_{db_name}(query) -> DataFrame (auto-transpiles SQL)")
                    else:
                        data_source_apis.append(f"- pd.read_sql(query, db_{db_name}) -> DataFrame")
                # NoSQL
                for db_name in self.schema_manager.nosql_connections.keys():
                    data_source_apis.append(f"- db_{db_name}: NoSQL connector")
                # File sources
                for db_name in self.schema_manager.file_connections.keys():
                    conn = self.schema_manager.file_connections[db_name]
                    if hasattr(conn, 'path'):
                        ext = str(conn.path).rsplit('.', 1)[-1].lower() if '.' in str(conn.path) else 'csv'
                        reader = {'csv': 'read_csv', 'json': 'read_json', 'parquet': 'read_parquet'}.get(ext, 'read_csv')
                        data_source_apis.append(f"- pd.{reader}(file_{db_name}) -> DataFrame")

            # APIs
            all_apis_for_desc = self.get_all_apis()
            if all_apis_for_desc:
                for api_name, api_config in all_apis_for_desc.items():
                    if api_config.type == "graphql":
                        data_source_apis.append(f"- api_{api_name}(query_string) -> dict (GraphQL 'data' portion)")
                    else:
                        data_source_apis.append(f"- api_{api_name}('GET /endpoint', {{params}}) -> data (REST)")

            # Build step code hints section (from exploratory session)
            step_hints_section = ""
            step_hints = getattr(self, '_proof_step_hints', [])
            if step_hints:
                relevant = []
                op_lower = (operation or "").lower()
                name_lower = (inf_name or "").lower()
                for step in step_hints:
                    goal_lower = (step.get("goal", "") or "").lower()
                    # Include step if goal overlaps with inference operation or name
                    if (any(word in goal_lower for word in name_lower.split() if len(word) > 3)
                            or any(word in goal_lower for word in op_lower.split() if len(word) > 3)):
                        relevant.append(step)
                if not relevant and step_hints:
                    # No keyword match — include all steps as general reference
                    relevant = step_hints
                if relevant:
                    hints = []
                    for step in relevant[:3]:  # Limit to 3 most relevant
                        goal = step.get("goal", f"Step {step.get('step_number', '?')}")
                        code_text = step.get("code", "")
                        if len(code_text) > 800:
                            code_text = code_text[:800] + "\n# ... (truncated)"
                        hints.append(f"# Step: {goal}\n{code_text}")
                    step_hints_section = (
                        "\n\nREFERENCE CODE from exploratory session (use as hints, adapt as needed):\n"
                        + "\n---\n".join(hints) + "\n"
                    )

            # Build codegen learnings for inference
            inference_learnings = ""
            try:
                inference_learnings = self._get_codegen_learnings(
                    f"inference {inf_id}: {operation}", TaskType.SQL_GENERATION
                )
            except Exception as e:
                logger.debug(f"Failed to get codegen learnings for inference: {e}")

            # Generate inference code
            inference_prompt = load_prompt("inference_prompt.md").format(
                inf_id=inf_id,
                inf_name=inf_name,
                operation=operation,
                explanation=explanation,
                scalars="\n".join(scalars) if scalars else "(none)",
                tables="\n".join(tables) if tables else "(none)",
                referenced_section=referenced_section,
                api_sources_section=api_sources_section,
                data_source_apis="\n".join(data_source_apis),
                table_name=table_name,
            ) + step_hints_section
            if inference_learnings:
                inference_prompt += f"\n\nLEARNINGS FROM PREVIOUS ERRORS:\n{inference_learnings}"

            max_retries = 7
            last_error = None
            code = None
            first_error = None  # Track for learning capture
            first_code = None
            _val_passed = []  # True assertions (structural + user-specified)
            _val_profile = []  # Data profile stats for human review

            for attempt in range(max_retries):
                prompt = inference_prompt
                if last_error:
                    prompt = f"PREVIOUS ERROR: {last_error}\n\n{inference_prompt}"

                code_result = self.router.execute(
                    task_type=TaskType.SQL_GENERATION,
                    system="Generate Python code. Return only executable code.",
                    user_message=prompt,
                    max_tokens=self.router.max_output_tokens,
                )

                code = code_result.content.strip()
                if code.startswith("```"):
                    code = re.sub(r'^```\w*\n?', '', code)
                    code = re.sub(r'\n?```$', '', code)

                logger.info(f"[INFERENCE_CODE] {inf_id} attempt {attempt+1}: code length={len(code)} chars")
                logger.debug(f"[INFERENCE_CODE] {inf_id} attempt {attempt+1} code:\n{code}")

                # Auto-fix DataFrame boolean errors
                code = re.sub(r'\bif\s+(df|result|data)\s*:', r'if not \1.empty:', code)
                code = re.sub(r'\bif\s+not\s+(df|result|data)\s*:', r'if \1.empty:', code)

                node.code = code

                import io
                import sys
                import numpy as np

                # Build full execution globals (databases, APIs, file sources)
                exec_globals = self._get_execution_globals()
                exec_globals["store"] = self.datastore
                exec_globals["pd"] = pd
                exec_globals["np"] = np
                exec_globals["llm_map"] = constat.llm.llm_map
                exec_globals["llm_classify"] = constat.llm.llm_classify
                exec_globals["llm_extract"] = constat.llm.llm_extract
                exec_globals["llm_summarize"] = constat.llm.llm_summarize
                exec_globals["llm_score"] = constat.llm.llm_score
                exec_globals["doc_read"] = self._create_doc_read_helper()
                self._inference_used_llm_map = False  # Reset per inference

                # Add resolved values to context
                for pid, fact in resolved_premises.items():
                    if fact and fact.value is not None:
                        val = fact.value
                        if isinstance(val, str) and "rows" in val:
                            try:
                                val = int(val.split()[0])
                            except (ValueError, IndexError):
                                pass
                        exec_globals[pid] = val

                for iid, ival in resolved_inferences.items():
                    if ival is not None:
                        val = ival
                        if isinstance(val, str) and "rows" in val:
                            try:
                                val = int(val.split()[0])
                            except (ValueError, IndexError):
                                pass
                        exec_globals[iid] = val

                captured = io.StringIO()
                old_stdout = sys.stdout
                try:
                    sys.stdout = captured
                    exec(code, exec_globals)

                    # --- Post-execution validation ---
                    _val_computed = exec_globals.get('_result')
                    _val_tables = [t['name'] for t in self.datastore.list_tables()] if self.datastore else []
                    _val_output = captured.getvalue().strip()
                    _val_error = None
                    _val_passed = []  # True assertions (structural + user-specified)
                    _val_profile = []  # Data profile stats for human review

                    # 1. Result exists (unless table saved or stdout produced)
                    if _val_computed is None and table_name not in _val_tables and not _val_output:
                        _val_error = f"No result produced. Set _result, save table '{table_name}', or print output."
                    else:
                        _val_passed.append("Result produced")

                    # 2. DataFrame not empty
                    if not _val_error and _val_computed is not None and hasattr(_val_computed, 'empty') and _val_computed.empty:
                        _val_error = f"Result DataFrame is empty (0 rows). Expected data in '{table_name}'."
                    elif not _val_error and _val_computed is not None and hasattr(_val_computed, 'empty'):
                        _val_passed.append(f"DataFrame has {len(_val_computed)} rows")

                    # 3. Saved table has rows
                    if not _val_error and table_name in _val_tables:
                        _row_ct = int(self.datastore.query(f'SELECT COUNT(*) FROM "{table_name}"').iloc[0, 0])
                        if _row_ct == 0:
                            _val_error = f"Table '{table_name}' saved but has 0 rows."
                        else:
                            _val_passed.append(f"Table '{table_name}' has {_row_ct} rows")

                    # 4. No all-null columns
                    if not _val_error and table_name in _val_tables:
                        try:
                            _null_df = self.datastore.query(
                                f'SELECT column_name FROM (SELECT * FROM "{table_name}" LIMIT 1000) '
                                f'UNPIVOT (value FOR column_name IN (*)) '
                                f'GROUP BY column_name HAVING COUNT(CASE WHEN value IS NOT NULL AND value != \'\' THEN 1 END) = 0'
                            )
                            if len(_null_df) > 0:
                                _val_error = f"Table '{table_name}' has all-NULL columns: {list(_null_df['column_name'])}. Data enrichment likely failed."
                            else:
                                _val_passed.append("No all-NULL columns")
                        except Exception:
                            pass

                    # 5. Column-level profiling → _val_profile (stats for human review)
                    if not _val_error and table_name in _val_tables:
                        try:
                            _cols_df = self.datastore.query(f'SELECT * FROM "{table_name}" LIMIT 0')
                            for _c in _cols_df.columns:
                                try:
                                    _stats = self.datastore.query(
                                        f'SELECT MIN("{_c}") as mn, MAX("{_c}") as mx, '
                                        f'COUNT("{_c}") as cnt, COUNT(*) as total '
                                        f'FROM "{table_name}"'
                                    )
                                    _mn, _mx = _stats.iloc[0]['mn'], _stats.iloc[0]['mx']
                                    _cnt = int(_stats.iloc[0]['cnt'])
                                    _tot = int(_stats.iloc[0]['total'])
                                    if isinstance(_mn, (int, float)) and isinstance(_mx, (int, float)):
                                        if _mn == _mx:
                                            _val_profile.append(f"{_c}: all values = {_mn}")
                                        else:
                                            _val_profile.append(f"{_c}: {_mn} to {_mx}")
                                    if _cnt < _tot:
                                        _val_profile.append(f"{_c}: {_cnt}/{_tot} non-null")
                                except Exception:
                                    pass
                        except Exception:
                            pass

                    # 6. Input-output row ratio vs dependencies
                    if not _val_error and table_name in _val_tables and node.dependencies:
                        try:
                            _out_ct = int(self.datastore.query(f'SELECT COUNT(*) FROM "{table_name}"').iloc[0, 0])
                            for _dep_name in node.dependencies:
                                _dep_node = dag.get_node(_dep_name)
                                if _dep_node and _dep_node.row_count and _dep_node.row_count > 1:
                                    if abs(_out_ct - _dep_node.row_count) <= max(1, _dep_node.row_count * 0.2):
                                        _val_profile.append(f"Row count matches {_dep_node.fact_id} ({_out_ct} vs {_dep_node.row_count})")
                                    elif _out_ct > _dep_node.row_count:
                                        _val_profile.append(f"Expanded from {_dep_node.fact_id}: {_dep_node.row_count} → {_out_ct} rows")
                                    else:
                                        _val_profile.append(f"Filtered from {_dep_node.fact_id}: {_dep_node.row_count} → {_out_ct} rows")
                        except Exception:
                            pass

                    # 7. User-specified validations (from query constraints)
                    _user_validations = getattr(self, '_proof_user_validations', [])
                    if not _val_error and _user_validations and table_name in _val_tables:
                        for _uv in _user_validations:
                            try:
                                _uv_result = self.datastore.query(_uv['sql'].format(table=table_name))
                                _uv_ok = bool(_uv_result.iloc[0, 0]) if len(_uv_result) > 0 else False
                                if _uv_ok:
                                    _val_passed.append(f"{_uv['label']}")
                                else:
                                    _val_error = f"User validation failed: {_uv['label']}"
                            except Exception as _uv_e:
                                logger.debug(f"User validation '{_uv.get('label', '?')}' skipped for {table_name}: {_uv_e}")

                    if _val_error:
                        last_error = _val_error
                        logger.warning(f"[INFERENCE_CODE] {inf_id} attempt {attempt+1} validation: {_val_error}")
                        if first_error is None:
                            first_error = last_error
                            first_code = code
                        continue
                    # --- End validation ---

                    last_error = None
                    # Capture learning if this was a successful retry
                    if attempt > 0 and first_error and self.learning_store:
                        try:
                            self._capture_error_learning(
                                context={
                                    "error_message": first_error,
                                    "original_code": first_code[:500] if first_code else "",
                                    "step_goal": f"inference {inf_id}: {operation}",
                                },
                                fixed_code=code,
                            )
                        except Exception as le:
                            logger.debug(f"Learning capture failed: {le}")
                    break
                except Exception as e:
                    last_error = str(e)
                    logger.warning(f"[INFERENCE_CODE] {inf_id} attempt {attempt+1}/{max_retries} failed: {last_error}")
                    logger.debug(f"[INFERENCE_CODE] {inf_id} attempt {attempt+1} code:\n{code}")
                    if first_error is None:
                        first_error = last_error
                        first_code = code
                finally:
                    sys.stdout = old_stdout

            if last_error:
                logger.error(f"[INFERENCE_CODE] {inf_id} all attempts failed: {last_error}")
                logger.debug(f"[INFERENCE_CODE] Failed code for {inf_id}:\n{code}")
                raise Exception(last_error)

            computed = exec_globals.get('_result')
            # Log execution results for debugging
            tables_after = [t['name'] for t in self.datastore.list_tables()] if self.datastore else []
            logger.debug(f"[INFERENCE_CODE] {inf_id} exec complete: _result={computed is not None}, tables={tables_after}, expected={table_name}")

            # Check if result was saved as table
            if self.datastore and table_name in [t['name'] for t in self.datastore.list_tables()]:
                count_df = self.datastore.query(f"SELECT COUNT(*) FROM {table_name}")
                row_count = int(count_df.iloc[0, 0])
                node.row_count = row_count
                resolved_inferences[inf_id] = f"{row_count} rows"
                result_value = f"{row_count} rows"

                # Warn if row count is suspiciously low compared to inputs
                if row_count <= 1 and node.dependencies:
                    # Check if any dependency had more rows
                    for dep_name in node.dependencies:
                        dep_node = dag.get_node(dep_name)
                        if dep_node and dep_node.row_count and dep_node.row_count > 5:
                            logger.warning(
                                f"[INFERENCE_CODE] {inf_id} produced only {row_count} row(s) but "
                                f"dependency '{dep_name}' had {dep_node.row_count} rows. "
                                f"This may indicate incorrect aggregation."
                            )
            elif computed is not None:
                resolved_inferences[inf_id] = computed
                result_value = computed
            else:
                output = captured.getvalue().strip()
                if output:
                    resolved_inferences[inf_id] = output
                    result_value = output
                else:
                    # No table created, no _result, no output - this is likely a failure
                    # Check if operation suggests a table should have been created
                    if any(kw in operation.lower() for kw in ['join', 'filter', 'merge', 'apply', 'calculate', 'select']):
                        logger.error(f"[INFERENCE_CODE] {inf_id} ({inf_name}) produced no output. Code:\n{code[:500]}...")
                        raise ValueError(f"Inference {inf_id} ({inf_name}) did not produce expected table '{table_name}'")
                    resolved_inferences[inf_id] = "completed"
                    result_value = "completed"

            # Persist inference code to disk (separate from step plan code)
            if self.history and code:
                output = captured.getvalue().strip()
                self.history.save_inference_code(
                    session_id=self.session_id,
                    inference_id=inf_id,
                    name=inf_name,
                    operation=operation,
                    code=code,
                    attempt=attempt + 1,
                    output=output or None,
                )

            # Reduce confidence if inference used LLM fuzzy mapping
            used_llm = getattr(self, '_inference_used_llm_map', False)

            # Safety net: detect hardcoded mapping dicts in generated code
            # LLM sometimes embeds its knowledge as a literal dict instead of calling llm_map()
            if not used_llm and code and '.map(' in code:
                # Check for dict literals with 3+ string key-value pairs followed by .map()
                import ast
                try:
                    tree = ast.parse(code)
                    for node_ast in ast.walk(tree):
                        if isinstance(node_ast, ast.Dict) and len(node_ast.keys) >= 3:
                            # Check if most keys are string constants
                            str_keys = sum(1 for k in node_ast.keys if isinstance(k, ast.Constant) and isinstance(k.value, str))
                            if str_keys >= 3:
                                used_llm = True
                                logger.info(f"[INFERENCE_CODE] {inf_id}: detected hardcoded mapping dict ({str_keys} string keys) — flagging as LLM knowledge")
                                break
                except Exception:
                    pass

            confidence = 0.65 if used_llm else 0.9
            source = FactSource.LLM_KNOWLEDGE if used_llm else FactSource.DERIVED

            self.fact_resolver.add_user_fact(
                fact_name=inf_name,
                value=result_value,
                reasoning=f"Computed: {operation}" + (" (includes LLM fuzzy mapping)" if used_llm else ""),
                source=source,
                context=f"Code:\n{code}" if code else None,
            )

            return result_value, confidence, "llm_knowledge" if used_llm else "derived", _val_passed, _val_profile

    @staticmethod
    def _find_skill_script(scripts_dir: Path) -> Path | None:
        """Find the first executable Python script in a skill's scripts directory."""
        if not scripts_dir.exists():
            return None
        for ext in ("*.py",):
            scripts = sorted(scripts_dir.glob(ext))
            if scripts:
                return scripts[0]
        return None

    def _ensure_session_datastore(self, problem: str) -> None:
        """Create history session_id and datastore if not yet initialized."""
        # session_id may be a server UUID (not a valid history directory) — check filesystem
        needs_history = not self.session_id
        if self.session_id:
            history_dir = self.history._session_dir(self.session_id)
            if not (history_dir / "session.json").exists():
                needs_history = True
        if needs_history:
            self.session_id = self.history.create_session(
                config_dict=self.config.model_dump(),
                databases=self.resources.database_names,
                apis=self.resources.api_names,
                documents=self.resources.document_names,
                server_session_id=self.server_session_id,
            )
        if not self.datastore:
            session_dir = self.history._session_dir(self.session_id)
            datastore_path = session_dir / "datastore.duckdb"
            tables_dir = session_dir / "tables"
            underlying_datastore = DataStore(db_path=datastore_path)
            self.datastore = RegistryAwareDataStore(
                datastore=underlying_datastore,
                registry=self.registry,
                user_id=self.user_id,
                session_id=self.session_id,
                tables_dir=tables_dir,
            )
            self.fact_resolver._datastore = self.datastore

    def _execute_skill_script(self, script_path: Path, problem: str) -> dict | None:
        """Execute a skill script directly, bypassing planning.

        Looks for a callable entry point (run_proof, run, main) in the script.
        Returns solve()-compatible result dict on success, None on failure.
        """
        import time
        start_time = time.time()
        skill_name = script_path.parent.parent.name

        self._emit_event(StepEvent(
            event_type="step_start",
            step_number=1,
            data={"goal": f"Executing skill: {skill_name}"}
        ))

        try:
            script_content = script_path.read_text()

            # Build execution globals with all session resources
            exec_globals = self._get_execution_globals()
            exec_globals['__file__'] = str(script_path)

            # Execute the script to define entry point functions
            exec(compile(script_content, str(script_path), 'exec'), exec_globals)

            # Find entry point: try common names
            entry_fn = None
            for fn_name in ('run_proof', 'run', 'main'):
                fn = exec_globals.get(fn_name)
                if callable(fn):
                    entry_fn = fn
                    break

            if entry_fn is None:
                logger.warning(f"[SKILL_EXEC] No entry point (run_proof/run/main) found in {script_path.name}")
                return None

            # Run the entry point
            results = entry_fn()

            if not isinstance(results, dict):
                logger.warning(f"[SKILL_EXEC] Entry point returned {type(results)}, expected dict")
                return None

            # Save result tables to datastore
            # run_proof() returns {name: parquet_path} or {name: DataFrame}
            import pandas as pd
            saved_tables = []
            for name, val in results.items():
                if name == '_result':
                    continue
                if isinstance(val, pd.DataFrame):
                    self.datastore.save_dataframe(name, val, step_number=1)
                    saved_tables.append(name)
                elif isinstance(val, str) and val.endswith('.parquet'):
                    df = pd.read_parquet(val)
                    self.datastore.save_dataframe(name, df, step_number=1)
                    saved_tables.append(name)

            # Save _result as the primary output
            primary_result = results.get('_result')
            if isinstance(primary_result, pd.DataFrame):
                self.datastore.save_dataframe('_result', primary_result, step_number=1)
            elif isinstance(primary_result, str) and primary_result.endswith('.parquet'):
                df = pd.read_parquet(primary_result)
                self.datastore.save_dataframe('_result', df, step_number=1)

            duration_ms = int((time.time() - start_time) * 1000)

            self._emit_event(StepEvent(
                event_type="step_complete",
                step_number=1,
                data={
                    "goal": f"Skill execution complete ({len(saved_tables)} tables)",
                    "duration_ms": duration_ms,
                }
            ))

            logger.info(f"[SKILL_EXEC] Success: {len(saved_tables)} tables saved in {duration_ms}ms")

            return {
                "steps": [{"step_number": 1, "goal": f"Skill: {skill_name}", "status": "success"}],
                "tables": saved_tables,
                "mode": "skill",
                "duration_ms": duration_ms,
            }

        except Exception as e:
            logger.warning(f"[SKILL_EXEC] Execution failed: {e}")
            self._emit_event(StepEvent(
                event_type="step_error",
                step_number=1,
                data={"error": str(e)}
            ))
            return None

    def _extract_user_validations(self, problem: str, inferences: list[dict]) -> list[dict]:
        """Extract user-specified validation constraints from the problem text.

        Looks for explicit constraints like "ensure no raise exceeds 15%",
        "verify total budget under $100k", etc. and converts them to SQL checks.

        Returns list of dicts with 'label' and 'sql' keys.
        The 'sql' value uses {table} placeholder for the target table name.
        """
        # Build inference context for the LLM
        inf_desc = "\n".join(
            f"- {inf.get('inference_id', '?')}: {inf.get('name', '')} = {inf.get('operation', '')}"
            for inf in inferences
        )

        prompt = f"""Analyze this user request for explicit validation constraints (ensure, verify, validate, must, should not exceed, at least, between, limit, cap, maximum, minimum, etc.):

USER REQUEST: {problem}

INFERENCES (output tables):
{inf_desc}

Extract ONLY explicitly stated constraints. Do NOT invent constraints that aren't in the request.

For each constraint, provide:
- label: Human-readable description of the check
- sql: A DuckDB SQL expression that returns TRUE if the constraint passes, using {{table}} as placeholder for the table name

Respond with ONLY valid JSON array. Empty array [] if no explicit constraints found.

Example:
[
  {{"label": "No raise exceeds 15%", "sql": "SELECT COUNT(*) = 0 FROM \\"{{table}}\\" WHERE raise_pct > 0.15"}},
  {{"label": "Total budget under $100k", "sql": "SELECT SUM(raise_amount) < 100000 FROM \\"{{table}}\\""}}
]

YOUR JSON RESPONSE:"""

        try:
            result = self.router.execute(
                task_type=TaskType.SQL_GENERATION,
                system="Extract validation constraints from user requests. Output ONLY valid JSON.",
                user_message=prompt,
            )
            content = result.content if hasattr(result, 'content') else str(result)
            # Strip markdown fences
            content = content.strip()
            if content.startswith("```"):
                content = content.split("\n", 1)[-1].rsplit("```", 1)[0].strip()

            import json
            validations = json.loads(content)
            if isinstance(validations, list):
                valid = [v for v in validations if isinstance(v, dict) and 'label' in v and 'sql' in v]
                if valid:
                    logger.info(f"[USER_VALIDATIONS] Extracted {len(valid)} constraints: {[v['label'] for v in valid]}")
                return valid
        except Exception as e:
            logger.debug(f"[USER_VALIDATIONS] Extraction failed: {e}")

        return []

    def add_user_validation(self, label: str, sql: str) -> None:
        """Add a user-specified validation constraint for proof inference checks.

        Args:
            label: Human-readable description (e.g., "No raise exceeds 15%")
            sql: DuckDB SQL that returns TRUE if valid. Use {table} placeholder.
        """
        if not hasattr(self, '_proof_user_validations'):
            self._proof_user_validations = []
        self._proof_user_validations.append({"label": label, "sql": sql})
        logger.info(f"[USER_VALIDATIONS] Added: {label}")

    def _profile_table(self, table_name: str) -> dict:
        """Profile a datastore table for data quality assessment.

        Returns dict with row_count, column_count, duplicate_rows, and per-column stats.
        """
        row_count_df = self.datastore.query(f"SELECT COUNT(*) as cnt FROM {table_name}")
        row_count = int(row_count_df.iloc[0, 0])

        schema_df = self.datastore.query(f"DESCRIBE {table_name}")
        columns = []
        for _, row in schema_df.iterrows():
            col_name = row.iloc[0]
            col_type = str(row.iloc[1]) if len(row) > 1 else "unknown"
            try:
                stats = self.datastore.query(
                    f"SELECT "
                    f"COUNT(DISTINCT \"{col_name}\") as distinct_count, "
                    f"SUM(CASE WHEN \"{col_name}\" IS NULL OR CAST(\"{col_name}\" AS VARCHAR) = '' THEN 1 ELSE 0 END) as null_count "
                    f"FROM {table_name}"
                )
                distinct = int(stats.iloc[0]['distinct_count'])
                null_count = int(stats.iloc[0]['null_count'])
            except Exception:
                distinct = 0
                null_count = row_count

            null_pct = (null_count / row_count * 100) if row_count > 0 else 0
            columns.append({
                "name": col_name,
                "type": col_type,
                "distinct": distinct,
                "null_count": null_count,
                "null_pct": null_pct,
                "all_null": null_count == row_count,
            })

        # Check for duplicate rows
        try:
            dup_df = self.datastore.query(
                f"SELECT COUNT(*) - COUNT(DISTINCT *) as dups FROM (SELECT * FROM {table_name})"
            )
            duplicate_rows = int(dup_df.iloc[0, 0])
        except Exception:
            duplicate_rows = 0

        return {
            "row_count": row_count,
            "column_count": len(columns),
            "columns": columns,
            "duplicate_rows": duplicate_rows,
        }

    def _make_json_serializable(self, obj: Any) -> Any:
        """Convert an object to be JSON-serializable.

        Handles pandas NA values (NAType), numpy types, and nested structures.
        """
        import numpy as np
        import pandas as pd

        if obj is None:
            return None
        # Handle pandas NA
        if pd.isna(obj):
            return None
        # Handle numpy types
        if isinstance(obj, (np.integer, np.floating)):
            return obj.item()
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        # Handle nested structures
        if isinstance(obj, dict):
            return {k: self._make_json_serializable(v) for k, v in obj.items()}
        if isinstance(obj, (list, tuple)):
            return [self._make_json_serializable(v) for v in obj]
        return obj

    def _get_system_prompt(self) -> str:
        """Get the system prompt with active role and skills appended.

        Returns:
            The config system prompt + active role prompt + active skills prompts
        """
        parts = []

        # Base system prompt from config
        base_prompt = self.config.system_prompt or ""
        if base_prompt:
            parts.append(base_prompt)

        # Role prompt (if active)
        role_prompt = self.role_manager.get_role_prompt()
        if role_prompt:
            parts.append(role_prompt)

        # Skills prompts (if any active)
        active_skill_objects = self.skill_manager.active_skill_objects
        if active_skill_objects:
            skill_parts = []
            for skill in active_skill_objects:
                skill_dir = self.skill_manager.skills_dir / skill.filename
                scripts_dir = skill_dir / "scripts"
                script_files = []
                if scripts_dir.exists():
                    script_files = sorted(
                        str(f) for f in scripts_dir.iterdir() if f.is_file()
                    )
                if script_files:
                    skill_parts.append(
                        f"## Skill: {skill.name} — EXECUTE, DO NOT REWRITE\n"
                        f"Scripts: {', '.join(script_files)}\n\n"
                        f"**Load the script and call its `run_proof()` function.** Do NOT reimplement the logic.\n\n"
                        f"Generate ONLY this exact pattern — no summary code, no column references, no extra logic:\n"
                        f"```python\n"
                        f"import pandas as pd\n\n"
                        f"# 1. Load script into its own namespace\n"
                        f"_ns = {{}}\n"
                        f"exec(open('{script_files[0]}').read(), _ns)\n\n"
                        f"# 2. Call run_proof() — returns dict[str, str] of Parquet file paths\n"
                        f"file_paths = _ns['run_proof']()\n\n"
                        f"# 3. Load and save each dataset to the session store\n"
                        f"for name, path in file_paths.items():\n"
                        f"    store.save_dataframe(name, pd.read_parquet(path))\n"
                        f"_result = pd.read_parquet(file_paths['_result'])\n"
                        f"```\n\n"
                        f"CRITICAL: Do NOT add print statements that reference specific column names.\n"
                        f"The skill script's output columns may differ from the documentation below.\n"
                        f"Just load, save, and assign `_result`. Nothing else.\n\n"
                        f"Skill documentation (for context only — do NOT hardcode column names from this):\n\n"
                        f"{skill.prompt}"
                    )
                else:
                    skill_parts.append(f"## Skill: {skill.name} (reference)\n{skill.prompt}")
            parts.append("# Active Skills\n\n" + "\n\n".join(skill_parts))

        return "\n\n".join(parts)

    def _build_source_context(self, include_user_facts: bool = True, query: str = None) -> dict:
        """Build context about available data sources (schema, APIs, documents, facts).

        Args:
            include_user_facts: Whether to include resolved user facts
            query: Optional natural language query for semantic source search.
                   When provided, finds relevant tables, documents, and APIs
                   via similarity search and includes targeted context.

        Returns:
            dict with keys: schema_overview, api_overview, doc_overview, user_facts
        """
        # When query is provided, use semantic search across ALL source types
        if query:
            sources = self.find_relevant_sources(query, table_limit=10, doc_limit=5, api_limit=5)
            schema_overview = self._format_relevant_tables(sources.get("tables", []))
            api_overview = self._format_relevant_apis(sources.get("apis", []))
            doc_overview = self._format_relevant_docs(sources.get("documents", []))
        else:
            # Schema overview - prefer preloaded hot tables over full listing
            if self._preloaded_context:
                schema_overview = self._preloaded_context
                schema_overview += "\n\nUse `find_relevant_tables(query)` or `get_table_schema(table)` for other tables."
            else:
                schema_overview = self._get_brief_schema_summary()
                schema_overview += "\n\nUse discovery tools to explore schemas: `find_relevant_tables(query)`, `get_table_schema(table)`"

            # API overview - use self.resources (single source of truth)
            api_overview = ""
            if self.resources.has_apis():
                api_lines = ["\n## Available APIs"]
                for name, api_info in self.resources.apis.items():
                    api_type = api_info.api_type.upper()
                    desc = api_info.description or f"{api_type} endpoint"
                    api_lines.append(f"- **{name}** ({api_type}): {desc}")
                api_overview = "\n".join(api_lines)

            # Document overview - use self.resources (single source of truth)
            doc_overview = ""
            if self.resources.has_documents():
                doc_lines = ["\n## Reference Documents"]
                for name, doc_info in self.resources.documents.items():
                    desc = doc_info.description or doc_info.doc_type
                    doc_lines.append(f"- **{name}**: {desc}")
                doc_overview = "\n".join(doc_lines)

        # User facts
        user_facts = ""
        if include_user_facts:
            try:
                all_facts = self.fact_resolver.get_all_facts()
                if all_facts:
                    fact_lines = ["\n## Known User Facts (use these values in code)"]
                    for name, fact in all_facts.items():
                        fact_lines.append(f"- **{name}**: {fact.value}")
                    user_facts = "\n".join(fact_lines)
            except Exception as e:
                logger.debug(f"Failed to get user facts for context: {e}")

        return {
            "schema_overview": schema_overview,
            "api_overview": api_overview,
            "doc_overview": doc_overview,
            "user_facts": user_facts,
        }

    def _build_available_sources_description(self) -> str:
        """Build a concise description of available data sources for Tier 2 assessment.

        This is used by the fact resolver's Tier 2 LLM assessment to understand
        what sources are available without providing full schema details.
        """
        lines = []

        # Databases (names + descriptions, not full schema)
        if self.config.databases:
            lines.append("Databases:")
            for name, db_config in self.config.databases.items():
                desc = db_config.description or db_config.type or "SQL database"
                lines.append(f"  - {name}: {desc}")

        # Documents
        if self.config.documents:
            lines.append("Documents:")
            for name, doc_config in self.config.documents.items():
                desc = doc_config.description or doc_config.type
                lines.append(f"  - {name}: {desc}")

        # APIs
        if self.config.apis:
            lines.append("APIs:")
            for name, api_config in self.config.apis.items():
                desc = api_config.description or f"{api_config.type} endpoint"
                lines.append(f"  - {name}: {desc}")

        # Config values
        if self.config.system_prompt:
            lines.append("Config: Domain context available in system prompt")

        # Hot tables (preloaded schema)
        if self._preloaded_context:
            lines.append("Hot tables (schema preloaded): See system context")

        # Tool calling documentation
        lines.append("Discovery tools available: find_relevant_tables(), get_table_schema(), search_documents()")

        return "\n".join(lines) if lines else "(no sources configured)"

    def _create_llm_ask_helper(self) -> callable:
        """Create a helper function for step code to query LLM for general knowledge."""
        def llm_ask(question: str) -> int | float | str:
            """
            Ask the LLM a general knowledge question.

            Use this for facts not available in the databases, such as:
            - Industry benchmarks and averages
            - General domain knowledge
            - Conversion factors or standard values
            - Definitions and explanations

            Args:
                question: The question to ask

            Returns:
                The value (number, string, or ISO date string)
            """
            return self._resolve_llm_knowledge(question)
        return llm_ask

    def _handle_llm_call_event(self, event) -> None:
        """Callback from constat.llm primitives — flags LLM knowledge usage."""
        self._inference_used_llm_map = True

    def _create_doc_read_helper(self) -> callable:
        """Create a helper to read reference documents at execution time."""
        def doc_read(name: str) -> str:
            """Read a reference document by name. Returns the document text content.

            Prints provenance info (source, path, modified date) to stdout so
            consumers can verify the document is current.

            Args:
                name: Document name as configured (e.g., 'compensation_policy', 'business_rules')

            Returns:
                Document text content.

            Raises:
                ValueError: If document not found or has no text content.
            """
            from datetime import datetime

            # Reload fresh to pick up any file changes
            self.doc_tools._load_document_with_mtime(name)

            result = self.doc_tools.get_document(name)
            if "error" in result:
                raise ValueError(f"Document '{name}' not found: {result['error']}")
            content = result.get("content")
            if not content:
                raise ValueError(f"Document '{name}' has no text content (may be binary)")

            # Build provenance info
            provenance_parts = [f"[DOC] Source: {name}"]
            path = result.get("path")
            if path:
                provenance_parts.append(f"Path: {path}")
            doc = self.doc_tools._loaded_documents.get(name)
            if doc and doc.file_mtime:
                mtime_str = datetime.fromtimestamp(doc.file_mtime).strftime("%Y-%m-%d %H:%M:%S")
                provenance_parts.append(f"Modified: {mtime_str}")
            if doc and doc.content_hash:
                provenance_parts.append(f"Hash: {doc.content_hash}")
            print(" | ".join(provenance_parts))

            return content

        return doc_read

    def _is_current_plan_sensitive(self) -> bool:
        """Check if the current plan involves sensitive data."""
        return self.plan is not None and self.plan.contains_sensitive_data

    def _create_publish_helper(self) -> callable:
        """Create a helper function for step code to publish artifacts.

        Published artifacts appear in the artifacts panel (consequential outputs).
        Unpublished artifacts are still accessible via inline links and /artifacts.
        """
        def publish(name: str, title: str = None, description: str = None) -> bool:
            """
            Mark an artifact as published for the artifacts panel.

            Call this for artifacts that are consequential outputs (final deliverables)
            rather than intermediate results. Published artifacts appear prominently
            in the artifacts panel.

            Args:
                name: The table or artifact name to publish
                title: Optional human-friendly display title
                description: Optional description

            Returns:
                True if published successfully, False if artifact not found

            Example:
                # After creating a final summary table
                store.save_dataframe('executive_summary', summary_df)
                publish('executive_summary', title='Executive Summary Report')
            """
            if not self.registry:
                return False

            # Try to publish as table first
            if self.registry.publish_table(
                self.user_id, self.session_id, name,
                is_published=True, title=title
            ):
                return True

            # Try as artifact
            if self.registry.publish_artifact(
                self.user_id, self.session_id, name,
                is_published=True, title=title
            ):
                return True

            return False

        return publish

    def _get_execution_globals(self) -> dict:
        """Get globals dict for code execution.

        Each step runs in isolation - only `store` (DuckDB) is shared.
        """
        def parse_number(val):
            """Parse string-formatted numbers, handling ranges, series, percentages, currency, and units.

            Returns a tuple of ALL extracted numbers, preserving order.
            Use min()/max() on the result for range bounds.

            Examples:
              "8-12%"           → (8.0, 12.0)
              "5%"              → (5.0,)
              "$1,200"          → (1200.0,)
              "8 to 12"         → (8.0, 12.0)
              "between 5 and 10"→ (5.0, 10.0)
              "1, 2, 3"         → (1.0, 2.0, 3.0)
              "1; 2; 3"         → (1.0, 2.0, 3.0)
              "10k"             → (10000.0,)
              "1.5M"            → (1500000.0,)
              "(5%)"            → (-5.0,)
              "up to 15%"       → (0.0, 15.0)
              None / NaN        → (0.0,)
            """
            import re as _re
            if val is None:
                return (0.0,)
            # Handle numeric passthrough
            if isinstance(val, (int, float)):
                if val != val:  # NaN check
                    return (0.0,)
                return (float(val),)

            s = str(val).strip()
            if not s:
                return (0.0,)

            # Detect accounting-style negatives: (5%) → -5
            is_accounting_neg = s.startswith('(') and s.endswith(')')
            if is_accounting_neg:
                s = s[1:-1].strip()

            # Strip currency symbols, normalize whitespace
            s = _re.sub(r'[£€¥₹]', '', s)
            s = s.replace('$', '').replace('%', '').replace(',', '').strip()

            # Expand unit suffixes: k, M, B, T
            _unit_mult = {'k': 1e3, 'K': 1e3, 'm': 1e6, 'M': 1e6, 'b': 1e9, 'B': 1e9, 't': 1e12, 'T': 1e12}
            def _apply_unit(num_str):
                num_str = num_str.strip()
                if num_str and num_str[-1] in _unit_mult:
                    try:
                        return float(num_str[:-1]) * _unit_mult[num_str[-1]]
                    except ValueError:
                        pass
                try:
                    return float(num_str)
                except ValueError:
                    return None

            # Split on range/list delimiters: -, –, —, to, and, /, ;, ,, |, or
            parts = _re.split(r'\s*[-–—/;,|]\s*|\s+to\s+|\s+and\s+|\s+or\s+', s, flags=_re.IGNORECASE)
            numbers = []
            for p in parts:
                p = p.strip()
                if not p:
                    continue
                v = _apply_unit(p)
                if v is not None:
                    numbers.append(v)

            # Fallback: regex extract all decimal numbers from original string
            if not numbers:
                for m in _re.finditer(r'-?\d+\.?\d*', s):
                    try:
                        numbers.append(float(m.group()))
                    except ValueError:
                        pass

            if not numbers:
                return (0.0,)

            if is_accounting_neg:
                numbers = [-n for n in numbers]

            # "up to X" → (0, X)
            if _re.match(r'up\s+to', str(val), _re.IGNORECASE) and len(numbers) == 1:
                return (0.0, numbers[0])

            return tuple(numbers)

        globals_dict = {
            "store": self.datastore,  # Persistent datastore - only shared state between steps
            "parse_number": parse_number,  # Parse string numbers/ranges: "8-12%" → (8.0, 12.0)
            "llm_ask": self._create_llm_ask_helper(),  # LLM query helper for general knowledge
            "send_email": create_send_email(
                self.config.email,
                is_sensitive=self._is_current_plan_sensitive,
            ),  # Email function - blocked if plan involves sensitive data
            "viz": create_viz_helper(
                datastore=self.datastore,
                print_file_refs=self.config.execution.print_file_refs,
                session_id=self.session_id,
                user_id=self.user_id,
                registry=self.registry,
                open_with_system_viewer=self.config.execution.open_with_system_viewer,
            ),  # Visualization/file output helper
            "publish": self._create_publish_helper(),  # Mark artifact as published for artifacts panel
            "facts": self._get_facts_dict(),  # Resolved facts as dict (loaded from _facts table)
            "llm_map": constat.llm.llm_map,
            "llm_classify": constat.llm.llm_classify,
            "llm_extract": constat.llm.llm_extract,
            "llm_summarize": constat.llm.llm_summarize,
            "llm_score": constat.llm.llm_score,
            "doc_read": self._create_doc_read_helper(),
        }

        # Provide database connections from config
        # SQL databases get both raw engine (db_<name>) and transpiling helper (sql_<name>)
        from constat.catalog.sql_transpiler import TranspilingConnection, create_sql_helper

        config_db_names = set()
        first_db = None
        for i, (db_name, db_config) in enumerate(self.config.databases.items()):
            config_db_names.add(db_name)
            conn = self.schema_manager.get_connection(db_name)

            # For SQL databases wrapped in TranspilingConnection:
            # - db_<name> = raw engine (for pd.read_sql compatibility)
            # - sql_<name> = helper function with auto-transpilation
            if isinstance(conn, TranspilingConnection):
                globals_dict[f"db_{db_name}"] = conn.engine
                globals_dict[f"sql_{db_name}"] = create_sql_helper(conn)
                if i == 0:
                    globals_dict["db"] = conn.engine
                    globals_dict["sql"] = create_sql_helper(conn)
                    first_db = conn.engine
            else:
                globals_dict[f"db_{db_name}"] = conn
                if i == 0:
                    globals_dict["db"] = conn
                    first_db = conn

        # Also include dynamically added databases (from projects) not in config
        for db_name in self.schema_manager.connections.keys():
            if db_name not in config_db_names:
                conn = self.schema_manager.connections[db_name]
                if isinstance(conn, TranspilingConnection):
                    globals_dict[f"db_{db_name}"] = conn.engine
                    globals_dict[f"sql_{db_name}"] = create_sql_helper(conn)
                    if first_db is None:
                        globals_dict["db"] = conn.engine
                        globals_dict["sql"] = create_sql_helper(conn)
                        first_db = conn.engine
                else:
                    globals_dict[f"db_{db_name}"] = conn
                    if first_db is None:
                        globals_dict["db"] = conn
                        first_db = conn
        for db_name in self.schema_manager.nosql_connections.keys():
            if db_name not in config_db_names:
                globals_dict[f"db_{db_name}"] = self.schema_manager.nosql_connections[db_name]
        for db_name in self.schema_manager.file_connections.keys():
            if db_name not in config_db_names:
                conn = self.schema_manager.file_connections[db_name]
                if hasattr(conn, 'path'):
                    globals_dict[f"file_{db_name}"] = conn.path

        # Provide API clients for GraphQL/REST APIs (config + project APIs)
        all_apis = self.get_all_apis()
        if all_apis:
            from constat.catalog.api_executor import APIExecutor
            # Create executor with merged config (config APIs + project APIs)
            api_executor = APIExecutor(self.config, project_apis=self._project_apis)
            for api_name, api_config in all_apis.items():
                if api_config.type == "graphql":
                    # Create a GraphQL query function
                    globals_dict[f"api_{api_name}"] = lambda query, variables=None, _name=api_name, _exec=api_executor: \
                        _exec.execute_graphql(_name, query, variables)
                else:
                    # Create a REST call function
                    globals_dict[f"api_{api_name}"] = lambda operation, params=None, _name=api_name, _exec=api_executor: \
                        _exec.execute_rest(_name, operation, params or {})

        return globals_dict

    def _get_facts_dict(self) -> dict:
        """Get resolved facts as a simple dict for use in generated code.

        Returns a dict mapping fact names to their values. This dict is injected
        into the execution globals so steps can reference facts['user_email'] etc.
        """
        facts_dict = {}
        try:
            all_facts = self.fact_resolver.get_all_facts()
            for name, fact in all_facts.items():
                if fact and fact.value is not None:
                    facts_dict[name] = fact.value
        except Exception as e:
            logger.debug(f"Error getting facts dict: {e}")
        return facts_dict

    def _materialize_facts_table(self) -> None:
        """Materialize resolved facts as a _facts table in the datastore.

        Creates a table with columns: name, value, source, description
        This table is used by downloaded scripts and for auditing.
        """
        import pandas as pd

        try:
            all_facts = self.fact_resolver.get_all_facts()
            if not all_facts:
                return

            rows = []
            for name, fact in all_facts.items():
                if fact and fact.value is not None:
                    rows.append({
                        "name": name,
                        "value": str(fact.value),
                        "source": fact.source.value if hasattr(fact.source, "value") else str(fact.source),
                        "description": fact.reasoning or "",
                    })

            if rows:
                facts_df = pd.DataFrame(rows)
                self.datastore.save_dataframe(
                    name="_facts",
                    df=facts_df,
                    step_number=0,  # Step 0 = pre-execution
                    description="Resolved facts for this analysis",
                )
                logger.debug(f"[FACTS] Materialized _facts table with {len(rows)} facts")

        except Exception as e:
            logger.warning(f"Failed to materialize _facts table: {e}")

    def _auto_save_results(self, namespace: dict, step_number: int) -> None:
        """
        Auto-save any DataFrames or lists found in the execution namespace.

        This ensures intermediate results are persisted even if the LLM
        forgot to explicitly save them.
        """
        import pandas as pd

        # Skip internal/injected variables
        skip_vars = {"store", "db", "pd", "np", "llm_ask", "send_email", "facts", "viz", "publish", "__builtins__"}
        skip_prefixes = ("db_", "_")

        # Already-saved tables (don't duplicate by name OR by data content)
        existing_tables = self.datastore.list_tables()
        existing_names = {t["name"] for t in existing_tables}

        # Get row counts of existing tables to detect duplicates by size
        # (cheap heuristic - if same rows, likely same data)
        existing_row_counts: dict[int, str] = {}
        for t in existing_tables:
            row_count = t.get("row_count")
            if row_count is not None:
                existing_row_counts[row_count] = t["name"]

        for var_name, value in namespace.items():
            # Skip internal variables
            if var_name in skip_vars or var_name.startswith(skip_prefixes):
                continue

            # Auto-save DataFrames
            if isinstance(value, pd.DataFrame) and var_name not in existing_names:
                # Skip if this data was already saved under a different name
                # (cheap heuristic: same row count = likely duplicate)
                if len(value) in existing_row_counts:
                    logger.debug(f"Skip auto-save of {var_name}: likely duplicate of {existing_row_counts[len(value)]}")
                    continue

                self.datastore.save_dataframe(
                    name=var_name,
                    df=value,
                    step_number=step_number,
                    description=f"Auto-saved from step {step_number}",
                    role_id=self._current_role_id,
                )

            # Auto-save lists (as state, since they might be useful)
            elif isinstance(value, (list, dict)) and len(value) > 0:
                # Check if already saved in state
                existing = self.datastore.get_state(var_name)
                if existing is None:
                    try:
                        # Convert pandas NA values to None for JSON serialization
                        clean_value = self._make_json_serializable(value)
                        self.datastore.set_state(var_name, clean_value, step_number)
                    except Exception as e:
                        logger.debug(f"Skip auto-save of {var_name}: not JSON-serializable: {e}")

    def _execute_step(self, step: Step) -> StepResult:
        """
        Execute a single step with retry on errors.

        Returns:
            StepResult with success/failure info
        """
        start_time = time.time()
        last_code = ""
        last_error = None
        pending_learning_context = None  # Track error for potential learning capture

        # Set current role context for this step (facts created will inherit this role_id)
        previous_role_id = self._current_role_id
        self._current_role_id = step.role_id

        self._emit_event(StepEvent(
            event_type="step_start",
            step_number=step.number,
            data={"goal": step.goal}
        ))

        max_attempts = self.session_config.max_retries_per_step
        for attempt in range(1, max_attempts + 1):
            # Emit detailed generating event
            self._emit_event(StepEvent(
                event_type="generating",
                step_number=step.number,
                data={
                    "attempt": attempt,
                    "max_attempts": max_attempts,
                    "is_retry": attempt > 1,
                    "goal": step.goal[:50] + "..." if len(step.goal) > 50 else step.goal,
                    "retry_reason": last_error[:100] if attempt > 1 and last_error else None,
                }
            ))

            # Use router with step's task_type for automatic model selection/escalation
            if attempt == 1:
                prompt = self._build_step_prompt(step)
                result = self.router.execute_code(
                    task_type=step.task_type,
                    system=STEP_SYSTEM_PROMPT,
                    user_message=prompt,
                    tools=self._get_schema_tools(),
                    tool_handlers=self._get_tool_handlers(),
                    complexity=step.complexity,
                )
            else:
                # Track error context for potential learning capture
                pending_learning_context = {
                    "error_message": last_error[:500] if last_error else "",
                    "original_code": last_code[:500] if last_code else "",
                    "step_goal": step.goal,
                    "attempt": attempt,
                }

                retry_prompt = RETRY_PROMPT_TEMPLATE.format(
                    error_details=last_error,
                    previous_code=last_code,
                )
                result = self.router.execute_code(
                    task_type=step.task_type,
                    system=STEP_SYSTEM_PROMPT,
                    user_message=retry_prompt,
                    tools=self._get_schema_tools(),
                    tool_handlers=self._get_tool_handlers(),
                    complexity=step.complexity,
                )

            if not result.success:
                # Router exhausted all models
                raise RuntimeError(f"Code generation failed: {result.content}")

            code = result.content

            step.code = code

            self._emit_event(StepEvent(
                event_type="executing",
                step_number=step.number,
                data={
                    "attempt": attempt,
                    "max_attempts": max_attempts,
                    "is_retry": attempt > 1,
                    "code_lines": len(code.split('\n')),
                }
            ))

            # Track tables before execution (name + version to detect updates)
            tables_before_list = self.datastore.list_tables() if self.datastore else []
            tables_before = set(t['name'] for t in tables_before_list)
            versions_before = {t['name']: t.get('version', 1) for t in tables_before_list}

            # Execute
            exec_globals = self._get_execution_globals()
            result = self.executor.execute(code, exec_globals)
            logger.debug(f"[Step {step.number}] Execution result (attempt {attempt}): success={result.success}, error={result.error_message()[:200] if not result.success else 'none'}")

            # Auto-save any DataFrames or lists created during execution
            if result.success and self.datastore:
                self._auto_save_results(result.namespace, step.number)

            # Record artifacts in datastore
            if self.datastore:
                self.datastore.add_artifact(step.number, attempt, "code", code, role_id=self._current_role_id)
                if result.stdout:
                    self.datastore.add_artifact(step.number, attempt, "output", result.stdout, role_id=self._current_role_id)

            if result.success:
                duration_ms = int((time.time() - start_time) * 1000)

                # Capture learning if this was a successful retry
                if attempt > 1 and pending_learning_context:
                    self._capture_error_learning(
                        context=pending_learning_context,
                        fixed_code=code,
                    )

                # Run post-validations
                validation_warnings: list[str] = []
                if step.post_validations:
                    logger.debug(f"[Step {step.number}] Running {len(step.post_validations)} post-validations (attempt {attempt})")
                    validation_warnings, failed_validation = self._run_post_validations(step, result.namespace)
                    logger.debug(f"[Step {step.number}] Post-validation result: failed={failed_validation is not None}, warnings={len(validation_warnings)}")

                    if failed_validation:
                        if failed_validation.on_fail == ValidationOnFail.CLARIFY:
                            clarify_response = self._ask_validation_clarification(
                                step, failed_validation
                            )
                            if clarify_response:
                                last_code = code
                                stdout_hint = f"\nCode stdout:\n{result.stdout[-2000:]}" if result.stdout else ""
                                last_error = f"Validation failed: {failed_validation.description}. User guidance: {clarify_response}{stdout_hint}"
                                continue
                            # User skipped — treat as warning
                            validation_warnings.append(f"Skipped: {failed_validation.description}")

                        elif failed_validation.on_fail == ValidationOnFail.RETRY:
                            last_code = code
                            stdout_context = ""
                            if result.stdout:
                                stdout_context = f"\nCode stdout (shows actual state):\n{result.stdout[-2000:]}\n"
                            last_error = (
                                f"Code executed without errors, but post-validation failed.\n"
                                f"Validation: {failed_validation.description}\n"
                                f"Expression: {failed_validation.expression}\n"
                                f"{stdout_context}"
                                f"The code must be fixed so this validation passes."
                            )
                            self._emit_event(StepEvent(
                                event_type="validation_retry",
                                step_number=step.number,
                                data={"validation": failed_validation.description}
                            ))
                            continue

                    if validation_warnings:
                        self._emit_event(StepEvent(
                            event_type="validation_warnings",
                            step_number=step.number,
                            data={"warnings": validation_warnings}
                        ))

                # Detect new AND updated tables
                tables_after_list = self.datastore.list_tables() if self.datastore else []
                tables_after = set(t['name'] for t in tables_after_list)
                versions_after = {t['name']: t.get('version', 1) for t in tables_after_list}
                new_tables = tables_after - tables_before
                updated_tables = {
                    name for name in tables_before & tables_after
                    if versions_after.get(name, 1) > versions_before.get(name, 1)
                    and not name.startswith('_')  # Skip internal tables
                }
                tables_created = list(new_tables | updated_tables)

                self._emit_event(StepEvent(
                    event_type="step_complete",
                    step_number=step.number,
                    data={
                        "goal": step.goal,
                        "code": code,
                        "stdout": result.stdout,
                        "attempts": attempt,
                        "duration_ms": duration_ms,
                        "tables_created": tables_created,
                    }
                ))

                # Persist step code to disk
                if self.session_id:
                    self.history.save_step_code(
                        session_id=self.session_id,
                        step_number=step.number,
                        goal=step.goal,
                        code=code,
                        output=result.stdout,
                    )

                # Restore previous role context
                self._current_role_id = previous_role_id

                return StepResult(
                    success=True,
                    stdout=result.stdout,
                    attempts=attempt,
                    duration_ms=duration_ms,
                    tables_created=tables_created,
                    code=code,
                    validation_warnings=validation_warnings,
                )

            # Prepare for retry
            last_code = code
            last_error = format_error_for_retry(result, code)

            # Log error for debugging
            logger.warning(f"[Step {step.number}] Execution error (attempt {attempt}/{max_attempts}): {result.error_message() or 'Unknown error'}")
            logger.debug(f"[Step {step.number}] Failed code:\n{code}")

            # Record error artifact
            if self.datastore:
                self.datastore.add_artifact(step.number, attempt, "error", last_error, role_id=self._current_role_id)

            # Determine error type for better status messages
            error_lower = last_error.lower() if last_error else ""
            if "sql" in error_lower or "query" in error_lower:
                error_type = "SQL error"
            elif "name" in error_lower and "not defined" in error_lower:
                error_type = "Variable not found"
            elif "type" in error_lower:
                error_type = "Type error"
            elif "key" in error_lower:
                error_type = "Key error"
            elif "index" in error_lower:
                error_type = "Index error"
            elif "timeout" in error_lower:
                error_type = "Timeout"
            else:
                error_type = "Runtime error"

            will_retry = attempt < max_attempts
            if will_retry:
                logger.debug(f"[Step {step.number}] Will retry ({max_attempts - attempt} attempts remaining)")
            self._emit_event(StepEvent(
                event_type="step_error",
                step_number=step.number,
                data={
                    "error": last_error,
                    "error_type": error_type,
                    "attempt": attempt,
                    "max_attempts": max_attempts,
                    "will_retry": will_retry,
                    "next_attempt": attempt + 1 if will_retry else None,
                }
            ))

        # Max retries exceeded - generate suggestions for alternative approaches
        logger.warning(f"[Step {step.number}] Failed after {max_attempts} attempts: {last_error[:200]}")
        duration_ms = int((time.time() - start_time) * 1000)
        suggestions = self._generate_failure_suggestions(step, last_error, last_code)

        # Emit step_failed event with suggestions
        self._emit_event(StepEvent(
            event_type="step_failed",
            step_number=step.number,
            data={
                "error": last_error,
                "attempts": self.session_config.max_retries_per_step,
                "suggestions": suggestions,
            }
        ))

        # Persist failed step code to disk (with error)
        if self.session_id and last_code:
            self.history.save_step_code(
                session_id=self.session_id,
                step_number=step.number,
                goal=step.goal,
                code=last_code,
                error=last_error,
            )

        # Restore previous role context
        self._current_role_id = previous_role_id

        return StepResult(
            success=False,
            stdout="",
            error=f"Failed after {self.session_config.max_retries_per_step} attempts. Last error: {last_error}",
            attempts=self.session_config.max_retries_per_step,
            duration_ms=duration_ms,
            suggestions=suggestions,
        )

    def _run_post_validations(
        self, step: Step, namespace: dict
    ) -> tuple[list[str], PostValidation | None]:
        """Run post-validations against step's execution namespace.

        Returns:
            (warnings, first_failing_validation)
            - warnings: list of warning messages from on_fail=WARN validations
            - first_failing_validation: first RETRY or CLARIFY validation that failed, or None
        """
        warnings: list[str] = []
        # Inject store tables so validations can reference them by name
        try:
            for t in self.datastore.list_tables():
                name = t["name"]
                if name not in namespace:
                    df = self.datastore.load_dataframe(name)
                    if df is not None:
                        namespace[name] = df
                        logger.debug(f"[Step {step.number}] Injected store table '{name}' into validation namespace")
        except Exception as e:
            logger.warning(f"[Step {step.number}] Failed to inject store tables for validation: {e}")
        for v in step.post_validations:
            try:
                result = eval(v.expression, {"__builtins__": __builtins__}, namespace)  # noqa: S307
                passed = bool(result)
            except Exception as e:
                logger.warning(f"[Step {step.number}] Post-validation expression error: {v.expression} -> {e}")
                passed = False

            if not passed:
                if v.on_fail == ValidationOnFail.WARN:
                    warnings.append(f"Validation warning: {v.description}")
                else:
                    # RETRY or CLARIFY — return immediately
                    return warnings, v
        return warnings, None

    def _ask_validation_clarification(self, step: Step, validation: PostValidation) -> str | None:
        """Ask user for clarification when a post-validation fails with on_fail=CLARIFY.

        Returns:
            User's response string, or None if skipped/unavailable.
        """
        if not self._clarification_callback:
            return None

        question = validation.clarify_question or f"Validation failed: {validation.description}. How should we proceed?"
        request = ClarificationRequest(
            original_question=step.goal,
            ambiguity_reason=f"Post-validation failed: {validation.description}",
            questions=[ClarificationQuestion(question=question)],
        )

        self._emit_event(StepEvent(
            event_type="clarification_needed",
            step_number=step.number,
            data={
                "reason": request.ambiguity_reason,
                "questions": request.questions,
            }
        ))

        response = self._clarification_callback(request)
        if response.skip:
            return None

        # Return first non-empty answer
        for answer in response.answers.values():
            if answer:
                return answer
        return None

    def _capture_error_learning(self, context: dict, fixed_code: str) -> None:
        """Capture a learning from a successful error fix.

        Args:
            context: Error context dict with error_message, original_code, step_goal
            fixed_code: The code that successfully fixed the error
        """
        try:
            # Determine category based on error type and context
            category = self._categorize_error(context)

            # Use LLM to generate a concise learning summary
            summary = self._summarize_error_fix(context, fixed_code)
            if not summary:
                # Fallback to a simple summary
                error_preview = context.get("error_message", "")[:100]
                summary = f"Fixed error: {error_preview}"

            # Add fixed code to context
            context["fixed_code"] = fixed_code[:500]

            # Save the learning
            self.learning_store.save_learning(
                category=category,
                context=context,
                correction=summary,
                source=LearningSource.AUTO_CAPTURE,
            )

            # Auto-compact if too many raw learnings
            stats = self.learning_store.get_stats()
            if stats["unpromoted"] >= 50:
                logger.info(f"[LEARNINGS] {stats['unpromoted']} unpromoted learnings — triggering compaction")
                try:
                    self._compact_learnings()
                except Exception as ce:
                    logger.warning(f"[LEARNINGS] Compaction failed (non-fatal): {ce}")
        except Exception as e:
            logger.debug(f"Learning capture failed (non-fatal): {e}")

    def _compact_learnings(self) -> None:
        """Compact raw learnings into rules using LLM to find patterns."""
        unpromoted = [
            l for l in self.learning_store.list_raw_learnings(limit=200, include_promoted=False)
            if not l.get("promoted_to")
        ]
        if len(unpromoted) < 20:
            return

        # Format learnings for LLM
        learning_texts = []
        for l in unpromoted:
            correction = l.get("correction", "")
            category = l.get("category", "")
            ctx = l.get("context", {})
            error = ctx.get("error_message", "")[:100]
            learning_texts.append(f"[{l['id']}] ({category}) {correction} | error: {error}")

        prompt = f"""Group these {len(learning_texts)} code generation learnings into reusable rules.
Each rule should capture a PATTERN that applies across multiple learnings.

Learnings:
{chr(10).join(learning_texts)}

Return JSON array of rules:
[
  {{"summary": "Rule description", "category": "codegen_error", "confidence": 0.85, "source_ids": ["learn_xxx", "learn_yyy"], "tags": ["sql", "duckdb"]}}
]

Guidelines:
- Only create a rule if 3+ learnings share the same pattern
- confidence = fraction of learnings in the group that clearly match
- Learnings not matching any pattern should be omitted (they stay as raw)
- Keep summaries actionable: "Use X instead of Y when Z"

Return ONLY the JSON array."""

        result = self.router.execute(
            task_type=TaskType.SYNTHESIS,
            system="You analyze code error patterns and extract reusable rules.",
            user_message=prompt,
            max_tokens=4096,
        )

        # Parse rules
        import json
        text = result.content.strip()
        if text.startswith("```"):
            text = re.sub(r'^```\w*\n?', '', text)
            text = re.sub(r'\n?```$', '', text)
        try:
            rules = json.loads(text)
        except json.JSONDecodeError:
            logger.warning("[LEARNINGS] Failed to parse compaction response as JSON")
            return

        promoted_count = 0
        for rule_data in rules:
            source_ids = rule_data.get("source_ids", [])
            if len(source_ids) < 3:
                continue
            category_str = rule_data.get("category", "codegen_error")
            try:
                category = LearningCategory(category_str)
            except ValueError:
                category = LearningCategory.CODEGEN_ERROR

            rule_id = self.learning_store.save_rule(
                summary=rule_data["summary"],
                category=category,
                confidence=rule_data.get("confidence", 0.8),
                source_learnings=source_ids,
                tags=rule_data.get("tags", []),
            )
            # Archive promoted learnings
            for lid in source_ids:
                self.learning_store.archive_learning(lid, rule_id)
                promoted_count += 1

        logger.info(f"[LEARNINGS] Compacted {promoted_count} learnings into {len(rules)} rules")

    def _categorize_error(self, context: dict) -> LearningCategory:
        """Categorize an error for learning storage.

        Categories:
        - HTTP_ERROR: 4xx/5xx errors from external API calls
        - EXTERNAL_API_ERROR: Other errors in API integration code
        - CODEGEN_ERROR: General code generation errors (default)
        """
        error_msg = context.get("error_message", "").lower()
        original_code = context.get("original_code", "").lower()
        step_goal = context.get("step_goal", "").lower()

        # Check for HTTP errors (4xx/5xx)
        http_error_patterns = [
            "status code 4", "status code 5",
            "status_code=4", "status_code=5",
            "http 4", "http 5",
            "400 ", "401 ", "403 ", "404 ", "405 ",
            "500 ", "502 ", "503 ", "504 ",
            "bad request", "unauthorized", "forbidden", "not found",
            "internal server error", "bad gateway", "service unavailable",
            "httperror", "httpstatuserror",
        ]
        for pattern in http_error_patterns:
            if pattern in error_msg:
                return LearningCategory.HTTP_ERROR

        # Check for external API integration errors (not HTTP status errors)
        api_indicators = [
            "requests.", "httpx.", "aiohttp.",
            "rest api", "graphql", "openapi",
            "api_client", "api_response",
            "response.json()", "response.text",
        ]
        is_api_code = any(ind in original_code for ind in api_indicators)

        # Also check step goal for API-related work
        api_goal_keywords = ["fetch", "call api", "api request", "rest", "graphql", "webhook"]
        is_api_goal = any(kw in step_goal for kw in api_goal_keywords)

        if is_api_code or is_api_goal:
            return LearningCategory.EXTERNAL_API_ERROR

        # Default to general code generation error
        return LearningCategory.CODEGEN_ERROR

    def _save_correction_as_learning(self, user_input: str) -> None:
        """Save a user correction as a reusable learning.

        Captures context from the current session to make the correction
        applicable to future similar queries.

        Args:
            user_input: The user's correction (e.g., "always use customer_id not cust_id")
        """
        if not self.learning_store:
            logger.debug("No learning store available, skipping correction capture")
            return

        try:
            # Build context from current session
            context = {}

            # Add current problem/query context
            if self.datastore:
                problem = self.datastore.get_session_meta("problem")
                if problem:
                    context["original_problem"] = problem[:500]

                # Add available tables for schema context
                tables = self.datastore.list_tables()
                if tables:
                    context["tables"] = [t.get("name", "") for t in tables[:10]]

                # Add any active schemas from catalog
                if hasattr(self, 'catalog') and self.catalog:
                    schemas = getattr(self.catalog, 'get_active_schemas', lambda: [])()
                    if schemas:
                        context["schemas"] = schemas[:5]

            # Add session ID for traceability
            if self.session_id:
                context["session_id"] = self.session_id

            # Save the correction as a learning
            learning_id = self.learning_store.save_learning(
                category=LearningCategory.USER_CORRECTION,
                context=context,
                correction=user_input,
                source=LearningSource.NL_DETECTION,
            )

            logger.info(f"Saved user correction as learning {learning_id}: {user_input[:50]}...")

            # Emit event so UI can acknowledge
            self._emit_event(StepEvent(
                event_type="correction_saved",
                step_number=0,
                data={"correction": user_input, "learning_id": learning_id}
            ))

        except Exception as e:
            logger.debug(f"Correction capture failed (non-fatal): {e}")

    def _generate_failure_suggestions(
        self, step: "Step", error: str, code: str
    ) -> list["FailureSuggestion"]:
        """Generate suggestions for alternative approaches when a step fails.

        Differentiates between:
        - Codegen failures: LLM couldn't produce working code (need different approach)
        - Runtime errors: Code ran but failed on data/environment (user can redirect)

        Args:
            step: The step that failed
            error: The last error message
            code: The last code that was attempted

        Returns:
            List of FailureSuggestion objects
        """
        from constat.core.models import FailureSuggestion

        suggestions = []
        error_lower = error.lower() if error else ""
        goal_lower = step.goal.lower() if step.goal else ""

        # Detect if this is a codegen failure vs runtime error
        codegen_indicators = [
            "syntax error", "invalid syntax", "unexpected token",
            "code generation failed", "could not generate",
            "parsing error", "indentation error"
        ]
        is_codegen_failure = any(ind in error_lower for ind in codegen_indicators)

        if is_codegen_failure:
            # Codegen failures - LLM couldn't produce working code
            suggestions.extend([
                FailureSuggestion(
                    id="break_down",
                    label="Break into smaller steps",
                    description="Split this step into simpler sub-steps that may be easier to generate",
                    action="break_down"
                ),
                FailureSuggestion(
                    id="simplify_goal",
                    label="Simplify the goal",
                    description="Rephrase the step goal in simpler terms",
                    action="rephrase"
                ),
                FailureSuggestion(
                    id="provide_code",
                    label="Provide code snippet",
                    description="Give a working code example or pattern to follow",
                    action="provide_code"
                ),
                FailureSuggestion(
                    id="report_issue",
                    label="Report this issue",
                    description="This appears to be a code generation bug - report it for investigation",
                    action="report"
                ),
            ])
        else:
            # Runtime errors - likely data/source issues user can redirect

            # Document/search related failures
            if any(term in error_lower or term in goal_lower for term in [
                "document", "search", "not found", "no results", "relevance", "policy", "guideline"
            ]):
                suggestions.extend([
                    FailureSuggestion(
                        id="rephrase_search",
                        label="Rephrase search query",
                        description="Try searching with different or broader terms",
                        action="rephrase"
                    ),
                    FailureSuggestion(
                        id="list_documents",
                        label="List available documents",
                        description="Show all documents so you can specify which one to use",
                        action="list_docs"
                    ),
                    FailureSuggestion(
                        id="load_full_doc",
                        label="Load full document",
                        description="Load an entire document instead of searching chunks",
                        action="load_doc"
                    ),
                ])

            # Database/query related failures
            if any(term in error_lower for term in [
                "table", "column", "sql", "query", "database", "no such", "does not exist"
            ]):
                suggestions.extend([
                    FailureSuggestion(
                        id="list_tables",
                        label="List available tables",
                        description="Show database schema to find correct table/column names",
                        action="list_tables"
                    ),
                    FailureSuggestion(
                        id="different_table",
                        label="Use different data source",
                        description="Specify which table or data source to use instead",
                        action="redirect"
                    ),
                ])

            # Data not found / empty results
            if any(term in error_lower for term in [
                "empty", "no data", "zero rows", "none", "missing"
            ]):
                suggestions.extend([
                    FailureSuggestion(
                        id="broaden_query",
                        label="Broaden the query",
                        description="Remove filters or expand date range to find data",
                        action="broaden"
                    ),
                    FailureSuggestion(
                        id="check_filters",
                        label="Check filter criteria",
                        description="The filters may be too restrictive",
                        action="check_filters"
                    ),
                ])

            # API/connection errors
            if any(term in error_lower for term in [
                "api", "connection", "timeout", "rate limit", "unauthorized", "403", "401", "500"
            ]):
                suggestions.extend([
                    FailureSuggestion(
                        id="use_cached",
                        label="Use cached/local data",
                        description="Check if we have data from a previous call or local source",
                        action="use_cache"
                    ),
                ])

        # Import/syntax errors (could be either codegen or config issue)
        if any(term in error_lower for term in [
            "import", "module", "not allowed"
        ]):
            suggestions.append(
                FailureSuggestion(
                    id="simplify_code",
                    label="Simplify approach",
                    description="Use a simpler method that doesn't require this import",
                    action="simplify"
                )
            )

        # Always offer these general options
        suggestions.extend([
            FailureSuggestion(
                id="modify_step",
                label="Modify this step",
                description="Change the goal or approach for this step",
                action="modify"
            ),
            FailureSuggestion(
                id="skip_step",
                label="Skip this step",
                description="Continue without this step (may affect later steps)",
                action="skip"
            ),
            FailureSuggestion(
                id="manual_input",
                label="Provide value manually",
                description="Enter the needed information directly",
                action="manual"
            ),
        ])

        return suggestions

    def _summarize_error_fix(self, context: dict, fixed_code: str) -> str:
        """Use LLM to generate a concise learning summary from an error fix.

        Args:
            context: Error context with error_message, original_code
            fixed_code: The code that fixed the error

        Returns:
            A concise summary of what was learned, or empty string on failure
        """
        try:
            prompt = f"""Summarize what was learned from this error fix in ONE sentence.

Error: {context.get('error_message', '')[:300]}
Original code snippet: {context.get('original_code', '')[:200]}
Fixed code snippet: {fixed_code[:200]}

Output ONLY a single sentence describing the lesson learned, e.g., "Always use X instead of Y when..."
Do not include any explanation or extra text."""

            response = self.llm.generate(
                system="You are a technical writer summarizing coding lessons learned.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )
            # generate() returns string directly
            return response.strip()
        except Exception as e:
            logger.debug(f"Failed to summarize learning (non-fatal): {e}")
            return ""

    def _request_approval(
        self,
        problem: str,
        planner_response: PlannerResponse,
    ) -> PlanApprovalResponse:
        """
        Request approval for a plan.

        If auto_approve is set or no callback is registered, auto-approves.
        Otherwise calls the registered callback.

        Args:
            problem: The original problem
            planner_response: The planner's response with plan and reasoning

        Returns:
            PlanApprovalResponse with user's decision
        """
        # Auto-approve if configured
        if self.session_config.auto_approve:
            return PlanApprovalResponse.approve()

        # No callback registered - auto-approve
        if not self._approval_callback:
            return PlanApprovalResponse.approve()

        # Build approval request
        steps = [
            {
                "number": step.number,
                "goal": step.goal,
                "inputs": step.expected_inputs,
                "outputs": step.expected_outputs,
                "role_id": step.role_id,
            }
            for step in planner_response.plan.steps
        ]

        request = PlanApprovalRequest(
            problem=problem,
            steps=steps,
            reasoning=planner_response.reasoning,
        )

        return self._approval_callback(request)

    def _ensure_enhance_updates_source(
        self,
        question: str,
        plan: Plan,
        existing_tables: list[dict],
    ) -> Plan:
        """Append an update step when an 'enhance' plan only creates a mapping table.

        The LLM planner consistently decomposes "enhance X with Y" into
        analyze → fetch reference → create mapping, but omits the final
        "apply mapping back to X" step. This method detects that pattern
        and appends the missing step.
        """
        import re

        # Only applies when there are steps and existing tables
        if not plan.steps or not existing_tables:
            return plan

        # Detect enhance intent
        enhance_re = re.compile(
            r'\b(?:enhance|enrich|extend|augment)\b|'
            r'\badd\s+(?:a\s+)?(?:column|field|the)\b',
            re.IGNORECASE,
        )
        if not enhance_re.search(question):
            return plan

        # Find candidate target table by matching table name fragments to question
        # Prefer higher step_number (most recent working dataset) to break ties
        question_lower = question.lower()
        candidates = [t for t in existing_tables if not t['name'].startswith('_')]

        target_table = None
        best_score = (0, -1)  # (word_overlap, step_number)
        for t in candidates:
            name = t['name']
            name_parts = [p for p in name.lower().replace('_', ' ').split() if len(p) > 3]
            overlap = sum(1 for p in name_parts if p in question_lower)
            step_num = t.get('step_number', 0) or 0
            score = (overlap, step_num)
            if score > best_score and overlap > 0:
                best_score = score
                target_table = name

        if not target_table:
            return plan

        # Check whether the last step already updates the target table
        last_step = plan.steps[-1]
        last_goal_lower = last_step.goal.lower()
        target_words = [p for p in target_table.lower().replace('_', ' ').split() if len(p) > 3]

        mentions_target = any(w in last_goal_lower for w in target_words)
        mentions_update = any(
            kw in last_goal_lower
            for kw in ['update', 'add column', 'enhance', 'enrich', 'modify', 'apply', 'save back']
        )

        if mentions_target and mentions_update:
            return plan  # Already correct

        # Check if ANY step already updates the target
        for step in plan.steps:
            goal_lower = step.goal.lower()
            if any(w in goal_lower for w in target_words) and any(
                kw in goal_lower
                for kw in ['update', 'add column', 'enhance', 'enrich', 'modify', 'apply', 'save back']
            ):
                return plan  # Some step already handles it

        # Append the missing update step
        logger.info(
            f"[PLAN_VALIDATION] Enhance plan missing update step for '{target_table}'. Appending step."
        )
        last = plan.steps[-1]
        update_step = Step(
            number=last.number + 1,
            goal=(
                f"Apply the mapping from previous steps to update `{target_table}` "
                f"by adding the new column(s) and saving it back with the same name"
            ),
            expected_inputs=[
                last.expected_outputs[0] if last.expected_outputs else "mapping",
                target_table,
            ],
            expected_outputs=[target_table],
            depends_on=[last.number],
            task_type=TaskType.PYTHON_ANALYSIS,
            complexity="low",
            role_id=last.role_id,
        )
        plan.steps.append(update_step)
        return plan

    def _build_plan_from_edited_steps(
        self,
        problem: str,
        edited_steps: list[dict],
        start_number: int = 1,
    ) -> Plan:
        """
        Build a Plan directly from user-edited steps, skipping the planner.

        Used when user edits the plan steps but provides no additional feedback,
        meaning they want exactly what they specified without re-interpretation.

        Args:
            problem: The original problem
            edited_steps: List of {"number": int, "goal": str} from user edits
            start_number: First step number (e.g. 4 for follow-up after 3 initial steps)

        Returns:
            Plan object ready for execution
        """
        from datetime import datetime

        # Get original plan steps for metadata preservation
        original_steps_by_goal = {}
        if self.plan:
            for step in self.plan.steps:
                # Key by normalized goal for fuzzy matching
                normalized = step.goal.lower().strip()
                original_steps_by_goal[normalized] = step

        # Build new steps with sequential numbering from start_number
        new_steps = []
        for i, edited in enumerate(edited_steps):
            step_num = start_number + i
            goal = edited.get("goal", "")
            original_number = edited.get("number", step_num)

            # Try to find original step for metadata (by goal similarity)
            normalized_goal = goal.lower().strip()
            original = original_steps_by_goal.get(normalized_goal)

            # Also try to find by original step number
            if not original and self.plan:
                original = self.plan.get_step(original_number)

            # Create step with preserved metadata where available
            step = Step(
                number=step_num,  # Renumber sequentially from start_number
                goal=goal,
                expected_inputs=original.expected_inputs if original else [],
                expected_outputs=original.expected_outputs if original else [],
                depends_on=[step_num - 1] if step_num > start_number else [],  # Sequential dependencies for safety
                step_type=original.step_type if original else StepType.PYTHON,
                task_type=original.task_type if original else TaskType.PYTHON_ANALYSIS,
                complexity=original.complexity if original else "medium",
                role_id=original.role_id if original else None,
                skill_ids=original.skill_ids if original else None,
            )
            new_steps.append(step)

        return Plan(
            problem=problem,
            steps=new_steps,
            created_at=datetime.now().isoformat(),
        )

    def _classify_question(self, problem: str) -> str:
        """
        Classify whether a question requires code execution or is a meta-question.

        Returns:
            QuestionType.DATA_ANALYSIS - needs code execution (queries, computation, actions)
            QuestionType.META_QUESTION - about system capabilities (what can you do?)

        Note: We route almost everything through code execution because:
        - Data questions need database queries
        - General knowledge questions can use llm_ask() + computation
        - Action requests (email, export) need code
        - Even "What is sqrt(8)?" benefits from actual computation
        """
        # Only meta-questions about the system bypass code execution
        if is_meta_question(problem):
            return QuestionType.META_QUESTION

        # Everything else goes through code execution
        # The generated code can use llm_ask() for general knowledge
        # and then compute/transform/act on the results
        return QuestionType.DATA_ANALYSIS

    def _try_show_existing_data(self, question: str) -> Optional[dict]:
        """Fast path for 'show me X' requests that display existing tables.

        Detects simple data display requests and handles them directly without
        going through LLM intent classification. This is much faster for
        simple lookups like 'show me the raise_recommendations table'.

        Args:
            question: The user's question

        Returns:
            Result dict if this was a show request, None otherwise
        """
        import re

        if not self.datastore:
            return None

        # Get available table names
        tables = self.datastore.list_tables()
        if not tables:
            return None

        table_names = {t['name'].lower(): t['name'] for t in tables}

        # Patterns for "show me X" type requests
        # Match: "show me X", "display X", "what's in X", "view X", "print X"
        show_patterns = [
            r"^(?:show\s+(?:me\s+)?(?:the\s+)?|display\s+(?:the\s+)?|view\s+(?:the\s+)?|print\s+(?:the\s+)?|what(?:'s| is)\s+in\s+(?:the\s+)?)(.+?)(?:\s+table)?(?:\s+data)?$",
            r"^(.+?)(?:\s+table|\s+data)$",  # "raise_recommendations table"
        ]

        question_lower = question.lower().strip()

        for pattern in show_patterns:
            match = re.match(pattern, question_lower, re.IGNORECASE)
            if match:
                candidate = match.group(1).strip()
                # Check if candidate matches a table name
                if candidate in table_names:
                    actual_table = table_names[candidate]
                    return self._quick_table_display(actual_table)

        # Also check if the entire question is just a table name
        if question_lower in table_names:
            return self._quick_table_display(table_names[question_lower])

        return None

    def _quick_table_display(self, table_name: str) -> dict:
        """Display a table quickly without going through planning.

        Args:
            table_name: The name of the table to display

        Returns:
            Result dict with table data
        """
        try:
            # Query the table
            df = self.datastore.query(f"SELECT * FROM {table_name} LIMIT 50")

            # Format output
            try:
                table_str = df.to_markdown(index=False)
            except Exception:
                table_str = df.to_string(index=False)

            row_count = len(df)
            total_rows = self.datastore.query(f"SELECT COUNT(*) as cnt FROM {table_name}").iloc[0]['cnt']

            output = f"**{table_name}** ({total_rows} rows)\n\n{table_str}"
            if row_count < total_rows:
                output += f"\n\n_Showing first {row_count} of {total_rows} rows_"

            self._emit_event(StepEvent(
                event_type="quick_display",
                step_number=0,
                data={"table": table_name, "rows": row_count}
            ))

            return {
                "success": True,
                "mode": "quick_display",
                "output": output,
                "datastore_tables": self.datastore.list_tables(),
                "suggestions": [
                    f"Run SQL query on {table_name}",
                    "Ask a question about this data",
                ],
            }
        except Exception as e:
            # If table query fails, return None to fall through to normal processing
            return None

    def _analyze_question(self, problem: str, previous_problem: str = None) -> QuestionAnalysis:
        """
        Analyze a question in a single LLM call: extract facts, classify type, check cached facts.

        This combines what were previously separate operations into one call for efficiency:
        1. Extract embedded facts (e.g., "my role as CFO" -> user_role: CFO)
        2. Classify question type (meta-question vs data analysis)
        3. Check if question can be answered from cached facts

        Args:
            problem: The question to analyze
            previous_problem: If this is a follow-up, the original problem for context

        Returns:
            QuestionAnalysis with question_type, extracted_facts, and optional cached_fact_answer
        """
        # First, use fast regex-based classification for obvious meta-questions
        # This avoids an LLM call for simple cases like "what can you do?"
        if is_meta_question(problem):
            # Return immediately for meta-questions - no LLM classification needed
            # Note: We skip fact extraction here for efficiency. Most meta-questions
            # like "how do you reason" don't contain extractable facts anyway.
            return QuestionAnalysis(
                question_type=QuestionType.META_QUESTION,
                extracted_facts=[],
                cached_fact_answer=None,
            )

        # Get cached facts for context
        cached_facts = self.fact_resolver.get_all_facts()
        fact_context = ""
        if cached_facts:
            fact_context = "Known facts:\n" + "\n".join(
                f"- {name}: {fact.display_value}" for name, fact in cached_facts.items()
            )

        # Build data source context from SessionResources (single source of truth)
        # This includes all databases/APIs/documents from config + active projects
        data_sources = []
        for name, desc in self.resources.get_database_descriptions():
            desc_line = desc.split('\n')[0] if desc else f"database '{name}'"
            data_sources.append(f"DATABASE '{name}': {desc_line}")
        for name, desc in self.resources.get_api_descriptions():
            desc_line = desc.split('\n')[0] if desc else f"API"
            data_sources.append(f"API '{name}': {desc_line}")
        for name, desc in self.resources.get_document_descriptions():
            desc_line = desc.split('\n')[0] if desc else "reference document"
            data_sources.append(f"DOCUMENT '{name}': {desc_line}")

        source_context = ""
        if data_sources:
            source_context = "\nAvailable data sources:\n" + "\n".join(f"- {s}" for s in data_sources)

        # Build follow-up context if this is a continuation of a previous analysis
        followup_context = ""
        if previous_problem:
            followup_context = f"""
This is a FOLLOW-UP to a previous analysis.
Previous question: "{previous_problem}"

CRITICAL INTENT RULES (apply in order):
1. RE-EXECUTION language ANYWHERE in message = REDO intent, NOT NEW_QUESTION:
   "redo", "again", "retry", "rerun", "try again", "this time", "instead", "once more", etc.

2. Requests to CHANGE how something is computed/calculated = STEER_PLAN + REDO:
   "change X to use Y", "use average instead", "compute it differently", "use the 2 most recent"

3. Simple value changes (literal numbers/strings) = MODIFY_FACT:
   "change age to 50", "use $100k instead", "set threshold to 10"

4. NEW_QUESTION is ONLY for completely unrelated questions about different topics.
   If the message references the previous analysis AT ALL, it is NOT NEW_QUESTION."""

        prompt = load_prompt("analyze_question.md").format(
            problem=problem,
            source_context=source_context,
            fact_context=fact_context,
            followup_context=followup_context,
        )

        try:
            result = self.router.execute(
                task_type=TaskType.INTENT_CLASSIFICATION,
                system="You analyze user questions efficiently. Be precise and concise.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            response = result.content.strip()

            # Debug logging for intent classification
            import logging
            _intent_logger = logging.getLogger(__name__)
            _intent_logger.debug(f"[INTENT CLASSIFICATION] Raw LLM response:\n{response}")

            # Parse response
            question_type = QuestionType.DATA_ANALYSIS
            extracted_facts = []
            cached_answer = None

            # Parse FACTS section
            if "FACTS:" in response:
                facts_section = response.split("FACTS:", 1)[1].split("---")[0].strip()
                if facts_section and facts_section != "NONE":
                    for line in facts_section.split("\n"):
                        line = line.strip().lstrip("-").strip()
                        if ":" in line and line.lower() != "none":
                            parts = line.split(":", 1)
                            fact_name = parts[0].strip()
                            value_part = parts[1].strip()

                            # Parse value and optional description (format: "value | description")
                            description = None
                            if "|" in value_part:
                                value_str, description = value_part.split("|", 1)
                                value_str = value_str.strip()
                                description = description.strip()
                            else:
                                value_str = value_part

                            # Try to parse as number
                            try:
                                value = float(value_str)
                                if value == int(value):
                                    value = int(value)
                            except ValueError:
                                value = value_str

                            # Add to fact resolver
                            fact = self.fact_resolver.add_user_fact(
                                fact_name=fact_name,
                                value=value,
                                reasoning=f"Extracted from question: {problem}",
                                description=description,
                            )
                            extracted_facts.append(fact)

            # Parse QUESTION_TYPE
            if "QUESTION_TYPE:" in response:
                type_line = response.split("QUESTION_TYPE:", 1)[1].split("\n")[0].strip()
                type_line = type_line.split("---")[0].strip().upper()
                if "META" in type_line:
                    question_type = QuestionType.META_QUESTION
                elif "GENERAL" in type_line:
                    question_type = QuestionType.GENERAL_KNOWLEDGE

            # Safety check: If classified as GENERAL_KNOWLEDGE but query mentions
            # configured data sources, override to DATA_ANALYSIS
            if question_type == QuestionType.GENERAL_KNOWLEDGE:
                problem_lower = problem.lower()
                # Check if query mentions any data source names or key terms from descriptions
                source_keywords = []
                if self.config.databases:
                    for name, db in self.config.databases.items():
                        source_keywords.append(name.lower())
                        if db.description:
                            # Extract key terms from description (nouns)
                            for word in db.description.lower().split():
                                if len(word) > 4 and word.isalpha():
                                    source_keywords.append(word)
                if self.config.documents:
                    for name, doc in self.config.documents.items():
                        source_keywords.append(name.lower().replace("_", " "))
                        source_keywords.append(name.lower())
                        if doc.description:
                            for word in doc.description.lower().split():
                                if len(word) > 4 and word.isalpha():
                                    source_keywords.append(word)

                # Check for matches
                for keyword in source_keywords:
                    if keyword in problem_lower:
                        logger.debug(f"Overriding GENERAL_KNOWLEDGE to DATA_ANALYSIS: query mentions '{keyword}'")
                        question_type = QuestionType.DATA_ANALYSIS
                        break

            # Parse CACHED_ANSWER
            if "CACHED_ANSWER:" in response:
                answer_section = response.split("CACHED_ANSWER:", 1)[1].split("---")[0].strip()
                if answer_section and answer_section.upper() != "NONE":
                    cached_answer = answer_section

            # Parse INTENTS (preserving order)
            _intent_logger.debug(f"[INTENT PARSING] Starting intent parsing, response length: {len(response)}")
            intents = []
            if "INTENTS:" in response:
                intents_section = response.split("INTENTS:", 1)[1].split("---")[0].strip()
                _intent_logger.debug(f"[INTENT PARSING] intents_section: '{intents_section}'")
                if intents_section and intents_section.upper() != "NONE":
                    for line in intents_section.split("\n"):
                        line = line.strip().lstrip("-").strip()
                        _intent_logger.debug(f"[INTENT PARSING] processing line: '{line}'")
                        if line and line.upper() != "NONE":
                            # Parse "INTENT_NAME | extracted_value" or just "INTENT_NAME"
                            if "|" in line:
                                intent_name, extracted = line.split("|", 1)
                                intent_name = intent_name.strip().upper()
                                extracted = extracted.strip()
                            else:
                                intent_name = line.strip().upper()
                                extracted = None
                            _intent_logger.debug(f"[INTENT PARSING] extracted intent: '{intent_name}'")
                            intents.append(DetectedIntent(
                                intent=intent_name,
                                confidence=0.8,
                                extracted_value=extracted,
                            ))
                else:
                    _intent_logger.debug(f"[INTENT PARSING] intents_section was empty or NONE")
            else:
                _intent_logger.debug(f"[INTENT PARSING] 'INTENTS:' not found in response")

            # Default to NEW_QUESTION if no intents detected
            if not intents:
                intents.append(DetectedIntent(intent="NEW_QUESTION", confidence=0.5))

            # Parse FACT_MODIFICATIONS
            fact_modifications = []
            if "FACT_MODIFICATIONS:" in response:
                mods_section = response.split("FACT_MODIFICATIONS:", 1)[1].split("---")[0].strip()
                if mods_section and mods_section.upper() != "NONE":
                    for line in mods_section.split("\n"):
                        line = line.strip().lstrip("-").strip()
                        if ":" in line and line.upper() != "NONE":
                            fact_name, new_value = line.split(":", 1)
                            fact_modifications.append({
                                "fact_name": fact_name.strip(),
                                "new_value": new_value.strip(),
                                "action": "modify",
                            })

            # Parse SCOPE_REFINEMENTS
            scope_refinements = []
            if "SCOPE_REFINEMENTS:" in response:
                scope_section = response.split("SCOPE_REFINEMENTS:", 1)[1].split("---")[0].strip()
                if scope_section and scope_section.upper() != "NONE":
                    for line in scope_section.split("\n"):
                        line = line.strip().lstrip("-").strip()
                        if line and line.upper() != "NONE":
                            scope_refinements.append(line)

            # Parse WANTS_BRIEF
            wants_brief = False
            if "WANTS_BRIEF:" in response:
                brief_section = response.split("WANTS_BRIEF:", 1)[1].split("---")[0].strip()
                wants_brief = brief_section.upper().startswith("YES")

            # Parse EXECUTION_MODE and MODE_REASON
            recommended_mode = "EXPLORATORY"  # Default
            mode_reasoning = None
            if "EXECUTION_MODE:" in response:
                mode_section = response.split("EXECUTION_MODE:", 1)[1].split("\n")[0].strip()
                mode_section = mode_section.split("---")[0].strip().upper()
                if "PROOF" in mode_section or "AUDITABLE" in mode_section:
                    recommended_mode = "PROOF"
                elif "EXPLORATORY" in mode_section:
                    recommended_mode = "EXPLORATORY"
            if "MODE_REASON:" in response:
                reason_section = response.split("MODE_REASON:", 1)[1].split("---")[0].strip()
                if reason_section and reason_section.upper() != "NONE":
                    mode_reasoning = reason_section.split("\n")[0].strip()

            _intent_logger.debug(f"[INTENT PARSING] Final parsed intents: {[i.intent for i in intents]}")

            # Store intent in datastore for debugging (if method exists)
            if self.datastore and hasattr(self.datastore, 'set_query_intent'):
                self.datastore.set_query_intent(
                    query_text=problem,
                    intents=[{"intent": i.intent, "confidence": i.confidence, "value": i.extracted_value} for i in intents],
                    is_followup=bool(self.session_id),
                )

            return QuestionAnalysis(
                question_type=question_type,
                extracted_facts=extracted_facts,
                cached_fact_answer=cached_answer,
                intents=intents,
                fact_modifications=fact_modifications,
                scope_refinements=scope_refinements,
                wants_brief=wants_brief,
                recommended_mode=recommended_mode,
                mode_reasoning=mode_reasoning,
            )

        except Exception as e:
            # On error, fall back to regex-based classification
            _intent_logger.exception(f"[INTENT PARSING] Exception during question analysis: {e}")
            return QuestionAnalysis(
                question_type=QuestionType.META_QUESTION if is_meta_question(problem) else QuestionType.DATA_ANALYSIS,
                extracted_facts=[],
                cached_fact_answer=None,
                intents=[DetectedIntent(intent="NEW_QUESTION", confidence=0.5)],
            )

    def _detect_ambiguity(self, problem: str, is_auditable_mode: bool = False, session_tables: Optional[list[dict]] = None) -> Optional[ClarificationRequest]:
        """
        Detect if a question is ambiguous and needs clarification before planning.

        Checks for missing parameters like:
        - Geographic scope ("how many bears" - where?)
        - Time period ("what were sales" - when?)
        - Threshold values ("top customers" - top how many?)
        - Category/segment ("product performance" - which products?)

        Args:
            problem: The user's question
            is_auditable_mode: If True, defer personal value questions to lazy resolution
            session_tables: List of existing session tables (from datastore.list_tables())

        Returns:
            ClarificationRequest if clarification needed, None otherwise
        """
        ctx = self._build_source_context()

        # Inject user correction and NL correction learnings so ambiguity detection
        # respects past corrections (e.g. "use business_rules not compensation_policy")
        learnings_text = ""
        if self.learning_store:
            try:
                rule_lines = []
                for category in [LearningCategory.USER_CORRECTION, LearningCategory.NL_CORRECTION]:
                    rules = self.learning_store.list_rules(
                        category=category,
                        min_confidence=0.6,
                    )
                    for rule in rules[:3]:
                        rule_lines.append(f"- {rule['summary']}")
                if rule_lines:
                    learnings_text = "\n## Learned Rules (respect these when suggesting options)\n" + "\n".join(rule_lines)
            except Exception:
                pass

        personal_values_guidance = (
            "NEVER ask for personal VALUES like age, salary, preferences - these will be requested later during fact resolution. The user explicitly referenced 'my age' means they intend to provide it - don't pre-ask."
            if is_auditable_mode else
            "For personal values mentioned (like 'my age'), you MAY ask since exploratory mode needs all values upfront."
        )
        # Format session tables so the LLM knows what datasets already exist
        session_tables_text = ""
        if session_tables:
            table_lines = ["\n## Session Tables (datasets created during this conversation)"]
            for t in session_tables:
                table_lines.append(f"- `{t['name']}`: {t.get('row_count', '?')} rows (step {t.get('step_number', '?')})")
            session_tables_text = "\n".join(table_lines)

        prompt = load_prompt("detect_ambiguity.md").format(
            problem=problem,
            schema_overview=ctx["schema_overview"],
            api_overview=ctx["api_overview"],
            doc_overview=ctx["doc_overview"],
            user_facts=ctx["user_facts"],
            learnings_text=learnings_text,
            personal_values_guidance=personal_values_guidance,
            session_tables=session_tables_text,
        )

        try:
            result = self.router.execute(
                task_type=TaskType.INTENT_CLASSIFICATION,
                system="You detect ambiguity in data analysis requests. Be practical - only flag truly ambiguous requests.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            response = result.content.strip()

            if response.startswith("CLEAR"):
                return None

            # Parse ambiguous response
            if "AMBIGUOUS" in response:
                lines = response.split("\n")
                reason = ""
                questions: list[ClarificationQuestion] = []
                current_question = None
                in_questions_section = False

                for line in lines:
                    line = line.strip()
                    if line.startswith("REASON:"):
                        reason = line[7:].strip()
                    elif line.upper().startswith("QUESTIONS"):
                        in_questions_section = True
                    elif line.startswith("SUGGESTIONS:") and current_question:
                        # Parse suggestions for current question
                        suggestions_text = line[12:].strip()
                        suggestions = [s.strip() for s in suggestions_text.split("|") if s.strip()]
                        current_question.suggestions = suggestions[:4]  # Max 4 suggestions
                    elif in_questions_section and line:
                        # Try to parse as a question in specific formats only:
                        # Q1: question, - question, 1. question, 1) question
                        # Do NOT capture arbitrary text as questions (could be LLM reasoning)
                        question_text = None

                        if line.startswith("Q") and ":" in line[:4]:
                            # Format: Q1: question text
                            question_text = line.split(":", 1)[1].strip()
                        elif line.startswith("- "):
                            # Format: - question text
                            question_text = line[2:].strip()
                        elif len(line) > 2 and line[0].isdigit() and line[1] in ".):":
                            # Format: 1. question or 1) question or 1: question
                            question_text = line[2:].strip()
                        elif len(line) > 3 and line[:2].isdigit() and line[2] in ".):":
                            # Format: 10. question (two digit number)
                            question_text = line[3:].strip()
                        # NOTE: We intentionally do NOT capture arbitrary text as questions
                        # The LLM sometimes adds explanatory text that shouldn't be treated as questions

                        # Only accept if it looks like a question (ends with ? or starts with question word)
                        if question_text and len(question_text) > 5:
                            is_question = (
                                question_text.endswith("?") or
                                question_text.lower().startswith(("what ", "which ", "how ", "when ", "where ", "who ", "should ", "do ", "does ", "is ", "are "))
                            )
                            if is_question:
                                # Save previous question and start new one
                                if current_question and current_question.text:
                                    questions.append(current_question)
                                current_question = ClarificationQuestion(text=question_text)

                # Don't forget the last question
                if current_question and current_question.text:
                    questions.append(current_question)

                if questions:
                    return ClarificationRequest(
                        original_question=problem,
                        ambiguity_reason=reason,
                        questions=questions[:3],  # Max 3 questions
                    )

            return None

        except Exception as e:
            # On error, proceed without clarification
            logger.debug(f"Clarification detection failed (proceeding without): {e}")
            return None

    def _request_clarification(self, request: ClarificationRequest) -> Optional[str]:
        """
        Request clarification from the user.

        Args:
            request: The clarification request

        Returns:
            Enhanced question with clarification, or None to skip
        """
        # Skip if disabled or no callback
        if self.session_config.skip_clarification:
            return None

        if not self._clarification_callback:
            return None

        # Emit event for UI
        self._emit_event(StepEvent(
            event_type="clarification_needed",
            step_number=0,
            data={
                "reason": request.ambiguity_reason,
                "questions": request.questions,
            }
        ))

        response = self._clarification_callback(request)

        if response.skip:
            return None

        # Build enhanced question with clarifications
        clarifications = []
        logger.debug(f"[CLARIFICATION] Response answers: {response.answers}")
        for question, answer in response.answers.items():
            logger.debug(f"[CLARIFICATION] Q: {question!r} -> A: {answer!r}")
            if answer:
                clarifications.append(f"{question}: {answer}")

        if clarifications:
            enhanced = f"{request.original_question}\n\nClarifications:\n" + "\n".join(clarifications)
            logger.debug(f"[CLARIFICATION] Enhanced problem:\n{enhanced}")
            return enhanced

        return None

    def _answer_general_question(self, problem: str) -> dict:
        """
        Answer a general knowledge question directly using LLM.
        """
        result = self.router.execute(
            task_type=TaskType.GENERAL,
            system="You are a helpful assistant. Answer the question directly and concisely.",
            user_message=problem,
        )

        return {
            "success": True,
            "meta_response": True,  # Reuse this flag to skip planning display
            "output": result.content,
            "plan": None,
        }

    def _answer_from_cached_facts(self, problem: str) -> Optional[dict]:
        """
        Try to answer a question from cached facts.

        Checks if the question references a fact already in the cache
        (e.g., "what is my role" -> user_role fact).

        Returns:
            Answer dict if fact found, None otherwise
        """
        cached_facts = self.fact_resolver.get_all_facts()
        if not cached_facts:
            return None

        # Create context about available facts (use display_value for table references)
        fact_context = "\n".join(
            f"- {name}: {fact.display_value}" for name, fact in cached_facts.items()
        )

        prompt = f"""Check if this question can be answered from these known facts:

Known facts:
{fact_context}

User question: {problem}

If the question asks about one of these facts, respond with:
FACT_MATCH: <fact_name>
ANSWER: <direct answer using the fact value>

If the question cannot be answered from these facts, respond with:
NO_MATCH

Examples:
- "what is my role" + fact user_role=CFO -> FACT_MATCH: user_role, ANSWER: Your role is CFO.
- "what's the target region" + fact target_region=US -> FACT_MATCH: target_region, ANSWER: The target region is US.
- "how many customers" + no matching fact -> NO_MATCH
"""

        try:
            result = self.router.execute(
                task_type=TaskType.GENERAL,
                system="You are a helpful assistant matching questions to known facts.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            response = result.content.strip()
            if "NO_MATCH" in response:
                return None

            # Extract the answer
            if "ANSWER:" in response:
                answer_start = response.index("ANSWER:") + 7
                answer = response[answer_start:].strip()

                return {
                    "success": True,
                    "meta_response": True,
                    "output": answer,
                    "plan": None,
                }
        except Exception as e:
            logger.debug(f"Meta question handling failed: {e}")

        return None

    def _explain_differentiators(self) -> dict:
        """Explain what makes Constat different from other AI tools."""
        explanation = load_prompt("explain_differentiators.md")

        return {
            "success": True,
            "meta_response": True,
            "output": explanation,
            "suggestions": [
                "What data is available?",
                "How do you reason about problems?",
            ],
            "plan": None,
        }

    def _explain_reasoning_methodology(self) -> dict:
        """Explain Constat's reasoning methodology."""
        explanation = load_prompt("explain_reasoning_methodology.md")

        return {
            "success": True,
            "meta_response": True,
            "output": explanation,
            "suggestions": [
                "What data is available?",
                "Show me an example analysis",
            ],
            "plan": None,
        }

    def _answer_personal_question(self) -> dict:
        """Answer personal questions about Vera."""
        explanation = load_prompt("answer_personal_question.md")

        return {
            "success": True,
            "meta_response": True,
            "output": explanation,
            "suggestions": [
                "What data is available?",
                "How do you reason about problems?",
                "What makes you different, Vera?",
            ],
            "plan": None,
        }

    def _answer_meta_question(self, problem: str) -> dict:
        """
        Answer meta-questions about capabilities without planning/execution.

        Uses schema overview and domain context to answer questions like
        "what questions can you answer" directly.
        """
        # Check if asking about reasoning methodology
        problem_lower = problem.lower()
        if any(phrase in problem_lower for phrase in [
            "how do you reason", "how do you think", "how do you work",
            "reasoning process", "methodology", "how does this work"
        ]):
            return self._explain_reasoning_methodology()

        # Check if asking what makes Constat/Vera different
        if any(phrase in problem_lower for phrase in [
            "what makes", "what's different", "how is .* different",
            "unique about", "special about", "why constat", "why use constat",
            "why vera"
        ]):
            return self._explain_differentiators()

        # Check if asking personal questions about Vera
        if any(phrase in problem_lower for phrase in [
            "who are you", "what are you", "your name", "about you",
            "how old", "your age", "are you a ", "are you an ",
            "who made you", "who created you", "who built you",
            "tell me about yourself", "introduce yourself", "vera"
        ]):
            return self._answer_personal_question()

        ctx = self._build_source_context(include_user_facts=False)
        domain_context = self._get_system_prompt()

        # Get user role if known
        user_role = None
        try:
            role_fact = self.fact_resolver.get_fact("user_role")
            if role_fact:
                user_role = role_fact.value
        except Exception as e:
            logger.debug(f"Failed to get user_role fact: {e}")

        role_context = f"\nThe user's role is: {user_role}" if user_role else ""

        prompt = f"""The user is asking about your capabilities. Answer based on the available data.

User question: {problem}{role_context}

Available databases and tables:
{ctx["schema_overview"]}
{ctx["api_overview"]}
{ctx["doc_overview"]}

Domain context:
{domain_context}

Provide a helpful summary tailored to the user's role (if known):
1. What data sources are relevant to their role (databases, APIs, and reference documents)
2. What types of analyses would be most valuable

Then provide 3-6 example questions the user could ask, each on its own line in quotes like:
"What is the revenue by region?"

Keep it concise and actionable."""

        result = self.router.execute(
            task_type=TaskType.SUMMARIZATION,
            system="You are a helpful assistant explaining data analysis capabilities.",
            user_message=prompt,
        )

        # Extract example questions from output to use as suggestions
        # Don't emit event here - let REPL display after output
        suggestions = self._extract_example_questions(result.content)

        # Strip the example questions section from output to avoid duplication
        output = self._strip_example_questions_section(result.content)

        return {
            "success": True,
            "meta_response": True,
            "output": output,
            "suggestions": suggestions,
            "plan": None,
        }

    def _extract_example_questions(self, text: str) -> list[str]:
        """
        Extract example questions from meta-response text.

        Looks for quoted questions in the text that the user could ask.
        """
        import re
        questions = []

        # Look for questions in quotes (single or double)
        # Pattern: "question?" or 'question?'
        quoted_pattern = r'["\u201c]([^"\u201d]+\?)["\u201d]'
        matches = re.findall(quoted_pattern, text)
        for match in matches:
            q = match.strip()
            if len(q) > 10 and q not in questions:  # Skip very short matches
                questions.append(q)

        # Limit to 6 suggestions
        return questions[:6]

    def _strip_example_questions_section(self, text: str) -> str:
        """
        Strip the example questions section from meta-response output.

        This avoids duplicating questions that will be shown as suggestions.
        """
        import re

        # Find where example questions section starts and remove from there
        # Match various header formats
        patterns = [
            r'\n*Example Questions[^\n]*:\s*\n',  # "Example Questions You Could Ask:"
            r'\n*#+\s*Example Questions?[^\n]*\n',  # Markdown header
            r'\n*\*\*Example Questions?[^\n]*\*\*\s*\n',  # Bold header
        ]

        for pattern in patterns:
            match = re.search(pattern, text, flags=re.IGNORECASE)
            if match:
                # Remove from the start of the example section to the end
                return text[:match.start()].rstrip()

        return text.rstrip()

    def _synthesize_answer(self, problem: str, step_outputs: str, artifacts: list[dict] | None = None) -> str:
        """
        Synthesize a final user-facing answer from step execution outputs.

        This takes the raw step outputs (which may be verbose technical details)
        and creates a clear, direct answer to the user's original question.
        """
        # Build artifact context if files were created
        artifact_context = ""
        if artifacts:
            artifact_lines = ["Files created:"]
            for a in artifacts:
                desc = a.get("description", "")
                uri = a.get("file_uri", "")
                if uri:
                    filename = uri.split("/")[-1]
                    artifact_lines.append(f"- `{filename}`: {desc}")
            artifact_context = "\n" + "\n".join(artifact_lines) + "\n"

        prompt = f"""Synthesize a final insight from this data analysis.

Question: {problem}

[Step outputs for context - DO NOT repeat these, user already saw them]
{step_outputs}
{artifact_context}
IMPORTANT: The user has ALREADY SEEN all the step-by-step output above in separate messages. Your job is ONLY to provide high-level analysis - NOT to summarize what each step did.

Write a brief insight (max 150 words) with:

1. **Answer**: Direct 1-sentence answer to their question
2. **Key Insight**: The most important finding or pattern (not a list of what steps did)
3. **Next Steps**: 2-3 follow-up questions they could ask (formatted as a numbered list)

WRONG (do not do this):
- "Step 1 loaded the data, Step 2 filtered it, Step 3 calculated..."
- "First I analyzed X, then I computed Y..."

RIGHT:
- "The top 3 customers account for 60% of revenue. This concentration suggests..."
- "Revenue grew 15% YoY, driven primarily by the Enterprise segment..."

Reference tables with backticks: `table_name`"""

        result = self.router.execute(
            task_type=TaskType.SUMMARIZATION,
            system="You are a data analyst presenting findings. Be clear and direct.",
            user_message=prompt,
            max_tokens=self.router.max_output_tokens,
        )

        if not result.success:
            logger.warning(f"Answer synthesis failed: {result.content}")
            return "Analysis completed but answer synthesis failed. See step outputs above."
        return result.content

    def _extract_facts_from_response(self, problem: str, answer: str) -> list:
        """
        Extract facts from the analysis response to cache for follow-up questions.

        For example, if the answer says "Total revenue was $2.4M", we cache
        the fact `total_revenue = 2400000` so follow-up questions like
        "How does that compare to last year?" can reference it.

        Returns:
            List of extracted Fact objects
        """
        prompt = load_prompt("extract_facts_from_response.md").format(
            problem=problem,
            answer=answer,
        )

        try:
            result = self.router.execute(
                task_type=TaskType.INTENT_CLASSIFICATION,
                system="You extract key facts and metrics from analysis results.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            response = result.content.strip()
            if "NO_FACTS" in response:
                return []

            extracted_facts = []
            for line in response.split("\n"):
                line = line.strip()
                if line == "---" or not line:
                    continue
                if ":" in line and not line.startswith("FACT"):
                    parts = line.split(":", 1)
                    fact_name = parts[0].strip().lower().replace(" ", "_")
                    value_part = parts[1].strip()

                    # Parse value and optional description (format: "value | description")
                    description = None
                    if "|" in value_part:
                        value_str, description = value_part.split("|", 1)
                        value_str = value_str.strip()
                        description = description.strip()
                    else:
                        value_str = value_part

                    # Try to parse as number
                    try:
                        # Handle currency (remove $, commas)
                        clean_value = value_str.replace("$", "").replace(",", "").replace("%", "")
                        value = float(clean_value)
                        if "%" in value_str:
                            value = value / 100  # Convert percentage
                        elif value == int(value):
                            value = int(value)
                    except ValueError:
                        value = value_str

                    # Add to fact resolver - source is DERIVED since synthesized from analysis
                    fact = self.fact_resolver.add_user_fact(
                        fact_name=fact_name,
                        value=value,
                        reasoning=f"Extracted from exploratory analysis. Run in auditable mode for full provenance.",
                        source=FactSource.DERIVED,
                        description=description,
                    )
                    extracted_facts.append(fact)

            return extracted_facts

        except Exception as e:
            logger.debug(f"Failed to extract facts from response: {e}")
            return []

    def _generate_suggestions(self, problem: str, answer: str, tables: list[dict]) -> list[str]:
        """
        Generate contextual follow-up suggestions based on the answer and available data.

        Args:
            problem: The original question
            answer: The synthesized answer
            tables: Available tables in the datastore

        Returns:
            List of 1-3 suggested follow-up questions
        """
        table_info = ", ".join(t["name"] for t in tables) if tables else "none"

        prompt = f"""Based on this completed analysis, suggest 1-3 actionable follow-up requests the user could make.

Original question: {problem}

Answer provided:
{answer}

Available data tables: {table_info}

Guidelines:
- Suggest ACTIONABLE REQUESTS that extend or build on the analysis (e.g., "Show a breakdown by region", "Compare this to last quarter")
- DO NOT ask clarifying questions back to the user (e.g., "Why did you need this?" or "What will you use this for?")
- Each suggestion should be something the system can execute
- Keep suggestions concise (under 12 words each)
- Consider: breakdowns, comparisons, visualizations, exports, time periods, rankings
- If the analysis seems complete, return just 1 suggestion or nothing

Return ONLY the suggestions, one per line, no numbering or bullets."""

        try:
            result = self.router.execute(
                task_type=TaskType.SUMMARIZATION,
                system="You suggest actionable follow-up analysis requests. Never ask clarifying questions.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            # Parse suggestions (one per line)
            suggestions = [
                s.strip().lstrip("0123456789.-) ")
                for s in result.content.strip().split("\n")
                if s.strip() and len(s.strip()) > 5
            ]
            return suggestions[:3]  # Max 3 suggestions
        except Exception as e:
            logger.debug(f"Failed to generate suggestions (non-fatal): {e}")
            return []

    def _run_post_synthesis_parallel(
        self,
        problem: str,
        final_answer: str,
        tables: list[dict],
    ) -> tuple[list, list[str]]:
        """
        Run facts extraction and suggestions generation in parallel.

        These two tasks are independent and can run concurrently to reduce
        total synthesis time from ~3 LLM calls to ~2 (synthesis + max(facts, suggestions)).

        Args:
            problem: The original question
            final_answer: The synthesized answer
            tables: Available tables for suggestions context

        Returns:
            Tuple of (extracted_facts, suggestions)
        """
        from concurrent.futures import ThreadPoolExecutor, as_completed

        response_facts = []
        suggestions = []

        def extract_facts():
            return self._extract_facts_from_response(problem, final_answer)

        def generate_suggestions():
            return self._generate_suggestions(problem, final_answer, tables)

        with ThreadPoolExecutor(max_workers=2) as executor:
            facts_future = executor.submit(extract_facts)
            suggestions_future = executor.submit(generate_suggestions)

            # Wait for both to complete
            for future in as_completed([facts_future, suggestions_future]):
                try:
                    if future == facts_future:
                        response_facts = future.result()
                    else:
                        suggestions = future.result()
                except Exception as e:
                    logger.debug(f"Post-synthesis task failed (non-fatal): {e}")

        return response_facts, suggestions

    def _replan_with_feedback(self, problem: str, feedback: str) -> PlannerResponse:
        """
        Generate a new plan incorporating user feedback.

        Feedback is appended verbatim to the problem, similar to how clarifications
        are handled. This preserves the exact user input for the planner.

        Args:
            problem: Original problem
            feedback: User's suggested changes (passed verbatim)

        Returns:
            New PlannerResponse with updated plan
        """
        # Append feedback verbatim, similar to clarifications format
        enhanced_problem = f"{problem}\n\nPlan Adjustments:\n{feedback}"

        self._sync_user_facts_to_planner()
        self._sync_available_roles_to_planner()
        return self.planner.plan(enhanced_problem)

    # =========================================================================
    # Phase 3: Intent Classification and Handler Methods
    # =========================================================================

    def _classify_turn_intent(self, user_input: str) -> TurnIntent:
        """
        Classify the user's input into a TurnIntent using the IntentClassifier.

        Builds context from current conversation state and delegates to the
        embedding-based classifier (with LLM fallback for low confidence).

        Args:
            user_input: The user's natural language input.

        Returns:
            TurnIntent with primary intent, optional sub-intent, and optional target.
        """
        # Build context dict for the classifier
        context = {
            "phase": self._conversation_state.phase,
            "has_plan": self._conversation_state.active_plan is not None,
        }

        # Delegate to the intent classifier
        return self._intent_classifier.classify(user_input, context)

    def _handle_query_intent(self, turn_intent: TurnIntent, user_input: str) -> dict:
        """
        Handle QUERY primary intent - answer from knowledge or current context.

        This handles sub-intents:
        - DETAIL: drill down into specific aspect
        - PROVENANCE: show proof chain
        - SUMMARY: condense results
        - LOOKUP: simple fact retrieval
        - Default: general answer using doc search + LLM fallback

        Args:
            turn_intent: The classified turn intent.
            user_input: The user's original input.

        Returns:
            Result dict with output, success, and other metadata.
        """
        sub_intent = turn_intent.sub
        target = turn_intent.target

        # Handle PROVENANCE sub-intent - show proof chain
        if sub_intent == SubIntent.PROVENANCE:
            return self._handle_provenance_query(target, user_input)

        # Handle DETAIL sub-intent - drill down into specific aspect
        if sub_intent == SubIntent.DETAIL:
            return self._handle_detail_query(target, user_input)

        # Handle SUMMARY sub-intent - condense results
        if sub_intent == SubIntent.SUMMARY:
            return self._handle_summary_query(user_input)

        # Handle LOOKUP sub-intent - simple fact retrieval
        if sub_intent == SubIntent.LOOKUP:
            result = self._handle_lookup_query(target, user_input)
            # Check if lookup found data sources that need planning
            if result.get("_route_to_planning"):
                return result  # Signal will be handled by solve()
            return result

        # Default: general answer using doc search + LLM fallback (KNOWLEDGE mode logic)
        return self._handle_general_query(user_input)

    def _handle_provenance_query(self, target: Optional[str], user_input: str) -> dict:
        """Handle provenance/proof chain query."""
        # Check if we have resolved facts with provenance
        all_facts = self.fact_resolver.get_all_facts()

        if not all_facts:
            return {
                "success": True,
                "output": "No facts have been resolved yet. Please run an analysis first to establish a proof chain.",
                "meta_response": True,
            }

        # Build provenance output
        provenance_lines = ["**Fact Provenance:**\n"]

        for name, fact in all_facts.items():
            # Check if this fact matches the target (if specified)
            if target and target.lower() not in name.lower():
                continue

            provenance_lines.append(f"**{name}**: {fact.display_value}")
            if hasattr(fact, "source") and fact.source:
                provenance_lines.append(f"  - Source: {fact.source}")
            if hasattr(fact, "reasoning") and fact.reasoning:
                provenance_lines.append(f"  - Reasoning: {fact.reasoning}")
            if hasattr(fact, "confidence") and fact.confidence is not None:
                provenance_lines.append(f"  - Confidence: {fact.confidence:.0%}")
            provenance_lines.append("")

        if len(provenance_lines) == 1:
            return {
                "success": True,
                "output": f"No facts found matching '{target}'." if target else "No facts available.",
                "meta_response": True,
            }

        return {
            "success": True,
            "output": "\n".join(provenance_lines),
            "meta_response": True,
        }

    def _handle_detail_query(self, target: Optional[str], user_input: str) -> dict:
        """Handle detail/drill-down query."""
        # If we have a datastore with results, try to get details from there
        if self.datastore:
            tables = self.datastore.list_tables()
            if tables:
                # Use LLM to generate a detail explanation
                table_info = "\n".join([f"- {t['name']}: {t['row_count']} rows" for t in tables])
                scratchpad_context = self.datastore.get_scratchpad_as_markdown()

                prompt = f"""The user wants more details about: {user_input}

Available data:
{table_info}

Previous analysis context:
{scratchpad_context}

Provide a detailed explanation or suggest what specific data to examine."""

                result = self.router.execute(
                    task_type=TaskType.SYNTHESIS,
                    system="You are a helpful data analyst providing detailed explanations.",
                    user_message=prompt,
                    max_tokens=self.router.max_output_tokens,
                )

                return {
                    "success": True,
                    "output": result.content,
                    "meta_response": True,
                }

        # Fallback to general query handling
        return self._handle_general_query(user_input)

    def _handle_summary_query(self, user_input: str) -> dict:
        """Handle summary/condensation query.

        Checks if the user is asking about a specific document first.
        If so, summarizes that document. Otherwise, summarizes previous analysis results.
        """
        # Check if the query mentions a document - search for relevant documents
        if self.doc_tools:
            try:
                doc_results = self.doc_tools.search_documents(user_input, limit=3)
                if doc_results and doc_results[0].get("similarity", 0) > 0.4:
                    # Found a relevant document - summarize it
                    doc_name = doc_results[0]["name"]
                    doc_content = self.doc_tools.get_document(doc_name)
                    if doc_content:
                        # Truncate if too long
                        max_chars = 15000
                        if len(doc_content) > max_chars:
                            doc_content = doc_content[:max_chars] + "\n\n[... truncated for length ...]"

                        prompt = f"""Summarize this document based on the user's request.

Document: {doc_name}
Content:
{doc_content}

User request: {user_input}

Provide a clear, structured summary focusing on what the user asked for."""

                        result = self.router.execute(
                            task_type=TaskType.SUMMARIZATION,
                            system="You are a document summarizer. Extract key concepts and structure your response clearly.",
                            user_message=prompt,
                            max_tokens=self.router.max_output_tokens,
                        )

                        return {
                            "success": True,
                            "output": result.content,
                            "meta_response": True,
                        }
            except Exception as e:
                logger.debug(f"Document search for summary failed: {e}")

        # Check if we have previous results to summarize
        if self.datastore:
            scratchpad_entries = self.datastore.get_scratchpad()
            if scratchpad_entries:
                # Combine all previous results
                context = "\n\n".join([
                    f"Step {e['step_number']}: {e['goal']}\n{e['narrative']}"
                    for e in scratchpad_entries
                ])

                prompt = f"""Summarize these analysis results concisely:

{context}

User request: {user_input}

Provide a brief, high-level summary of the key findings."""

                result = self.router.execute(
                    task_type=TaskType.SUMMARIZATION,
                    system="You are a concise summarizer. Focus on key insights.",
                    user_message=prompt,
                    max_tokens=self.router.max_output_tokens,
                )

                return {
                    "success": True,
                    "output": result.content,
                    "meta_response": True,
                }

        return {
            "success": True,
            "output": "No previous analysis results to summarize. Please run an analysis first.",
            "meta_response": True,
        }

    def _handle_lookup_query(self, target: Optional[str], user_input: str) -> dict:
        """Handle simple fact lookup query.

        Checks all available sources: cached facts, APIs, databases, and documents.
        Returns a signal to route to planning if data sources are found.
        """
        # First, try to answer from cached facts
        cached_result = self._answer_from_cached_facts(user_input)
        if cached_result:
            return cached_result

        # Check if any data sources (APIs, tables) can answer this query
        # Use a reasonable similarity threshold
        sources = self.find_relevant_sources(
            user_input,
            table_limit=3,
            doc_limit=3,
            api_limit=3,
            min_similarity=0.3,
        )
        logger.debug(f"[LOOKUP] find_relevant_sources returned: apis={len(sources.get('apis', []))}, tables={len(sources.get('tables', []))}, docs={len(sources.get('documents', []))}")

        # If we found relevant APIs or tables, signal to route to planning
        has_data_sources = bool(sources.get("apis") or sources.get("tables"))
        if has_data_sources:
            # Log what we found for debugging
            api_names = [a["name"] for a in sources.get("apis", [])]
            table_names = [t["name"] for t in sources.get("tables", [])]
            logger.debug(
                f"[LOOKUP] Found data sources for query - APIs: {api_names}, Tables: {table_names}"
            )
            # Return signal to route to planning
            return {"_route_to_planning": True}

        # If no data sources match, fall back to document search + LLM synthesis
        return self._handle_general_query(user_input)

    def _handle_general_query(self, user_input: str) -> dict:
        """
        Handle general query using document search + LLM fallback.

        This uses document lookup and LLM synthesis for knowledge/explanation queries.
        """
        # Use the existing _solve_knowledge method which does doc search + LLM synthesis
        return self._solve_knowledge(user_input)

    def _handle_plan_new_intent(self, turn_intent: TurnIntent, user_input: str) -> dict:
        """
        Handle PLAN_NEW primary intent - enhance problem based on sub-intent.

        NOTE: This method is called from solve() BEFORE the planning flow.
        It enhances the problem statement based on sub-intent (COMPARE, PREDICT).
        The actual planning is done by solve() after this returns.

        This handles sub-intents:
        - COMPARE: evaluate alternatives
        - PREDICT: what-if / forecast
        - Default: pass through unchanged

        Args:
            turn_intent: The classified turn intent.
            user_input: The user's original input.

        Returns:
            Dict with enhanced_problem for solve() to use.
        """
        sub_intent = turn_intent.sub
        enhanced_problem = user_input

        if sub_intent == SubIntent.COMPARE:
            # Add comparison context to the problem
            enhanced_problem = f"Compare and evaluate: {user_input}\n\nProvide a comparative analysis highlighting differences, pros/cons, and recommendations."

        elif sub_intent == SubIntent.PREDICT:
            # Add forecasting context to the problem
            enhanced_problem = f"Forecast/What-if analysis: {user_input}\n\nProvide predictive analysis with assumptions clearly stated."

        return {"enhanced_problem": enhanced_problem, "sub_intent": sub_intent}

    def _handle_plan_continue_intent(self, turn_intent: TurnIntent, user_input: str) -> dict:
        """
        Handle PLAN_CONTINUE primary intent - refine or extend the active plan.

        Uses the user's message as context for replanning.

        Args:
            turn_intent: The classified turn intent.
            user_input: The user's original input (used as modification context).

        Returns:
            Result dict from the replanning flow.
        """
        # Check for CORRECTION sub-intent - save as reusable learning
        if turn_intent.sub == SubIntent.CORRECTION:
            self._save_correction_as_learning(user_input)

        # Transition phase to PLANNING
        self._apply_phase_transition("plan_new")  # Returns to planning state

        # Check if there's a previous problem to continue from
        previous_problem = None
        if self.datastore:
            previous_problem = self.datastore.get_session_meta("problem")

        if not previous_problem:
            # No previous context - treat as a new plan
            return self._handle_plan_new_intent(turn_intent, user_input)

        # Build enhanced problem with user's modification context (verbatim like clarifications)
        enhanced_problem = f"{previous_problem}\n\nPlan Adjustments:\n{user_input}"

        # Delegate to existing follow_up() method which handles replanning
        result = self.follow_up(user_input, auto_classify=False)

        # Update conversation state based on result
        if result.get("success"):
            if result.get("plan"):
                self._conversation_state.active_plan = result["plan"]
            self._apply_phase_transition("complete")
        else:
            self._apply_phase_transition("fail")
            self._conversation_state.failure_context = result.get("error")

        return result

    def _handle_control_intent(self, turn_intent: TurnIntent, user_input: str) -> dict:
        """
        Handle CONTROL primary intent - system/session commands.

        This handles sub-intents:
        - RESET: clear session state
        - REDO_CMD: re-execute last plan
        - HELP: show available commands
        - STATUS: show current state
        - EXIT: end session (returns signal to caller)
        - CANCEL: stop execution (if executing)
        - REPLAN: stop execution, return to planning

        Args:
            turn_intent: The classified turn intent.
            user_input: The user's original input.

        Returns:
            Result dict with output, success, and control signals.
        """
        sub_intent = turn_intent.sub
        target = turn_intent.target

        # RESET: clear session state
        if sub_intent == SubIntent.RESET:
            return self._handle_reset()

        # REDO_CMD: re-execute last plan
        if sub_intent == SubIntent.REDO_CMD:
            return self._handle_redo()

        # HELP: show available commands
        if sub_intent == SubIntent.HELP:
            return self._handle_help()

        # STATUS: show current state
        if sub_intent == SubIntent.STATUS:
            return self._handle_status()

        # EXIT: end session
        if sub_intent == SubIntent.EXIT:
            return {
                "success": True,
                "exit": True,
                "output": "Ending session.",
                "meta_response": True,
            }

        # CANCEL: stop execution
        if sub_intent == SubIntent.CANCEL:
            return self._handle_cancel()

        # REPLAN: stop execution and return to planning
        if sub_intent == SubIntent.REPLAN:
            return self._handle_replan(user_input)

        # Default: unknown control command
        return {
            "success": True,
            "output": f"Unknown control command. Use /help to see available commands.",
            "meta_response": True,
        }

    def _handle_slash_command(self, command_text: str) -> dict:
        """Handle explicit slash commands via the command registry.

        This is a fast path for slash commands that bypasses intent classification.
        Commands like /tables, /show, /help are routed directly to their handlers.

        Args:
            command_text: The full command text (e.g., "/tables" or "/show orders")

        Returns:
            Result dict with output, success, and other metadata.
        """
        from constat.commands.registry import execute_command, parse_command
        from constat.commands.base import (
            TableResult,
            ListResult,
            TextResult,
            ErrorResult,
        )

        # Parse the command
        cmd, args = parse_command(command_text)

        # Execute via registry
        result = execute_command(self, cmd, args)

        # Convert CommandResult to dict format expected by solve()
        if isinstance(result, TableResult):
            # Format table as markdown
            lines = []
            if result.title:
                lines.append(f"**{result.title}**\n")
            if result.columns and result.rows:
                # Header
                lines.append("| " + " | ".join(str(c) for c in result.columns) + " |")
                lines.append("| " + " | ".join("---" for _ in result.columns) + " |")
                # Rows
                for row in result.rows:
                    lines.append("| " + " | ".join(str(v) for v in row) + " |")
            if result.footer:
                lines.append(f"\n*{result.footer}*")

            return {
                "success": result.success,
                "output": "\n".join(lines) if lines else result.footer or "No data.",
                "meta_response": True,
            }

        elif isinstance(result, ListResult):
            # Format list
            lines = []
            if result.title:
                lines.append(f"**{result.title}**\n")
            if result.items:
                for item in result.items:
                    if isinstance(item, dict):
                        # Code block
                        if "code" in item:
                            lines.append(f"### Step {item.get('step', '?')}: {item.get('goal', '')}")
                            lines.append(f"```{item.get('language', 'python')}")
                            lines.append(item["code"])
                            lines.append("```\n")
                        # Artifact-style item (has name, type, step)
                        elif "name" in item and "type" in item:
                            name = item.get("name", "")
                            atype = item.get("type", "")
                            step = item.get("step", "-")
                            title = item.get("title", "")
                            # Format: - **name** (type) - title if available
                            if title:
                                lines.append(f"- **{name}** ({atype}, step {step}) - {title}")
                            else:
                                lines.append(f"- **{name}** ({atype}, step {step})")
                        # Generic dict - format key: value pairs
                        else:
                            name = item.get("name", item.get("id", "Item"))
                            lines.append(f"- **{name}**")
                            for key, value in item.items():
                                if key not in ("name", "id") and value is not None:
                                    lines.append(f"  - {key}: {value}")
                    else:
                        lines.append(f"- {item}")
            elif result.empty_message:
                lines.append(result.empty_message)

            return {
                "success": result.success,
                "output": "\n".join(lines) if lines else result.empty_message or "No items.",
                "meta_response": True,
            }

        elif isinstance(result, TextResult):
            return {
                "success": result.success,
                "output": result.content,
                "meta_response": True,
            }

        elif isinstance(result, ErrorResult):
            return {
                "success": False,
                "output": f"Error: {result.error}" + (f"\n{result.details}" if result.details else ""),
                "meta_response": True,
            }

        else:
            # Generic fallback
            return {
                "success": getattr(result, "success", True),
                "output": str(result),
                "meta_response": True,
            }

    def _handle_reset(self) -> dict:
        """Handle reset control command - clear session state."""
        # Clear conversation state
        self._conversation_state = ConversationState(
            phase=Phase.IDLE,
        )

        # Clear fact resolver
        self.fact_resolver.clear_all_facts()

        # Clear plan
        self.plan = None

        # Clear scratchpad
        self.scratchpad = Scratchpad()

        # Clear session_id to indicate fresh start
        old_session_id = self.session_id
        self.session_id = None
        self.datastore = None

        self._apply_phase_transition("abandon")

        return {
            "success": True,
            "output": "Session reset. All facts and context cleared. Ready for a new question.",
            "meta_response": True,
            "reset": True,
        }

    def _handle_redo(self) -> dict:
        """Handle redo control command - re-execute last plan."""
        if not self.datastore:
            return {
                "success": False,
                "output": "No previous session to redo. Please run an analysis first.",
                "meta_response": True,
            }

        # Get the original problem
        problem = self.datastore.get_session_meta("problem")
        if not problem:
            return {
                "success": False,
                "output": "No previous problem found to redo.",
                "meta_response": True,
            }

        # Use replay to re-execute stored code
        try:
            result = self.replay(problem)
            return result
        except ValueError as e:
            return {
                "success": False,
                "output": f"Cannot redo: {e}",
                "meta_response": True,
            }

    def _handle_help(self) -> dict:
        """Handle help control command - show available commands."""
        # Use centralized help text from HELP_COMMANDS
        return {
            "success": True,
            "output": get_help_markdown(),
            "meta_response": True,
        }

    def _handle_status(self) -> dict:
        """Handle status control command - show current state."""
        state = self._conversation_state

        status_lines = [
            "**Current Session State:**",
            "",
            f"**Mode:** {state.mode.value.upper()}",
            f"**Phase:** {state.phase.value}",
        ]

        if state.active_plan:
            status_lines.append(f"**Active Plan:** {len(state.active_plan.steps)} steps")

        if state.session_facts:
            status_lines.append(f"**Session Facts:** {len(state.session_facts)} facts")

        # Add cached facts info
        all_facts = self.fact_resolver.get_all_facts()
        if all_facts:
            status_lines.append(f"**Resolved Facts:** {len(all_facts)}")
            for name, fact in list(all_facts.items())[:5]:  # Show first 5
                status_lines.append(f"  - {name}: {fact.display_value}")
            if len(all_facts) > 5:
                status_lines.append(f"  ... and {len(all_facts) - 5} more")

        if state.failure_context:
            status_lines.append(f"**Last Error:** {state.failure_context}")

        if self.datastore:
            tables = self.datastore.list_tables()
            if tables:
                status_lines.append(f"**Data Tables:** {len(tables)}")
                for t in tables[:5]:
                    status_lines.append(f"  - {t['name']}: {t['row_count']} rows")

        return {
            "success": True,
            "output": "\n".join(status_lines),
            "meta_response": True,
        }

    def _handle_cancel(self) -> dict:
        """Handle cancel control command - stop execution.

        Uses the Phase 4 cancel_execution() method to set the cancellation
        flag which will be checked between steps. Completed facts are preserved.
        """
        # Request cancellation (signals the execution loop to stop)
        self.cancel_execution()

        # Transition to IDLE
        self._apply_phase_transition("abandon")

        # Clear any queued intents
        cleared = self.clear_intent_queue()

        output = "Execution cancelled. Returned to idle state."
        if cleared > 0:
            output += f" ({cleared} queued intent(s) cleared.)"

        return {
            "success": True,
            "output": output,
            "meta_response": True,
            "cancelled": True,
        }

    def _handle_replan(self, user_input: str) -> dict:
        """Handle replan control command - stop and revise the plan.

        Uses the Phase 4 cancel_execution() method to stop execution
        and preserve completed facts.
        """
        # Request cancellation
        self.cancel_execution()

        # Transition to PLANNING
        self._apply_phase_transition("replan")

        return {
            "success": True,
            "output": "Ready to revise the plan. Please provide your modifications or ask a new question.",
            "meta_response": True,
            "replan": True,
        }

    def _apply_phase_transition(self, trigger: str) -> None:
        """
        Apply a phase transition based on the trigger.

        Valid triggers and their effects:
        - plan_new: IDLE -> PLANNING (or any -> PLANNING for plan_continue)
        - plan_ready: PLANNING -> AWAITING_APPROVAL
        - approve: AWAITING_APPROVAL -> EXECUTING
        - reject: AWAITING_APPROVAL -> PLANNING
        - complete: EXECUTING -> IDLE
        - fail: EXECUTING -> FAILED
        - retry: FAILED -> EXECUTING
        - replan: FAILED/EXECUTING -> PLANNING
        - abandon: any -> IDLE

        Args:
            trigger: The transition trigger name.
        """
        current_phase = self._conversation_state.phase

        if trigger == "plan_new":
            self._conversation_state.phase = Phase.PLANNING

        elif trigger == "plan_ready":
            if current_phase == Phase.PLANNING:
                self._conversation_state.phase = Phase.AWAITING_APPROVAL

        elif trigger == "approve":
            if current_phase == Phase.AWAITING_APPROVAL:
                self._conversation_state.phase = Phase.EXECUTING

        elif trigger == "reject":
            if current_phase == Phase.AWAITING_APPROVAL:
                self._conversation_state.phase = Phase.PLANNING

        elif trigger == "complete":
            self._conversation_state.phase = Phase.IDLE
            self._conversation_state.failure_context = None

        elif trigger == "fail":
            self._conversation_state.phase = Phase.FAILED

        elif trigger == "retry":
            if current_phase == Phase.FAILED:
                self._conversation_state.phase = Phase.EXECUTING

        elif trigger == "replan":
            if current_phase in (Phase.FAILED, Phase.EXECUTING):
                self._conversation_state.phase = Phase.PLANNING

        elif trigger == "abandon":
            self._conversation_state.phase = Phase.IDLE
            self._conversation_state.failure_context = None
            self._conversation_state.active_plan = None

        else:
            logger.warning(f"Unknown phase transition trigger: {trigger}")

        logger.debug(f"Phase transition: {current_phase.value} --({trigger})--> {self._conversation_state.phase.value}")

    def get_conversation_state(self) -> ConversationState:
        """
        Get the current conversation state.

        Returns:
            The current ConversationState object.
        """
        return self._conversation_state

    # =========================================================================
    # Phase 4: Execution Control
    # =========================================================================

    def cancel_execution(self) -> None:
        """
        Cancel the current execution.

        Sets the cancellation flag which will be checked between steps.
        Completed facts and results are preserved; only pending steps are cancelled.

        This method is thread-safe and can be called from another thread
        (e.g., in response to Ctrl+C or user typing "stop").
        """
        self._cancelled = True
        self._execution_context.cancel()

        # Emit cancellation event
        self._emit_event(StepEvent(
            event_type="execution_cancelled",
            step_number=0,
            data={"message": "Execution cancelled by user"}
        ))

        logger.debug("Execution cancellation requested")

    def is_cancelled(self) -> bool:
        """
        Check if cancellation has been requested.

        Returns:
            True if cancellation has been requested.
        """
        return self._cancelled

    def reset_cancellation(self) -> None:
        """
        Reset the cancellation flag.

        Called at the start of a new execution to clear any previous
        cancellation state.
        """
        self._cancelled = False
        self._execution_context.reset()

    def queue_intent(self, intent: TurnIntent, user_input: str) -> bool:
        """
        Queue an intent for processing after current execution completes.

        Queue behavior depends on intent type:
        - plan_new: Queue 1, latest wins (new request replaces queued)
        - control: Queue in order, process after execution
        - query: Not queued - should be handled in parallel immediately

        Args:
            intent: The TurnIntent to queue.
            user_input: The original user input (stored with intent for later processing).

        Returns:
            True if the intent was queued, False if it should be handled immediately.
        """
        if intent.primary == PrimaryIntent.QUERY:
            # Query intents are not queued - they're answered in parallel
            return False

        if intent.primary == PrimaryIntent.PLAN_NEW:
            # Queue 1, latest wins - replace any existing queued plan_new
            self._intent_queue = [
                (i, inp) for i, inp in self._intent_queue
                if i.primary != PrimaryIntent.PLAN_NEW
            ]
            self._intent_queue.append((intent, user_input))
            logger.debug(f"Queued plan_new intent (latest wins), queue size: {len(self._intent_queue)}")
            return True

        if intent.primary == PrimaryIntent.CONTROL:
            # Control intents queue in order
            self._intent_queue.append((intent, user_input))
            logger.debug(f"Queued control intent, queue size: {len(self._intent_queue)}")
            return True

        if intent.primary == PrimaryIntent.PLAN_CONTINUE:
            # Plan continue triggers replan - cancel current and queue
            self.cancel_execution()
            self._intent_queue.append((intent, user_input))
            logger.debug(f"Queued plan_continue intent (triggers replan), queue size: {len(self._intent_queue)}")
            return True

        return False

    def get_queued_intents_count(self) -> int:
        """
        Get the number of queued intents.

        Returns:
            The number of intents waiting to be processed.
        """
        return len(self._intent_queue)

    def process_queued_intents(self) -> list[dict]:
        """
        Process all queued intents after execution completes.

        Intents are processed in queue order. The queue is cleared
        after processing.

        Returns:
            List of result dicts from processing each queued intent.
        """
        if not self._intent_queue:
            return []

        results = []
        queued = list(self._intent_queue)
        self._intent_queue = []

        logger.debug(f"Processing {len(queued)} queued intents")

        for intent, user_input in queued:
            try:
                if intent.primary == PrimaryIntent.PLAN_NEW:
                    result = self._handle_plan_new_intent(intent, user_input)
                elif intent.primary == PrimaryIntent.PLAN_CONTINUE:
                    result = self._handle_plan_continue_intent(intent, user_input)
                elif intent.primary == PrimaryIntent.CONTROL:
                    result = self._handle_control_intent(intent, user_input)
                else:
                    result = {
                        "success": False,
                        "error": f"Unexpected queued intent type: {intent.primary}",
                    }
                results.append(result)
            except Exception as e:
                logger.error(f"Error processing queued intent: {e}")
                results.append({
                    "success": False,
                    "error": str(e),
                })

        return results

    def clear_intent_queue(self) -> int:
        """
        Clear all queued intents.

        Returns:
            The number of intents that were cleared.
        """
        count = len(self._intent_queue)
        self._intent_queue = []
        logger.debug(f"Cleared {count} queued intents")
        return count

    def get_execution_context(self) -> ExecutionContext:
        """
        Get the execution context for external monitoring.

        Returns:
            The ExecutionContext that can be used to check/request cancellation.
        """
        return self._execution_context

    def is_executing(self) -> bool:
        """
        Check if execution is currently in progress.

        Returns:
            True if the session is in EXECUTING phase.
        """
        return self._conversation_state.phase == Phase.EXECUTING

    def solve(self, problem: str) -> dict:
        """
        Solve a problem with intent-based routing and multi-step execution.

        Workflow:
        1. Classify intent (QUERY, PLAN_NEW, PLAN_CONTINUE, CONTROL)
        2. Route QUERY and CONTROL intents to handlers (no planning needed)
        3. For PLAN_NEW/PLAN_CONTINUE:
           a. Determine execution mode (exploratory vs proof)
           b. Generate plan
           c. Request user approval (if require_approval is True)
           d. Execute steps in parallel waves
           e. Synthesize answer and generate follow-up suggestions

        Args:
            problem: Natural language problem to solve

        Returns:
            Dict with plan, results, and summary
        """
        # Fast path 1: Handle slash commands directly via command registry
        # These are explicit commands like /tables, /show, /help etc.
        # Must be checked FIRST since they're the most explicit user intent
        if problem.strip().startswith("/"):
            return self._handle_slash_command(problem.strip())

        # Fast path 2: Check for meta-questions (no intent classification needed)
        # These are questions about capabilities, available data, etc.
        if is_meta_question(problem):
            self._emit_event(StepEvent(
                event_type="progress",
                step_number=0,
                data={"message": "Reviewing available data sources..."}
            ))
            return self._answer_meta_question(problem)

        # PARALLEL OPTIMIZATION: Run intent, analysis, ambiguity, and planning ALL in parallel
        # Most queries need planning, so we speculatively start it while classifying intent.
        # If planning isn't needed (CONTROL intent, META question), we discard the speculative plan.
        self._emit_event(StepEvent(
            event_type="progress",
            step_number=0,
            data={"message": "Analyzing your question..."}
        ))

        from concurrent.futures import ThreadPoolExecutor, as_completed
        import time

        def run_intent():
            return self._classify_turn_intent(problem)

        def run_analysis():
            return self._analyze_question(problem)

        def run_ambiguity():
            existing_tables = self.datastore.list_tables()
            return self._detect_ambiguity(problem, is_auditable_mode=True, session_tables=existing_tables)

        def run_dynamic_context():
            # Match skills and roles dynamically based on query
            try:
                return self.get_dynamic_context(problem)
            except Exception as e:
                logger.debug(f"[PARALLEL] Dynamic context matching failed: {e}")
                return None

        def run_planning():
            # Speculative planning - may be discarded if intent doesn't need it
            try:
                self._sync_user_facts_to_planner()
                self._sync_available_roles_to_planner()
                return self.planner.plan(problem)
            except Exception as e:
                logger.debug(f"[PARALLEL] Speculative planning failed: {e}")
                return None

        # Determine which tasks to run
        tasks = {
            "intent": run_intent,
            "analysis": run_analysis,
            "dynamic_context": run_dynamic_context,
        }
        if self.session_config.ask_clarifications and self._clarification_callback:
            tasks["ambiguity"] = run_ambiguity
        # Always run speculative planning in parallel
        tasks["planning"] = run_planning

        # Run all tasks in parallel
        parallel_start = time.time()
        results = {}
        with ThreadPoolExecutor(max_workers=len(tasks)) as executor:
            futures = {executor.submit(fn): name for name, fn in tasks.items()}
            for future in as_completed(futures):
                name = futures[future]
                try:
                    results[name] = future.result()
                except Exception as e:
                    logger.warning(f"[PARALLEL] Task {name} failed: {e}")
                    results[name] = None

        parallel_duration = time.time() - parallel_start
        logger.debug(f"[PARALLEL] All tasks completed in {parallel_duration:.2f}s")

        turn_intent = results.get("intent")
        analysis = results.get("analysis")
        clarification_request = results.get("ambiguity")
        speculative_plan = results.get("planning")
        dynamic_context = results.get("dynamic_context")

        # Emit dynamic context event (role and skills matched for this query)
        logger.info(f"[DYNAMIC_CONTEXT] dynamic_context={dynamic_context}")
        if dynamic_context:
            # Activate matched skills so they flow into planner and codegen prompts
            matched_skills = dynamic_context.get("skills", [])
            if matched_skills:
                skill_names = [s["name"] for s in matched_skills]
                activated = self.skill_manager.set_active_skills(skill_names)
                logger.info(f"[DYNAMIC_CONTEXT] Activated skills: {activated}")
                if activated:
                    speculative_plan = None
                    logger.debug("[PARALLEL] Skills activated, discarding speculative plan")

            event_data = {
                "role": dynamic_context.get("role"),
                "skills": matched_skills,
                "role_source": dynamic_context.get("role_source"),
            }
            logger.info(f"[DYNAMIC_CONTEXT] Emitting event with data: role={event_data.get('role')}, skills={event_data.get('skills')}")
            self._emit_event(StepEvent(
                event_type="dynamic_context",
                step_number=0,
                data=event_data,
            ))

        # Route based on primary intent (may discard speculative plan)
        route_to_planning = False

        if turn_intent and turn_intent.primary == PrimaryIntent.QUERY:
            # QUERY intent - answer from knowledge or current context
            result = self._handle_query_intent(turn_intent, problem)
            if not result.get("_route_to_planning"):
                logger.debug("[PARALLEL] QUERY handled without planning (speculative plan discarded)")
                return result
            logger.info("[ROUTING] QUERY/LOOKUP found data sources, routing to planning")
            route_to_planning = True

        if turn_intent and turn_intent.primary == PrimaryIntent.CONTROL and not route_to_planning:
            logger.debug("[PARALLEL] CONTROL intent (speculative plan discarded)")
            return self._handle_control_intent(turn_intent, problem)

        # PLAN_NEW, PLAN_CONTINUE, or re-routed QUERY - continue with planning flow
        self._apply_phase_transition("plan_new")

        # Apply sub-intent enhancements for PLAN_NEW (COMPARE, PREDICT)
        enhanced_problem = problem
        if turn_intent and turn_intent.primary == PrimaryIntent.PLAN_NEW:
            enhancement = self._handle_plan_new_intent(turn_intent, problem)
            enhanced_problem = enhancement.get("enhanced_problem", problem)
            # If problem was enhanced, speculative plan may be stale - replan
            if enhanced_problem != problem:
                logger.debug("[PARALLEL] Problem enhanced, replanning...")
                speculative_plan = None
        problem = enhanced_problem

        # Emit facts if any were extracted
        if analysis.extracted_facts:
            self._emit_event(StepEvent(
                event_type="facts_extracted",
                step_number=0,
                data={
                    "facts": [f.to_dict() for f in analysis.extracted_facts],
                    "source": "question",
                }
            ))

        # Return cached fact answer if question was about a known fact
        if analysis.cached_fact_answer:
            return {
                "success": True,
                "meta_response": True,
                "output": analysis.cached_fact_answer,
                "plan": None,
            }

        question_type = analysis.question_type

        if question_type == QuestionType.META_QUESTION:
            logger.debug("[PARALLEL] META_QUESTION (speculative plan discarded)")
            self._emit_event(StepEvent(
                event_type="progress",
                step_number=0,
                data={"message": "Reviewing available data sources..."}
            ))
            return self._answer_meta_question(problem)
        elif question_type == QuestionType.GENERAL_KNOWLEDGE:
            logger.debug("[PARALLEL] GENERAL_KNOWLEDGE (speculative plan discarded)")
            self._emit_event(StepEvent(
                event_type="progress",
                step_number=0,
                data={"message": "Generating response..."}
            ))
            return self._answer_general_question(problem)

        # Check for clarification (clarification_request already computed in parallel above)
        if clarification_request:
            enhanced_problem = self._request_clarification(clarification_request)
            if enhanced_problem:
                problem = enhanced_problem
                # Re-analyze with clarified problem - speculative plan is stale
                logger.debug("[PARALLEL] Problem clarified, replanning...")
                speculative_plan = None
                analysis = self._analyze_question(problem)

        # Create session + datastore (idempotent — may already exist from skill execution)
        self._ensure_session_datastore(problem)

        # Initialize session state
        self.scratchpad = Scratchpad(initial_context=f"Problem: {problem}")

        # Save problem statement to datastore (for UI restoration)
        self.datastore.set_session_meta("problem", problem)
        self.datastore.set_session_meta("status", "planning")

        # Log the initial query
        self.history.log_user_input(self.session_id, problem, "query")

        # Check for unclear/garbage input before processing
        if self._is_unclear_input(problem):
            return {
                "success": True,
                "meta_response": True,
                "output": "I'm not sure I understand that input. Could you rephrase your question?\n\n"
                          "You can ask me things like:\n"
                          "- Questions about your data (e.g., \"What's our total revenue?\")\n"
                          "- Verification requests (e.g., \"Prove that sales increased\")\n"
                          "- Explanations (e.g., \"How do you reason about problems?\")\n\n"
                          "Type /help for available commands.",
                "suggestions": [
                    "What data is available?",
                    "How can you help me?",
                ],
            }

        # All queries use exploratory mode by default
        # Use /prove command to generate auditable proofs when needed

        # Generate plan with approval loop
        current_problem = problem
        display_problem = problem  # What to show in UI (just feedback on replan)
        replan_attempt = 0

        while replan_attempt <= self.session_config.max_replan_attempts:
            # Use speculative plan if available (from parallel execution), otherwise generate new plan
            if speculative_plan is not None and replan_attempt == 0:
                logger.debug("[PARALLEL] Using speculative plan (saved ~1 LLM call)")
                planner_response = speculative_plan
                self.plan = planner_response.plan
                # Emit planning events for UI consistency
                self._emit_event(StepEvent(
                    event_type="planning_start",
                    step_number=0,
                    data={"message": "Plan ready..."}
                ))
            else:
                # Emit planning start event
                self._emit_event(StepEvent(
                    event_type="planning_start",
                    step_number=0,
                    data={"message": "Analyzing data sources and creating plan..."}
                ))

                # Sync user facts and roles to planner before generating plan
                self._sync_user_facts_to_planner()
                self._sync_available_roles_to_planner()

                # Generate plan
                planner_response = self.planner.plan(current_problem)
                self.plan = planner_response.plan

            # Emit planning complete event
            self._emit_event(StepEvent(
                event_type="planning_complete",
                step_number=0,
                data={"steps": len(self.plan.steps)}
            ))

            # Record plan to plan directory
            self.history.save_plan_data(
                self.session_id,
                raw_response=planner_response.raw_response or None,
                parsed_plan={
                    "steps": [
                        {
                            "number": s.number,
                            "goal": s.goal,
                            "inputs": s.expected_inputs,
                            "outputs": s.expected_outputs,
                            "depends_on": s.depends_on,
                            "task_type": s.task_type.value if s.task_type else None,
                            "role_id": s.role_id,
                        }
                        for s in self.plan.steps
                    ],
                },
                reasoning=planner_response.reasoning or None,
                iteration=replan_attempt,
            )

            # Request approval if required
            if self.session_config.require_approval:
                # Use display_problem for UI (just feedback on replan, full problem initially)
                approval = self._request_approval(display_problem, planner_response)

                if approval.decision == PlanApproval.REJECT:
                    # User rejected the plan
                    self.history.save_plan_data(
                        self.session_id,
                        approval_decision="rejected",
                        user_feedback=approval.reason,
                        iteration=replan_attempt,
                    )
                    self.datastore.set_session_meta("status", "rejected")
                    self.history.complete_session(self.session_id, status="rejected")
                    return {
                        "success": False,
                        "rejected": True,
                        "plan": self.plan,
                        "reason": approval.reason,
                        "message": "Plan was rejected by user.",
                    }

                elif approval.decision == PlanApproval.COMMAND:
                    # User entered a slash command - pass back to REPL
                    return {
                        "success": False,
                        "command": approval.command,
                        "message": "Slash command entered during approval.",
                    }

                elif approval.decision == PlanApproval.SUGGEST:
                    # User wants changes - check if we can skip replanning
                    suggestion_text = (approval.suggestion or "").strip()
                    has_edited_steps = bool(approval.edited_steps)
                    has_meaningful_feedback = bool(suggestion_text) and suggestion_text not in ("", "Edited plan")

                    # Record suggestion
                    self.history.save_plan_data(
                        self.session_id,
                        approval_decision="suggest",
                        user_feedback=suggestion_text or None,
                        edited_steps=approval.edited_steps,
                        iteration=replan_attempt,
                    )

                    # If user edited steps but provided no additional feedback, use edited plan directly
                    if has_edited_steps and not has_meaningful_feedback:
                        logger.info("[REPLAN] User edited steps with no feedback - using edited plan directly")

                        # Log the revision (the edited plan itself)
                        edited_summary = "; ".join(f"{s['number']}. {s['goal'][:50]}" for s in approval.edited_steps)
                        self.history.log_user_input(self.session_id, f"[Edited plan] {edited_summary}", "revision")

                        # Build Plan directly from edited steps
                        self.plan = self._build_plan_from_edited_steps(problem, approval.edited_steps)

                        # Create a synthetic planner_response for the plan_ready event
                        planner_response = PlannerResponse(
                            plan=self.plan,
                            reasoning="User-edited plan (approved without replanning)",
                        )

                        # Emit event for UI consistency
                        self._emit_event(StepEvent(
                            event_type="planning_complete",
                            step_number=0,
                            data={"steps": len(self.plan.steps), "edited": True}
                        ))

                        break  # Skip replanning, proceed to execution

                    # User has meaningful feedback - replan with feedback
                    replan_attempt += 1
                    if replan_attempt > self.session_config.max_replan_attempts:
                        self.datastore.set_session_meta("status", "max_replans_exceeded")
                        self.history.complete_session(self.session_id, status="failed")
                        return {
                            "success": False,
                            "plan": self.plan,
                            "error": f"Maximum replan attempts ({self.session_config.max_replan_attempts}) exceeded.",
                        }

                    # Log the revision
                    self.history.log_user_input(self.session_id, suggestion_text, "revision")

                    # Emit replan event
                    self._emit_event(StepEvent(
                        event_type="replanning",
                        step_number=0,
                        data={
                            "attempt": replan_attempt,
                            "feedback": suggestion_text,
                        }
                    ))

                    # Build replan prompt with original query + edited plan structure
                    # The edited plan is what the user approved/modified
                    if has_edited_steps:
                        edited_plan_text = "\n".join(
                            f"{step['number']}. {step['goal']}"
                            for step in approval.edited_steps
                        )
                        current_problem = f"""{problem}

**Requested plan structure (follow this exactly):**
{edited_plan_text}

**User notes:** {suggestion_text}"""
                    else:
                        # No edited steps provided - use original problem + feedback
                        current_problem = f"{problem}\n\n**User Revision (takes precedence):** {suggestion_text}"

                    display_problem = suggestion_text
                    continue  # Go back to planning

                # APPROVE - record and apply any edits/deletions, then proceed
                self.history.save_plan_data(
                    self.session_id,
                    approval_decision="approved",
                    edited_steps=approval.edited_steps if approval.edited_steps else None,
                    iteration=replan_attempt,
                )
                if approval.edited_steps:
                    self.plan = self._build_plan_from_edited_steps(problem, approval.edited_steps)
                elif approval.deleted_steps:
                    deleted_set = set(approval.deleted_steps)
                    self.plan.steps = [s for s in self.plan.steps if s.number not in deleted_set]
                    for i, step in enumerate(self.plan.steps, 1):
                        step.number = i
                break
            else:
                # No approval required - auto-approved
                self.history.save_plan_data(
                    self.session_id,
                    approval_decision="auto_approved",
                    iteration=replan_attempt,
                )
                break

        # Save plan to datastore (for UI restoration)
        self.datastore.set_session_meta("status", "executing")
        self.datastore.set_session_meta("mode", "exploratory")
        for step in self.plan.steps:
            self.datastore.save_plan_step(
                step_number=step.number,
                goal=step.goal,
                expected_inputs=step.expected_inputs,
                expected_outputs=step.expected_outputs,
                status="pending",
            )

        # Emit plan_ready event BEFORE execution starts
        self._emit_event(StepEvent(
            event_type="plan_ready",
            step_number=0,
            data={
                "steps": [
                    {"number": s.number, "goal": s.goal, "depends_on": s.depends_on, "role_id": s.role_id}
                    for s in self.plan.steps
                ],
                "reasoning": planner_response.reasoning,
                "is_followup": False,
            }
        ))

        # Materialize facts table before execution starts
        self._materialize_facts_table()

        # Execute steps in parallel waves based on dependencies
        # Phase 4: Reset cancellation state before starting execution
        self.reset_cancellation()
        all_results = []

        # Debug: Log full plan structure before computing waves
        logger.debug(f"[EXECUTION] Plan has {len(self.plan.steps)} steps:")
        for step in self.plan.steps:
            logger.debug(f"[EXECUTION]   Step {step.number}: status={step.status.value}, "
                         f"depends_on={step.depends_on}, goal='{step.goal[:60]}...'")

        execution_waves = self.plan.get_execution_order()
        logger.debug(f"[EXECUTION] execution_waves: {execution_waves}, total steps: {len(self.plan.steps)}")
        cancelled = False

        for wave_num, wave_step_nums in enumerate(execution_waves):
            logger.debug(f"[EXECUTION] Starting wave {wave_num + 1}, steps: {wave_step_nums}")
            # Phase 4: Check for cancellation before starting each wave
            if self.is_cancelled():
                cancelled = True
                self._emit_event(StepEvent(
                    event_type="execution_cancelled",
                    step_number=0,
                    data={
                        "message": "Execution cancelled between waves",
                        "wave": wave_num,
                        "completed_steps": len(all_results),
                    }
                ))
                break

            # Get steps for this wave
            wave_steps = [self.plan.get_step(num) for num in wave_step_nums]
            wave_steps = [s for s in wave_steps if s is not None]

            # Execute all steps in this wave in parallel
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=len(wave_steps)) as executor:
                # Submit all steps in wave
                future_to_step = {}
                for step in wave_steps:
                    # Phase 4: Check for cancellation before starting each step
                    if self.is_cancelled():
                        cancelled = True
                        break

                    step.status = StepStatus.RUNNING
                    self.datastore.update_plan_step(step.number, status="running")
                    self._emit_event(StepEvent(
                        event_type="wave_step_start",
                        step_number=step.number,
                        data={"wave": wave_num + 1, "goal": step.goal}
                    ))
                    future = executor.submit(self._execute_step, step)
                    future_to_step[future] = step

                if cancelled:
                    # Cancel any pending futures
                    for future in future_to_step:
                        future.cancel()
                    break

                # Collect results as they complete
                wave_failed = False
                for future in concurrent.futures.as_completed(future_to_step):
                    # Phase 4: Check for cancellation while collecting results
                    # Note: We still collect completed results even if cancelled
                    step = future_to_step[future]
                    try:
                        result = future.result()
                    except concurrent.futures.CancelledError:
                        # Step was cancelled before it started
                        result = StepResult(
                            success=False,
                            stdout="",
                            error="Step cancelled",
                            attempts=0,
                        )
                        cancelled = True
                    except Exception as e:
                        logger.error(f"[EXECUTION] Step {step.number} raised unhandled exception: {e}", exc_info=True)
                        result = StepResult(
                            success=False,
                            stdout="",
                            error=str(e),
                            attempts=1,
                        )

                    logger.debug(f"[EXECUTION] Step {step.number} result: success={result.success}, attempts={result.attempts}, error={result.error[:200] if result.error else 'none'}")
                    if result.success:
                        self.plan.mark_step_completed(step.number, result)
                        self.scratchpad.add_step_result(
                            step_number=step.number,
                            goal=step.goal,
                            result=result.stdout,
                            tables_created=result.tables_created,
                        )
                        if self.datastore:
                            self.datastore.add_scratchpad_entry(
                                step_number=step.number,
                                goal=step.goal,
                                narrative=result.stdout,
                                tables_created=result.tables_created,
                                code=result.code,
                            )
                            self.datastore.update_plan_step(
                                step.number,
                                status="completed",
                                code=step.code,
                                attempts=result.attempts,
                                duration_ms=result.duration_ms,
                            )
                        all_results.append(result)
                    else:
                        self.plan.mark_step_failed(step.number, result)
                        if self.datastore:
                            self.datastore.update_plan_step(
                                step.number,
                                status="failed" if not cancelled else "cancelled",
                                code=step.code,
                                error=result.error,
                                attempts=result.attempts,
                                duration_ms=result.duration_ms,
                            )
                        if not cancelled:  # Only mark as failed if not cancelled
                            wave_failed = True
                        all_results.append(result)

                if cancelled:
                    logger.debug(f"[EXECUTION] Wave {wave_num + 1} cancelled, breaking out of loop")
                    break

                logger.debug(f"[EXECUTION] Wave {wave_num + 1} completed: "
                             f"wave_failed={wave_failed}, all_results={len(all_results)}, "
                             f"completed_steps={self.plan.completed_steps}")

                # If any step in wave failed, stop execution
                if wave_failed:
                    self.datastore.set_session_meta("status", "failed")
                    failed_result = next(r for r in all_results if not r.success)
                    self.history.record_query(
                        session_id=self.session_id,
                        question=problem,
                        success=False,
                        attempts=failed_result.attempts,
                        duration_ms=failed_result.duration_ms,
                        error=failed_result.error,
                    )
                    self.history.complete_session(self.session_id, status="failed")
                    return {
                        "success": False,
                        "plan": self.plan,
                        "error": failed_result.error,
                        "completed_steps": self.plan.completed_steps,
                    }

        # Log exit reason
        logger.debug(f"[EXECUTION] Wave loop finished: cancelled={cancelled}, "
                     f"all_results={len(all_results)}, completed_steps={self.plan.completed_steps}")

        # Phase 4: Handle cancellation - return with completed results preserved
        if cancelled:
            self.datastore.set_session_meta("status", "cancelled")
            self.history.complete_session(self.session_id, status="cancelled")

            # Combine output from completed steps
            completed_output = ""
            if all_results:
                completed_output = "\n\n".join([
                    f"Step {i+1}: {self.plan.steps[i].goal}\n{r.stdout}"
                    for i, r in enumerate(all_results) if r.success
                ])

            # Process any queued intents
            queued_results = self.process_queued_intents()

            return {
                "success": False,
                "cancelled": True,
                "plan": self.plan,
                "completed_steps": self.plan.completed_steps,
                "partial_output": completed_output,
                "queued_intent_results": queued_results,
                "message": f"Execution cancelled. {len(self.plan.completed_steps)} step(s) completed.",
            }

        # Record successful completion
        total_duration = sum(r.duration_ms for r in all_results)
        total_attempts = sum(r.attempts for r in all_results)

        # Combine all step outputs
        combined_output = "\n\n".join([
            f"Step {i+1}: {self.plan.steps[i].goal}\n{r.stdout}"
            for i, r in enumerate(all_results)
        ])

        # Auto-publish final step artifacts (they appear in artifacts panel)
        # Find the last step that actually created tables (may not be the last step if it's just a summary)
        if all_results and self.registry:
            # Collect tables from last 2 steps that created tables (final outputs often span multiple steps)
            final_tables = []
            steps_with_tables = 0
            for result in reversed(all_results):
                if result.success and result.tables_created:
                    final_tables.extend(result.tables_created)
                    steps_with_tables += 1
                    if steps_with_tables >= 2:
                        break

            # Filter to most important tables (avoid publishing every intermediate table)
            # Prioritize: tables with "final", "report", "result", "recommendation" in name
            important_keywords = ("final", "report", "result", "recommendation", "summary", "output")
            important_tables = [t for t in final_tables if any(kw in t.lower() for kw in important_keywords)]

            # If no important tables found, use tables from last step only (max 5)
            if not important_tables:
                for result in reversed(all_results):
                    if result.success and result.tables_created:
                        important_tables = result.tables_created[:5]
                        break

            # Limit to max 8 published tables to avoid clutter
            tables_to_publish = important_tables[:8] if important_tables else []

            # Also find markdown artifacts from final step (highest priority for View Result)
            final_artifacts = []
            if self.datastore:
                all_artifacts = self.datastore.list_artifacts()
                last_step_num = len(self.plan.steps) if self.plan else 0
                # Get markdown artifacts from last 2 steps
                final_artifacts = [
                    a["name"] for a in all_artifacts
                    if a.get("step_number", 0) >= last_step_num - 1
                    and a.get("type") in ("markdown", "md", "html")
                ]

            if tables_to_publish or final_artifacts:
                self.registry.mark_final_step(
                    user_id=self.user_id,
                    session_id=self.session_id,
                    table_names=tables_to_publish if tables_to_publish else None,
                    artifact_names=final_artifacts if final_artifacts else None,
                )
                logger.debug(f"Auto-published final step: tables={tables_to_publish}, artifacts={final_artifacts}")

        # Note: Facts created during role-scoped steps are tagged with role_id
        # for provenance but remain globally accessible. No promotion needed.

        # Emit raw results first (so user can see them immediately)
        self._emit_event(StepEvent(
            event_type="raw_results_ready",
            step_number=0,
            data={"output": combined_output}
        ))

        # Check for created artifacts (to mention in synthesis)
        from constat.visualization.output import peek_pending_outputs
        pending_artifacts = peek_pending_outputs()

        # Check if insights are enabled (config or per-query brief detection via LLM)
        skip_insights = not self.session_config.enable_insights or analysis.wants_brief
        suggestions = []  # Initialize for brief mode (no suggestions)

        if skip_insights:
            # Use raw output as final answer
            final_answer = combined_output
            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer, "brief": True}
            ))
        else:
            # Synthesize final answer from step results
            self._emit_event(StepEvent(
                event_type="synthesizing",
                step_number=0,
                data={"message": "Synthesizing final answer..."}
            ))

            final_answer = self._synthesize_answer(problem, combined_output, pending_artifacts)

            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer}
            ))

            # Run facts extraction and suggestions generation in parallel
            tables = self.datastore.list_tables() if self.datastore else []
            response_facts, suggestions = self._run_post_synthesis_parallel(
                problem, final_answer, tables
            )

            if response_facts:
                self._emit_event(StepEvent(
                    event_type="facts_extracted",
                    step_number=0,
                    data={
                        "facts": [f.to_dict() for f in response_facts],
                        "source": "response",
                    }
                ))

            if suggestions:
                self._emit_event(StepEvent(
                    event_type="suggestions_ready",
                    step_number=0,
                    data={"suggestions": suggestions}
                ))

        self.history.record_query(
            session_id=self.session_id,
            question=problem,
            success=True,
            attempts=total_attempts,
            duration_ms=total_duration,
            answer=final_answer,
        )
        self.history.complete_session(self.session_id, status="completed")

        # Mark session as completed in datastore (for UI restoration)
        if self.datastore:
            self.datastore.set_session_meta("status", "completed")

        # Auto-compact if context is too large
        self._auto_compact_if_needed()

        # Ensure execution history is available as a queryable table
        if self.datastore:
            self.datastore.ensure_execution_history_table()

        return {
            "success": True,
            "plan": self.plan,
            "results": all_results,
            "output": combined_output,
            "final_answer": final_answer,
            "suggestions": suggestions,
            "scratchpad": self.scratchpad.to_markdown(),
            "datastore_tables": self.datastore.list_tables() if self.datastore else [],
            "datastore_path": str(self.datastore.db_path) if self.datastore and self.datastore.db_path else None,
        }

    def resume(self, session_id: str) -> bool:
        """
        Resume a previous session, loading its datastore and context.

        Args:
            session_id: The session ID to resume

        Returns:
            True if successfully resumed, False if session not found
        """
        # Check if session exists
        session_detail = self.history.get_session(session_id)
        if not session_detail:
            return False

        self.session_id = session_id

        # Load the datastore (contains tables, state, scratchpad, artifacts)
        session_dir = self.history._session_dir(session_id)
        datastore_path = session_dir / "datastore.duckdb"
        tables_dir = session_dir / "tables"

        # Create underlying datastore
        if datastore_path.exists():
            underlying_datastore = DataStore(db_path=datastore_path)
        else:
            # No datastore file - create empty one
            underlying_datastore = DataStore(db_path=datastore_path)

        # Wrap with registry-aware datastore
        self.datastore = RegistryAwareDataStore(
            datastore=underlying_datastore,
            registry=self.registry,
            user_id=self.user_id,
            session_id=session_id,
            tables_dir=tables_dir,
        )

        if datastore_path.exists():
            # Rebuild scratchpad from datastore
            scratchpad_entries = self.datastore.get_scratchpad()
            if scratchpad_entries:
                # Get the original problem from the first query
                if session_detail.queries:
                    initial_context = f"Problem: {session_detail.queries[0].question}"
                else:
                    initial_context = ""
                self.scratchpad = Scratchpad(initial_context=initial_context)

                # Add each step result
                for entry in scratchpad_entries:
                    self.scratchpad.add_step_result(
                        step_number=entry["step_number"],
                        goal=entry["goal"],
                        result=entry["narrative"],
                        tables_created=entry.get("tables_created", []),
                    )

        # Update fact resolver's datastore reference (for storing large facts as tables)
        self.fact_resolver._datastore = self.datastore

        return True

    def classify_follow_up_intent(self, user_text: str) -> dict:
        """
        Classify the intent of a follow-up message.

        This helps determine how to handle user input that could be:
        - Providing facts (e.g., "There were 1 million people")
        - Revising the request (e.g., "Use $50k threshold instead")
        - Making a new request (e.g., "Show me sales by region")
        - A combination of the above

        Args:
            user_text: The user's follow-up message

        Returns:
            Dict with:
                - intent: PRIMARY intent (PROVIDE_FACTS, REVISE, NEW_REQUEST, MIXED)
                - facts: List of any facts detected
                - revision: Description of any revision detected
                - new_request: The new request if detected
        """
        # Check for unresolved facts
        unresolved = self.fact_resolver.get_unresolved_facts()
        unresolved_names = [f.name for f in unresolved]

        prompt = f"""Analyze this user follow-up message and classify its intent.

User message: "{user_text}"

Context:
- There are {len(unresolved)} unresolved facts: {unresolved_names if unresolved else 'none'}

Classify the PRIMARY intent as one of:
- PROVIDE_FACTS: User is providing factual information (numbers, values, definitions)
- REVISE: User wants to modify/refine the previous request
- NEW_REQUEST: User is making an unrelated new request
- MIXED: Combination of the above

Also extract any facts, revisions, or new requests detected.

Respond in this exact format:
INTENT: <one of PROVIDE_FACTS, REVISE, NEW_REQUEST, MIXED>
FACTS: <comma-separated list of fact=value pairs, or NONE>
REVISION: <description of revision, or NONE>
NEW_REQUEST: <the new request, or NONE>
"""

        try:
            response = self.llm.generate(
                system="You are an intent classifier. Analyze user messages precisely.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            result = {
                "intent": "NEW_REQUEST",  # Default
                "facts": [],
                "revision": None,
                "new_request": None,
            }

            for line in response.split("\n"):
                line = line.strip()
                if line.startswith("INTENT:"):
                    intent = line.split(":", 1)[1].strip().upper()
                    if intent in ("PROVIDE_FACTS", "REVISE", "NEW_REQUEST", "MIXED"):
                        result["intent"] = intent
                elif line.startswith("FACTS:"):
                    facts_str = line.split(":", 1)[1].strip()
                    if facts_str != "NONE":
                        result["facts"] = [f.strip() for f in facts_str.split(",")]
                elif line.startswith("REVISION:"):
                    rev = line.split(":", 1)[1].strip()
                    if rev != "NONE":
                        result["revision"] = rev
                elif line.startswith("NEW_REQUEST:"):
                    req = line.split(":", 1)[1].strip()
                    if req != "NONE":
                        result["new_request"] = req

            return result

        except Exception as e:
            # Default to treating as new request
            logger.debug(f"Conversational intent classification failed: {e}")
            return {
                "intent": "NEW_REQUEST",
                "facts": [],
                "revision": None,
                "new_request": user_text,
            }

    def _classify_premise_response(self, user_response: str, fact_name: str, question: str) -> dict:
        """
        Classify a user's response to a premise clarification question.

        Determines if the response is:
        - VALUE: An actual value/answer (e.g., "5% for rating 4, 3% for rating 3")
        - STEER: Guidance on where/how to find the answer (e.g., "Look in the HR policy document")

        Args:
            user_response: The user's response text
            fact_name: The name of the fact being resolved
            question: The original clarification question asked

        Returns:
            Dict with:
                - type: "VALUE" or "STEER"
                - value: The parsed value if VALUE type
                - steer: The guidance/direction if STEER type
        """
        prompt = f"""Classify this user response to a data clarification question.

Question asked: "{question}"
Fact needed: {fact_name}
User response: "{user_response}"

Is this response:
A) VALUE - A direct answer providing the actual data/value/information requested
   Examples: "5%", "use 10000 as threshold", "rating 5 gets 10%, rating 4 gets 6%"

B) STEER - Guidance on WHERE or HOW to find the answer (not the answer itself)
   Examples: "Look in the HR policy", "Check the business_rules document",
   "Use the performance review guidelines", "It should be in the config"

Respond in this exact format:
TYPE: <VALUE or STEER>
CONTENT: <the value if VALUE, or the guidance/direction if STEER>
"""

        try:
            response = self.llm.generate(
                system="You classify user responses precisely. Distinguish between direct answers and guidance about where to find answers.",
                user_message=prompt,
                max_tokens=self.router.max_output_tokens,
            )

            result = {"type": "VALUE", "value": user_response, "steer": None}

            for line in response.split("\n"):
                line = line.strip()
                if line.startswith("TYPE:"):
                    resp_type = line.split(":", 1)[1].strip().upper()
                    if resp_type == "STEER":
                        result["type"] = "STEER"
                elif line.startswith("CONTENT:"):
                    content = line.split(":", 1)[1].strip()
                    if result["type"] == "VALUE":
                        result["value"] = content
                    else:
                        result["steer"] = content

            return result

        except Exception as e:
            # Default to treating as value
            logger.debug(f"Premise response classification failed: {e}")
            return {"type": "VALUE", "value": user_response, "steer": None}

    def follow_up(self, question: str, auto_classify: bool = True) -> dict:
        """
        Ask a follow-up question that builds on the current session's context.

        The follow-up has access to all tables and state from previous steps.
        If there are unresolved facts, the system will first try to extract
        facts from the user's message.

        Automatically detects if the question suggests auditable mode (verify,
        validate, etc.) and uses the fact resolver for formal verification.

        Args:
            question: The follow-up question
            auto_classify: If True, classify intent and handle accordingly

        Returns:
            Dict with plan, results, and summary (same format as solve())
        """
        if not self.session_id:
            raise ValueError("No active session. Call solve() or resume() first.")

        if not self.datastore:
            raise ValueError("No datastore available. Session may not have been properly initialized.")

        # Fast path: Handle slash commands directly via command registry
        if question.strip().startswith("/"):
            return self._handle_slash_command(question.strip())

        # Log the follow-up query
        self.history.log_user_input(self.session_id, question, "followup")

        # Detect corrections/hints and save as learnings
        from constat.api.detection.correction import detect_nl_correction
        nl_correction = detect_nl_correction(question)
        if nl_correction.detected:
            self._save_correction_as_learning(question)

        # Fast path: check if this is a simple "show me X" request for existing data
        # This avoids expensive LLM classification for simple data lookups
        show_result = self._try_show_existing_data(question)
        if show_result:
            return show_result

        # Get previous problem for follow-up context
        previous_problem = self.datastore.get_session_meta("problem")

        # Store follow-up question for /prove command
        follow_ups_json = self.datastore.get_session_meta("follow_ups")
        try:
            follow_ups = json.loads(follow_ups_json) if follow_ups_json else []
        except json.JSONDecodeError:
            follow_ups = []
        follow_ups.append(question)
        self.datastore.set_session_meta("follow_ups", json.dumps(follow_ups))
        logger.debug(f"[follow_up] Stored follow-up #{len(follow_ups)}: {question[:50]}...")

        # Use LLM to analyze the question and detect intents (with follow-up context)
        # This single call determines intent, facts, AND execution mode
        analysis = self._analyze_question(question, previous_problem=previous_problem)

        # Check for validation constraint addition (before REDO check)
        intent_names = [i.intent for i in analysis.intents]
        _validation_keywords = ("validate", "verify", "ensure", "assert", "check that", "confirm that")
        session_mode = self.datastore.get_session_meta("mode") if self.datastore else "exploratory"
        if session_mode == "auditable" and any(kw in question.lower() for kw in _validation_keywords):
            # Extract validation from follow-up and re-run proof with it
            new_validations = self._extract_user_validations(question, [])
            if new_validations:
                for v in new_validations:
                    self.add_user_validation(v['label'], v['sql'])
                return self.prove_conversation(guidance=f"Re-run with added validation: {question}")

        # Check for REDO intent — apply fact modifications and re-execute
        if "REDO" in intent_names:
            # Apply any fact modifications first
            for mod in analysis.fact_modifications:
                self.fact_resolver.add_user_fact(
                    fact_name=mod["fact_name"],
                    value=mod["new_value"],
                    reasoning=f"User correction: {question}",
                )
                logger.debug(f"[follow_up] Applied fact modification: {mod['fact_name']}={mod['new_value']}")
            self.fact_resolver.add_user_facts_from_text(question)

            # Determine mode: use LLM recommendation, fall back to session's current mode
            session_mode = self.datastore.get_session_meta("mode") if self.datastore else "exploratory"
            use_proof = analysis.recommended_mode == "PROOF" or session_mode == "auditable"
            logger.info(f"[follow_up] REDO intent detected, mode={session_mode}, recommended={analysis.recommended_mode}, use_proof={use_proof}")

            if use_proof:
                return self.prove_conversation(guidance=question)
            else:
                # Re-run exploratory: get original problem, re-solve with updated facts
                original_problem = self.datastore.get_session_meta("problem") if self.datastore else question
                return self.solve(original_problem)

        # Check for ambiguity and request clarification if needed
        if self.session_config.ask_clarifications and self._clarification_callback:
            existing_tables = self.datastore.list_tables()
            clarification_request = self._detect_ambiguity(question, session_tables=existing_tables)
            if clarification_request:
                enhanced_question = self._request_clarification(clarification_request)
                if enhanced_question:
                    question = enhanced_question
                    # Re-analyze with clarified question
                    logger.debug("[follow_up] Question clarified, re-analyzing...")
                    analysis = self._analyze_question(question, previous_problem=previous_problem)

        # All follow-ups use exploratory mode (planning + execution)
        # Use /prove command to generate auditable proofs when needed
        # Check for unresolved facts and try to extract facts from user message
        unresolved = self.fact_resolver.get_unresolved_facts()
        extracted_facts = []

        if auto_classify and (unresolved or "=" in question or any(c.isdigit() for c in question)):
            # Try to extract facts from the message
            extracted_facts = self.fact_resolver.add_user_facts_from_text(question)

            if extracted_facts:
                # Clear unresolved status to allow re-resolution
                self.fact_resolver.clear_unresolved()

        # Get context from previous work
        existing_tables = self.datastore.list_tables()
        existing_state = self.datastore.get_all_state()
        scratchpad_context = self.datastore.get_scratchpad_as_markdown()

        # Ensure execution history is available as a queryable table
        # This includes step goals, code, and outputs
        self.datastore.ensure_execution_history_table()
        existing_tables = self.datastore.list_tables()  # Refresh after adding history table

        # Calculate next step number
        existing_scratchpad = self.datastore.get_scratchpad()
        next_step_number = max((e["step_number"] for e in existing_scratchpad), default=0) + 1

        # Generate a plan for the follow-up, providing context
        existing_tables_list = (
            chr(10).join(f'  - `{t["name"]}`: {t.get("row_count", "?")} rows (step {t.get("step_number", "?")})' for t in existing_tables)
            if existing_tables else '(none)'
        )
        first_table_name = existing_tables[0]["name"] if existing_tables else "final_answer"
        context_prompt = load_prompt("followup_context.md").format(
            scratchpad_context=scratchpad_context,
            existing_tables_list=existing_tables_list,
            existing_state=existing_state if existing_state else '(none)',
            first_table_name=first_table_name,
            question=question,
        )
        # Emit planning start event
        self._emit_event(StepEvent(
            event_type="planning_start",
            step_number=0,
            data={"message": "Planning follow-up analysis..."}
        ))

        # Sync user facts and roles to planner before generating plan
        self._sync_user_facts_to_planner()
        self._sync_available_roles_to_planner()

        # Generate plan for follow-up
        planner_response = self.planner.plan(context_prompt)
        follow_up_plan = planner_response.plan

        # Validate: ensure "enhance" plans include a step that updates the source table
        follow_up_plan = self._ensure_enhance_updates_source(
            question, follow_up_plan, existing_tables,
        )

        # Emit planning complete event
        self._emit_event(StepEvent(
            event_type="planning_complete",
            step_number=0,
            data={"steps": len(follow_up_plan.steps)}
        ))

        # Renumber steps to continue from where we left off
        for i, step in enumerate(follow_up_plan.steps):
            step.number = next_step_number + i

        # Emit plan_ready event for display
        self._emit_event(StepEvent(
            event_type="plan_ready",
            step_number=0,
            data={
                "steps": [
                    {"number": s.number, "goal": s.goal, "depends_on": s.depends_on, "role_id": s.role_id}
                    for s in follow_up_plan.steps
                ],
                "reasoning": planner_response.reasoning,
                "is_followup": True,
            }
        ))

        # Request approval if required (same as solve())
        if self.session_config.require_approval:
            approval = self._request_approval(question, planner_response)

            if approval.decision == PlanApproval.REJECT:
                return {
                    "success": False,
                    "rejected": True,
                    "plan": follow_up_plan,
                    "reason": approval.reason,
                    "message": "Follow-up plan was rejected by user.",
                }

            elif approval.decision == PlanApproval.COMMAND:
                # User entered a slash command - pass back to REPL
                return {
                    "success": False,
                    "command": approval.command,
                    "message": "Slash command entered during approval.",
                }

            elif approval.decision == PlanApproval.SUGGEST:
                # Check if user edited steps but provided no meaningful feedback
                suggestion_text = (approval.suggestion or "").strip()
                has_edited_steps = bool(approval.edited_steps)
                has_meaningful_feedback = bool(suggestion_text) and suggestion_text not in ("", "Edited plan")

                if has_edited_steps and not has_meaningful_feedback:
                    # User edited steps directly — use edited plan without replanning
                    logger.info("[follow_up REPLAN] User edited steps with no feedback - using edited plan directly")
                    follow_up_plan = self._build_plan_from_edited_steps(
                        question, approval.edited_steps, start_number=next_step_number
                    )
                else:
                    # User has meaningful feedback — replan
                    if has_edited_steps:
                        edited_plan_text = "\n".join(
                            f"{step['number']}. {step['goal']}" for step in approval.edited_steps
                        )
                        context_prompt_with_feedback = f"""{context_prompt}

**Requested plan structure (follow this exactly):**
{edited_plan_text}

**User notes:** {suggestion_text}"""
                    else:
                        context_prompt_with_feedback = f"""{context_prompt}

User feedback: {suggestion_text}
"""

                    # Emit replanning event
                    self._emit_event(StepEvent(
                        event_type="replanning",
                        step_number=0,
                        data={"feedback": suggestion_text}
                    ))

                    self._sync_user_facts_to_planner()
                    self._sync_available_roles_to_planner()
                    planner_response = self.planner.plan(context_prompt_with_feedback)
                    follow_up_plan = planner_response.plan

                    # Validate: ensure "enhance" plans include a step that updates the source table
                    follow_up_plan = self._ensure_enhance_updates_source(
                        question, follow_up_plan, existing_tables,
                    )

                    # Renumber steps to continue from where we left off
                    for i, step in enumerate(follow_up_plan.steps):
                        step.number = next_step_number + i

                    # Emit updated plan
                    self._emit_event(StepEvent(
                        event_type="plan_ready",
                        step_number=0,
                        data={
                            "steps": [
                                {"number": s.number, "goal": s.goal, "depends_on": s.depends_on, "role_id": s.role_id}
                                for s in follow_up_plan.steps
                            ],
                            "reasoning": planner_response.reasoning,
                            "is_followup": True,
                        }
                    ))

                    # Request approval again
                    approval = self._request_approval(question, planner_response)
                    if approval.decision == PlanApproval.REJECT:
                        return {
                            "success": False,
                            "rejected": True,
                            "plan": follow_up_plan,
                            "reason": approval.reason,
                            "message": "Follow-up plan was rejected by user.",
                        }
                    elif approval.decision == PlanApproval.COMMAND:
                        return {
                            "success": False,
                            "command": approval.command,
                            "message": "Slash command entered during approval.",
                        }

            # APPROVE — apply any edits/deletions to the follow-up plan
            if approval.edited_steps:
                follow_up_plan = self._build_plan_from_edited_steps(
                    question, approval.edited_steps, start_number=next_step_number
                )
            elif approval.deleted_steps:
                deleted_set = set(approval.deleted_steps)
                follow_up_plan.steps = [s for s in follow_up_plan.steps if s.number not in deleted_set]
                # Renumber remaining steps sequentially from next_step_number
                for i, step in enumerate(follow_up_plan.steps):
                    step.number = next_step_number + i

        # Materialize facts table before execution starts
        self._materialize_facts_table()

        # Execute each step
        # Phase 4: Reset cancellation state before starting execution
        self.reset_cancellation()
        all_results = []
        cancelled = False

        for step in follow_up_plan.steps:
            # Phase 4: Check for cancellation before starting each step
            if self.is_cancelled():
                cancelled = True
                self._emit_event(StepEvent(
                    event_type="execution_cancelled",
                    step_number=step.number,
                    data={
                        "message": "Execution cancelled",
                        "completed_steps": len([r for r in all_results if r.success]),
                    }
                ))
                break

            step.status = StepStatus.RUNNING

            result = self._execute_step(step)

            if result.success:
                follow_up_plan.mark_step_completed(step.number, result)
                self.scratchpad.add_step_result(
                    step_number=step.number,
                    goal=step.goal,
                    result=result.stdout,
                    tables_created=result.tables_created,
                )
                if self.datastore:
                    self.datastore.add_scratchpad_entry(
                        step_number=step.number,
                        goal=step.goal,
                        narrative=result.stdout,
                        tables_created=result.tables_created,
                        code=result.code,
                    )
            else:
                follow_up_plan.mark_step_failed(step.number, result)
                self.history.record_query(
                    session_id=self.session_id,
                    question=question,
                    success=False,
                    attempts=result.attempts,
                    duration_ms=result.duration_ms,
                    error=result.error,
                )
                return {
                    "success": False,
                    "plan": follow_up_plan,
                    "error": result.error,
                    "completed_steps": follow_up_plan.completed_steps,
                }

            all_results.append(result)

        # Phase 4: Handle cancellation - return with completed results preserved
        if cancelled:
            # Combine output from completed steps
            completed_output = ""
            if all_results:
                completed_output = "\n\n".join([
                    f"Step {step.number}: {step.goal}\n{r.stdout}"
                    for step, r in zip(follow_up_plan.steps, all_results) if r.success
                ])

            # Process any queued intents
            queued_results = self.process_queued_intents()

            return {
                "success": False,
                "cancelled": True,
                "plan": follow_up_plan,
                "completed_steps": follow_up_plan.completed_steps,
                "partial_output": completed_output,
                "queued_intent_results": queued_results,
                "message": f"Execution cancelled. {len(follow_up_plan.completed_steps)} step(s) completed.",
            }

        # Record successful follow-up
        total_duration = sum(r.duration_ms for r in all_results)
        total_attempts = sum(r.attempts for r in all_results)

        combined_output = "\n\n".join([
            f"Step {step.number}: {step.goal}\n{r.stdout}"
            for step, r in zip(follow_up_plan.steps, all_results)
        ])

        # Auto-publish final step artifacts (they appear in artifacts panel)
        # Find the last step that actually created tables (may not be the last step if it's just a summary)
        if all_results and self.registry:
            # Collect tables from last 2 steps that created tables
            final_tables = []
            steps_with_tables = 0
            for result in reversed(all_results):
                if result.success and result.tables_created:
                    final_tables.extend(result.tables_created)
                    steps_with_tables += 1
                    if steps_with_tables >= 2:
                        break

            # Filter to most important tables
            important_keywords = ("final", "report", "result", "recommendation", "summary", "output")
            important_tables = [t for t in final_tables if any(kw in t.lower() for kw in important_keywords)]

            if not important_tables:
                for result in reversed(all_results):
                    if result.success and result.tables_created:
                        important_tables = result.tables_created[:5]
                        break

            tables_to_publish = important_tables[:8] if important_tables else []

            if tables_to_publish:
                self.registry.mark_final_step(
                    user_id=self.user_id,
                    session_id=self.session_id,
                    table_names=tables_to_publish,
                )
                logger.debug(f"Auto-published follow-up final step tables: {tables_to_publish}")

        # Emit raw results first (so user can see them immediately)
        self._emit_event(StepEvent(
            event_type="raw_results_ready",
            step_number=0,
            data={"output": combined_output}
        ))

        # Analyze follow-up question for brief output preference (LLM-based, not keywords)
        follow_up_analysis = self._analyze_question(question)

        # Check for created artifacts (to mention in synthesis)
        from constat.visualization.output import peek_pending_outputs
        pending_artifacts = peek_pending_outputs()  # Peek without clearing

        # Check if insights are enabled (config or per-query brief detection via LLM)
        skip_insights = (
            not self.session_config.enable_insights
            or follow_up_analysis.wants_brief
        )
        suggestions = []  # Initialize for brief mode (no suggestions)

        if skip_insights:
            # Use raw output as final answer
            final_answer = combined_output
            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer, "brief": True}
            ))
        else:
            # Synthesize final answer
            self._emit_event(StepEvent(
                event_type="synthesizing",
                step_number=0,
                data={"message": "Synthesizing final answer..."}
            ))

            final_answer = self._synthesize_answer(question, combined_output, pending_artifacts)

            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer}
            ))

            # Run facts extraction and suggestions generation in parallel
            tables = self.datastore.list_tables() if self.datastore else []
            response_facts, suggestions = self._run_post_synthesis_parallel(
                question, final_answer, tables
            )

            if response_facts:
                self._emit_event(StepEvent(
                    event_type="facts_extracted",
                    step_number=0,
                    data={
                        "facts": [f.to_dict() for f in response_facts],
                        "source": "response",
                    }
                ))

            if suggestions:
                self._emit_event(StepEvent(
                    event_type="suggestions_ready",
                    step_number=0,
                    data={"suggestions": suggestions}
                ))

        self.history.record_query(
            session_id=self.session_id,
            question=question,
            success=True,
            attempts=total_attempts,
            duration_ms=total_duration,
            answer=final_answer,
        )

        # Auto-compact if context is too large
        self._auto_compact_if_needed()

        return {
            "success": True,
            "plan": follow_up_plan,
            "results": all_results,
            "output": combined_output,
            "final_answer": final_answer,
            "suggestions": suggestions,
            "scratchpad": self.scratchpad.to_markdown(),
            "datastore_tables": self.datastore.list_tables() if self.datastore else [],
        }

    def _solve_auditable_with_steer_handling(
        self, problem: str, cached_fact_hints: list[dict] = None
    ) -> dict:
        """
        Wrapper for _solve_auditable that handles user steers.

        If a user provides guidance (steer) instead of a value when asked for
        a premise, this re-plans with the steer added to the problem context.

        Args:
            problem: The original problem
            cached_fact_hints: Optional cached fact hints

        Returns:
            Result dict from _solve_auditable
        """
        max_steer_attempts = 3
        steer_attempt = 0
        augmented_problem = problem

        while steer_attempt < max_steer_attempts:
            result = self._solve_auditable(augmented_problem, cached_fact_hints)

            # Check if we need to re-plan due to user steer
            if result.get("status") == "replan_needed":
                steer_attempt += 1
                steer = result.get("steer", "")
                fact_name = result.get("fact_name", "")
                logger.debug(f"Re-planning with steer: {steer}")

                # Augment the problem with the user's guidance
                augmented_problem = f"{problem}\n\nAdditional context: For {fact_name}, {steer}"

                # Clear cached plan to force re-planning
                self._last_plan = None

                # Emit event so user sees the message
                self._emit_event(StepEvent(
                    event_type="steer_detected",
                    step_number=0,
                    data={"message": result.get("message", "Re-planning with your guidance...")}
                ))
                continue

            return result

        # Max attempts reached
        return {"error": "Max re-planning attempts reached", "status": "failed"}

    def _solve_auditable(self, problem: str, cached_fact_hints: list[dict] = None) -> dict:
        """
        Solve a problem in auditable mode using fact-based derivation.

        Instead of generating a stepwise execution plan, this method:
        1. Identifies the question to answer
        2. Decomposes into required premises (facts from sources)
        3. Shows a derivation plan for approval
        4. Resolves facts with provenance tracking
        5. Generates an auditable derivation trace

        Args:
            problem: The problem/question to solve
            cached_fact_hints: Optional list of cached facts from previous run (for redo)
                               Each dict has 'name', 'value', 'value_type' keys

        Returns:
            Dict with derivation trace and verification result
        """
        import time
        start_time = time.time()

        # Save mode to datastore for follow-up handling
        if self.datastore:
            self.datastore.set_session_meta("mode", "auditable")

        # Step 1: Generate fact-based plan (identify required facts)
        self._emit_event(StepEvent(
            event_type="planning_start",
            step_number=0,
            data={"message": "Identifying required facts for verification..."}
        ))

        # Get source context for the planner
        # Proof planner is a single LLM call (no tool use), so it needs full table schemas
        # _build_source_context may only return database names without tables
        ctx = self._build_source_context(include_user_facts=False)
        # Override schema_overview with full schema if brief summary lacks table names
        if self.schema_manager and "Table:" not in ctx.get("schema_overview", ""):
            ctx["schema_overview"] = self.schema_manager.get_overview()

        # Build hint about cached facts for redo (helps LLM use consistent names)
        cached_facts_hint = ""
        if cached_fact_hints:
            # Separate tables from scalar facts
            cached_tables = [f for f in cached_fact_hints if f.get("value_type") == "table"]
            cached_scalars = [f for f in cached_fact_hints if f.get("value_type") != "table"]

            hint_lines = [
                "CACHED DATA AVAILABLE AS INPUT - Use these as premises with [source: cache]:",
            ]
            if cached_tables:
                hint_lines.append("\nCached tables:")
                for fact in cached_tables:
                    name = fact.get("name", "")
                    desc = fact.get("description", "")
                    hint_lines.append(f"  - \"{name}\" ({fact.get('row_count', '?')} rows) {desc}")
            if cached_scalars:
                hint_lines.append("\nCached values:")
                for fact in cached_scalars:
                    name = fact.get("name", "")
                    hint_lines.append(f"  - \"{name}\" = {fact.get('value', '?')}")
            hint_lines.append("")
            hint_lines.append("Use [source: cache] with EXACT names above. Do NOT rename them.")
            hint_lines.append("CRITICAL: You must still build a COMPLETE derivation chain that PROVES the answer.")
            hint_lines.append("CRITICAL: Do NOT just verify cached data exists - you must show HOW the answer is derived.")
            hint_lines.append("")
            cached_facts_hint = "\n".join(hint_lines) + "\n"

        # Extract document sources used during exploratory session
        exploratory_docs_hint = ""
        step_hints = getattr(self, '_proof_step_hints', [])
        if step_hints:
            doc_names = set()
            for step in step_hints:
                code = step.get("code", "")
                for m in re.findall(r"""doc_read\(\s*['"]([^'"]+)['"]\s*\)""", code):
                    doc_names.add(m)
            if doc_names:
                names_str = ", ".join(f'"{n}"' for n in sorted(doc_names))
                exploratory_docs_hint = f"DOCUMENT CONSTRAINT: The exploratory analysis used these documents: {names_str}. Use the SAME document sources for consistency.\n\n"

        fact_plan_prompt = f"""Construct a logical derivation to answer this question with full provenance.

Question: {problem}

{cached_facts_hint}{exploratory_docs_hint}Available databases:
{ctx["schema_overview"]}
{ctx["doc_overview"]}
{ctx["api_overview"]}

Build a formal derivation with EXACTLY this format:

QUESTION: <restate the question>

PREMISES:
P1: <fact_name> = ? (<what data to retrieve>) [source: database:<db_name>]
P2: <fact_name> = ? (<description>) [source: api:<api_name>]
P3: <fact_name> = ? (<description>) [source: document:<doc_name>]
P4: <fact_name> = <known_value> (<description>) [source: llm_knowledge]

PREMISE RULES:
- Premises are DATA only (tables, records, values) - NOT functions or operations
- Every premise MUST be referenced by at least one inference
- Use "cache" for data already in cache (PREFERRED - fastest)
- Use "database" for SQL queries to configured databases
- Use "api" for external API data (GraphQL or REST endpoints)
- Use "document:<doc_name>" for reference documents, policies, and guidelines (e.g., [source: document:business_rules])
- Use "llm_knowledge" for universal facts (mathematical constants, scientific facts) and well-established reference data (ISO codes, country info, currency codes)
- For known universal values, embed directly: P2: pi_value = 3.14159 (Pi constant) [source: llm_knowledge]
- NEVER ASSUME personal values (age, location, preferences) - use [source: user] and leave value as ?
- Example: P4: my_age = ? (User's age) [source: user]
- IMPORTANT: If cached data is available for what you need, ALWAYS use [source: cache] instead of fetching from database/api again.
- IMPORTANT: If the question mentions clarifications or user preferences (like "use guidelines from X"), treat these as DATA to be retrieved, NOT as embedded values. Always use = ? and resolve from the appropriate source.
- IMPORTANT: If a configured API can provide the data, use [source: api:<name>] instead of llm_knowledge.
- IMPORTANT: Extract numeric constraints from the question as premises. Example: "top 5 results" becomes P2: limit_count = 5 (Requested limit) [source: user]

INFERENCE:
I1: <result_name> = <operation>(P1, P2) -- <explanation>
I2: <result_name> = <operation>(I1) -- <explanation>

INFERENCE RULES:
- Each inference must reference at least one premise (P1/P2/etc) or prior inference (I1/I2/etc)
- CRITICAL: ALL premises MUST be used in the inference chain. Never define a premise that isn't referenced.
- CRITICAL: Each inference result_name MUST be GLOBALLY UNIQUE. NEVER reuse any name. BAD: two inferences both named "data_verified". GOOD: "validation_result" then "final_verification".
- CRITICAL: The final inference(s) should COMPUTE THE ACTUAL ANSWER, not just verify data exists. If the user asks for recommendations, calculate them. If they ask for comparisons, compute them.
- CRITICAL: Only ONE verify_exists() at the very end, referencing the computed answer. Do NOT add validate() or verify() steps before it.
- Operations like date extraction, filtering, grouping belong HERE, not in premises
- Keep operations simple: filter, join, group_sum, count, apply_rules, calculate, etc.
- For FUZZY MAPPING (e.g., free-text country names → ISO codes): plan the inference to first attempt mapping via data sources (APIs, databases), then use llm_map() as fallback for unmatched values. This reduces confidence but enables mapping when no exact data source exists.

CONCLUSION:
C: <final sentence describing what the final inference contains - use ENGLISH NAMES not I1/I2 references>
IMPORTANT: In the conclusion, ALWAYS use the English result_name (e.g., "raise_recommendations") NOT the ID (e.g., "I4")

EXAMPLE 1 - "What is revenue multiplied by Pi?":

PREMISES:
P1: orders = ? (All orders with amounts) [source: database:sales_db]
P2: pi_value = 3.14159 (Mathematical constant) [source: llm_knowledge]

INFERENCE:
I1: total_revenue = sum(P1.amount) -- Sum all order amounts
I2: adjusted_revenue = multiply(I1, P2) -- Multiply by Pi

CONCLUSION:
C: The revenue multiplied by Pi is provided in adjusted_revenue, calculated by multiplying total_revenue by pi_value.

EXAMPLE 2 - "Monthly revenue trend for last 12 months":

PREMISES:
P1: orders = ? (All orders with date and amount) [source: database:sales_db]

INFERENCE:
I1: recent_orders = filter(P1, last_12_months) -- Filter to last 12 months
I2: monthly_revenue = group_sum(I1, month, amount) -- Group by month, sum amounts
I3: trend = analyze(I2) -- Calculate trend direction

CONCLUSION:
C: Monthly revenue trend is provided in trend, showing direction based on monthly_revenue analysis.

EXAMPLE 3 - "Recommend raises based on performance reviews and guidelines":

PREMISES:
P1: employees = ? (All employees with current salary) [source: database:hr]
P2: performance_reviews = ? (Performance review records with ratings) [source: database:hr]
P3: raise_guidelines = ? (Business rules for raise percentages by rating) [source: document]

INFERENCE:
I1: recent_reviews = filter(P2, most_recent_per_employee) -- Get most recent review per employee
I2: employee_data = join(P1, I1, employee_id) -- Join employees with their reviews
I3: raises_with_rules = apply_guidelines(I2, P3) -- Apply raise guidelines based on rating (NOTE: P3 is USED here)
I4: raise_recommendations = calculate(I3, salary * raise_percentage) -- Calculate actual raise amounts

CONCLUSION:
C: Raise recommendations with calculated amounts are provided in raise_recommendations, derived by applying raise_guidelines to employee performance ratings.

Now generate the derivation. Use P1:, P2:, I1:, I2: prefixes EXACTLY as shown.
Premises are DATA. Operations (filter, extract, group, apply_guidelines) go in INFERENCE.
IMPORTANT: ALL premises must appear in at least one inference. The final inference must compute the answer.
"""

        result = self.router.execute(
            task_type=TaskType.INTENT_CLASSIFICATION,
            system="You analyze questions and decompose them into premises and inferences for auditable answers.",
            user_message=fact_plan_prompt,
            max_tokens=self.router.max_output_tokens,
        )
        fact_plan_text = result.content

        # Parse the proof structure
        claim = ""
        premises = []  # P1, P2, ... - base facts from sources
        inferences = []  # I1, I2, ... - derived facts
        conclusion = ""

        lines = fact_plan_text.split("\n")
        current_section = None
        for line in lines:
            line = line.strip()
            if line.startswith("QUESTION:"):
                claim = line.split("QUESTION:", 1)[1].strip()
            elif line.startswith("PREMISES:"):
                current_section = "premises"
            elif line.startswith("INFERENCE:"):
                current_section = "inference"
            elif line.startswith("CONCLUSION:"):
                current_section = "conclusion"
            elif current_section == "premises" and re.match(r'^P\d+:', line):
                # Parse: P1: fact_name = ? (description) [source: xxx]
                # Also handle: P1: fact_name = ? (description) without source
                match = re.match(r'^(P\d+):\s*(.+?)\s*=\s*\?\s*\(([^)]+)\)\s*(?:\[source:\s*([^\]]+)\])?', line)
                if match:
                    premises.append({
                        "id": match.group(1),
                        "name": match.group(2).strip(),
                        "description": match.group(3).strip(),
                        "source": match.group(4).strip() if match.group(4) else "database",
                    })
                else:
                    # Try format with embedded value: P1: fact_name = 8 (description) [source: llm_knowledge]
                    value_match = re.match(r'^(P\d+):\s*(.+?)\s*=\s*([^\s(]+)\s*\(([^)]+)\)\s*(?:\[source:\s*([^\]]+)\])?', line)
                    if value_match:
                        # Include the value in the name so embedded value extraction works
                        fact_name = value_match.group(2).strip()
                        embedded_val = value_match.group(3).strip()
                        premises.append({
                            "id": value_match.group(1),
                            "name": f"{fact_name} = {embedded_val}",  # Include value for extraction
                            "description": value_match.group(4).strip(),
                            "source": value_match.group(5).strip() if value_match.group(5) else "knowledge",
                        })
                    else:
                        # Try simpler format: P1: fact_name (description)
                        simple_match = re.match(r'^(P\d+):\s*(.+?)\s*\(([^)]+)\)', line)
                        if simple_match:
                            premises.append({
                                "id": simple_match.group(1),
                                "name": simple_match.group(2).strip().rstrip('=?').strip(),
                                "description": simple_match.group(3).strip(),
                                "source": "database",
                            })
            elif current_section == "inference" and re.match(r'^I\d+:', line):
                # Parse: I1: derived_fact = operation(inputs) -- explanation
                match = re.match(r'^(I\d+):\s*(.+?)\s*=\s*(.+?)\s*--\s*(.+)$', line)
                if match:
                    inferences.append({
                        "id": match.group(1),
                        "name": match.group(2).strip(),
                        "operation": match.group(3).strip(),
                        "explanation": match.group(4).strip(),
                    })
                else:
                    # Simpler format without operation details
                    simple_match = re.match(r'^(I\d+):\s*(.+)$', line)
                    if simple_match:
                        inferences.append({
                            "id": simple_match.group(1),
                            "name": "",
                            "operation": simple_match.group(2).strip(),
                            "explanation": "",
                        })
            elif current_section == "conclusion" and line:
                if line.startswith("C:"):
                    conclusion = line.split("C:", 1)[1].strip()
                elif not conclusion:
                    conclusion = line

        # Detect collection-oriented queries (reports, lists, etc.)
        # These don't have explicit scalar comparisons - the implicit goal is "data exists"
        comparison_keywords = [
            "greater than", "less than", "equal to", "equals",
            ">", "<", ">=", "<=", "==", "!=",
            "compare", "check if", "verify that", "prove that",
            "is positive", "is negative", "is zero",
        ]
        is_collection_query = True
        combined_text = (problem + " " + conclusion).lower()
        for keyword in comparison_keywords:
            if keyword.lower() in combined_text:
                is_collection_query = False
                break

        # Also check if the last inference is a comparison operation
        if inferences:
            last_op = inferences[-1].get("operation", "").lower()
            if any(kw in last_op for kw in ["compare", "check", "verify", ">", "<"]):
                is_collection_query = False

        # For collection queries, add implicit verification inference
        if is_collection_query and inferences:
            last_inf_id = inferences[-1]["id"]
            verify_id = f"I{len(inferences) + 1}"
            # Use a unique name that won't conflict with LLM-generated names
            existing_names = {inf.get("name", "") for inf in inferences}
            verify_name = "final_data_verification"
            suffix = 1
            while verify_name in existing_names:
                verify_name = f"final_data_verification_{suffix}"
                suffix += 1
            inferences.append({
                "id": verify_id,
                "name": verify_name,
                "operation": f"verify_exists({last_inf_id})",
                "explanation": "Verify result has data (count > 0) to confirm derivation succeeded",
            })
            # Update conclusion to emphasize provenance
            if not conclusion.lower().startswith("the data"):
                conclusion = f"The data is verified to exist with provenance: {conclusion}"

        # Validate plan structure BEFORE execution
        # This catches invalid references, unused premises, duplicates, etc.
        from constat.execution.dag import validate_proof_plan

        max_validation_retries = 3
        for validation_attempt in range(1, max_validation_retries + 1):
            # Emit validation start event
            self._emit_event(StepEvent(
                event_type="plan_validating",
                step_number=0,
                data={
                    "attempt": validation_attempt,
                    "max_attempts": max_validation_retries,
                    "premises_count": len(premises),
                    "inferences_count": len(inferences),
                }
            ))

            validation_result = validate_proof_plan(premises, inferences)

            if validation_result.valid:
                # Emit validation success
                self._emit_event(StepEvent(
                    event_type="plan_validated",
                    step_number=0,
                    data={
                        "attempt": validation_attempt,
                        "premises_count": len(premises),
                        "inferences_count": len(inferences),
                    }
                ))
                break  # Plan is valid, proceed

            # Plan has validation errors - emit detailed error event
            error_feedback = validation_result.format_for_retry()
            error_types = list(set(e.error_type for e in validation_result.errors))
            error_summary = ", ".join(error_types)

            self._emit_event(StepEvent(
                event_type="plan_validation_failed",
                step_number=0,
                data={
                    "attempt": validation_attempt,
                    "max_attempts": max_validation_retries,
                    "error_count": len(validation_result.errors),
                    "error_types": error_types,
                    "error_summary": error_summary,
                    "errors": [{"type": e.error_type, "fact_id": e.fact_id, "message": e.message}
                               for e in validation_result.errors],
                    "will_retry": validation_attempt < max_validation_retries,
                }
            ))

            if validation_attempt >= max_validation_retries:
                # Max retries reached - raise with helpful message
                raise ValueError(
                    f"Plan validation failed after {max_validation_retries} attempts.\n\n"
                    f"{error_feedback}\n\n"
                    f"Try rephrasing your question or reducing complexity."
                )

            # Emit retry event before regenerating
            self._emit_event(StepEvent(
                event_type="plan_regenerating",
                step_number=0,
                data={
                    "attempt": validation_attempt + 1,
                    "max_attempts": max_validation_retries,
                    "reason": error_summary,
                    "fixing_errors": error_types,
                }
            ))

            # Retry with explicit feedback about the errors
            retry_prompt = f"""{fact_plan_prompt}

{error_feedback}

REMEMBER:
- Each premise (P1, P2) and inference (I1, I2) must have a UNIQUE name
- ALL premises must be referenced in at least one inference operation
- Inference operations must only reference EXISTING premises (P1, P2, ...) or PRIOR inferences (I1 before I2)
- Do NOT reference facts that don't exist (e.g., P5 when only P1-P3 are defined)
"""
            retry_result = self.router.execute(
                task_type=TaskType.INTENT_CLASSIFICATION,
                system="You analyze questions and decompose them into premises and inferences for auditable answers. CRITICAL: Ensure all fact references are valid.",
                user_message=retry_prompt,
                max_tokens=self.router.max_output_tokens,
            )

            # Re-parse the retried plan
            fact_plan_text = retry_result.content
            claim = ""
            premises = []
            inferences = []
            conclusion = ""

            lines = fact_plan_text.split("\n")
            current_section = None
            for line in lines:
                line = line.strip()
                if line.startswith("QUESTION:"):
                    claim = line.split("QUESTION:", 1)[1].strip()
                elif line.startswith("PREMISES:"):
                    current_section = "premises"
                elif line.startswith("INFERENCE:"):
                    current_section = "inference"
                elif line.startswith("CONCLUSION:"):
                    current_section = "conclusion"
                elif current_section == "premises" and re.match(r'^P\d+:', line):
                    match = re.match(r'^(P\d+):\s*(.+?)\s*=\s*\?\s*\(([^)]+)\)\s*(?:\[source:\s*([^\]]+)\])?', line)
                    if match:
                        premises.append({
                            "id": match.group(1),
                            "name": match.group(2).strip(),
                            "description": match.group(3).strip(),
                            "source": match.group(4).strip() if match.group(4) else "database",
                        })
                    else:
                        value_match = re.match(r'^(P\d+):\s*(.+?)\s*=\s*([^\s(]+)\s*\(([^)]+)\)\s*(?:\[source:\s*([^\]]+)\])?', line)
                        if value_match:
                            fact_name = value_match.group(2).strip()
                            embedded_val = value_match.group(3).strip()
                            premises.append({
                                "id": value_match.group(1),
                                "name": f"{fact_name} = {embedded_val}",
                                "description": value_match.group(4).strip(),
                                "source": value_match.group(5).strip() if value_match.group(5) else "knowledge",
                            })
                        else:
                            simple_match = re.match(r'^(P\d+):\s*(.+?)\s*\(([^)]+)\)', line)
                            if simple_match:
                                premises.append({
                                    "id": simple_match.group(1),
                                    "name": simple_match.group(2).strip().rstrip('=?').strip(),
                                    "description": simple_match.group(3).strip(),
                                    "source": "database",
                                })
                elif current_section == "inference" and re.match(r'^I\d+:', line):
                    match = re.match(r'^(I\d+):\s*(.+?)\s*=\s*(.+?)\s*--\s*(.+)$', line)
                    if match:
                        inferences.append({
                            "id": match.group(1),
                            "name": match.group(2).strip(),
                            "operation": match.group(3).strip(),
                            "explanation": match.group(4).strip(),
                        })
                    else:
                        simple_match = re.match(r'^(I\d+):\s*(.+)$', line)
                        if simple_match:
                            inferences.append({
                                "id": simple_match.group(1),
                                "name": "",
                                "operation": simple_match.group(2).strip(),
                                "explanation": "",
                            })
                elif current_section == "conclusion" and line:
                    if line.startswith("C:"):
                        conclusion = line.split("C:", 1)[1].strip()
                    elif not conclusion:
                        conclusion = line

        # Emit planning complete
        total_steps = len(premises) + len(inferences) + 1  # +1 for conclusion
        self._emit_event(StepEvent(
            event_type="planning_complete",
            step_number=0,
            data={"steps": total_steps}
        ))

        # Build proof steps for display
        # Structure: Premises (resolve from sources) → Inferences (derive) → Conclusion
        proof_steps = []
        step_num = 1

        # Add premises as steps (these need to be resolved from sources)
        for p in premises:
            # Format: fact_name = ? (description) [source: xxx]
            proof_steps.append({
                "number": step_num,
                "goal": f"{p['name']} = ? ({p['description']}) [source: {p['source']}]",
                "depends_on": [],
                "type": "premise",
                "fact_id": p['id'],  # Keep P1/P2 id for execution reference
            })
            step_num += 1

        # Add inferences as steps (these depend on premises/prior inferences)
        premise_count = len(premises)
        for inf in inferences:
            # Format: derived_fact = operation -- explanation
            goal = inf['operation']
            if inf.get('name'):
                goal = f"{inf['name']} = {inf['operation']}"
            if inf.get('explanation'):
                goal += f" -- {inf['explanation']}"
            proof_steps.append({
                "number": step_num,
                "goal": goal,
                "depends_on": list(range(1, premise_count + 1)),  # Depends on all premises
                "type": "inference",
                "fact_id": inf['id'],  # Keep I1/I2 id for execution reference
            })
            step_num += 1

        # Add conclusion as final step
        all_prior_steps = list(range(1, step_num))
        proof_steps.append({
            "number": step_num,
            "goal": conclusion,
            "depends_on": all_prior_steps,
            "type": "conclusion",
        })

        # Store parsed derivation for later use
        self._current_proof = {
            "question": claim,  # The question being answered
            "premises": premises,
            "inferences": inferences,
            "conclusion": conclusion,
        }

        # Request approval if required
        # For auditable mode, we call the approval callback directly with proof_steps
        # that preserve the type and fact_id fields for proper P1:/I1:/C: display
        if self.session_config.require_approval:
            from constat.execution.mode import PlanApprovalRequest, PlanApprovalResponse, PlanApproval

            logger.debug(f"[_solve_auditable] require_approval=True, auto_approve={self.session_config.auto_approve}, has_callback={self._approval_callback is not None}")

            # Auto-approve if configured
            if self.session_config.auto_approve:
                logger.debug("[_solve_auditable] Auto-approving (auto_approve=True)")
                approval = PlanApprovalResponse.approve()
            elif not self._approval_callback:
                logger.debug("[_solve_auditable] Auto-approving (no callback)")
                approval = PlanApprovalResponse.approve()
            else:
                logger.debug("[_solve_auditable] Calling approval callback...")
                # Build approval request with full proof structure (preserves type, fact_id)
                request = PlanApprovalRequest(
                    problem=problem,
                    steps=proof_steps,  # Includes type, fact_id for proper display
                    reasoning=f"Question: {claim}",
                )
                approval = self._approval_callback(request)

            if approval.decision == PlanApproval.REJECT:
                self.datastore.set_session_meta("status", "rejected")
                return {
                    "success": False,
                    "rejected": True,
                    "reason": approval.reason,
                    "message": "Verification plan was rejected by user.",
                }

            elif approval.decision == PlanApproval.COMMAND:
                # User entered a slash command - pass back to REPL
                return {
                    "success": False,
                    "command": approval.command,
                    "message": "Slash command entered during approval.",
                }

            elif approval.decision == PlanApproval.SUGGEST:
                # Replan with feedback - for now, just include feedback in context
                problem = f"{problem}\n\nUser guidance: {approval.suggestion}"

            # Filter out deleted steps if any
            if approval.deleted_steps:
                deleted_set = set(approval.deleted_steps)
                # Filter premises by step number (P1, P2, etc. -> 1, 2, etc.)
                premises = [p for p in premises if p.get('number', 0) not in deleted_set]
                # Filter inferences by step number
                inferences = [i for i in inferences if i.get('number', 0) not in deleted_set]
                logger.info(f"Filtered out {len(deleted_set)} deleted steps: {approval.deleted_steps}")

        # Step 2: Execute plan using DAG-based parallel resolution
        # Start proof tree display for auditable mode (will print at end)
        self._emit_event(StepEvent(
            event_type="proof_start",
            step_number=0,
            data={
                "conclusion_fact": "answer",
                "conclusion_description": conclusion,
            }
        ))

        self._emit_event(StepEvent(
            event_type="verifying",
            step_number=0,
            data={"message": f"Resolving facts for: {claim or problem}"}
        ))

        try:
            from constat.execution.dag import parse_plan_to_dag, DAGExecutor, NodeStatus
            from constat.execution.fact_resolver import Fact, FactSource

            # Parse plan into DAG
            dag = parse_plan_to_dag(premises, inferences)

            # Emit fact_start for ALL nodes upfront so UI can show complete DAG
            for node in dag.nodes.values():
                # Determine if premise or inference (use fact_id, not dict key which is node.name)
                is_premise = node.fact_id.startswith("P")
                self._emit_event(StepEvent(
                    event_type="fact_start",
                    step_number=0,
                    data={
                        "fact_name": f"{node.fact_id}: {node.name}",
                        "fact_id": node.fact_id,
                        "fact_description": node.description if hasattr(node, 'description') else None,
                        "dependencies": [f"{dag.nodes[dep].fact_id}: {dep}" for dep in node.dependencies if dep in dag.nodes] if node.dependencies else [],
                        "is_premise": is_premise,
                    }
                ))

            # Get schema for SQL generation
            detailed_schema = self.schema_manager.get_overview()

            # Shared state for node execution
            resolved_premises = {}
            resolved_inferences = {}
            inference_names = {}
            derivation_lines = ["**Premise Resolution:**", ""]

            # Define node executor that calls back to Session
            def execute_node(node):
                return self._execute_dag_node(
                    node=node,
                    dag=dag,
                    problem=problem,
                    detailed_schema=detailed_schema,
                    premises=premises,
                    inferences=inferences,
                    resolved_premises=resolved_premises,
                    resolved_inferences=resolved_inferences,
                    inference_names=inference_names,
                )

            # Define event callback for progress reporting
            def dag_event_callback(event_type, data):
                node_name = data.get("name", "")
                fact_id = data.get("fact_id", "")
                level = data.get("level", 0)

                if event_type == "node_running":
                    # Determine if premise or inference
                    is_premise = fact_id.startswith("P") if fact_id else level == 0

                    # INVERTED TREE: Find what this node FEEDS INTO (not what it depends on)
                    # This shows derivation flow: premises -> inferences -> answer
                    def find_consumer(source_name: str) -> str | None:
                        """Find the first node that uses this node as input."""
                        for other_node in dag.nodes.values():
                            if source_name in other_node.dependencies:
                                return other_node.name
                        return None

                    # Emit fact_executing for DAG visualization
                    self._emit_event(StepEvent(
                        event_type="fact_executing",
                        step_number=level + 1,
                        data={
                            "fact_name": f"{fact_id}: {node_name}",
                            "fact_id": fact_id,
                        }
                    ))

                    if is_premise:
                        # Parent = what inference uses this premise
                        consumer = find_consumer(node_name)
                        logger.debug(f"[DAG] {fact_id} (premise) feeds into: {consumer}")
                        self._emit_event(StepEvent(
                            event_type="premise_resolving",
                            step_number=level + 1,
                            data={
                                "fact_name": f"{fact_id}: {node_name}",
                                "step": level + 1,
                                "parent": consumer,  # What this premise feeds into
                            }
                        ))
                    else:
                        # Parent = what inference uses this inference (or root if terminal)
                        consumer = find_consumer(node_name)
                        logger.debug(f"[DAG] {fact_id} (inference) feeds into: {consumer}")
                        self._emit_event(StepEvent(
                            event_type="inference_executing",
                            step_number=level + 1,
                            data={
                                "inference_id": fact_id,
                                "operation": node_name,
                                "parent": consumer,  # What this inference feeds into
                            }
                        ))

                elif event_type == "node_resolved":
                    value = data.get("value")
                    confidence = data.get("confidence", 0.9)
                    source = data.get("source", "")
                    is_premise = fact_id.startswith("P") if fact_id else level == 0

                    # Get dependencies for this node (dag.nodes keyed by name, not fact_id)
                    node_obj = dag.nodes.get(node_name)
                    deps = [f"{dag.nodes[dep].fact_id}: {dep}" for dep in node_obj.dependencies if dep in dag.nodes] if node_obj and node_obj.dependencies else []

                    # Emit fact_resolved for DAG visualization
                    fact_resolved_data = {
                        "fact_name": f"{fact_id}: {node_name}",
                        "fact_id": fact_id,
                        "value": value,
                        "confidence": confidence,
                        "source": source,
                        "dependencies": deps,
                    }
                    validations = data.get("validations")
                    if validations:
                        fact_resolved_data["validations"] = validations
                    self._emit_event(StepEvent(
                        event_type="fact_resolved",
                        step_number=level + 1,
                        data=fact_resolved_data,
                    ))

                    if is_premise:
                        derivation_lines.append(f"- {fact_id}: {node_name} = {str(value)[:100]} (confidence: {confidence:.0%})")
                        self._emit_event(StepEvent(
                            event_type="premise_resolved",
                            step_number=level + 1,
                            data={"fact_name": f"{fact_id}: {node_name}", "value": value, "confidence": confidence, "source": source}
                        ))
                    else:
                        self._emit_event(StepEvent(
                            event_type="inference_complete",
                            step_number=level + 1,
                            data={"inference_id": fact_id, "inference_name": node_name, "result": value}
                        ))

                elif event_type == "node_failed":
                    error = data.get("error", "Unknown error")
                    logger.error(f"{fact_id} ({node_name}) failed: {error}")

                    # Emit fact_failed for DAG visualization
                    self._emit_event(StepEvent(
                        event_type="fact_failed",
                        step_number=level + 1,
                        data={
                            "fact_name": f"{fact_id}: {node_name}",
                            "fact_id": fact_id,
                            "reason": error,
                        }
                    ))

                    self._emit_event(StepEvent(
                        event_type="premise_resolved" if fact_id.startswith("P") else "inference_failed",
                        step_number=level + 1,
                        data={"fact_name": f"{fact_id}: {node_name}", "error": error}
                    ))

                elif event_type == "node_blocked":
                    blocked_by = data.get("blocked_by", "dependency failed")
                    logger.info(f"{fact_id} ({node_name}) blocked by {blocked_by}")

                    self._emit_event(StepEvent(
                        event_type="fact_blocked",
                        step_number=level + 1,
                        data={
                            "fact_name": f"{fact_id}: {node_name}",
                            "fact_id": fact_id,
                            "reason": f"blocked by {blocked_by}",
                        }
                    ))

                elif event_type == "node_started":
                    # Log actual thread start for parallelism diagnosis
                    start_time = data.get("start_time_ms", 0)
                    logger.debug(f"DAG node {fact_id} STARTED at {start_time}ms")

                elif event_type == "node_timing":
                    # Log timing for parallelism analysis
                    start_ms = data.get("start_ms", 0)
                    end_ms = data.get("end_ms", 0)
                    duration_ms = data.get("duration_ms", 0)
                    failed = data.get("failed", False)
                    status = "FAILED" if failed else "COMPLETED"
                    logger.info(f"DAG timing: {fact_id} {status} - start:{start_ms} end:{end_ms} duration:{duration_ms}ms")

            # Phase 4: Extract user-specified validation constraints from the problem
            self._proof_user_validations = self._extract_user_validations(problem, inferences)

            # Reset cancellation state before starting execution
            self.reset_cancellation()

            # Execute DAG with parallel resolution
            executor = DAGExecutor(
                dag=dag,
                node_executor=execute_node,
                max_workers=min(10, len(premises) + len(inferences)),
                event_callback=dag_event_callback,
                fail_fast=True,
                execution_context=self._execution_context,
            )

            # Build pre-resolved info for nodes with embedded values (e.g., "breed_limit = 10")
            pre_resolved = {}
            for node in dag.nodes.values():
                if node.status == NodeStatus.RESOLVED and node.is_leaf:
                    pre_resolved[node.fact_id] = {
                        "value": node.value,
                        "confidence": node.confidence,
                    }

            # Emit event to start live plan display (includes pre-resolved info)
            self._emit_event(StepEvent(
                event_type="dag_execution_start",
                step_number=0,
                data={"premises": premises, "inferences": inferences, "pre_resolved": pre_resolved}
            ))

            # Phase 4: Check for cancellation before starting DAG execution
            if self.is_cancelled():
                self._emit_event(StepEvent(
                    event_type="execution_cancelled",
                    step_number=0,
                    data={"message": "Execution cancelled before DAG execution started"}
                ))
                return {
                    "success": False,
                    "cancelled": True,
                    "message": "Execution cancelled before fact resolution started.",
                    "queued_intent_results": self.process_queued_intents(),
                }

            logger.info(f"Executing DAG with {len(premises)} premises and {len(inferences)} inferences")

            # Temporarily disable fact_resolver events during DAG execution
            # Session.py already emits all necessary events with consistent naming (e.g., "P1: employees")
            # fact_resolver.resolve_tiered() emits duplicate events with different naming (e.g., "employees")
            saved_callback = self.fact_resolver._event_callback
            self.fact_resolver._event_callback = None
            try:
                result = executor.execute()
            finally:
                self.fact_resolver._event_callback = saved_callback

            # Emit event to stop live plan display
            self._emit_event(StepEvent(
                event_type="dag_execution_complete",
                step_number=0,
                data={"success": result.success, "failed_nodes": result.failed_nodes, "cancelled": result.cancelled}
            ))

            # Check for cancellation
            if result.cancelled:
                self._emit_event(StepEvent(
                    event_type="execution_cancelled",
                    step_number=0,
                    data={"message": "Execution cancelled during DAG execution"}
                ))
                return {
                    "success": False,
                    "cancelled": True,
                    "message": "Execution cancelled during fact resolution.",
                    "queued_intent_results": self.process_queued_intents(),
                }

            if not result.success:
                failed = ", ".join(result.failed_nodes)
                raise Exception(f"Plan execution failed. Could not resolve: {failed}")

            # Build inference lines from results
            inference_lines = ["", "**Inference Execution:**", ""]
            for inf in inferences:
                inf_id = inf['id']
                inf_name = inf.get('name', inf_id)
                if inf_id in resolved_inferences:
                    val = resolved_inferences[inf_id]
                    inference_lines.append(f"- {inf_id}: {inf_name} = {val} ✓")
                else:
                    # Check DAG node for value
                    node = dag.get_node(inf_name)
                    if node and node.value is not None:
                        inference_lines.append(f"- {inf_id}: {inf_name} = {node.value} ✓")
                        resolved_inferences[inf_id] = node.value

            derivation_lines.extend(inference_lines)
            # Step 4: Synthesize answer from resolved premises and inferences
            self._emit_event(StepEvent(
                event_type="synthesizing",
                step_number=0,
                data={"message": "Synthesizing answer from resolved facts..."}
            ))

            # Build synthesis context - use English names for clarity
            # Truncate values to prevent token overflow
            def truncate_value(val, max_chars=500):
                s = str(val)
                return s[:max_chars] + "..." if len(s) > max_chars else s

            resolved_context = "\n".join([
                f"- {pid} ({premises[int(pid[1:])-1]['name']}): {truncate_value(p.value)}"
                for pid, p in resolved_premises.items() if p and p.value
            ])
            # Use English variable names (e.g., "budget_validated_raises") not IDs (e.g., "I6")
            inference_context = "\n".join([
                f"- {inf_id} ({inference_names.get(inf_id, inf_id)}): {truncate_value(result, 200)}"
                for inf_id, result in resolved_inferences.items()
            ])

            # Collect artifact table names for synthesis (don't include table data)
            artifact_tables = []
            if inferences and self.datastore:
                available_tables = {t['name'] for t in self.datastore.list_tables()}
                for inf in inferences:
                    inf_id = inf['id']
                    table_name = inference_names.get(inf_id, inf_id.lower())
                    result = resolved_inferences.get(inf_id, "")
                    # Track tables that were created (skip verification steps)
                    if "rows" in str(result) and "verified" not in str(result).lower() and "FAILED" not in str(result):
                        if table_name in available_tables:
                            artifact_tables.append(table_name)

            # Build artifact reference for synthesis
            artifact_reference = ""
            if artifact_tables:
                artifact_reference = f"\n\nArtifact tables created: {', '.join(artifact_tables)}"
                artifact_reference += "\n(User can view these via /tables command)"

            synthesis_prompt = f"""Based on the resolved premises and inference plan, provide the answer.

Question: {claim}

Resolved Premises:
{resolved_context if resolved_context else "(no premises resolved)"}

Inference Steps:
{inference_context}

Conclusion to derive: {conclusion}
{artifact_reference}

IMPORTANT INSTRUCTIONS:
1. Always refer to data by its English variable name (e.g., "budget_validated_raises"), NEVER by ID (e.g., "I6")
2. Do NOT display tables inline - data is available as clickable artifacts
3. ALWAYS use backticks when referencing artifacts: `table_name` (these become clickable links)
4. For tables, include row count: `budget_validated_raises` (15 rows)
5. ONLY reference artifacts that were actually created - do not invent table names
6. Focus on the key findings and conclusions, not on showing raw data
7. If the user asked for recommendations/suggestions, summarize them with key values
8. EVALUATE GOAL COMPLETENESS: Re-read the original question carefully. Explicitly assess whether EVERY aspect of the question has been addressed. If any goal was partially or not addressed, state what is missing and why.

Provide a concise, clear answer with inline artifact references."""

            synthesis_result = self.router.execute(
                task_type=TaskType.SYNTHESIS,
                system="You synthesize answers from resolved facts with full provenance.",
                user_message=synthesis_prompt,
                max_tokens=self.router.max_output_tokens,
            )

            if not synthesis_result.success:
                logger.warning(f"Synthesis failed: {synthesis_result.content}")
                answer = f"Verification completed but answer synthesis failed. See derivation below."
            else:
                answer = synthesis_result.content
            confidence = sum(p.confidence for p in resolved_premises.values() if p) / max(len(resolved_premises), 1)
            derivation_trace = "\n".join(derivation_lines)

            verify_result = {
                "answer": answer,
                "confidence": confidence,
                "derivation": derivation_trace,
                "sources": [{"type": p["source"], "description": p["description"]} for p in premises],
            }

            # Generate insights if enabled
            insights = ""
            skip_insights = not self.session_config.enable_insights
            if not skip_insights:
                self._emit_event(StepEvent(
                    event_type="generating_insights",
                    step_number=0,
                    data={"message": "Generating insights..."}
                ))

                # Build source summary
                source_types = set()
                for p in resolved_premises.values():
                    if p and p.source:
                        source_types.add(p.source.value if hasattr(p.source, 'value') else str(p.source))

                insights_prompt = f"""Provide a brief summary of this analysis.

Original question: {claim}

Resolved premises:
{resolved_context}

Inference results:
{inference_context}

Conclusion: {conclusion}

Sources used: {', '.join(source_types) if source_types else 'various'}

Write a SHORT summary (2-3 sentences max) in plain prose explaining what this analysis shows.
Do NOT use bullet points or numbered lists. Just a brief paragraph.
Focus on the key finding and its significance.
If the original question had multiple goals or sub-questions, note whether all were addressed."""

                try:
                    insights_result = self.router.execute(
                        task_type=TaskType.SYNTHESIS,
                        system="You analyze proofs and provide actionable insights.",
                        user_message=insights_prompt,
                        max_tokens=self.router.max_output_tokens,
                    )
                    insights = insights_result.content
                except Exception as e:
                    logger.debug(f"Failed to generate insights (non-fatal): {e}")
                    insights = ""

            duration_ms = int((time.time() - start_time) * 1000)

            # Format output
            answer = verify_result.get("answer", "")
            confidence = verify_result.get("confidence", 0.0)
            derivation_trace = verify_result.get("derivation", "")
            sources = verify_result.get("sources", [])

            # Build final output (derivation details shown during execution, not in final answer)
            output_parts = [
                f"**Verification Result** (confidence: {confidence:.0%})",
                "",
                answer,
            ]

            if insights:
                output_parts.extend([
                    "",
                    "**Summary:**",
                    insights,
                ])

            final_output = "\n".join(output_parts)

            self._emit_event(StepEvent(
                event_type="verification_complete",
                step_number=0,
                data={
                    "answer": answer,
                    "confidence": confidence,
                    "has_derivation": bool(derivation_trace),
                }
            ))

            # Record in history
            self.history.record_query(
                session_id=self.session_id,
                question=problem,
                success=True,
                attempts=1,
                duration_ms=duration_ms,
                answer=final_output,
            )

            # Save resolved facts for redo operations
            if self.datastore:
                import json
                cached_facts = self.fact_resolver.export_cache()
                self.datastore.set_session_meta("resolved_facts", json.dumps(cached_facts))

            # Generate and save Data Flow Diagram (DFD) as published artifact
            dfd_text = ""
            try:
                from constat.visualization.box_dag import generate_proof_dfd
                dfd_text = generate_proof_dfd(proof_steps, max_width=80, max_name_len=10)
                if dfd_text and self.datastore:
                    from pathlib import Path
                    artifacts_dir = Path(".constat") / self.user_id / "sessions" / self.session_id / "artifacts"
                    artifacts_dir.mkdir(parents=True, exist_ok=True)
                    dfd_path = artifacts_dir / "data_flow.txt"
                    dfd_path.write_text(dfd_text)

                    # Register as published artifact
                    self.registry.register_artifact(
                        user_id=self.user_id,
                        session_id=self.session_id,
                        name="data_flow",
                        file_path=str(dfd_path.resolve()),
                        artifact_type="diagram",
                        size_bytes=len(dfd_text.encode('utf-8')),
                        description="Data flow diagram showing proof dependencies",
                        is_published=True,
                        title="Data Flow Diagram",
                    )
                    logger.debug(f"Saved DFD artifact: {dfd_path}")
            except Exception as e:
                logger.warning(f"Failed to generate DFD: {e}")

            # Build proof nodes for summary generation
            proof_nodes = []
            for p in premises:
                pid = p['id']
                resolved = resolved_premises.get(pid)
                proof_nodes.append({
                    "id": pid,
                    "name": p['name'],
                    "value": resolved.value if resolved else None,
                    "source": resolved.source.value if resolved and hasattr(resolved.source, 'value') else str(resolved.source) if resolved else p.get('source'),
                    "confidence": resolved.confidence if resolved else None,
                    "dependencies": [],
                })
            for inf in inferences:
                iid = inf['id']
                value = resolved_inferences.get(iid)
                # Get actual confidence and reasoning from DAG node
                inf_name = inf.get('name', iid)
                dag_node = dag.get_node(inf_name)
                node_confidence = dag_node.confidence if dag_node else (1.0 if value is not None else 0.0)
                # Get reasoning from the fact resolver if available
                node_reasoning = None
                if self.fact_resolver:
                    fact = self.fact_resolver.get_fact(inf_name)
                    if fact and hasattr(fact, 'reasoning'):
                        node_reasoning = fact.reasoning
                deps = [p['id'] for p in premises]  # Inferences depend on premises
                proof_nodes.append({
                    "id": iid,
                    "name": inf_name,
                    "value": value,
                    "source": "derived",
                    "confidence": node_confidence,
                    "dependencies": deps,
                    "reasoning": node_reasoning,
                })

            return {
                "success": True,
                "mode": Mode.PROOF.value,
                "output": final_output,
                "final_answer": answer,
                "confidence": confidence,
                "derivation": derivation_trace,
                "derivation_chain": derivation_trace,  # Alias for UI
                "sources": sources,
                "proof_nodes": proof_nodes,  # For summary generation
                "problem": problem,  # Original problem text
                "suggestions": [
                    "Show me the supporting data for this verification",
                    "What assumptions were made in this analysis?",
                ],
                "datastore_tables": self.datastore.list_tables() if self.datastore else [],
            }

        except Exception as e:
            duration_ms = int((time.time() - start_time) * 1000)

            self._emit_event(StepEvent(
                event_type="verification_error",
                step_number=0,
                data={"error": str(e)}
            ))

            return {
                "success": False,
                "mode": Mode.PROOF.value,
                "error": str(e),
                "output": f"Verification failed: {e}",
            }

    def _solve_knowledge(self, problem: str) -> dict:
        """
        Solve a problem in knowledge mode using document lookup + LLM synthesis.

        This mode is for explanation/knowledge requests that don't need data analysis.
        It searches configured documents and synthesizes an explanation.

        Args:
            problem: The question/request to answer

        Returns:
            Dict with synthesized explanation and sources
        """
        start_time = time.time()

        # Step 1: Search documents for relevant content
        self._emit_event(StepEvent(
            event_type="searching_documents",
            step_number=0,
            data={"message": "Searching reference documents..."}
        ))

        sources = []
        doc_context = ""

        if self.doc_tools and self.config.documents:
            # Search for relevant document excerpts
            search_results = self.doc_tools.search_documents(problem, limit=5)

            if search_results:
                doc_lines = ["Relevant document excerpts:"]
                for i, result in enumerate(search_results, 1):
                    doc_name = result.get("document", "unknown")
                    excerpt = result.get("excerpt", "")
                    relevance = result.get("relevance", 0)
                    section = result.get("section", "")

                    source_info = {
                        "document": doc_name,
                        "section": section,
                        "relevance": relevance,
                    }
                    sources.append(source_info)

                    doc_lines.append(f"\n[{i}] From '{doc_name}'" + (f" - {section}" if section else ""))
                    doc_lines.append(excerpt)

                doc_context = "\n".join(doc_lines)

        # Step 2: Build prompt for LLM synthesis
        self._emit_event(StepEvent(
            event_type="synthesizing",
            step_number=0,
            data={"message": "Synthesizing explanation..."}
        ))

        # System prompt for knowledge/explanation queries
        system_prompt = """You are a knowledgeable assistant for explanation and lookup requests.

Answer questions using configured reference documents and your general knowledge.
Be accurate and cite your sources when referencing specific documents.
If you don't have enough information, say so rather than guessing."""

        # Add context about the configuration (including active role)
        config_prompt = self._get_system_prompt()
        if config_prompt:
            system_prompt = f"{system_prompt}\n\n{config_prompt}"

        # Build user message with document context
        if doc_context:
            user_message = f"""Question: {problem}

{doc_context}

Please provide a clear, accurate explanation based on the documents above and your general knowledge.
Cite specific documents when referencing them."""
        else:
            user_message = f"""Question: {problem}

No reference documents are configured. Please provide an explanation based on your general knowledge.
If you don't have enough information, say so rather than guessing."""

        # Step 3: Generate response
        try:
            result = self.router.execute(
                task_type=TaskType.SYNTHESIS,
                system=system_prompt,
                user_message=user_message,
                max_tokens=self.router.max_output_tokens,
            )

            answer = result.content
            duration_ms = int((time.time() - start_time) * 1000)

            self._emit_event(StepEvent(
                event_type="knowledge_complete",
                step_number=0,
                data={
                    "has_documents": bool(sources),
                    "source_count": len(sources),
                }
            ))

            # Build final output
            output_parts = [answer]

            if sources:
                output_parts.extend([
                    "",
                    "**Sources consulted:**",
                ])
                for src in sources:
                    src_line = f"- {src['document']}"
                    if src.get('section'):
                        src_line += f" ({src['section']})"
                    output_parts.append(src_line)

            final_output = "\n".join(output_parts)

            # Record in history (only if session exists)
            if self.session_id:
                self.history.record_query(
                    session_id=self.session_id,
                    question=problem,
                    success=True,
                    attempts=1,
                    duration_ms=duration_ms,
                    answer=final_output,
                )

            return {
                "success": True,
                "meta_response": True,  # Display as meta-response (no tables)
                "mode": "knowledge",
                "output": final_output,
                "sources": sources,
                "plan": None,  # No plan in knowledge mode
                "suggestions": [
                    "Tell me more about a specific aspect",
                    "What data is available to analyze?",
                ],
            }

        except Exception as e:
            duration_ms = int((time.time() - start_time) * 1000)

            self._emit_event(StepEvent(
                event_type="knowledge_error",
                step_number=0,
                data={"error": str(e)}
            ))

            return {
                "success": False,
                "mode": "knowledge",
                "error": str(e),
                "output": f"Failed to generate explanation: {e}",
            }

    def prove_conversation(self, guidance: str | None = None) -> dict:
        """
        Re-run the original request in proof/auditable mode.

        Takes the original problem from the session and runs it through
        the auditable solver, reusing existing facts and tables from
        the exploratory session.

        Args:
            guidance: Optional user guidance for the proof (e.g., "focus on X", "use table Y")

        Returns:
            Dict with proof results (same format as _solve_auditable)
        """
        import json

        logger.debug("[prove_conversation] Starting prove_conversation")

        # Check if we have a session to prove
        if not self.session_id:
            return {"error": "No active session to prove"}

        if not self.datastore:
            return {"error": "No datastore available"}

        # Get the original problem
        original_problem = self.datastore.get_session_meta("problem")
        if not original_problem:
            return {"no_claims": True, "error": "No conversation to prove"}

        # Get follow-up questions to include in proof
        follow_ups_json = self.datastore.get_session_meta("follow_ups")
        follow_ups = []
        if follow_ups_json:
            try:
                follow_ups = json.loads(follow_ups_json)
            except json.JSONDecodeError:
                pass

        # Build combined problem statement
        if follow_ups:
            combined_problem = f"""Original request: {original_problem}

Follow-up requests:
{chr(10).join(f'- {q}' for q in follow_ups)}

Prove all of the above claims and provide a complete audit trail."""
            logger.debug(f"[prove_conversation] Combined problem with {len(follow_ups)} follow-ups")
        else:
            combined_problem = original_problem

        if guidance:
            combined_problem += f"\n\nAdditional guidance for this proof: {guidance}"
            logger.debug(f"[prove_conversation] Added guidance: {guidance[:100]}")

        logger.debug(f"[prove_conversation] Running proof for: {combined_problem[:150]}...")

        # Gather step codes from exploratory session as hints for inference generation
        self._proof_step_hints = []
        if self.history and self.session_id:
            try:
                step_codes = self.history.list_step_codes(self.session_id)
                if step_codes:
                    self._proof_step_hints = step_codes
                    logger.info(f"[prove_conversation] Loaded {len(step_codes)} step code hints for proof")
            except Exception as e:
                logger.debug(f"[prove_conversation] Could not load step codes: {e}")

        # Clear old inference codes from previous proof runs
        if self.history and self.session_id:
            self.history.clear_inferences(self.session_id)

        # Emit proof_start event so UI shows "Generating proof..." instead of "Planning..."
        self._emit_event(StepEvent(
            event_type="proof_start",
            step_number=0,
            data={"problem": combined_problem[:100]}
        ))

        # For proof mode, we do NOT pass cached/derived facts as hints.
        # Proof must derive from GROUND TRUTH sources only (databases, APIs, documents).
        # This ensures the proof is independent and verifiable.
        logger.debug("[prove_conversation] Proof will derive from ground truth sources only (no cached facts)")

        # Auto-approve during /prove
        original_auto_approve = self.session_config.auto_approve
        self.session_config.auto_approve = True

        try:
            # Run the combined problem through the auditable solver
            # No cached_fact_hints - proof derives from ground truth only
            result = self._solve_auditable(combined_problem)

            self._emit_event(StepEvent(
                event_type="proof_complete",
                step_number=0,
                data={
                    "success": result.get("success", False),
                    "confidence": result.get("confidence", 0.0),
                }
            ))

            # Generate proof summary asynchronously
            proof_nodes = result.get("proof_nodes", [])
            logger.info(f"[prove_conversation] Generating summary: success={result.get('success')}, proof_nodes count={len(proof_nodes)}")
            if result.get("success") and proof_nodes:
                try:
                    from constat.api.summarization import summarize_proof
                    logger.info(f"[prove_conversation] Calling summarize_proof with {len(proof_nodes)} nodes")
                    summary_result = summarize_proof(
                        problem=result.get("problem", combined_problem),
                        proof_nodes=proof_nodes,
                        llm=self.router,
                    )
                    logger.info(f"[prove_conversation] summarize_proof returned success={summary_result.success}, has_summary={bool(summary_result.summary)}, error={summary_result.error}")
                    if summary_result.success and summary_result.summary:
                        # Save as artifact (optional - don't fail if this fails)
                        if self.datastore:
                            try:
                                self.datastore.add_artifact(
                                    step_number=0,
                                    attempt=1,
                                    artifact_type="markdown",
                                    content=f"# Proof Summary\n\n{summary_result.summary}",
                                    name="proof_summary",
                                    title="Proof Summary",
                                )
                            except Exception as ae:
                                logger.warning(f"[prove_conversation] Failed to save summary artifact: {ae}")
                        # Emit event that summary is ready (always emit if summary generated)
                        self._emit_event(StepEvent(
                            event_type="proof_summary_ready",
                            step_number=0,
                            data={"summary": summary_result.summary}
                        ))
                        logger.info(f"[prove_conversation] Proof summary generated and emitted")
                    else:
                        logger.warning(f"[prove_conversation] summarize_proof failed: {summary_result.error}")
                except Exception as e:
                    logger.warning(f"Failed to generate proof summary: {e}", exc_info=True)
            else:
                logger.warning(f"[prove_conversation] Skipping summary: success={result.get('success')}, has_nodes={bool(proof_nodes)}")

            self.last_proof_result = result
            return result

        except Exception as e:
            logger.error(f"[prove_conversation] Error: {e}")
            return {"error": str(e), "success": False}

        finally:
            self.session_config.auto_approve = original_auto_approve
            self._proof_step_hints = []

    def replay(self, problem: str) -> dict:
        """
        Replay a previous session by re-executing stored code without LLM codegen.

        This loads the stored code from the scratchpad and re-executes it,
        then synthesizes a new answer (which still uses the LLM).

        Useful for demos, debugging, or re-running with modified data.

        Args:
            problem: The original problem (used for answer synthesis)

        Returns:
            Dict with results (same format as solve())
        """
        if not self.datastore:
            raise ValueError("No datastore available for replay")

        # Load stored scratchpad entries
        entries = self.datastore.get_scratchpad()
        if not entries:
            raise ValueError("No stored steps to replay")

        # Emit planning complete (we're using stored plan)
        self._emit_event(StepEvent(
            event_type="plan_ready",
            step_number=0,
            data={
                "steps": [
                    {"number": e["step_number"], "goal": e["goal"], "depends_on": []}
                    for e in entries
                ],
                "reasoning": "Replaying stored execution",
                "is_followup": False,
            }
        ))

        all_results = []
        for entry in entries:
            step_number = entry["step_number"]
            goal = entry["goal"]
            code = entry["code"]

            if not code:
                raise ValueError(f"Step {step_number} has no stored code to replay")

            self._emit_event(StepEvent(
                event_type="step_start",
                step_number=step_number,
                data={"goal": goal}
            ))

            self._emit_event(StepEvent(
                event_type="executing",
                step_number=step_number,
                data={"attempt": 1, "code": code}
            ))

            start_time = time.time()

            # Track tables before execution (name + version to detect updates)
            tables_before_list = self.datastore.list_tables()
            tables_before = set(t['name'] for t in tables_before_list)
            versions_before = {t['name']: t.get('version', 1) for t in tables_before_list}

            # Execute stored code
            exec_globals = self._get_execution_globals()
            result = self.executor.execute(code, exec_globals)

            # Auto-save any DataFrames
            if result.success:
                self._auto_save_results(result.namespace, step_number)

            duration_ms = int((time.time() - start_time) * 1000)
            tables_after_list = self.datastore.list_tables()
            tables_after = set(t['name'] for t in tables_after_list)
            versions_after = {t['name']: t.get('version', 1) for t in tables_after_list}
            new_tables = tables_after - tables_before
            updated_tables = {
                name for name in tables_before & tables_after
                if versions_after.get(name, 1) > versions_before.get(name, 1)
                and not name.startswith('_')
            }
            tables_created = list(new_tables | updated_tables)

            if result.success:
                self._emit_event(StepEvent(
                    event_type="step_complete",
                    step_number=step_number,
                    data={
                        "goal": goal,
                        "code": code,
                        "stdout": result.stdout,
                        "attempts": 1,
                        "duration_ms": duration_ms,
                        "tables_created": tables_created,
                    }
                ))

                all_results.append(StepResult(
                    success=True,
                    stdout=result.stdout,
                    attempts=1,
                    duration_ms=duration_ms,
                    tables_created=tables_created,
                    code=code,
                ))
            else:
                self._emit_event(StepEvent(
                    event_type="step_error",
                    step_number=step_number,
                    data={"error": result.stderr or "Execution failed", "attempt": 1}
                ))
                return {
                    "success": False,
                    "error": result.stderr or "Replay execution failed",
                    "step_number": step_number,
                }

        # Synthesize final answer (respects insights config)
        combined_output = "\n\n".join([
            f"Step {entry['step_number']}: {entry['goal']}\n{r.stdout}"
            for entry, r in zip(entries, all_results)
        ])

        # Emit raw results first
        self._emit_event(StepEvent(
            event_type="raw_results_ready",
            step_number=0,
            data={"output": combined_output}
        ))

        # Check for created artifacts (to mention in synthesis)
        from constat.visualization.output import peek_pending_outputs
        pending_artifacts = peek_pending_outputs()

        # Check if insights are enabled
        skip_insights = not self.session_config.enable_insights

        if skip_insights:
            final_answer = combined_output
            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer, "brief": True}
            ))
        else:
            self._emit_event(StepEvent(
                event_type="synthesizing",
                step_number=0,
                data={"message": "Synthesizing final answer..."}
            ))

            final_answer = self._synthesize_answer(problem, combined_output, pending_artifacts)

            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer}
            ))

        total_duration = sum(r.duration_ms for r in all_results)

        return {
            "success": True,
            "results": all_results,
            "output": combined_output,
            "final_answer": final_answer,
            "datastore_tables": self.datastore.list_tables(),
            "duration_ms": total_duration,
            "replay": True,
        }

    def get_state(self) -> dict:
        """Get current session state for inspection or resumption."""
        return {
            "session_id": self.session_id,
            "plan": self.plan,
            "scratchpad": self.scratchpad.to_markdown(),
            "state": self.datastore.get_all_state() if self.datastore else {},
            "completed_steps": self.plan.completed_steps if self.plan else [],
            "datastore_tables": self.datastore.list_tables() if self.datastore else [],
        }

    # --- Session Data Sources ---

    def add_database(
        self,
        name: str,
        db_type: str,
        uri: str,
        description: str = "",
    ) -> bool:
        """Add a database to the current session.

        The database will be available as `db_<name>` in code execution.
        The schema is introspected and table/column names are added as
        session entities for entity extraction.

        Args:
            name: Database name (used as db_<name> variable)
            db_type: Database type (sql, csv, json, parquet, mongodb, etc.)
            uri: Connection URI or file path
            description: Human-readable description

        Returns:
            True if added successfully
        """
        self.session_databases[name] = {
            "type": db_type,
            "uri": uri,
            "description": description,
        }

        # Update schema entities list for NER recognition in future document indexing
        # Entities are created via NER extraction from description chunks, not from bare names
        if self.doc_tools and db_type in ("sql", "sqlite", "postgresql", "mysql"):
            try:
                import duckdb

                conn = duckdb.connect(":memory:")
                if db_type == "sqlite" or uri.endswith(".db") or uri.endswith(".sqlite"):
                    conn.execute(f"ATTACH '{uri}' AS session_db (TYPE SQLITE)")
                    tables = conn.execute(
                        "SELECT table_name FROM information_schema.tables WHERE table_schema = 'session_db'"
                    ).fetchall()
                else:
                    conn.close()
                    tables = []

                if tables:
                    table_names = [t[0] for t in tables]
                    column_names = []
                    for table in table_names:
                        try:
                            cols = conn.execute(
                                f"SELECT column_name FROM information_schema.columns WHERE table_name = '{table}'"
                            ).fetchall()
                            column_names.extend([c[0] for c in cols])
                        except Exception:
                            pass

                    # Update schema entities for NER to recognize in documents
                    current_entities = self.doc_tools._schema_entities or []
                    new_entities = list(set(current_entities + table_names + column_names))
                    self.doc_tools._schema_entities = new_entities

                conn.close()
            except Exception:
                pass  # Non-fatal - database still added to session

        return True

    def add_file(
        self,
        name: str,
        uri: str,
        auth: str = "",
        description: str = "",
    ) -> bool:
        """Add a file to the current session.

        The file will be available as `file_<name>` in code execution.
        For local files, this is a Path. For HTTP files, content is fetched on-demand.
        Document files (md, txt, pdf, docx) are also indexed in the vector store.

        Args:
            name: File name (used as file_<name> variable)
            uri: File URI (file:// or http://)
            auth: Auth header for HTTP (e.g., "Bearer token123")
            description: Human-readable description

        Returns:
            True if added successfully
        """
        self.session_files[name] = {
            "uri": uri,
            "auth": auth,
            "description": description,
        }

        # Index document files in the vector store
        if self.doc_tools:
            doc_extensions = {'.md', '.txt', '.pdf', '.docx', '.html', '.htm', '.pptx'}
            from pathlib import Path

            # Handle file:// URIs
            file_path = uri
            if uri.startswith("file://"):
                file_path = uri[7:]

            path = Path(file_path)
            if path.suffix.lower() in doc_extensions and path.exists():
                try:
                    # Read file content based on type
                    if path.suffix.lower() == '.pdf':
                        from pypdf import PdfReader
                        reader = PdfReader(path)
                        content = "\n\n".join(
                            page.extract_text() for page in reader.pages if page.extract_text()
                        )
                    elif path.suffix.lower() == '.docx':
                        from docx import Document
                        doc = Document(path)
                        content = "\n\n".join(para.text for para in doc.paragraphs if para.text)
                    elif path.suffix.lower() == '.pptx':
                        from pptx import Presentation
                        prs = Presentation(path)
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    texts.append(shape.text)
                        content = "\n\n".join(texts)
                    else:
                        content = path.read_text()

                    # Detect format
                    format_map = {'.md': 'markdown', '.txt': 'text', '.html': 'html', '.htm': 'html'}
                    doc_format = format_map.get(path.suffix.lower(), 'text')

                    # Add as session document
                    self.doc_tools._add_document_internal(
                        name=f"session:{name}",
                        content=content,
                        doc_format=doc_format,
                        description=description,
                    )
                    logger.debug(f"Indexed document: session:{name}")
                except Exception as e:
                    logger.warning(f"Failed to index document {name}: {e}")

        return True

    def get_all_databases(self) -> dict[str, dict]:
        """Get all databases (config + session-added).

        Returns:
            Dict of name -> {type, uri, description, source}
        """
        from constat.storage.bookmarks import BookmarkStore

        result = {}

        # Config databases
        for name, db_config in self.config.databases.items():
            result[name] = {
                "type": db_config.type or "sql",
                "uri": db_config.uri or db_config.path or "",
                "description": db_config.description or "",
                "source": "config",
            }

        # Bookmarked databases
        bookmarks = BookmarkStore()
        for name, bm in bookmarks.list_databases().items():
            if name not in result:  # Don't override config
                result[name] = {
                    "type": bm["type"],
                    "uri": bm["uri"],
                    "description": bm["description"],
                    "source": "bookmark",
                }

        # Session databases
        for name, db in self.session_databases.items():
            result[name] = {
                "type": db["type"],
                "uri": db["uri"],
                "description": db["description"],
                "source": "session",
            }

        return result

    def get_all_files(self) -> dict[str, dict]:
        """Get all files (config documents + file sources + bookmarks + session).

        Returns:
            Dict of name -> {uri, description, auth, source, file_type}
        """
        from constat.storage.bookmarks import BookmarkStore

        result = {}

        # Config documents
        if self.config.documents:
            for name, doc_config in self.config.documents.items():
                uri = ""
                if doc_config.path:
                    uri = f"file://{doc_config.path}"
                elif doc_config.url:
                    uri = doc_config.url
                result[name] = {
                    "uri": uri,
                    "description": doc_config.description or "",
                    "auth": "",
                    "source": "config",
                    "file_type": "document",
                }

        # Config file-type databases (csv, json, parquet)
        for name, db_config in self.config.databases.items():
            if db_config.type in ("csv", "json", "jsonl", "parquet", "arrow", "feather"):
                path = db_config.path or db_config.uri or ""
                result[name] = {
                    "uri": f"file://{path}" if not path.startswith(("file://", "http")) else path,
                    "description": db_config.description or "",
                    "auth": "",
                    "source": "config",
                    "file_type": db_config.type,
                }

        # Bookmarked files
        bookmarks = BookmarkStore()
        for name, bm in bookmarks.list_files().items():
            if name not in result:  # Don't override config
                result[name] = {
                    "uri": bm["uri"],
                    "description": bm["description"],
                    "auth": bm.get("auth", ""),
                    "source": "bookmark",
                    "file_type": "file",
                }

        # Session files
        for name, f in self.session_files.items():
            result[name] = {
                "uri": f["uri"],
                "description": f["description"],
                "auth": f.get("auth", ""),
                "source": "session",
                "file_type": "file",
            }

        return result

    # --- Context Management ---

    def get_context_stats(self) -> Optional[ContextStats]:
        """
        Get statistics about context size.

        Returns:
            ContextStats with token estimates and breakdown, or None if no datastore
        """
        if not self.datastore:
            return None

        estimator = ContextEstimator(self.datastore)
        return estimator.estimate()

    def compact_context(
        self,
        summarize_scratchpad: bool = True,
        sample_tables: bool = True,
        clear_old_state: bool = False,
        keep_recent_steps: int = 3,
    ) -> Optional[CompactionResult]:
        """
        Compact session context to reduce token usage.

        This is useful for long-running sessions where context grows too large.

        Args:
            summarize_scratchpad: Truncate old scratchpad narratives
            sample_tables: Sample large tables down to max rows
            clear_old_state: Clear state variables from old steps
            keep_recent_steps: Number of recent steps to preserve intact

        Returns:
            CompactionResult with details, or None if no datastore
        """
        if not self.datastore:
            return None

        compactor = ContextCompactor(self.datastore)
        return compactor.compact(
            summarize_scratchpad=summarize_scratchpad,
            sample_tables=sample_tables,
            clear_old_state=clear_old_state,
            keep_recent_steps=keep_recent_steps,
        )

    def _auto_compact_if_needed(self) -> Optional[CompactionResult]:
        """
        Automatically compact context if it exceeds critical threshold.

        This is called after step execution to prevent context from growing
        too large for the LLM context window.

        Returns:
            CompactionResult if compaction was performed, None otherwise
        """
        if not self.datastore:
            return None

        stats = self.get_context_stats()
        if not stats or not stats.is_critical:
            return None

        # Context is critical - auto-compact
        self._emit_event(StepEvent(
            event_type="progress",
            step_number=0,
            data={"message": f"Auto-compacting context ({stats.total_tokens:,} tokens)..."}
        ))

        result = self.compact_context(
            summarize_scratchpad=True,
            sample_tables=True,
            clear_old_state=False,  # Conservative - don't clear state
            keep_recent_steps=5,    # Keep more steps for auto-compact
        )

        if result:
            self._emit_event(StepEvent(
                event_type="progress",
                step_number=0,
                data={"message": f"Context compacted: {result.tokens_before:,} → {result.tokens_after:,} tokens"}
            ))

        return result

    def reset_context(self) -> Optional[CompactionResult]:
        """
        Fully reset session context (clear all state).

        WARNING: This clears all scratchpad entries, tables, state variables,
        and artifacts. Use with caution.

        Returns:
            CompactionResult with details, or None if no datastore
        """
        if not self.datastore:
            return None

        compactor = ContextCompactor(self.datastore)
        result = compactor.clear_all()

        # Also reset in-memory scratchpad
        self.scratchpad = Scratchpad()
        self.plan = None

        return result

    # --- Saved Plans ---

    CONSTAT_BASE_DIR = Path(".constat")
    DEFAULT_USER_ID = "default"

    @classmethod
    def _get_user_plans_file(cls, user_id: str) -> Path:
        """Get path to user-scoped saved plans file."""
        return cls.CONSTAT_BASE_DIR / user_id / "saved_plans.json"

    @classmethod
    def _get_shared_plans_file(cls) -> Path:
        """Get path to shared plans file."""
        return cls.CONSTAT_BASE_DIR / "shared" / "saved_plans.json"

    def save_plan(self, name: str, problem: str, user_id: Optional[str] = None, shared: bool = False) -> None:
        """
        Save the current session's plan and code for future replay.

        Args:
            name: Name for the saved plan
            problem: The original problem (for replay context)
            user_id: User ID (defaults to DEFAULT_USER_ID)
            shared: If True, save as shared plan accessible to all users
        """
        if not self.datastore:
            raise ValueError("No datastore available")

        entries = self.datastore.get_scratchpad()
        if not entries:
            raise ValueError("No steps to save")

        user_id = user_id or self.DEFAULT_USER_ID

        plan_data = {
            "problem": problem,
            "created_by": user_id,
            "steps": [
                {
                    "step_number": e["step_number"],
                    "goal": e["goal"],
                    "code": e["code"],
                }
                for e in entries
            ],
        }

        if shared:
            plans = self._load_shared_plans()
            plans[name] = plan_data
            self._save_shared_plans(plans)
        else:
            plans = self._load_user_plans(user_id)
            plans[name] = plan_data
            self._save_user_plans(user_id, plans)

    @classmethod
    def load_saved_plan(cls, name: str, user_id: Optional[str] = None) -> dict:
        """
        Load a saved plan by name.

        Searches user's plans first, then shared plans.

        Args:
            name: Name of the saved plan
            user_id: User ID (defaults to DEFAULT_USER_ID)

        Returns:
            Dict with problem and steps
        """
        user_id = user_id or cls.DEFAULT_USER_ID

        # Check user's plans first
        user_plans = cls._load_user_plans(user_id)
        if name in user_plans:
            return user_plans[name]

        # Check shared plans
        shared_plans = cls._load_shared_plans()
        if name in shared_plans:
            return shared_plans[name]

        raise ValueError(f"No saved plan named '{name}'")

    @classmethod
    def list_saved_plans(cls, user_id: Optional[str] = None, include_shared: bool = True) -> list[dict]:
        """
        List saved plans accessible to the user.

        Args:
            user_id: User ID (defaults to DEFAULT_USER_ID)
            include_shared: Include shared plans in the list

        Returns:
            List of dicts with name, problem, shared flag
        """
        user_id = user_id or cls.DEFAULT_USER_ID
        result = []

        # User's plans
        user_plans = cls._load_user_plans(user_id)
        for name, data in user_plans.items():
            result.append({
                "name": name,
                "problem": data.get("problem", ""),
                "shared": False,
                "steps": len(data.get("steps", [])),
            })

        # Shared plans
        if include_shared:
            shared_plans = cls._load_shared_plans()
            for name, data in shared_plans.items():
                result.append({
                    "name": name,
                    "problem": data.get("problem", ""),
                    "shared": True,
                    "created_by": data.get("created_by", "unknown"),
                    "steps": len(data.get("steps", [])),
                })

        return result

    @classmethod
    def delete_saved_plan(cls, name: str, user_id: Optional[str] = None) -> bool:
        """Delete a saved plan by name (only user's own plans)."""
        user_id = user_id or cls.DEFAULT_USER_ID
        user_plans = cls._load_user_plans(user_id)

        if name not in user_plans:
            return False

        del user_plans[name]
        cls._save_user_plans(user_id, user_plans)
        return True

    @classmethod
    def share_plan_with(cls, name: str, target_user: str, from_user: Optional[str] = None) -> bool:
        """
        Share a plan with a specific user (copy to their plans).

        Args:
            name: Name of the plan to share
            target_user: User ID to share with
            from_user: Source user ID (defaults to DEFAULT_USER_ID)

        Returns:
            True if shared successfully
        """
        from_user = from_user or cls.DEFAULT_USER_ID

        # Find the plan (check user's plans first, then shared)
        source_plans = cls._load_user_plans(from_user)
        if name in source_plans:
            plan_data = source_plans[name].copy()
        else:
            shared_plans = cls._load_shared_plans()
            if name in shared_plans:
                plan_data = shared_plans[name].copy()
            else:
                return False

        # Copy to target user's plans
        target_plans = cls._load_user_plans(target_user)
        plan_data["shared_by"] = from_user
        target_plans[name] = plan_data
        cls._save_user_plans(target_user, target_plans)
        return True

    @classmethod
    def _load_user_plans(cls, user_id: str) -> dict:
        """Load saved plans for a specific user."""
        plans_file = cls._get_user_plans_file(user_id)
        if not plans_file.exists():
            return {}
        try:
            return json.loads(plans_file.read_text())
        except (json.JSONDecodeError, OSError) as e:
            logger.debug(f"Could not load user plans from {plans_file}: {e}")
            return {}

    @classmethod
    def _save_user_plans(cls, user_id: str, plans: dict) -> None:
        """Save plans to user-scoped file."""
        plans_file = cls._get_user_plans_file(user_id)
        plans_file.parent.mkdir(parents=True, exist_ok=True)
        plans_file.write_text(json.dumps(plans, indent=2))

    @classmethod
    def _load_shared_plans(cls) -> dict:
        """Load shared plans accessible to all users."""
        plans_file = cls._get_shared_plans_file()
        if not plans_file.exists():
            return {}
        try:
            return json.loads(plans_file.read_text())
        except (json.JSONDecodeError, OSError) as e:
            logger.debug(f"Could not load shared plans from {plans_file}: {e}")
            return {}

    @classmethod
    def _save_shared_plans(cls, plans: dict) -> None:
        """Save shared plans."""
        plans_file = cls._get_shared_plans_file()
        plans_file.parent.mkdir(parents=True, exist_ok=True)
        plans_file.write_text(json.dumps(plans, indent=2))

    def replay_saved(self, name: str, user_id: Optional[str] = None) -> dict:
        """
        Replay a saved plan by name.

        Args:
            name: Name of the saved plan
            user_id: User ID for plan lookup (defaults to DEFAULT_USER_ID)

        Returns:
            Dict with results (same format as solve())
        """
        plan_data = self.load_saved_plan(name, user_id=user_id)

        if not self.datastore:
            raise ValueError("No datastore available for replay")

        # Clear existing scratchpad and load saved steps
        # (We'll execute fresh but use stored code)
        problem = plan_data["problem"]
        steps = plan_data["steps"]

        # Emit plan ready
        self._emit_event(StepEvent(
            event_type="plan_ready",
            step_number=0,
            data={
                "steps": [
                    {"number": s["step_number"], "goal": s["goal"], "depends_on": []}
                    for s in steps
                ],
                "reasoning": f"Replaying saved plan: {name}",
                "is_followup": False,
            }
        ))

        all_results = []
        for step_data in steps:
            step_number = step_data["step_number"]
            goal = step_data["goal"]
            code = step_data["code"]

            if not code:
                raise ValueError(f"Step {step_number} has no stored code")

            self._emit_event(StepEvent(
                event_type="step_start",
                step_number=step_number,
                data={"goal": goal}
            ))

            self._emit_event(StepEvent(
                event_type="executing",
                step_number=step_number,
                data={"attempt": 1, "code": code}
            ))

            start_time = time.time()
            tables_before_list = self.datastore.list_tables()
            tables_before = set(t['name'] for t in tables_before_list)
            versions_before = {t['name']: t.get('version', 1) for t in tables_before_list}

            exec_globals = self._get_execution_globals()
            result = self.executor.execute(code, exec_globals)

            if result.success:
                self._auto_save_results(result.namespace, step_number)

            duration_ms = int((time.time() - start_time) * 1000)
            tables_after_list = self.datastore.list_tables()
            tables_after = set(t['name'] for t in tables_after_list)
            versions_after = {t['name']: t.get('version', 1) for t in tables_after_list}
            new_tables = tables_after - tables_before
            updated_tables = {
                name for name in tables_before & tables_after
                if versions_after.get(name, 1) > versions_before.get(name, 1)
                and not name.startswith('_')
            }
            tables_created = list(new_tables | updated_tables)

            if result.success:
                self._emit_event(StepEvent(
                    event_type="step_complete",
                    step_number=step_number,
                    data={
                        "goal": goal,
                        "code": code,
                        "stdout": result.stdout,
                        "attempts": 1,
                        "duration_ms": duration_ms,
                        "tables_created": tables_created,
                    }
                ))

                all_results.append(StepResult(
                    success=True,
                    stdout=result.stdout,
                    attempts=1,
                    duration_ms=duration_ms,
                    tables_created=tables_created,
                    code=code,
                ))
            else:
                self._emit_event(StepEvent(
                    event_type="step_error",
                    step_number=step_number,
                    data={"error": result.stderr or "Execution failed", "attempt": 1}
                ))
                return {
                    "success": False,
                    "error": result.stderr or "Replay execution failed",
                    "step_number": step_number,
                }

        # Synthesize answer (respects insights config)
        combined_output = "\n\n".join([
            f"Step {s['step_number']}: {s['goal']}\n{r.stdout}"
            for s, r in zip(steps, all_results)
        ])

        # Emit raw results first
        self._emit_event(StepEvent(
            event_type="raw_results_ready",
            step_number=0,
            data={"output": combined_output}
        ))

        # Check for created artifacts (to mention in synthesis)
        from constat.visualization.output import peek_pending_outputs
        pending_artifacts = peek_pending_outputs()

        # Check if insights are enabled
        skip_insights = not self.session_config.enable_insights

        if skip_insights:
            final_answer = combined_output
            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer, "brief": True}
            ))
        else:
            self._emit_event(StepEvent(
                event_type="synthesizing",
                step_number=0,
                data={"message": "Synthesizing final answer..."}
            ))

            final_answer = self._synthesize_answer(problem, combined_output, pending_artifacts)

            self._emit_event(StepEvent(
                event_type="answer_ready",
                step_number=0,
                data={"answer": final_answer}
            ))

        total_duration = sum(r.duration_ms for r in all_results)

        return {
            "success": True,
            "results": all_results,
            "output": combined_output,
            "final_answer": final_answer,
            "datastore_tables": self.datastore.list_tables(),
            "duration_ms": total_duration,
            "replay": True,
            "plan_name": name,
        }

    def get_unresolved_facts(self) -> list[dict]:
        """Get list of facts that could not be resolved."""
        return [f.to_dict() for f in self.fact_resolver.get_unresolved_facts()]

    def get_unresolved_summary(self) -> str:
        """Get human-readable summary of unresolved facts."""
        return self.fact_resolver.get_unresolved_summary()

    def provide_facts(self, user_text: str) -> dict:
        """
        Extract facts from user text and add to resolver cache.

        This is used in auditable mode when facts could not be resolved.
        The user provides facts in natural language, and the LLM extracts
        them into structured facts that can be used for re-resolution.

        Example:
            session.provide_facts("There were 1 million people at the march")
            # Extracts: march_attendance = 1000000

        Args:
            user_text: Natural language text containing facts

        Returns:
            Dict with:
                - extracted_facts: List of facts extracted and added
                - unresolved_remaining: List of still-unresolved facts
        """
        # Extract facts from user text
        extracted = self.fact_resolver.add_user_facts_from_text(user_text)

        # Clear unresolved facts to allow re-resolution
        self.fact_resolver.clear_unresolved()

        return {
            "extracted_facts": [f.to_dict() for f in extracted],
            "unresolved_remaining": [f.to_dict() for f in self.fact_resolver.get_unresolved_facts()],
        }

    def add_fact(self, fact_name: str, value, reasoning: str = None, **params) -> dict:
        """
        Explicitly add a fact to the resolver cache.

        This is a more direct way to provide facts than provide_facts(),
        useful when you know the exact fact name and value.

        Args:
            fact_name: Name of the fact (e.g., "march_attendance")
            value: The value to set
            reasoning: Optional explanation
            **params: Additional parameters for the fact

        Returns:
            Dict with the created fact
        """
        # Use current role context if not explicitly provided
        if "role_id" not in params and self._current_role_id:
            params["role_id"] = self._current_role_id

        fact = self.fact_resolver.add_user_fact(
            fact_name=fact_name,
            value=value,
            reasoning=reasoning,
            **params,
        )
        return fact.to_dict()

    # =========================================================================
    # Dynamic Role and Skill Selection
    # =========================================================================

    def match_role_for_query(self, query: str) -> Optional[str]:
        """Match a query to the best-fitting role using semantic similarity.

        Args:
            query: User's natural language query

        Returns:
            Role name if matched, None if no match (use shared context)
        """

        match = self.role_matcher.match(query)
        if match:
            self._current_role_id = match.role.name
            logger.info(f"Matched role '{match.role.name}' for query (similarity: {match.similarity:.2f})")
            return match.role.name

        self._current_role_id = None
        return None

    def match_skills_for_query(self, query: str) -> list[str]:
        """Match a query to relevant skills using semantic similarity.

        Args:
            query: User's natural language query

        Returns:
            List of matched skill names (may be empty)
        """
        matches = self.skill_matcher.match(query)
        return [m.skill.name for m in matches]

    def get_dynamic_context(self, query: str) -> dict:
        """Get the dynamically selected role and skills for a query.

        This is the main entry point for dynamic selection. It:
        1. Matches the query to skills (multiple selection)
        2. Checks if any skill specifies a role/agent context
        3. If skill specifies role → use that role
        4. Otherwise → matches the query to a role (single selection)
        5. Returns the combined context for prompt building

        Args:
            query: User's natural language query

        Returns:
            Dict with:
                - role: Optional[dict] with name, description, similarity
                - skills: List of dicts with name, prompt, description
                - role_prompt: Combined role prompt content
                - skills_prompt: Combined skills prompt content
                - role_source: "skill" if role came from skill, "query" if from query match
        """
        # Step 1: Match skills first
        skill_matches = self.skill_matcher.match(query)
        skills_info = []
        skills_prompts = []
        skill_specified_role = None

        for match in skill_matches:
            skills_info.append({
                "name": match.skill.name,
                "description": match.skill.description,
                "similarity": match.similarity,
            })
            skills_prompts.append(f"## {match.skill.name}\n{match.skill.prompt}")

            # Check if skill specifies a role/agent
            # Uses the 'agent' field from SKILL.md frontmatter
            if match.skill.agent and not skill_specified_role:
                skill_specified_role = match.skill.agent
                logger.info(f"Skill '{match.skill.name}' specifies agent/role: {skill_specified_role}")

        if skill_matches:
            skill_names = [m.skill.name for m in skill_matches]
            similarities = [f"{m.skill.name}({m.similarity:.2f})" for m in skill_matches]
            logger.info(f"[CONTEXT] Selected skills: {similarities}")
        else:
            logger.info("[CONTEXT] No skills matched for query")

        # Step 2: Determine role
        role_info = None
        role_prompt = ""
        role_source = None

        if skill_specified_role:
            # Skill specified a role - try to use it
            role = self.role_manager.get_role(skill_specified_role)
            if role:
                self._current_role_id = role.name
                role_info = {
                    "name": role.name,
                    "description": role.description,
                    "similarity": 1.0,  # Explicit specification = full match
                }
                role_prompt = role.prompt
                role_source = "skill"
                logger.info(f"[CONTEXT] Selected role: {role.name} (specified by skill)")
            else:
                logger.warning(f"Skill specified role '{skill_specified_role}' not found, falling back to query match")

        if not role_info:
            # No skill-specified role, match based on query
            role_match = self.role_matcher.match(query)
            if role_match:
                self._current_role_id = role_match.role.name
                role_info = {
                    "name": role_match.role.name,
                    "description": role_match.role.description,
                    "similarity": role_match.similarity,
                }
                role_prompt = role_match.role.prompt
                role_source = "query"
                logger.info(f"[CONTEXT] Selected role: {role_match.role.name} (similarity: {role_match.similarity:.2f})")
            else:
                self._current_role_id = None
                logger.info("[CONTEXT] No role matched for query")

        # Step 3: Merge role-declared skills
        # If the selected role declares explicit skills, add them if not already matched
        if role_info:
            role_obj = self.role_manager.get_role(role_info["name"])
            if role_obj and role_obj.skills:
                matched_skill_names = {s["name"] for s in skills_info}
                for skill_name in role_obj.skills:
                    if skill_name not in matched_skill_names:
                        skill_obj = self.skill_manager.get_skill(skill_name)
                        if skill_obj:
                            skills_info.append({
                                "name": skill_obj.name,
                                "description": skill_obj.description,
                                "similarity": 1.0,
                                "source": "role",
                            })
                            skills_prompts.append(
                                f"## {skill_obj.name} (required by role: {role_info['name']})\n{skill_obj.prompt}"
                            )
                            logger.info(f"[CONTEXT] Added role-declared skill: {skill_name}")
                        else:
                            logger.warning(f"Role '{role_info['name']}' declares skill '{skill_name}' but it was not found")

        return {
            "role": role_info,
            "skills": skills_info,
            "role_prompt": role_prompt,
            "skills_prompt": "\n\n".join(skills_prompts),
            "role_source": role_source,
        }

    @property
    def current_role_id(self) -> Optional[str]:
        """Get the current role ID for this query context."""
        return self._current_role_id

    def set_current_role(self, role_name: Optional[str]) -> bool:
        """Manually set the current role (override dynamic selection).

        Args:
            role_name: Role name or None to clear

        Returns:
            True if successful, False if role not found
        """
        if role_name is None:
            self._current_role_id = None
            self.role_manager.set_active_role(None)
            return True

        if self.role_manager.set_active_role(role_name):
            self._current_role_id = role_name
            return True
        return False


def create_session(config_path: str, session_id: str) -> Session:
    """Create a session from a config file path.

    Args:
        config_path: Path to config YAML file
        session_id: Client-provided session ID (required)
    """
    config = Config.from_yaml(config_path)
    return Session(config, session_id=session_id)
