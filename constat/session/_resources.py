# Copyright (c) 2025 Kenneth Stott
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Resources mixin: database/file management, context, agents, facts."""

from __future__ import annotations

import logging
from typing import Optional

from constat.context import ContextEstimator, ContextCompactor, ContextStats, CompactionResult
from constat.execution.scratchpad import Scratchpad
from constat.session._types import StepEvent

logger = logging.getLogger(__name__)


# noinspection PyUnresolvedReferences
class ResourcesMixin:

    def add_database(
        self,
        name: str,
        db_type: str,
        uri: str,
        description: str = "",
    ) -> bool:
        """Add a database to the current session.

        The database will be available as `db_<name>` in code execution.
        The schema is introspected and table/column names are added as
        session entities for entity extraction.

        Args:
            name: Database name (used as db_<name> variable)
            db_type: Database type (sql, csv, json, parquet, mongodb, etc.)
            uri: Connection URI or file path
            description: Human-readable description

        Returns:
            True if added successfully
        """
        self.session_databases[name] = {
            "type": db_type,
            "uri": uri,
            "description": description,
        }

        # Update schema entities list for NER recognition in future document indexing
        # Entities are created via NER extraction from description chunks, not from bare names
        if self.doc_tools and db_type in ("sql", "sqlite", "postgresql", "mysql"):
            try:
                import duckdb

                conn = duckdb.connect(":memory:")
                if db_type == "sqlite" or uri.endswith(".db") or uri.endswith(".sqlite"):
                    conn.execute(f"ATTACH '{uri}' AS session_db (TYPE SQLITE)")
                    tables = conn.execute(
                        "SELECT table_name FROM information_schema.tables WHERE table_schema = 'session_db'"
                    ).fetchall()
                else:
                    conn.close()
                    tables = []

                if tables:
                    table_names = [t[0] for t in tables]
                    column_names = []
                    for table in table_names:
                        try:
                            cols = conn.execute(
                                f"SELECT column_name FROM information_schema.columns WHERE table_name = '{table}'"
                            ).fetchall()
                            column_names.extend([c[0] for c in cols])
                        except duckdb.Error:
                            pass

                    # Update schema entities for NER to recognize in documents
                    current_entities = self.doc_tools._schema_entities or []
                    new_entities = list(set(current_entities + table_names + column_names))
                    self.doc_tools._schema_entities = new_entities

                conn.close()
            # noinspection PyBroadException - duckdb import is conditional, can't narrow
            except Exception:
                pass  # Non-fatal - database still added to session

        return True

    def add_file(
        self,
        name: str,
        uri: str,
        auth: str = "",
        description: str = "",
    ) -> bool:
        """Add a file to the current session.

        The file will be available as `file_<name>` in code execution.
        For local files, this is a Path. For HTTP files, content is fetched on-demand.
        Document files (md, txt, pdf, docx) are also indexed in the vector store.

        Args:
            name: File name (used as file_<name> variable)
            uri: File URI (file:// or http://)
            auth: Auth header for HTTP (e.g., "Bearer token123")
            description: Human-readable description

        Returns:
            True if added successfully
        """
        self.session_files[name] = {
            "uri": uri,
            "auth": auth,
            "description": description,
        }

        # Index document files in the vector store
        if self.doc_tools:
            doc_extensions = {'.md', '.txt', '.pdf', '.docx', '.html', '.htm', '.pptx'}
            from pathlib import Path

            # Handle file:// URIs
            file_path = uri
            if uri.startswith("file://"):
                file_path = uri[7:]

            path = Path(file_path)
            if path.suffix.lower() in doc_extensions and path.exists():
                try:
                    # Read file content based on type
                    if path.suffix.lower() == '.pdf':
                        from pypdf import PdfReader
                        reader = PdfReader(str(path))
                        content = "\n\n".join(
                            page.extract_text() for page in reader.pages if page.extract_text()
                        )
                    elif path.suffix.lower() == '.docx':
                        from docx import Document
                        doc = Document(str(path))
                        content = "\n\n".join(para.text for para in doc.paragraphs if para.text)
                    elif path.suffix.lower() == '.pptx':
                        from pptx import Presentation
                        prs = Presentation(str(path))
                        texts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text:
                                    texts.append(shape.text)
                        content = "\n\n".join(texts)
                    else:
                        content = path.read_text()

                    # Detect format
                    format_map = {'.md': 'markdown', '.txt': 'text', '.html': 'html', '.htm': 'html'}
                    doc_format = format_map.get(path.suffix.lower(), 'text')

                    # Add as session document
                    self.doc_tools._add_document_internal(
                        name=f"session:{name}",
                        content=content,
                        doc_format=doc_format,
                        description=description,
                    )
                    logger.debug(f"Indexed document: session:{name}")
                except Exception as e:
                    logger.warning(f"Failed to index document {name}: {e}")

        return True

    def get_all_databases(self) -> dict[str, dict]:
        """Get all databases (config + session-added).

        Returns:
            Dict of name -> {type, uri, description, source}
        """
        from constat.storage.bookmarks import BookmarkStore

        result = {}

        # Config databases
        for name, db_config in self.config.databases.items():
            result[name] = {
                "type": db_config.type or "sql",
                "uri": db_config.uri or db_config.path or "",
                "description": db_config.description or "",
                "source": "config",
            }

        # Bookmarked databases
        bookmarks = BookmarkStore()
        for name, bm in bookmarks.list_databases().items():
            if name not in result:  # Don't override config
                result[name] = {
                    "type": bm["type"],
                    "uri": bm["uri"],
                    "description": bm["description"],
                    "source": "bookmark",
                }

        # Session databases
        for name, db in self.session_databases.items():
            result[name] = {
                "type": db["type"],
                "uri": db["uri"],
                "description": db["description"],
                "source": "session",
            }

        return result

    def get_all_files(self) -> dict[str, dict]:
        """Get all files (config documents + file sources + bookmarks + session).

        Returns:
            Dict of name -> {uri, description, auth, source, file_type}
        """
        from constat.storage.bookmarks import BookmarkStore

        result = {}

        # Config documents
        if self.config.documents:
            for name, doc_config in self.config.documents.items():
                uri = ""
                if doc_config.path:
                    uri = f"file://{doc_config.path}"
                elif doc_config.url:
                    uri = doc_config.url
                result[name] = {
                    "uri": uri,
                    "description": doc_config.description or "",
                    "auth": "",
                    "source": "config",
                    "file_type": "document",
                }

        # Config file-type databases (csv, json, parquet)
        for name, db_config in self.config.databases.items():
            if db_config.type in ("csv", "json", "jsonl", "parquet", "arrow", "feather"):
                path = db_config.path or db_config.uri or ""
                result[name] = {
                    "uri": f"file://{path}" if not path.startswith(("file://", "http")) else path,
                    "description": db_config.description or "",
                    "auth": "",
                    "source": "config",
                    "file_type": db_config.type,
                }

        # Bookmarked files
        bookmarks = BookmarkStore()
        for name, bm in bookmarks.list_files().items():
            if name not in result:  # Don't override config
                result[name] = {
                    "uri": bm["uri"],
                    "description": bm["description"],
                    "auth": bm.get("auth", ""),
                    "source": "bookmark",
                    "file_type": "file",
                }

        # Session files
        for name, f in self.session_files.items():
            result[name] = {
                "uri": f["uri"],
                "description": f["description"],
                "auth": f.get("auth", ""),
                "source": "session",
                "file_type": "file",
            }

        return result

    # --- Context Management ---

    def get_context_stats(self) -> Optional[ContextStats]:
        """
        Get statistics about context size.

        Returns:
            ContextStats with token estimates and breakdown, or None if no datastore
        """
        if not self.datastore:
            return None

        estimator = ContextEstimator(self.datastore)
        return estimator.estimate()

    def compact_context(
        self,
        summarize_scratchpad: bool = True,
        sample_tables: bool = True,
        clear_old_state: bool = False,
        keep_recent_steps: int = 3,
    ) -> Optional[CompactionResult]:
        """
        Compact session context to reduce token usage.

        This is useful for long-running sessions where context grows too large.

        Args:
            summarize_scratchpad: Truncate old scratchpad narratives
            sample_tables: Sample large tables down to max rows
            clear_old_state: Clear state variables from old steps
            keep_recent_steps: Number of recent steps to preserve intact

        Returns:
            CompactionResult with details, or None if no datastore
        """
        if not self.datastore:
            return None

        compactor = ContextCompactor(self.datastore)
        return compactor.compact(
            summarize_scratchpad=summarize_scratchpad,
            sample_tables=sample_tables,
            clear_old_state=clear_old_state,
            keep_recent_steps=keep_recent_steps,
        )

    def _auto_compact_if_needed(self) -> Optional[CompactionResult]:
        """
        Automatically compact context if it exceeds critical threshold.

        This is called after step execution to prevent context from growing
        too large for the LLM context window.

        Returns:
            CompactionResult if compaction was performed, None otherwise
        """
        if not self.datastore:
            return None

        stats = self.get_context_stats()
        if not stats or not stats.is_critical:
            return None

        # Context is critical - auto-compact
        self._emit_event(StepEvent(
            event_type="progress",
            step_number=0,
            data={"message": f"Auto-compacting context ({stats.total_tokens:,} tokens)..."}
        ))

        result = self.compact_context(
            summarize_scratchpad=True,
            sample_tables=True,
            clear_old_state=False,  # Conservative - don't clear state
            keep_recent_steps=5,    # Keep more steps for auto-compact
        )

        if result:
            self._emit_event(StepEvent(
                event_type="progress",
                step_number=0,
                data={"message": f"Context compacted: {result.tokens_before:,} → {result.tokens_after:,} tokens"}
            ))

        return result

    def reset_context(self) -> Optional[CompactionResult]:
        """
        Fully reset session context (clear all state).

        WARNING: This clears all scratchpad entries, tables, state variables,
        and artifacts. Use with caution.

        Returns:
            CompactionResult with details, or None if no datastore
        """
        if not self.datastore:
            return None

        compactor = ContextCompactor(self.datastore)
        result = compactor.clear_all()

        # Also reset in-memory scratchpad
        # noinspection PyAttributeOutsideInit
        self.scratchpad = Scratchpad()
        # noinspection PyAttributeOutsideInit
        self.plan = None

        return result

    # --- Facts ---

    def get_unresolved_facts(self) -> list[dict]:
        """Get list of facts that could not be resolved."""
        return [f.to_dict() for f in self.fact_resolver.get_unresolved_facts()]

    def get_unresolved_summary(self) -> str:
        """Get human-readable summary of unresolved facts."""
        return self.fact_resolver.get_unresolved_summary()

    def provide_facts(self, user_text: str) -> dict:
        """
        Extract facts from user text and add to resolver cache.

        This is used in auditable mode when facts could not be resolved.
        The user provides facts in natural language, and the LLM extracts
        them into structured facts that can be used for re-resolution.

        Example:
            session.provide_facts("There were 1 million people at the march")
            # Extracts: march_attendance = 1000000

        Args:
            user_text: Natural language text containing facts

        Returns:
            Dict with:
                - extracted_facts: List of facts extracted and added
                - unresolved_remaining: List of still-unresolved facts
        """
        # Extract facts from user text
        extracted = self.fact_resolver.add_user_facts_from_text(user_text)

        # Clear unresolved facts to allow re-resolution
        self.fact_resolver.clear_unresolved()

        return {
            "extracted_facts": [f.to_dict() for f in extracted],
            "unresolved_remaining": [f.to_dict() for f in self.fact_resolver.get_unresolved_facts()],
        }

    def add_fact(self, fact_name: str, value, reasoning: str = None, **params) -> dict:
        """
        Explicitly add a fact to the resolver cache.

        This is a more direct way to provide facts than provide_facts(),
        useful when you know the exact fact name and value.

        Args:
            fact_name: Name of the fact (e.g., "march_attendance")
            value: The value to set
            reasoning: Optional explanation
            **params: Additional parameters for the fact

        Returns:
            Dict with the created fact
        """
        # Use current agent context if not explicitly provided
        if "role_id" not in params and self._current_agent_id:
            params["role_id"] = self._current_agent_id

        fact = self.fact_resolver.add_user_fact(
            fact_name=fact_name,
            value=value,
            reasoning=reasoning,
            **params,
        )
        return fact.to_dict()

    # =========================================================================
    # Dynamic Agent and Skill Selection
    # =========================================================================

    def match_agent_for_query(self, query: str) -> Optional[str]:
        """Match a query to the best-fitting agent using semantic similarity.

        Args:
            query: User's natural language query

        Returns:
            Agent name if matched, None if no match (use shared context)
        """

        match = self.agent_matcher.match(query)
        if match:
            # noinspection PyAttributeOutsideInit
            self._current_agent_id = match.agent.name
            logger.info(f"Matched agent '{match.agent.name}' for query (similarity: {match.similarity:.2f})")
            return match.agent.name

        # noinspection PyAttributeOutsideInit
        self._current_agent_id = None
        return None

    def match_skills_for_query(self, query: str) -> list[str]:
        """Match a query to relevant skills using semantic similarity.

        Args:
            query: User's natural language query

        Returns:
            List of matched skill names (maybe empty)
        """
        matches = self.skill_matcher.match(query)
        return [m.skill.name for m in matches]

    def get_dynamic_context(self, query: str) -> dict:
        """Get the dynamically selected agent and skills for a query.

        This is the main entry point for dynamic selection. It:
        1. Matches the query to skills (multiple selection)
        2. Checks if any skill specifies an agent context
        3. If skill specifies agent -> use that agent
        4. Otherwise -> matches the query to an agent (single selection)
        5. Returns the combined context for prompt building

        Args:
            query: User's natural language query

        Returns:
            Dict with:
                - agent: Optional[dict] with name, description, similarity
                - skills: List of dicts with name, prompt, description
                - agent_prompt: Combined agent prompt content
                - skills_prompt: Combined skills prompt content
                - agent_source: "skill" if agent came from skill, "query" if from query match
        """
        # Step 1: Match skills first
        skill_matches = self.skill_matcher.match(query)
        skills_info = []
        skills_prompts = []
        skill_specified_agent = None

        for match in skill_matches:
            skills_info.append({
                "name": match.skill.name,
                "description": match.skill.description,
                "similarity": match.similarity,
            })
            skills_prompts.append(f"## {match.skill.name}\n{match.skill.prompt}")

            # Check if skill specifies an agent
            # Uses the 'agent' field from SKILL.md frontmatter
            if match.skill.agent and not skill_specified_agent:
                skill_specified_agent = match.skill.agent
                logger.info(f"Skill '{match.skill.name}' specifies agent: {skill_specified_agent}")

        if skill_matches:
            similarities = [f"{m.skill.name}({m.similarity:.2f})" for m in skill_matches]
            logger.info(f"[CONTEXT] Selected skills: {similarities}")
        else:
            logger.info("[CONTEXT] No skills matched for query")

        # Step 2: Determine agent
        agent_info = None
        agent_prompt = ""
        agent_source = None

        if skill_specified_agent:
            # Skill specified an agent - try to use it
            agent = self.agent_manager.get_agent(skill_specified_agent)
            if agent:
                self._current_agent_id = agent.name
                agent_info = {
                    "name": agent.name,
                    "description": agent.description,
                    "similarity": 1.0,  # Explicit specification = full match
                }
                agent_prompt = agent.prompt
                agent_source = "skill"
                logger.info(f"[CONTEXT] Selected agent: {agent.name} (specified by skill)")
            else:
                logger.warning(f"Skill specified agent '{skill_specified_agent}' not found, falling back to query match")

        if not agent_info:
            # No skill-specified agent, match based on query
            agent_match = self.agent_matcher.match(query)
            if agent_match:
                self._current_agent_id = agent_match.agent.name
                agent_info = {
                    "name": agent_match.agent.name,
                    "description": agent_match.agent.description,
                    "similarity": agent_match.similarity,
                }
                agent_prompt = agent_match.agent.prompt
                agent_source = "query"
                logger.info(f"[CONTEXT] Selected agent: {agent_match.agent.name} (similarity: {agent_match.similarity:.2f})")
            else:
                # noinspection PyAttributeOutsideInit
                self._current_agent_id = None
                logger.info("[CONTEXT] No agent matched for query")

        # Step 3: Merge agent-declared skills
        # If the selected agent declares explicit skills, add them if not already matched
        if agent_info:
            agent_obj = self.agent_manager.get_agent(agent_info["name"])
            if agent_obj and agent_obj.skills:
                matched_skill_names = {s["name"] for s in skills_info}
                for skill_name in agent_obj.skills:
                    if skill_name not in matched_skill_names:
                        skill_obj = self.skill_manager.get_skill(skill_name)
                        if skill_obj:
                            skills_info.append({
                                "name": skill_obj.name,
                                "description": skill_obj.description,
                                "similarity": 1.0,
                                "source": "agent",
                            })
                            skills_prompts.append(
                                f"## {skill_obj.name} (required by agent: {agent_info['name']})\n{skill_obj.prompt}"
                            )
                            logger.info(f"[CONTEXT] Added agent-declared skill: {skill_name}")
                        else:
                            logger.warning(f"Agent '{agent_info['name']}' declares skill '{skill_name}' but it was not found")

        return {
            "agent": agent_info,
            "skills": skills_info,
            "agent_prompt": agent_prompt,
            "skills_prompt": "\n\n".join(skills_prompts),
            "agent_source": agent_source,
        }

    @property
    def current_agent_id(self) -> Optional[str]:
        """Get the current agent ID for this query context."""
        return self._current_agent_id

    def set_current_agent(self, agent_name: Optional[str]) -> bool:
        """Manually set the current agent (override dynamic selection).

        Args:
            agent_name: Agent name or None to clear

        Returns:
            True if successful, False if agent not found
        """
        if agent_name is None:
            # noinspection PyAttributeOutsideInit
            self._current_agent_id = None
            self.agent_manager.set_active_agent(None)
            return True

        if self.agent_manager.set_active_agent(agent_name):
            # noinspection PyAttributeOutsideInit
            self._current_agent_id = agent_name
            return True
        return False
