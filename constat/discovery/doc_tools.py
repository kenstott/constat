# Copyright (c) 2025 Kenneth Stott
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""Document discovery tools for reference documents.

These tools allow the LLM to discover and search reference documents
on-demand rather than loading everything into the system prompt.
"""

import glob as glob_module
import hashlib
import logging
import threading
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

from constat.core.config import Config, DocumentConfig
from constat.embedding_loader import EmbeddingModelLoader
from constat.discovery.models import (
    DocumentChunk,
    LoadedDocument,
    StructuredFileSchema,
    ChunkEntity,
)
from constat.discovery.vector_store import (
    VectorStoreBackend,
    create_vector_store,
)
from constat.discovery.entity_extractor import EntityExtractor


def _is_glob_pattern(path: str) -> bool:
    """Check if a path contains glob pattern characters."""
    return any(c in path for c in ["*", "?", "[", "]"])


def _expand_file_paths(path: str) -> list[tuple[str, Path]]:
    """
    Expand a path that may be a glob pattern, directory, or single file.

    Args:
        path: File path, glob pattern, or directory path

    Returns:
        List of (display_name, resolved_path) tuples
    """
    p = Path(path)

    # Case 1: Glob pattern
    if _is_glob_pattern(path):
        matches = sorted(glob_module.glob(path, recursive=True))
        return [(Path(m).name, Path(m)) for m in matches if Path(m).is_file()]

    # Case 2: Directory - list all files
    if p.is_dir():
        files = []
        for f in sorted(p.iterdir()):
            if f.is_file() and not f.name.startswith("."):
                files.append((f.name, f))
        return files

    # Case 3: Single file
    if p.exists():
        return [(p.name, p)]

    # Path doesn't exist yet - return as-is for later error handling
    return [(p.name, p)]



def _infer_csv_schema(filepath: Path, sample_rows: int = 100) -> StructuredFileSchema:
    """Infer schema from a CSV file."""
    import csv

    columns = []
    row_count = 0

    with open(filepath, 'r', newline='', encoding='utf-8', errors='replace') as f:
        reader = csv.reader(f)
        try:
            headers = next(reader)
        except StopIteration:
            return StructuredFileSchema(
                filename=filepath.name,
                filepath=str(filepath),
                file_format="csv",
                row_count=0,
                columns=[],
            )

        # Initialize column info
        col_data = {h: {'name': h, 'values': []} for h in headers}

        # Sample rows to infer types and collect sample values
        for i, row in enumerate(reader):
            row_count += 1
            if i < sample_rows:
                for j, val in enumerate(row):
                    if j < len(headers):
                        col_data[headers[j]]['values'].append(val)

        # Count remaining rows
        for _ in reader:
            row_count += 1

    # Infer types and get sample values
    for header in headers:
        values = col_data[header]['values']
        col_type = _infer_column_type(values)
        unique_values = list(set(v for v in values if v))[:10]

        columns.append({
            'name': header,
            'type': col_type,
            'sample_values': unique_values,
        })

    return StructuredFileSchema(
        filename=filepath.name,
        filepath=str(filepath),
        file_format="csv",
        row_count=row_count,
        columns=columns,
    )


def _infer_json_schema(filepath: Path, sample_docs: int = 100) -> StructuredFileSchema:
    """Infer schema from a JSON file (array of objects or single object)."""
    import json

    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        try:
            data = json.load(f)
        except json.JSONDecodeError:
            return StructuredFileSchema(
                filename=filepath.name,
                filepath=str(filepath),
                file_format="json",
                row_count=0,
                columns=[],
            )

    # Handle array of objects vs single object
    if isinstance(data, list):
        docs = data[:sample_docs]
        row_count = len(data)
    elif isinstance(data, dict):
        docs = [data]
        row_count = 1
    else:
        return StructuredFileSchema(
            filename=filepath.name,
            filepath=str(filepath),
            file_format="json",
            row_count=1,
            columns=[],
        )

    # Collect all keys and their values
    key_values: dict[str, list] = {}
    for doc in docs:
        if isinstance(doc, dict):
            for key, val in doc.items():
                if key not in key_values:
                    key_values[key] = []
                key_values[key].append(val)

    # Build column info
    columns = []
    for key, values in key_values.items():
        col_type = _infer_json_value_type(values)
        # Get sample values (only for simple types)
        sample_values = []
        for v in values[:10]:
            if isinstance(v, (str, int, float, bool)) and v is not None:
                sample_values.append(str(v) if not isinstance(v, str) else v)
        unique_samples = list(set(sample_values))[:10]

        columns.append({
            'name': key,
            'type': col_type,
            'sample_values': unique_samples,
        })

    return StructuredFileSchema(
        filename=filepath.name,
        filepath=str(filepath),
        file_format="json",
        row_count=row_count,
        columns=columns,
    )


def _infer_column_type(values: list[str]) -> str:
    """Infer column type from string values."""
    if not values:
        return "unknown"

    # Check for numeric types
    int_count = 0
    float_count = 0
    date_count = 0

    for v in values:
        if not v:
            continue
        try:
            int(v)
            int_count += 1
            continue
        except ValueError:
            pass
        try:
            float(v)
            float_count += 1
            continue
        except ValueError:
            pass
        # Simple date check
        if len(v) == 10 and v[4:5] == '-' and v[7:8] == '-':
            date_count += 1

    non_empty = len([v for v in values if v])
    if non_empty == 0:
        return "unknown"

    if int_count == non_empty:
        return "integer"
    if int_count + float_count == non_empty:
        return "float"
    if date_count > non_empty * 0.8:
        return "date"
    return "string"


def _infer_json_value_type(values: list) -> str:
    """Infer type from JSON values."""
    if not values:
        return "unknown"

    types = set()
    for v in values:
        if v is None:
            continue
        elif isinstance(v, bool):
            types.add("boolean")
        elif isinstance(v, int):
            types.add("integer")
        elif isinstance(v, float):
            types.add("float")
        elif isinstance(v, str):
            types.add("string")
        elif isinstance(v, list):
            types.add("array")
        elif isinstance(v, dict):
            types.add("object")

    if len(types) == 0:
        return "null"
    if len(types) == 1:
        return types.pop()
    if types == {"integer", "float"}:
        return "float"
    return "mixed"


def _infer_structured_schema(filepath: Path, description: Optional[str] = None) -> Optional[StructuredFileSchema]:
    """Infer schema for a structured file based on its extension."""
    suffix = filepath.suffix.lower()

    if suffix == ".csv":
        schema = _infer_csv_schema(filepath)
    elif suffix == ".json":
        schema = _infer_json_schema(filepath)
    elif suffix == ".jsonl":
        # JSON Lines - read first N lines as separate JSON objects
        schema = _infer_jsonl_schema(filepath)
    else:
        return None

    schema.description = description
    return schema


def _infer_jsonl_schema(filepath: Path, sample_lines: int = 100) -> StructuredFileSchema:
    """Infer schema from a JSON Lines file."""
    import json

    docs = []
    row_count = 0

    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        for i, line in enumerate(f):
            row_count += 1
            if i < sample_lines and line.strip():
                try:
                    docs.append(json.loads(line))
                except json.JSONDecodeError:
                    pass

    # Collect all keys and their values
    key_values: dict[str, list] = {}
    for doc in docs:
        if isinstance(doc, dict):
            for key, val in doc.items():
                if key not in key_values:
                    key_values[key] = []
                key_values[key].append(val)

    # Build column info
    columns = []
    for key, values in key_values.items():
        col_type = _infer_json_value_type(values)
        sample_values = []
        for v in values[:10]:
            if isinstance(v, (str, int, float, bool)) and v is not None:
                sample_values.append(str(v) if not isinstance(v, str) else v)
        unique_samples = list(set(sample_values))[:10]

        columns.append({
            'name': key,
            'type': col_type,
            'sample_values': unique_samples,
        })

    return StructuredFileSchema(
        filename=filepath.name,
        filepath=str(filepath),
        file_format="jsonl",
        row_count=row_count,
        columns=columns,
    )


class DocumentDiscoveryTools:
    """Tools for discovering and searching reference documents on-demand.

    Supports incremental updates - only reloads documents that have changed
    based on file modification times and content hashes.
    """

    EMBEDDING_MODEL = "BAAI/bge-large-en-v1.5"
    # bge-large-en-v1.5 has max_seq_length=512 tokens (~2048 chars)
    # Use 1500 chars to stay within limit while maximizing context
    CHUNK_SIZE = 1500
    CACHE_FILENAME = "doc_index_cache.json"

    def __init__(
        self,
        config: Config,
        cache_dir: Optional[Path] = None,
        vector_store: Optional[VectorStoreBackend] = None,
        schema_entities: Optional[list[str]] = None,
        allowed_documents: Optional[set[str]] = None,
        skip_auto_index: bool = False,
    ):
        """Initialize document discovery tools.

        Args:
            config: Config with document definitions
            cache_dir: Directory for caching document metadata
            vector_store: Vector store backend for semantic search
            schema_entities: Schema entities for entity extraction
            allowed_documents: Set of allowed document names. If None, all documents
                are visible. If empty set, no documents are visible.
            skip_auto_index: If True, skip automatic indexing of config documents
        """
        self._skip_auto_index = skip_auto_index
        self.config = config
        self.allowed_documents = allowed_documents
        self._loaded_documents: dict[str, LoadedDocument] = {}

        # Use shared embedding model loader (may already be loading in background)
        self._model = EmbeddingModelLoader.get_instance().get_model()
        # Lock for thread-safe access to embedding model (not thread-safe for concurrent encode)
        self._model_lock = threading.Lock()

        # Schema entities for entity extraction (table names, column names)
        self._schema_entities: list[str] = schema_entities or []

        # OpenAPI entities (operations, schemas)
        self._openapi_operations: list[str] = []
        self._openapi_schemas: list[str] = []

        # GraphQL entities (types, fields)
        self._graphql_types: list[str] = []
        self._graphql_fields: list[str] = []

        # Cache directory for persisting document metadata
        if cache_dir:
            self._cache_dir = cache_dir
        else:
            self._cache_dir = Path.cwd() / ".constat"
        self._cache_file = self._cache_dir / self.CACHE_FILENAME

        # Initialize vector store
        if vector_store:
            self._vector_store = vector_store
        else:
            self._vector_store = self._create_vector_store()

        # Active project IDs for automatic search filtering
        # Set this when projects are loaded so searches automatically include project docs
        self._active_project_ids: list[str] = []

        # Note: No cleanup needed - data is scoped by project_id/session_id

        # Index documents that aren't already indexed (incremental)
        # Skip if caller will handle indexing (e.g., warmup with hash-based invalidation)
        print(f"[DOC_INIT] skip_auto_index={self._skip_auto_index}, documents={list(self.config.documents.keys()) if self.config.documents else []}")
        if not self._skip_auto_index and self.config.documents:
            unindexed_docs = self._get_unindexed_documents()
            print(f"[DOC_INIT] unindexed_docs={unindexed_docs}")
            if unindexed_docs:
                logger.info(f"[DOC_INIT] Indexing {len(unindexed_docs)} documents...")
                for name in unindexed_docs:
                    try:
                        self._load_document(name)
                    except Exception as e:
                        logger.warning(f"[DOC_INIT] Failed to load {name}: {e}")
                # Incrementally add new documents (don't clear existing chunks)
                self._index_loaded_documents(self._schema_entities)
                logger.info(f"[DOC_INIT] Indexing complete, count={self._vector_store.count()}")

    def _is_document_allowed(self, doc_name: str) -> bool:
        """Check if a document is allowed based on permissions."""
        if self.allowed_documents is None:
            return True  # No filtering
        return doc_name in self.allowed_documents

    def _get_unindexed_documents(self) -> list[str]:
        """Get list of configured documents that are not yet indexed.

        Returns:
            List of document names that need indexing
        """
        if not self.config.documents:
            return []

        # Get document names that already have chunks in the vector store
        try:
            indexed_docs = set()
            result = self._vector_store._conn.execute("""
                SELECT DISTINCT document_name
                FROM embeddings
                WHERE source = 'document'
            """).fetchall()
            indexed_docs = {row[0] for row in result}
        except Exception:
            # If query fails, assume nothing is indexed
            indexed_docs = set()

        # Return documents in config that are not yet indexed
        return [name for name in self.config.documents if name not in indexed_docs]

    def _index_loaded_documents(self, schema_entities: Optional[list[str]] = None) -> None:
        """Add document chunks to the vector store (incremental).

        This is the core indexing method for the DOCUMENT resource type.
        It adds chunks for all loaded documents without clearing existing data.

        Vector store chunk naming:
        - Documents: document_name = "<doc_name>" (e.g., "business_rules")
        - Schema (from SchemaManager): document_name = "schema:<db>.<table>"
        - APIs (from APISchemaManager): document_name = "api:<api_name>.<endpoint>"

        Args:
            schema_entities: Optional list of known schema entity names for entity extraction
        """
        print(f"[DOC_INDEX] _loaded_documents={list(self._loaded_documents.keys())}")
        chunks = []
        for name, doc in self._loaded_documents.items():
            doc_chunks = self._chunk_document(name, doc.content)
            print(f"[DOC_INDEX] {name}: {len(doc_chunks)} chunks")
            chunks.extend(doc_chunks)

        if not chunks:
            print("[DOC_INDEX] No chunks to index!")
            return

        # Generate embeddings
        texts = [chunk.content for chunk in chunks]
        with self._model_lock:
            embeddings = self._model.encode(texts, convert_to_numpy=True)

        # Add to vector store
        print(f"[DOC_INDEX] Adding {len(chunks)} chunks with source='document'")
        self._vector_store.add_chunks(chunks, embeddings, source="document")

        # Extract and store entities
        self._extract_and_store_entities(chunks, schema_entities)

    def _create_vector_store(self) -> VectorStoreBackend:
        """Create vector store based on config."""
        storage_config = self.config.storage
        vs_config = storage_config.vector_store if storage_config else None

        backend = vs_config.backend if vs_config else "duckdb"
        db_path = vs_config.db_path if vs_config else None

        return create_vector_store(backend=backend, db_path=db_path)

    def set_schema_entities(self, entities: set[str] | list[str]) -> None:
        """Set database schema entities (table names, column names) for pattern matching.

        When schema entities change, entity extraction is re-run on existing
        documents to link schema terms to document references.

        Args:
            entities: Set or list of entity names
        """
        new_entities = list(entities) if isinstance(entities, set) else entities

        # Check if entities actually changed
        if set(new_entities) == set(self._schema_entities or []):
            logger.debug(f"set_schema_entities: no change, skipping (have {len(new_entities)} entities)")
            return

        logger.info(f"set_schema_entities: updating from {len(self._schema_entities or [])} to {len(new_entities)} entities")
        logger.debug(f"set_schema_entities: new entities include: {list(new_entities)[:10]}...")
        self._schema_entities = new_entities

        # Re-extract entities from existing documents if we have indexed chunks
        chunk_count = self._vector_store.count()
        if chunk_count > 0 and hasattr(self._vector_store, 'add_entities'):
            logger.info(f"set_schema_entities: re-extracting entities from {chunk_count} chunks")
            # Clear existing entity links (but keep entities from other sources)
            if hasattr(self._vector_store, 'clear_chunk_entity_links'):
                self._vector_store.clear_chunk_entity_links()

            # Get all chunks and re-extract entities
            chunks = self._vector_store.get_chunks()
            if chunks:
                self._extract_and_store_entities(chunks, self._schema_entities)
                logger.info(f"set_schema_entities: extraction complete")
        else:
            logger.debug(f"set_schema_entities: no chunks to re-extract ({chunk_count} chunks)")

    def extract_entities_for_session(
        self,
        session_id: str,
        project_ids: list[str],
        schema_entities: list[str],
        api_entities: list[str] | None = None,
    ) -> int:
        """Run entity extraction for a session's visible documents.

        Called at session creation to build chunk-entity links using the
        session's entity catalog. Links are stored with session_id.

        Args:
            session_id: Session ID for storing links
            project_ids: List of loaded project IDs
            schema_entities: Schema entity names (tables, columns)
            api_entities: API entity names (endpoints, schemas)

        Returns:
            Number of chunk-entity links created
        """
        if not hasattr(self._vector_store, 'add_entities'):
            return 0

        # Update internal entity lists for extraction
        self._schema_entities = schema_entities or []
        if api_entities:
            self._openapi_operations = api_entities
            self._openapi_schemas = api_entities

        # Get chunks visible to this session (base + projects)
        # Base chunks have project_id='__base__' or NULL
        # Project chunks have project_id in project_ids
        logger.info(f"extract_entities_for_session({session_id}): looking for chunks with project_ids={project_ids}")
        chunks = self._get_session_visible_chunks(project_ids)
        if not chunks:
            logger.warning(f"extract_entities_for_session({session_id}): no visible chunks found!")
            # Debug: check what's in the database
            if hasattr(self._vector_store, '_conn'):
                try:
                    count = self._vector_store._conn.execute("SELECT COUNT(*) FROM embeddings").fetchone()[0]
                    by_proj = self._vector_store._conn.execute("SELECT project_id, COUNT(*) FROM embeddings GROUP BY project_id").fetchall()
                    logger.warning(f"extract_entities_for_session: total embeddings={count}, by_project={by_proj}")
                except Exception as e:
                    logger.warning(f"extract_entities_for_session: failed to check embeddings: {e}")
            return 0

        logger.info(f"extract_entities_for_session({session_id}): extracting from {len(chunks)} chunks")

        # Combine API terms from all sources
        api_terms = list(set(
            (self._openapi_operations or []) +
            (self._openapi_schemas or []) +
            (self._graphql_types or []) +
            (self._graphql_fields or [])
        ))

        # Create extractor with session's entity catalog
        extractor = EntityExtractor(
            session_id=session_id,
            schema_terms=self._schema_entities,
            api_terms=api_terms if api_terms else None,
        )

        all_links: list[ChunkEntity] = []
        for chunk in chunks:
            extractions = extractor.extract(chunk)
            for entity, link in extractions:
                all_links.append(link)

        # Store entities - Entity model now has semantic_type instead of metadata
        entities = extractor.get_all_entities()
        if entities:
            # Add all entities to vector store (session_id is required)
            self._vector_store.add_entities(entities, session_id=session_id)

        # Store links WITH session_id
        if all_links:
            # Deduplicate links by (chunk_id, entity_id)
            unique_links = {}
            for link in all_links:
                key = (link.chunk_id, link.entity_id)
                if key not in unique_links:
                    unique_links[key] = link
                else:
                    existing = unique_links[key]
                    unique_links[key] = ChunkEntity(
                        chunk_id=link.chunk_id,
                        entity_id=link.entity_id,
                        mention_count=existing.mention_count + link.mention_count,
                        confidence=max(existing.confidence, link.confidence),
                        mention_text=existing.mention_text or link.mention_text,
                    )
            self._vector_store.link_chunk_entities(list(unique_links.values()))
            logger.info(f"extract_entities_for_session({session_id}): created {len(unique_links)} links")
            return len(unique_links)

        return 0

    def _get_session_visible_chunks(self, project_ids: list[str]) -> list[DocumentChunk]:
        """Get chunks visible to a session (base + loaded projects).

        Args:
            project_ids: List of loaded project IDs

        Returns:
            List of DocumentChunk objects
        """
        if not hasattr(self._vector_store, '_conn'):
            return self._vector_store.get_chunks()

        # Query chunks where project_id is NULL, '__base__', or in project_ids
        conditions = ["project_id IS NULL", "project_id = '__base__'"]
        params = []

        if project_ids:
            placeholders = ",".join(["?" for _ in project_ids])
            conditions.append(f"project_id IN ({placeholders})")
            params.extend(project_ids)

        where_clause = " OR ".join(conditions)

        result = self._vector_store._conn.execute(
            f"""
            SELECT chunk_id, document_name, content, section, chunk_index
            FROM embeddings
            WHERE {where_clause}
            """,
            params,
        ).fetchall()

        chunks = []
        for row in result:
            chunk_id, doc_name, content, section, chunk_idx = row
            chunk = DocumentChunk(
                document_name=doc_name,
                content=content,
                section=section,
                chunk_index=chunk_idx,
            )
            # Store chunk_id for linking (hacky but needed for entity extraction)
            chunk._chunk_id = chunk_id
            chunks.append(chunk)

        return chunks

    def process_metadata_through_ner(
        self,
        metadata_texts: list[tuple[str, str]],
        source_type: str = "schema",
    ) -> None:
        """Process schema/API metadata through NER for cross-datasource entity linking.

        Creates pseudo-chunks from metadata text (names, descriptions) and runs
        entity extraction to find entities that appear across datasources.

        Args:
            metadata_texts: List of (source_name, text) tuples
            source_type: Source type for the chunks ("schema" or "api")
        """
        if not metadata_texts:
            return

        if not hasattr(self._vector_store, 'add_entities'):
            return

        logger.info(f"Processing {len(metadata_texts)} {source_type} metadata items through NER")

        # Create pseudo-chunks from metadata
        chunks = []
        for source_name, text in metadata_texts:
            if text and text.strip():
                chunk = DocumentChunk(
                    document_name=f"__{source_type}_metadata__",
                    content=text,
                    section=source_name,
                    chunk_index=len(chunks),
                )
                chunks.append(chunk)

        if not chunks:
            return

        # Combine API terms
        api_terms = list(set(
            (self._openapi_operations or []) +
            (self._openapi_schemas or []) +
            (self._graphql_types or []) +
            (self._graphql_fields or [])
        ))

        # Run entity extraction on metadata chunks
        # Use "__metadata__" as session_id for metadata processing
        extractor = EntityExtractor(
            session_id="__metadata__",
            schema_terms=self._schema_entities,
            api_terms=api_terms if api_terms else None,
        )

        all_links: list[ChunkEntity] = []
        for chunk in chunks:
            extractions = extractor.extract(chunk)
            for entity, link in extractions:
                all_links.append(link)

        entities = extractor.get_all_entities()
        logger.debug(f"Metadata NER: {len(entities)} entities, {len(all_links)} links from {len(chunks)} metadata items")

        if entities:
            # Add all entities to vector store
            self._vector_store.add_entities(entities, session_id="__metadata__")

        if all_links:
            # Deduplicate links by (chunk_id, entity_id)
            unique_links = {}
            for link in all_links:
                key = (link.chunk_id, link.entity_id)
                if key not in unique_links:
                    unique_links[key] = link
                else:
                    existing = unique_links[key]
                    unique_links[key] = ChunkEntity(
                        chunk_id=link.chunk_id,
                        entity_id=link.entity_id,
                        mention_count=existing.mention_count + link.mention_count,
                        confidence=max(existing.confidence, link.confidence),
                        mention_text=existing.mention_text or link.mention_text,
                    )
            self._vector_store.link_chunk_entities(list(unique_links.values()))

    def set_openapi_entities(
        self,
        operations: list[str],
        schemas: list[str],
    ) -> None:
        """Set OpenAPI entities for pattern matching in documents.

        Args:
            operations: List of operation/endpoint names
            schemas: List of schema definition names
        """
        self._openapi_operations = operations
        self._openapi_schemas = schemas

    def set_graphql_entities(
        self,
        types: list[str],
        fields: list[str],
    ) -> None:
        """Set GraphQL entities for pattern matching in documents.

        Args:
            types: List of type names
            fields: List of field/operation names
        """
        self._graphql_types = types
        self._graphql_fields = fields

    def _add_document_internal(
        self,
        name: str,
        content: str,
        doc_format: str = "text",
        description: str = "",
        project_id: str | None = None,
        session_id: str | None = None,
        skip_entity_extraction: bool = False,
    ) -> bool:
        """Internal method to add a document with full control over flags.

        Args:
            name: Document name
            content: Document content
            doc_format: Format (text, markdown, etc.)
            description: Optional description
            project_id: Project this document belongs to (for project filtering)
            session_id: Session this document was added in (for session filtering)
            skip_entity_extraction: If True, skip NER (done later at session creation)

        Returns:
            True if indexed successfully
        """
        from datetime import datetime

        # Remove existing document with same name to avoid duplicate key errors
        if name in self._loaded_documents:
            self.remove_document(name)
        # Always try to delete from vector store in case chunks exist
        if hasattr(self._vector_store, 'delete_by_document'):
            self._vector_store.delete_by_document(name)

        # Create a minimal config for the document
        doc_config = DocumentConfig(
            type="inline",
            content=content,
            description=description or "",
            format=doc_format,
        )

        # Extract sections for markdown
        sections = []
        if doc_format in ("markdown", "md"):
            for line in content.split("\n"):
                if line.startswith("#"):
                    sections.append(line.lstrip("#").strip())

        # Store the loaded document
        self._loaded_documents[name] = LoadedDocument(
            name=name,
            config=doc_config,
            content=content,
            format=doc_format,
            sections=sections,
            loaded_at=datetime.now().isoformat(),
        )

        # Chunk and embed the document
        chunks = self._chunk_document(name, content)
        logger.debug(f"Document '{name}': {len(chunks)} chunks generated")
        if not chunks:
            return True

        # Generate embeddings
        texts = [chunk.content for chunk in chunks]
        embeddings = self._model.encode(texts, convert_to_numpy=True)

        # Add to vector store with project_id/session_id for filtering
        if hasattr(self._vector_store, 'add_chunks'):
            self._vector_store.add_chunks(
                chunks,
                embeddings,
                source="document",
                session_id=session_id,
                project_id=project_id,
            )

        # Extract entities with appropriate scope (unless skipped for later session-level extraction)
        if not skip_entity_extraction:
            if session_id:
                # Session-added documents get session_id for session filtering
                self._extract_and_store_entities_session(chunks, session_id)
            elif project_id:
                # Project documents get project_id for project filtering
                self._extract_and_store_entities_project(chunks, project_id)
            else:
                # Base documents (no project_id, no session_id) - permanent
                self._extract_and_store_entities(chunks, self._schema_entities)

        return True

    def add_document_from_file(
        self,
        file_path: str,
        name: str | None = None,
        description: str = "",
        project_id: str | None = None,
        session_id: str | None = None,
        skip_entity_extraction: bool = False,
    ) -> tuple[bool, str]:
        """Add a document from a file path.

        Args:
            file_path: Path to the document file
            name: Optional name (defaults to filename without extension)
            description: Optional description
            project_id: Optional project ID for project-scoped documents
            session_id: Optional session ID for session-scoped documents
            skip_entity_extraction: If True, skip NER (done later at session creation)

        Returns:
            Tuple of (success, message)
        """
        path = Path(file_path)
        if not path.exists():
            return False, f"File not found: {file_path}"

        # Use filename as name if not provided
        if not name:
            name = path.stem.replace(" ", "_").replace("-", "_").lower()

        # Determine format and extract content
        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                content = self._extract_pdf_text(path)
                doc_format = "text"
            elif suffix == ".docx":
                content = self._extract_docx_text(path)
                doc_format = "text"
            elif suffix == ".xlsx":
                content = self._extract_xlsx_text(path)
                doc_format = "text"
            elif suffix == ".pptx":
                content = self._extract_pptx_text(path)
                doc_format = "text"
            elif suffix in (".md", ".markdown"):
                content = path.read_text()
                doc_format = "markdown"
            elif suffix in (".json",):
                content = path.read_text()
                doc_format = "json"
            elif suffix in (".yaml", ".yml"):
                content = path.read_text()
                doc_format = "yaml"
            else:
                content = path.read_text()
                doc_format = "text"
        except Exception as e:
            return False, f"Failed to read file: {e}"

        if not content.strip():
            return False, "File is empty"

        # Add document with specified scope
        success = self._add_document_internal(
            name=name,
            content=content,
            doc_format=doc_format,
            description=description or f"Document from {path.name}",
            project_id=project_id,
            session_id=session_id,
            skip_entity_extraction=skip_entity_extraction,
        )

        if success:
            return True, f"Added document '{name}' ({len(content):,} chars)"
        return False, "Failed to index document"

    def _extract_and_store_entities_session(
        self,
        chunks: list[DocumentChunk],
        session_id: str,
    ) -> None:
        """Extract entities from chunks and store them with session_id.

        Uses spaCy NER for named entity extraction plus pattern matching
        for database, OpenAPI, and GraphQL schemas.

        Args:
            chunks: Document chunks to extract entities from
            session_id: Session ID for session-scoped entities
        """
        if not hasattr(self._vector_store, 'add_entities'):
            return

        logger.debug(f"Extracting entities from {len(chunks)} chunks for session {session_id}")

        # Combine API terms
        api_terms = list(set(
            (self._openapi_operations or []) +
            (self._openapi_schemas or []) +
            (self._graphql_types or []) +
            (self._graphql_fields or [])
        ))

        # Create extractor with all known schema entities
        extractor = EntityExtractor(
            session_id=session_id,
            schema_terms=self._schema_entities,
            api_terms=api_terms if api_terms else None,
        )

        all_links: list[ChunkEntity] = []

        # Extract entities from all chunks using spaCy NER
        for chunk in chunks:
            extractions = extractor.extract(chunk)
            logger.debug(f"[ENTITY] Chunk '{chunk.section}' -> {len(extractions)} entities")

            for entity, link in extractions:
                all_links.append(link)

        entities = extractor.get_all_entities()
        logger.debug(f"Extracted {len(entities)} entities, {len(all_links)} links")
        if entities:
            self._vector_store.add_entities(entities, session_id=session_id)

        if all_links:
            # Deduplicate links by (chunk_id, entity_id)
            unique_links = {}
            for link in all_links:
                key = (link.chunk_id, link.entity_id)
                if key not in unique_links:
                    unique_links[key] = link
                else:
                    existing = unique_links[key]
                    unique_links[key] = ChunkEntity(
                        chunk_id=link.chunk_id,
                        entity_id=link.entity_id,
                        mention_count=existing.mention_count + link.mention_count,
                        confidence=max(existing.confidence, link.confidence),
                        mention_text=existing.mention_text or link.mention_text,
                    )
            self._vector_store.link_chunk_entities(list(unique_links.values()))

    def _extract_and_store_entities_project(
        self,
        chunks: list[DocumentChunk],
        project_id: str,
    ) -> None:
        """Extract entities from chunks and store them with project_id.

        Uses spaCy NER for named entity extraction plus pattern matching
        for database, OpenAPI, and GraphQL schemas.

        Args:
            chunks: Document chunks to extract entities from
            project_id: Project ID for project-scoped entities
        """
        if not hasattr(self._vector_store, 'add_entities'):
            return

        logger.debug(f"Extracting entities from {len(chunks)} chunks for project {project_id}")

        # Combine API terms
        api_terms = list(set(
            (self._openapi_operations or []) +
            (self._openapi_schemas or []) +
            (self._graphql_types or []) +
            (self._graphql_fields or [])
        ))

        # Create extractor with all known schema entities
        # Use project_id as session_id for project-scoped extraction
        extractor = EntityExtractor(
            session_id=project_id,
            project_id=project_id,
            schema_terms=self._schema_entities,
            api_terms=api_terms if api_terms else None,
        )

        all_links: list[ChunkEntity] = []

        # Extract entities from all chunks using spaCy NER
        for chunk in chunks:
            extractions = extractor.extract(chunk)
            logger.debug(f"[ENTITY] Chunk '{chunk.section}' -> {len(extractions)} entities")

            for entity, link in extractions:
                all_links.append(link)

        entities = extractor.get_all_entities()
        logger.debug(f"Extracted {len(entities)} entities, {len(all_links)} links")
        if entities:
            self._vector_store.add_entities(entities, session_id=project_id)

        if all_links:
            # Deduplicate links by (chunk_id, entity_id)
            unique_links = {}
            for link in all_links:
                key = (link.chunk_id, link.entity_id)
                if key not in unique_links:
                    unique_links[key] = link
                else:
                    existing = unique_links[key]
                    unique_links[key] = ChunkEntity(
                        chunk_id=link.chunk_id,
                        entity_id=link.entity_id,
                        mention_count=existing.mention_count + link.mention_count,
                        confidence=max(existing.confidence, link.confidence),
                        mention_text=existing.mention_text or link.mention_text,
                    )
            self._vector_store.link_chunk_entities(list(unique_links.values()))

    def remove_document(self, name: str) -> bool:
        """Remove a document and its vectors from the index.

        Args:
            name: Document name to remove

        Returns:
            True if removed, False if not found
        """
        if name not in self._loaded_documents:
            return False

        del self._loaded_documents[name]

        # Remove from vector store
        if hasattr(self._vector_store, 'delete_by_document'):
            self._vector_store.delete_by_document(name)

        return True

    def get_session_documents(self) -> dict[str, dict]:
        """Get documents added during this session (not in config).

        Returns:
            Dict of {name: {format, char_count, loaded_at}} for session docs
        """
        # Get names of documents from config (documents is a dict keyed by name)
        config_doc_names = set(self.config.documents.keys()) if self.config.documents else set()

        # Find documents that were added but not in config
        session_docs = {}
        for name, doc in self._loaded_documents.items():
            if name not in config_doc_names:
                session_docs[name] = {
                    "format": doc.format,
                    "char_count": len(doc.content),
                    "loaded_at": doc.loaded_at,
                }

        return session_docs

    def refresh(self, force_full: bool = False) -> dict:
        """Refresh documents, using incremental update by default.

        Args:
            force_full: If True, force full rebuild of DOCUMENT chunks only
                       (preserves schema and API chunks from other managers)

        Returns:
            Dict with refresh statistics: {added, updated, removed, unchanged}
        """
        if force_full:
            self._loaded_documents.clear()

            # Clear only DOCUMENT chunks, preserve schema:* and api:* from other managers
            try:
                self._vector_store._conn.execute("""
                    DELETE FROM embeddings
                    WHERE document_name NOT LIKE 'schema:%'
                      AND document_name NOT LIKE 'api:%'
                """)
            except Exception:
                # Fallback for non-DuckDB backends - this will clear everything
                logger.warning("refresh: falling back to full clear (non-DuckDB backend)")
                self._vector_store.clear()

            # Clear document-sourced entities only
            if hasattr(self._vector_store, 'clear_entities'):
                self._vector_store.clear_entities(source='document')

            # Reload all documents and rebuild index
            stats = {"added": 0, "updated": 0, "removed": 0, "unchanged": 0, "mode": "full_rebuild"}
            if self.config.documents:
                for name in self.config.documents:
                    try:
                        self._load_document(name)
                        stats["added"] += 1
                    except Exception as e:
                        logger.warning(f"Failed to load {name}: {e}")
                self._index_loaded_documents(self._schema_entities)
            return stats

        return self._refresh_incremental()

    def _refresh_incremental(self) -> dict:
        """Incrementally update document index based on file changes.

        Returns:
            Dict with refresh statistics
        """
        stats = {"added": 0, "updated": 0, "removed": 0, "unchanged": 0}

        # Get current document list (expanded for globs/directories)
        current_docs = self._get_all_document_paths()

        # Track which documents to reload
        docs_to_reload: list[str] = []
        docs_to_remove: list[str] = []

        # Check each configured document
        for doc_name, doc_path in current_docs.items():
            if doc_path and doc_path.exists():
                current_mtime = doc_path.stat().st_mtime

                if doc_name in self._loaded_documents:
                    # Check if file has changed
                    loaded_doc = self._loaded_documents[doc_name]
                    if loaded_doc.file_mtime != current_mtime:
                        docs_to_reload.append(doc_name)
                        stats["updated"] += 1
                    else:
                        stats["unchanged"] += 1
                else:
                    # New document
                    docs_to_reload.append(doc_name)
                    stats["added"] += 1
            elif doc_name in self._loaded_documents:
                # Document was removed
                docs_to_remove.append(doc_name)
                stats["removed"] += 1

        # Check for documents that no longer exist in config
        for doc_name in list(self._loaded_documents.keys()):
            if doc_name not in current_docs:
                docs_to_remove.append(doc_name)
                stats["removed"] += 1

        # Remove deleted documents
        for doc_name in docs_to_remove:
            if doc_name in self._loaded_documents:
                del self._loaded_documents[doc_name]

        # Reload changed documents
        for doc_name in docs_to_reload:
            try:
                self._load_document_with_mtime(doc_name)
            except Exception:
                pass  # Skip documents that fail to load

        # Rebuild index if anything changed
        if docs_to_reload or docs_to_remove:
            self._build_index(self._schema_entities)

        return stats

    def _get_all_document_paths(self) -> dict[str, Optional[Path]]:
        """Get all document names mapped to their file paths.

        Returns:
            Dict mapping document names to Path objects (None for non-file docs)
        """
        result = {}

        for doc_name, doc_config in self.config.documents.items():
            if doc_config.type == "file" and doc_config.path:
                expanded = _expand_file_paths(doc_config.path)
                if len(expanded) > 1:
                    # Multiple files - each gets its own entry
                    for filename, filepath in expanded:
                        full_name = f"{doc_name}:{filename}"
                        result[full_name] = filepath
                elif len(expanded) == 1:
                    _, filepath = expanded[0]
                    result[doc_name] = filepath
                else:
                    result[doc_name] = None
            else:
                # Non-file documents (inline, http, etc.)
                result[doc_name] = None

        return result

    def _load_document_with_mtime(self, name: str) -> None:
        """Load a document and record its modification time."""
        # Get the file path
        doc_paths = self._get_all_document_paths()
        filepath = doc_paths.get(name)

        # Load the document
        if ":" in name:
            # Expanded glob/directory entry
            parent_name, filename = name.split(":", 1)
            if parent_name in self.config.documents:
                doc_config = self.config.documents[parent_name]
                if filepath and filepath.exists():
                    self._load_file_directly(name, filepath, doc_config)
        else:
            self._load_document(name)

        # Record mtime and content hash
        if name in self._loaded_documents:
            doc = self._loaded_documents[name]
            if filepath and filepath.exists():
                doc.file_mtime = filepath.stat().st_mtime
            doc.content_hash = hashlib.sha256(doc.content.encode()).hexdigest()[:16]

    def list_documents(self) -> list[dict]:
        """
        List all configured reference documents with descriptions.

        For file-type documents with glob patterns or directory paths,
        expands to show each individual file.

        Returns:
            List of document info dicts with name, type, description, tags, path
        """
        results = []

        for doc_name, doc_config in self.config.documents.items():
            # Skip documents not allowed by permissions
            if not self._is_document_allowed(doc_name):
                continue
            # For file types, check if it's a glob/directory that needs expansion
            if doc_config.type == "file" and doc_config.path:
                expanded = _expand_file_paths(doc_config.path)

                if len(expanded) > 1:
                    # Multiple files - list each one with collection context
                    collection_desc = doc_config.description or f"Collection: {doc_name}"
                    for filename, filepath in expanded:
                        full_name = f"{doc_name}:{filename}"
                        results.append({
                            "name": full_name,
                            "type": doc_config.type,
                            "description": f"File in collection. Collection description: {collection_desc}",
                            "format": doc_config.format or self._detect_format(filepath.suffix),
                            "tags": doc_config.tags,
                            "path": str(filepath),
                            "collection": doc_name,
                            "collection_description": collection_desc,
                            "loaded": full_name in self._loaded_documents,
                        })
                elif len(expanded) == 1:
                    # Single file
                    filename, filepath = expanded[0]
                    results.append({
                        "name": doc_name,
                        "type": doc_config.type,
                        "description": doc_config.description or f"File: {filename}",
                        "format": doc_config.format or self._detect_format(filepath.suffix),
                        "tags": doc_config.tags,
                        "path": str(filepath),
                        "loaded": doc_name in self._loaded_documents,
                    })
                else:
                    # No files matched - still list it (will error on load)
                    results.append({
                        "name": doc_name,
                        "type": doc_config.type,
                        "description": doc_config.description or f"Document: {doc_name}",
                        "format": doc_config.format,
                        "tags": doc_config.tags,
                        "path": doc_config.path,
                        "loaded": doc_name in self._loaded_documents,
                    })
            else:
                # Non-file types (inline, http, etc.)
                entry = {
                    "name": doc_name,
                    "type": doc_config.type,
                    "description": doc_config.description or f"Document: {doc_name}",
                    "format": doc_config.format,
                    "tags": doc_config.tags,
                    "loaded": doc_name in self._loaded_documents,
                }
                # Include source location where applicable
                if doc_config.type == "http" and doc_config.url:
                    entry["url"] = doc_config.url
                elif doc_config.type == "inline":
                    entry["source"] = "inline content"
                results.append(entry)

        return results

    def get_document(self, name: str) -> dict:
        """
        Get the full content of a document.

        Args:
            name: Document name. For expanded glob/directory entries,
                  use format "config_name:filename" (e.g., "data_files:sales.csv")

        Returns:
            Dict with document content and metadata
        """
        # Check if already loaded
        if name in self._loaded_documents:
            doc = self._loaded_documents[name]
            result = {
                "name": doc.name,
                "content": doc.content,
                "format": doc.format,
                "sections": doc.sections,
                "loaded_at": doc.loaded_at,
            }
            if hasattr(doc.config, 'path') and doc.config.path:
                result["path"] = doc.config.path
            return result

        # Handle expanded names from glob/directory (format: "parent:filename")
        # But skip API-style document names (api:*, Database:*, Table:*, etc.)
        api_prefixes = ("api:", "Database:", "Table:", "__")
        if ":" in name and not any(name.startswith(p) for p in api_prefixes):
            parent_name, filename = name.split(":", 1)
            # Check permissions on parent name
            if not self._is_document_allowed(parent_name):
                return {"error": f"Access denied to document: {parent_name}"}

            # Try base config documents first
            doc_config = self.config.documents.get(parent_name)

            # Try project documents if not in base
            if not doc_config:
                for project in self.config.projects.values():
                    if project.documents and parent_name in project.documents:
                        doc_config = project.documents[parent_name]
                        break

            if not doc_config:
                return {"error": f"Document config not found: {parent_name}"}

            if doc_config.type != "file" or not doc_config.path:
                return {"error": f"Document {parent_name} is not a file type"}

            # Find the specific file in the expanded paths
            expanded = _expand_file_paths(doc_config.path)
            matching = [(fn, fp) for fn, fp in expanded if fn == filename]

            if not matching:
                return {"error": f"File '{filename}' not found in {parent_name}"}

            _, filepath = matching[0]

            # Load this specific file
            try:
                self._load_file_directly(name, filepath, doc_config)
            except Exception as e:
                return {"error": f"Failed to load file: {str(e)}"}
        else:
            # System-generated documents (API, Database, Table) skip permission checks
            # and go straight to vector store reconstruction
            system_prefixes = ("api:", "Database:", "Table:", "__")
            if any(name.startswith(p) for p in system_prefixes):
                return self._reconstruct_from_chunks(name)

            # Standard document name - check permissions
            if not self._is_document_allowed(name):
                return {"error": f"Access denied to document: {name}"}

            # Try base config documents
            doc_config = self.config.documents.get(name)

            # Try project documents if not in base
            if not doc_config:
                for proj_name, project in self.config.projects.items():
                    if project.documents and name in project.documents:
                        doc_config = project.documents[name]
                        break

            # Fallback: reconstruct from vector store chunks
            if not doc_config:
                return self._reconstruct_from_chunks(name)

            # Load from config
            try:
                result = self._load_document(name)
                # Binary files (PDF, Office) return dict directly instead of storing
                if isinstance(result, dict):
                    return result
            except Exception as e:
                return {"error": f"Failed to load document: {str(e)}"}

        doc = self._loaded_documents.get(name)
        if not doc:
            return {"error": f"Document not loaded: {name}"}

        result = {
            "name": doc.name,
            "content": doc.content,
            "format": doc.format,
            "sections": doc.sections,
            "loaded_at": doc.loaded_at,
        }

        # Include path for file-based documents
        if hasattr(doc.config, 'path') and doc.config.path:
            result["path"] = doc.config.path

        return result

    def _reconstruct_from_chunks(self, name: str) -> dict:
        """Reconstruct document content from vector store chunks.

        This is a fallback for documents that were indexed but aren't in config
        (e.g., session-uploaded documents, API metadata).

        Args:
            name: Document name to reconstruct

        Returns:
            Dict with document content or error
        """
        try:
            # Query chunks for this document
            chunks = self._vector_store._conn.execute(
                """
                SELECT content, section, chunk_index
                FROM embeddings
                WHERE document_name = ?
                ORDER BY chunk_index
                """,
                [name],
            ).fetchall()

            if not chunks:
                logger.warning(f"Document '{name}' not found in embeddings table")
                return {"error": f"Document not found: {name}"}

            # Reconstruct content from chunks
            sections = []
            content_parts = []
            for content, section, _ in chunks:
                content_parts.append(content)
                if section and section not in sections:
                    sections.append(section)

            return {
                "name": name,
                "content": "\n\n".join(content_parts),
                "format": "text",  # Assume plain text for reconstructed docs
                "sections": sections,
                "loaded_at": None,
                "reconstructed": True,
            }
        except Exception as e:
            logger.warning(f"Failed to reconstruct document {name} from chunks: {e}")
            return {"error": f"Document not found: {name}"}

    def search_documents(
        self,
        query: str,
        limit: int = 5,
        session_id: str | None = None,
    ) -> list[dict]:
        """
        Search across all documents for relevant content using semantic search.

        Uses this instance's active_project_ids (set when projects are loaded into the session).

        Args:
            query: Natural language query
            limit: Maximum results to return
            session_id: Session ID to include (for session-specific documents)

        Returns:
            List of relevant document excerpts with relevance scores
        """
        if self._model is None:
            logger.debug(f"[SEARCH] No embedding model")
            return []

        # Use lock for thread-safe access - both embedding model AND DuckDB connection
        # are not thread-safe for concurrent operations
        with self._model_lock:
            # Embed the query
            query_embedding = self._model.encode([query], convert_to_numpy=True)

            # Search using vector store with filtering (also protected by lock for consistency)
            search_results = self._vector_store.search(
                query_embedding,
                limit=limit,
                project_ids=self._active_project_ids,
                session_id=session_id,
            )

        logger.debug(f"[SEARCH] Found {len(search_results)} results for '{query}'")

        results = []
        for chunk_id, similarity, chunk in search_results:
            logger.debug(f"[SEARCH] Result: doc={chunk.document_name}, section={chunk.section}, score={similarity:.3f}")
            results.append({
                "document": chunk.document_name,
                # Return full chunk content (already sized at CHUNK_SIZE=800)
                # Truncating at 500 chars can cut content mid-sentence
                "excerpt": chunk.content,
                "relevance": round(similarity, 3),
                "section": chunk.section,
            })

        return results

    def search_documents_enriched(
        self,
        query: str,
        limit: int = 5,
        session_id: str | None = None,
    ) -> list[dict]:
        """Search documents with entity enrichment.

        Like search_documents, but includes entities mentioned in each chunk.
        Useful for understanding what concepts are discussed in relevant chunks.

        Uses this instance's active_project_ids (set when projects are loaded into the session).

        Args:
            query: Natural language query
            limit: Maximum results to return
            session_id: Session ID to include (for session-specific documents)

        Returns:
            List of dicts with document, excerpt, relevance, section, and entities
        """
        if self._model is None:
            return []

        # Use enriched search if available
        if hasattr(self._vector_store, 'search_enriched'):
            # Use lock for thread-safe access - both embedding model AND DuckDB connection
            with self._model_lock:
                query_embedding = self._model.encode([query], convert_to_numpy=True)
                enriched_results = self._vector_store.search_enriched(
                    query_embedding,
                    limit=limit,
                    project_ids=self._active_project_ids,
                    session_id=session_id,
                )

            results = []
            for enriched in enriched_results:
                results.append({
                    "document": enriched.chunk.document_name,
                    "excerpt": enriched.chunk.content,
                    "relevance": round(enriched.score, 3),
                    "section": enriched.chunk.section,
                    "entities": [
                        {"name": e.name, "type": e.type}
                        for e in enriched.entities
                    ],
                })
            return results

        # Fall back to regular search
        return self.search_documents(query, limit, session_id)

    def explore_entity(
        self,
        entity_name: str,
        limit: int = 5,
        session_id: str | None = None,
    ) -> list[dict]:
        """Find chunks mentioning the given entity.

        Use when the LLM notices a relevant entity and wants more context.
        Returns chunks ordered by relevance (mention count, recency).

        Uses this instance's active_project_ids (set when projects are loaded into the session).

        This is designed to be exposed as an LLM tool:
        {
            "name": "explore_entity",
            "description": "Find additional context about an entity (table, concept, term)
                           mentioned in the retrieved chunks. Use when you need more
                           information about something specific.",
            "parameters": {
                "entity_name": {
                    "type": "string",
                    "description": "Name of the entity to explore"
                }
            }
        }

        Args:
            entity_name: Name of the entity to explore
            limit: Maximum number of chunks to return
            session_id: Session ID to include (for session-specific documents)

        Returns:
            List of dicts with document, excerpt, section, mention_count, confidence
            Empty list if entity not found
        """
        if not hasattr(self._vector_store, 'find_entity_by_name'):
            return []

        entity = self._vector_store.find_entity_by_name(entity_name, project_ids=self._active_project_ids, session_id=session_id)
        if not entity:
            return []

        chunks = self._vector_store.get_chunks_for_entity(
            entity.id,
            limit=limit,
            project_ids=self._active_project_ids,
            session_id=session_id,
        )

        results = []
        for chunk_id, chunk, mention_count, confidence in chunks:
            results.append({
                "document": chunk.document_name,
                "excerpt": chunk.content,
                "section": chunk.section,
                "entity": {
                    "name": entity.name,
                    "type": entity.type,
                },
                "mention_count": mention_count,
                "confidence": round(confidence, 3),
            })

        return results

    def get_document_section(self, name: str, section: str) -> dict:
        """
        Get a specific section from a document.

        Args:
            name: Document name
            section: Section header/title

        Returns:
            Dict with section content
        """
        doc = self.get_document(name)
        if "error" in doc:
            return doc

        content = doc["content"]

        # Try to find the section
        section_lower = section.lower()
        lines = content.split("\n")

        in_section = False
        section_lines = []
        section_level = 0

        for line in lines:
            # Check for markdown headers
            if line.startswith("#"):
                header_level = len(line) - len(line.lstrip("#"))
                header_text = line.lstrip("#").strip().lower()

                if section_lower in header_text:
                    in_section = True
                    section_level = header_level
                    section_lines.append(line)
                elif in_section and header_level <= section_level:
                    # End of section
                    break
                elif in_section:
                    section_lines.append(line)
            elif in_section:
                section_lines.append(line)

        if section_lines:
            return {
                "document": name,
                "section": section,
                "content": "\n".join(section_lines),
            }
        else:
            return {"error": f"Section '{section}' not found in document '{name}'"}

    def _load_document(self, name: str) -> dict | None:
        """Load a document from its configured source.

        Returns:
            None for text documents (stored in _loaded_documents)
            dict for binary files (PDF, Office) to be returned directly
        """
        from datetime import datetime

        # Check base config first
        doc_config = self.config.documents.get(name)

        # Check project documents if not in base
        if not doc_config:
            for project in self.config.projects.values():
                if project.documents and name in project.documents:
                    doc_config = project.documents[name]
                    break

        if not doc_config:
            configured = list(self.config.documents.keys())
            raise ValueError(
                f"Document not configured: {name}. "
                f"Configured documents: {configured}. "
                f"doc_read() is ONLY for configured reference documents. "
                f"If the data is in a datastore table, use store.query() instead."
            )
        content = ""
        doc_format = doc_config.format

        if doc_config.type == "inline":
            content = doc_config.content or ""
            if doc_format == "auto":
                doc_format = "text"

        elif doc_config.type == "file":
            if doc_config.path:
                path = Path(doc_config.path)
                # Resolve relative paths from config directory
                if not path.is_absolute() and self.config.config_dir:
                    path = (Path(self.config.config_dir) / doc_config.path).resolve()
                if path.exists():
                    suffix = path.suffix.lower()

                    # Binary document formats - extract text content
                    if suffix == ".pdf":
                        content = self._extract_pdf_text(path)
                        doc_format = "text"
                    elif suffix == ".docx":
                        content = self._extract_docx_text(path)
                        doc_format = "text"
                    elif suffix == ".xlsx":
                        content = self._extract_xlsx_text(path)
                        doc_format = "text"
                    elif suffix == ".pptx":
                        content = self._extract_pptx_text(path)
                        doc_format = "text"
                    else:
                        # Check for structured data files - use schema metadata
                        schema = _infer_structured_schema(path, doc_config.description)
                        if schema:
                            content = schema.to_metadata_doc()
                            doc_format = schema.file_format
                        # Text-based files - return content for rendering
                        else:
                            content = path.read_text()
                            if doc_format == "auto":
                                doc_format = self._detect_format(suffix)
                else:
                    raise FileNotFoundError(f"Document file not found: {doc_config.path}")

        elif doc_config.type == "http":
            if doc_config.url:
                import requests
                headers = doc_config.headers or {}
                response = requests.get(doc_config.url, headers=headers, timeout=30)
                response.raise_for_status()

                # Check content type and URL extension for binary formats
                content_type = response.headers.get("content-type", "")
                url_lower = doc_config.url.lower() if doc_config.url else ""

                if "pdf" in content_type or url_lower.endswith(".pdf"):
                    content = self._extract_pdf_text_from_bytes(response.content)
                    doc_format = "text"
                elif "wordprocessingml" in content_type or url_lower.endswith(".docx"):
                    content = self._extract_docx_text_from_bytes(response.content)
                    doc_format = "text"
                elif "spreadsheetml" in content_type or url_lower.endswith(".xlsx"):
                    content = self._extract_xlsx_text_from_bytes(response.content)
                    doc_format = "text"
                elif "presentationml" in content_type or url_lower.endswith(".pptx"):
                    content = self._extract_pptx_text_from_bytes(response.content)
                    doc_format = "text"
                else:
                    content = response.text
                    if doc_format == "auto":
                        doc_format = self._detect_format_from_content_type(content_type)

        elif doc_config.type == "pdf":
            # Direct PDF type - load from path or url
            if doc_config.path:
                path = Path(doc_config.path)
                if path.exists():
                    content = self._extract_pdf_text(path)
                    doc_format = "text"
                else:
                    raise FileNotFoundError(f"PDF file not found: {doc_config.path}")
            elif doc_config.url:
                import requests
                headers = doc_config.headers or {}
                response = requests.get(doc_config.url, headers=headers, timeout=30)
                response.raise_for_status()
                content = self._extract_pdf_text_from_bytes(response.content)
                doc_format = "text"

        elif doc_config.type == "docx":
            # Word document - load from path or url
            if doc_config.path:
                path = Path(doc_config.path)
                if path.exists():
                    content = self._extract_docx_text(path)
                    doc_format = "text"
                else:
                    raise FileNotFoundError(f"Word document not found: {doc_config.path}")
            elif doc_config.url:
                import requests
                headers = doc_config.headers or {}
                response = requests.get(doc_config.url, headers=headers, timeout=30)
                response.raise_for_status()
                content = self._extract_docx_text_from_bytes(response.content)
                doc_format = "text"

        elif doc_config.type == "xlsx":
            # Excel spreadsheet - load from path or url
            if doc_config.path:
                path = Path(doc_config.path)
                if path.exists():
                    content = self._extract_xlsx_text(path)
                    doc_format = "text"
                else:
                    raise FileNotFoundError(f"Excel file not found: {doc_config.path}")
            elif doc_config.url:
                import requests
                headers = doc_config.headers or {}
                response = requests.get(doc_config.url, headers=headers, timeout=30)
                response.raise_for_status()
                content = self._extract_xlsx_text_from_bytes(response.content)
                doc_format = "text"

        elif doc_config.type == "pptx":
            # PowerPoint presentation - load from path or url
            if doc_config.path:
                path = Path(doc_config.path)
                if path.exists():
                    content = self._extract_pptx_text(path)
                    doc_format = "text"
                else:
                    raise FileNotFoundError(f"PowerPoint file not found: {doc_config.path}")
            elif doc_config.url:
                import requests
                headers = doc_config.headers or {}
                response = requests.get(doc_config.url, headers=headers, timeout=30)
                response.raise_for_status()
                content = self._extract_pptx_text_from_bytes(response.content)
                doc_format = "text"

        # TODO: Implement confluence, notion loaders
        else:
            raise NotImplementedError(f"Document type not yet implemented: {doc_config.type}")

        # Extract sections for markdown
        sections = []
        if doc_format in ("markdown", "md"):
            for line in content.split("\n"):
                if line.startswith("#"):
                    sections.append(line.lstrip("#").strip())

        self._loaded_documents[name] = LoadedDocument(
            name=name,
            config=doc_config,
            content=content,
            format=doc_format,
            sections=sections,
            loaded_at=datetime.now().isoformat(),
        )

    def _load_file_directly(self, name: str, filepath: Path, doc_config: DocumentConfig) -> None:
        """Load a file directly from a path (for expanded glob/directory entries)."""
        from datetime import datetime

        suffix = filepath.suffix.lower()
        doc_format = doc_config.format

        # Check if it's a structured data file - use schema metadata instead of raw content
        schema = _infer_structured_schema(filepath, doc_config.description)
        if schema:
            # For structured files, index the metadata, not the raw data
            content = schema.to_metadata_doc()
            doc_format = schema.file_format
            sections = ["Schema", "Columns"]

            # Store schema for later reference
            file_config = DocumentConfig(
                type="file",
                path=str(filepath),
                description=doc_config.description,
                format=doc_format,
                tags=doc_config.tags,
            )

            self._loaded_documents[name] = LoadedDocument(
                name=name,
                config=file_config,
                content=content,
                format=doc_format,
                sections=sections,
                loaded_at=datetime.now().isoformat(),
            )
            return

        # Handle other file types
        if suffix == ".pdf":
            content = self._extract_pdf_text(filepath)
            doc_format = "text"
        elif suffix == ".docx":
            content = self._extract_docx_text(filepath)
            doc_format = "text"
        elif suffix == ".xlsx":
            content = self._extract_xlsx_text(filepath)
            doc_format = "text"
        elif suffix == ".pptx":
            content = self._extract_pptx_text(filepath)
            doc_format = "text"
        else:
            content = filepath.read_text()
            if doc_format == "auto" or not doc_format:
                doc_format = self._detect_format(suffix)

        # Extract sections for markdown
        sections = []
        if doc_format in ("markdown", "md"):
            for line in content.split("\n"):
                if line.startswith("#"):
                    sections.append(line.lstrip("#").strip())

        # Create a modified config with the actual path
        file_config = DocumentConfig(
            type="file",
            path=str(filepath),
            description=doc_config.description,
            format=doc_format,
            tags=doc_config.tags,
        )

        self._loaded_documents[name] = LoadedDocument(
            name=name,
            config=file_config,
            content=content,
            format=doc_format,
            sections=sections,
            loaded_at=datetime.now().isoformat(),
        )

    def _is_structured_data_format(self, doc_format: str) -> bool:
        """Check if format is structured data that shouldn't be semantically indexed."""
        # These formats are data to be queried, not text to be searched
        structured_formats = {"csv", "json", "jsonl", "parquet", "xml", "yaml", "yml"}
        return doc_format.lower() in structured_formats

    def _extract_pdf_text(self, path) -> str:
        """Extract text content from a PDF file.

        Args:
            path: Path to the PDF file

        Returns:
            Extracted text content with page markers
        """
        from pypdf import PdfReader

        reader = PdfReader(path)
        pages = []

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i}]\n{text.strip()}")

        return "\n\n".join(pages)

    def _extract_pdf_text_from_bytes(self, pdf_bytes: bytes) -> str:
        """Extract text content from PDF bytes.

        Args:
            pdf_bytes: Raw PDF file content

        Returns:
            Extracted text content with page markers
        """
        from io import BytesIO
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(pdf_bytes))
        pages = []

        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i}]\n{text.strip()}")

        return "\n\n".join(pages)

    def _extract_docx_text(self, path) -> str:
        """Extract text content from a Word document.

        Args:
            path: Path to the .docx file

        Returns:
            Extracted text content with paragraph separation
        """
        from docx import Document

        doc = Document(path)
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Check if it's a heading
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        level_num = int(level)
                        paragraphs.append(f"{'#' * level_num} {text}")
                    except ValueError:
                        paragraphs.append(text)
                else:
                    paragraphs.append(text)

        # Also extract text from tables
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            if table_rows:
                paragraphs.append("\n".join(table_rows))

        return "\n\n".join(paragraphs)

    def _extract_docx_text_from_bytes(self, docx_bytes: bytes) -> str:
        """Extract text content from Word document bytes.

        Args:
            docx_bytes: Raw .docx file content

        Returns:
            Extracted text content
        """
        from io import BytesIO
        from docx import Document

        doc = Document(BytesIO(docx_bytes))
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                if para.style and para.style.name.startswith("Heading"):
                    level = para.style.name.replace("Heading ", "")
                    try:
                        level_num = int(level)
                        paragraphs.append(f"{'#' * level_num} {text}")
                    except ValueError:
                        paragraphs.append(text)
                else:
                    paragraphs.append(text)

        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            if table_rows:
                paragraphs.append("\n".join(table_rows))

        return "\n\n".join(paragraphs)

    def _extract_xlsx_text(self, path) -> str:
        """Extract text content from an Excel spreadsheet.

        Args:
            path: Path to the .xlsx file

        Returns:
            Extracted text content with sheet and cell markers
        """
        from openpyxl import load_workbook

        wb = load_workbook(path, data_only=True)
        sheets = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows():
                cells = []
                for cell in row:
                    if cell.value is not None:
                        cells.append(str(cell.value))
                    else:
                        cells.append("")
                # Only include rows that have some content
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))

            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        return "\n\n".join(sheets)

    def _extract_xlsx_text_from_bytes(self, xlsx_bytes: bytes) -> str:
        """Extract text content from Excel spreadsheet bytes.

        Args:
            xlsx_bytes: Raw .xlsx file content

        Returns:
            Extracted text content
        """
        from io import BytesIO
        from openpyxl import load_workbook

        wb = load_workbook(BytesIO(xlsx_bytes), data_only=True)
        sheets = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows():
                cells = []
                for cell in row:
                    if cell.value is not None:
                        cells.append(str(cell.value))
                    else:
                        cells.append("")
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))

            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        return "\n\n".join(sheets)

    def _extract_pptx_text(self, path) -> str:
        """Extract text content from a PowerPoint presentation.

        Args:
            path: Path to the .pptx file

        Returns:
            Extracted text content with slide markers
        """
        from pptx import Presentation

        prs = Presentation(path)
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables in slides
                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_rows.append(" | ".join(cells))
                    if table_rows:
                        slide_text.append("\n".join(table_rows))

            if slide_text:
                slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))

        return "\n\n".join(slides)

    def _extract_pptx_text_from_bytes(self, pptx_bytes: bytes) -> str:
        """Extract text content from PowerPoint presentation bytes.

        Args:
            pptx_bytes: Raw .pptx file content

        Returns:
            Extracted text content
        """
        from io import BytesIO
        from pptx import Presentation

        prs = Presentation(BytesIO(pptx_bytes))
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                if shape.has_table:
                    table_rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        table_rows.append(" | ".join(cells))
                    if table_rows:
                        slide_text.append("\n".join(table_rows))

            if slide_text:
                slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))

        return "\n\n".join(slides)

    def _detect_format(self, suffix: str) -> str:
        """Detect document format from file extension."""
        format_map = {
            ".md": "markdown",
            ".markdown": "markdown",
            ".txt": "text",
            ".html": "html",
            ".htm": "html",
            ".pdf": "pdf",
            ".docx": "docx",
            ".xlsx": "xlsx",
            ".pptx": "pptx",
        }
        return format_map.get(suffix.lower(), "text")

    def _detect_format_from_content_type(self, content_type: str) -> str:
        """Detect document format from HTTP content-type."""
        if "markdown" in content_type:
            return "markdown"
        if "html" in content_type:
            return "html"
        if "pdf" in content_type:
            return "pdf"
        return "text"

    def _build_index(self, schema_entities: Optional[list[str]] = None) -> None:
        """Full rebuild of DOCUMENT chunks in the vector store.

        Resource Type: DOCUMENT
        - Clears all document chunks (document_name NOT LIKE 'schema:%' AND NOT LIKE 'api:%')
        - Re-indexes all loaded documents
        - Preserves schema and API chunks from other managers

        The vector store is shared across 3 resource types:
        - DOCUMENT: managed by DocumentDiscoveryTools (this class)
        - SCHEMA: managed by SchemaManager (document_name LIKE 'schema:%')
        - API: managed by APISchemaManager (document_name LIKE 'api:%')

        Each manager is responsible for its own resource type's cache lifecycle.

        Args:
            schema_entities: Optional list of known schema entity names for entity extraction
        """
        # Clear DOCUMENT chunks only (preserve schema:* and api:* from other managers)
        try:
            self._vector_store._conn.execute("""
                DELETE FROM embeddings
                WHERE document_name NOT LIKE 'schema:%'
                  AND document_name NOT LIKE 'api:%'
            """)
        except Exception:
            # Fallback for non-DuckDB backends
            self._vector_store.clear()

        # Clear document-sourced entities only
        if hasattr(self._vector_store, 'clear_entities'):
            self._vector_store.clear_entities(source='document')

        # Index all loaded documents
        self._index_loaded_documents(schema_entities)

    def _extract_and_store_entities(
        self,
        chunks: list[DocumentChunk],
        schema_entities: Optional[list[str]] = None,
    ) -> None:
        """Extract entities from chunks and store them using spaCy NER.

        Args:
            chunks: Document chunks to extract entities from
            schema_entities: Known schema entity names for matching
        """
        # Skip if vector store doesn't support entities
        if not hasattr(self._vector_store, 'add_entities'):
            return

        # Combine API terms
        api_terms = list(set(
            (self._openapi_operations or []) +
            (self._openapi_schemas or []) +
            (self._graphql_types or []) +
            (self._graphql_fields or [])
        ))

        # Create extractor with schema entities and spaCy NER
        # Use "__document__" as session_id for general document extraction
        extractor = EntityExtractor(
            session_id="__document__",
            schema_terms=schema_entities,
            api_terms=api_terms if api_terms else None,
        )

        # Extract entities from each chunk
        all_links: list[ChunkEntity] = []

        for chunk in chunks:
            extractions = extractor.extract(chunk)
            for entity, link in extractions:
                all_links.append(link)

        # Store all unique entities
        entities = extractor.get_all_entities()
        logger.debug(f"Entity extraction: {len(entities)} unique entities, {len(all_links)} links from {len(chunks)} chunks")
        if entities:
            self._vector_store.add_entities(entities, session_id="__document__")

        # Store all chunk-entity links (deduplicated by chunk_id + entity_id)
        if all_links:
            # Deduplicate links - same entity in same chunk should only have one link
            unique_links = {}
            for link in all_links:
                key = (link.chunk_id, link.entity_id)
                if key not in unique_links:
                    unique_links[key] = link
                else:
                    # Merge mention counts if duplicate
                    existing = unique_links[key]
                    unique_links[key] = ChunkEntity(
                        chunk_id=link.chunk_id,
                        entity_id=link.entity_id,
                        mention_count=existing.mention_count + link.mention_count,
                        confidence=max(existing.confidence, link.confidence),
                        mention_text=existing.mention_text or link.mention_text,
                    )
            self._vector_store.link_chunk_entities(list(unique_links.values()))

    def _chunk_document(self, name: str, content: str) -> list[DocumentChunk]:
        """Split a document into chunks for embedding.

        Chunks are split only on paragraph/line boundaries - never mid-paragraph.
        Paragraphs are combined until hitting CHUNK_SIZE, then a new chunk starts.
        A chunk may exceed CHUNK_SIZE if a single paragraph is larger (we don't split it).
        """
        chunks = []
        current_section = None

        # Determine paragraph separator based on content
        # Use double newline if present, otherwise single newline
        if "\n\n" in content:
            paragraphs = content.split("\n\n")
            separator = "\n\n"
        else:
            paragraphs = content.split("\n")
            separator = "\n"

        chunk_index = 0
        current_chunk = ""

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue

            # Track sections from markdown headers
            if para.startswith("#"):
                current_section = para.lstrip("#").strip()

            # Check if adding this paragraph would exceed chunk size
            potential_chunk = (current_chunk + separator + para).strip() if current_chunk else para

            if len(potential_chunk) <= self.CHUNK_SIZE:
                # Fits in current chunk - add it
                current_chunk = potential_chunk
            else:
                # Would exceed chunk size
                if current_chunk:
                    # Save current chunk first
                    chunks.append(DocumentChunk(
                        document_name=name,
                        content=current_chunk,
                        section=current_section,
                        chunk_index=chunk_index,
                    ))
                    chunk_index += 1
                # Start new chunk with this paragraph (even if it exceeds CHUNK_SIZE)
                current_chunk = para

        # Save final chunk
        if current_chunk:
            chunks.append(DocumentChunk(
                document_name=name,
                content=current_chunk,
                section=current_section,
                chunk_index=chunk_index,
            ))

        return chunks


# Tool schemas for LLM
DOC_TOOL_SCHEMAS = [
    {
        "name": "list_documents",
        "description": "List all available reference documents with descriptions and tags. Use this to see what documentation is available.",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "get_document",
        "description": "Get the full content of a reference document. Use this to read business rules, data dictionaries, or other documentation.",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {
                    "type": "string",
                    "description": "Name of the document to retrieve",
                },
            },
            "required": ["name"],
        },
    },
    {
        "name": "search_documents",
        "description": "Search across all documents for relevant content using semantic search. Use this to find specific information without reading entire documents.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "Natural language query describing what information you're looking for",
                },
                "limit": {
                    "type": "integer",
                    "description": "Maximum number of results to return",
                    "default": 5,
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "get_document_section",
        "description": "Get a specific section from a document by header/title. Use this when you know which section contains the information you need.",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {
                    "type": "string",
                    "description": "Name of the document",
                },
                "section": {
                    "type": "string",
                    "description": "Section header/title to retrieve",
                },
            },
            "required": ["name", "section"],
        },
    },
    {
        "name": "explore_entity",
        "description": "Find all document chunks that mention a specific entity (table, column, API endpoint, concept, or business term). Use this to gather additional context about an entity discovered in search results. This follows entity links to find related documentation across all sources.",
        "input_schema": {
            "type": "object",
            "properties": {
                "entity_name": {
                    "type": "string",
                    "description": "Name of the entity to explore (e.g., 'customers', 'revenue', 'OrderAPI')",
                },
                "limit": {
                    "type": "integer",
                    "description": "Maximum number of chunks to return",
                    "default": 5,
                },
            },
            "required": ["entity_name"],
        },
    },
]
