# Copyright (c) 2025 Kenneth Stott
#
# This source code is licensed under the Business Source License 1.1
# found in the LICENSE file in the root directory of this source tree.
#
# NOTICE: Use of this software for training artificial intelligence or
# machine learning models is strictly prohibited without explicit written
# permission from the copyright holder.

"""File content extraction functions for PDF, DOCX, XLSX, PPTX."""

import hashlib
import io
import logging
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

MIN_IMAGE_BYTES = 2048
MIN_IMAGE_DIMENSION = 64
_MAX_VISION_CALLS_PER_DOC = 50

# Skip vector metafile formats (not rasterizable without system deps)
_SKIP_MIMES = {"image/x-emf", "image/x-wmf", "image/emf", "image/wmf"}
_SKIP_EXTS = {".emf", ".wmf"}


@dataclass
class ExtractedImage:
    name: str
    data: bytes
    mime_type: str
    page: int | None
    index: int


def _guess_mime(data: bytes) -> str:
    """Guess MIME type from magic bytes."""
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if data[:2] == b"\xff\xd8":
        return "image/jpeg"
    if data[:4] == b"GIF8":
        return "image/gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    if data[:4] in (b"II\x2a\x00", b"MM\x00\x2a"):
        return "image/tiff"
    if data[:2] == b"BM":
        return "image/bmp"
    return "image/png"


def _ext_from_mime(mime: str) -> str:
    """Get file extension from MIME type."""
    mapping = {
        "image/png": ".png",
        "image/jpeg": ".jpg",
        "image/gif": ".gif",
        "image/webp": ".webp",
        "image/tiff": ".tiff",
        "image/bmp": ".bmp",
    }
    return mapping.get(mime, ".png")


def _is_emf_wmf(name: str, mime: str) -> bool:
    """Check if image is EMF/WMF format."""
    lower_name = name.lower()
    return (
        mime.lower() in _SKIP_MIMES
        or any(lower_name.endswith(ext) for ext in _SKIP_EXTS)
    )


def _passes_size_filter(data: bytes) -> bool:
    """Check if image meets minimum size requirements."""
    if len(data) < MIN_IMAGE_BYTES:
        return False
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        w, h = img.size
        if w < MIN_IMAGE_DIMENSION and h < MIN_IMAGE_DIMENSION:
            return False
    except Exception:
        return False
    return True


def _filter_and_dedup(images: list[ExtractedImage]) -> list[ExtractedImage]:
    """Apply size filter, EMF/WMF skip, and SHA-256 dedup."""
    seen: set[str] = set()
    result: list[ExtractedImage] = []
    for img in images:
        if _is_emf_wmf(img.name, img.mime_type):
            logger.warning("Skipping EMF/WMF image: %s", img.name)
            continue
        if not _passes_size_filter(img.data):
            continue
        h = hashlib.sha256(img.data).hexdigest()[:16]
        if h in seen:
            continue
        seen.add(h)
        result.append(img)
    return result


def _fetch_linked_image(ref: str, config_dir: str | None = None) -> bytes | None:
    """Fetch a linked/external image by reference.

    Supports: http(s), file:/, or bare path (treated as file:/).
    Returns None with warning on failure.
    """
    try:
        if ref.startswith(("http://", "https://")):
            from ._transport import _get_http_session
            resp = _get_http_session().get(ref, timeout=15)
            resp.raise_for_status()
            return resp.content
        else:
            path = ref.removeprefix("file://")
            p = Path(path)
            if not p.is_absolute() and config_dir:
                p = (Path(config_dir) / path).resolve()
            return p.read_bytes()
    except Exception as e:
        logger.warning("Failed to fetch linked image %s: %s", ref, e)
        return None


def _extract_pdf_images(path: Path | None, data: bytes | None) -> list[ExtractedImage]:
    """Extract embedded images from a PDF."""
    from pypdf import PdfReader

    reader = PdfReader(path or io.BytesIO(data))
    images: list[ExtractedImage] = []
    for page_num, page in enumerate(reader.pages, 1):
        for img_idx, image in enumerate(page.images, 1):
            mime = _guess_mime(image.data)
            ext = _ext_from_mime(mime)
            images.append(ExtractedImage(
                name=f"page_{page_num}_img_{img_idx}{ext}",
                data=image.data,
                mime_type=mime,
                page=page_num,
                index=img_idx,
            ))
    return _filter_and_dedup(images)


def _extract_docx_images(
    path: Path | None, data: bytes | None, config_dir: str | None = None,
) -> list[ExtractedImage]:
    """Extract embedded and linked images from a DOCX."""
    from docx import Document
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    doc = Document(path or io.BytesIO(data))
    images: list[ExtractedImage] = []
    idx = 0
    for rel in doc.part.rels.values():
        if rel.reltype != RT.IMAGE:
            continue
        idx += 1
        if rel.is_external:
            img_data = _fetch_linked_image(rel.target_ref, config_dir)
            if img_data is None:
                continue
            mime = _guess_mime(img_data)
            ext = _ext_from_mime(mime)
            images.append(ExtractedImage(
                name=f"linked_{idx}{ext}",
                data=img_data,
                mime_type=mime,
                page=None,
                index=idx,
            ))
        else:
            image_part = rel.target_part
            filename = image_part.partname.split("/")[-1]
            images.append(ExtractedImage(
                name=filename,
                data=image_part.blob,
                mime_type=image_part.content_type,
                page=None,
                index=idx,
            ))
    return _filter_and_dedup(images)


def _extract_pptx_images(
    path: Path | None, data: bytes | None, config_dir: str | None = None,
) -> list[ExtractedImage]:
    """Extract embedded images from a PPTX."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path or io.BytesIO(data))
    images: list[ExtractedImage] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_idx += 1
                image = shape.image
                ext = "." + image.content_type.split("/")[-1]
                if ext == ".jpeg":
                    ext = ".jpg"
                images.append(ExtractedImage(
                    name=f"slide_{slide_num}_img_{img_idx}{ext}",
                    data=image.blob,
                    mime_type=image.content_type,
                    page=slide_num,
                    index=img_idx,
                ))
    return _filter_and_dedup(images)


def _extract_xlsx_images(
    path: Path | None, data: bytes | None, config_dir: str | None = None,
) -> list[ExtractedImage]:
    """Extract embedded images from an XLSX."""
    from openpyxl import load_workbook

    wb = load_workbook(path or io.BytesIO(data), data_only=True)
    images: list[ExtractedImage] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        for img_idx, image in enumerate(ws._images, 1):
            img_data = image._data()
            mime = _guess_mime(img_data)
            ext = _ext_from_mime(mime)
            images.append(ExtractedImage(
                name=f"sheet_{sheet_name}_img_{img_idx}{ext}",
                data=img_data,
                mime_type=mime,
                page=None,
                index=img_idx,
            ))
    return _filter_and_dedup(images)


def _extract_images_from_document(
    path: Path | None,
    data: bytes | None,
    doc_type: str,
    config_dir: str | None = None,
) -> list[ExtractedImage]:
    """Dispatch to format-specific image extractor."""
    try:
        if doc_type == "pdf":
            return _extract_pdf_images(path, data)
        elif doc_type == "docx":
            return _extract_docx_images(path, data, config_dir)
        elif doc_type == "pptx":
            return _extract_pptx_images(path, data, config_dir)
        elif doc_type == "xlsx":
            return _extract_xlsx_images(path, data, config_dir)
    except Exception as e:
        logger.warning("Image extraction failed for %s: %s", doc_type, e)
    return []


def _is_structured_data_format(doc_format: str) -> bool:
    """Check if format is structured data that shouldn't be semantically indexed."""
    # These formats are data to be queried, not text to be searched
    structured_formats = {"csv", "json", "jsonl", "parquet", "xml", "yaml", "yml"}
    return doc_format.lower() in structured_formats


def _extract_pdf_text(path) -> str:
    """Extract text content from a PDF file.

    Args:
        path: Path to the PDF file

    Returns:
        Extracted text content with page markers
    """
    from pypdf import PdfReader

    reader = PdfReader(path)
    pages = []

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i}]\n{text.strip()}")

    return "\n\n".join(pages)


def _extract_pdf_text_from_bytes(pdf_bytes: bytes) -> str:
    """Extract text content from PDF bytes.

    Args:
        pdf_bytes: Raw PDF file content

    Returns:
        Extracted text content with page markers
    """
    from io import BytesIO
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(pdf_bytes))
    pages = []

    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if text and text.strip():
            pages.append(f"[Page {i}]\n{text.strip()}")

    return "\n\n".join(pages)


def _extract_docx_content(doc) -> str:
    """Extract text content from a python-docx Document object.

    Args:
        doc: python-docx Document instance

    Returns:
        Extracted text content with paragraph separation
    """
    paragraphs = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Check if it's a heading
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "")
                try:
                    level_num = int(level)
                    paragraphs.append(f"{'#' * level_num} {text}")
                except ValueError:
                    paragraphs.append(text)
            else:
                paragraphs.append(text)

    # Also extract text from tables
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            table_rows.append(" | ".join(cells))
        if table_rows:
            paragraphs.append("\n".join(table_rows))

    return "\n\n".join(paragraphs)


def _extract_docx_text(path) -> str:
    """Extract text content from a Word document.

    Args:
        path: Path to the .docx file

    Returns:
        Extracted text content with paragraph separation
    """
    from docx import Document

    return _extract_docx_content(Document(path))


def _extract_docx_text_from_bytes(docx_bytes: bytes) -> str:
    """Extract text content from Word document bytes.

    Args:
        docx_bytes: Raw .docx file content

    Returns:
        Extracted text content
    """
    from io import BytesIO
    from docx import Document

    return _extract_docx_content(Document(BytesIO(docx_bytes)))


def _extract_xlsx_content(wb) -> str:
    """Extract text content from an openpyxl Workbook object.

    Args:
        wb: openpyxl Workbook instance

    Returns:
        Extracted text content with sheet and cell markers
    """
    sheets = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        rows = []

        for row in sheet.iter_rows():
            cells = []
            for cell in row:
                if cell.value is not None:
                    cells.append(str(cell.value))
                else:
                    cells.append("")
            # Only include rows that have some content
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))

        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    return "\n\n".join(sheets)


def _extract_xlsx_text(path) -> str:
    """Extract text content from an Excel spreadsheet.

    Args:
        path: Path to the .xlsx file

    Returns:
        Extracted text content with sheet and cell markers
    """
    from openpyxl import load_workbook

    return _extract_xlsx_content(load_workbook(path, data_only=True))


def _extract_xlsx_text_from_bytes(xlsx_bytes: bytes) -> str:
    """Extract text content from Excel spreadsheet bytes.

    Args:
        xlsx_bytes: Raw .xlsx file content

    Returns:
        Extracted text content
    """
    from io import BytesIO
    from openpyxl import load_workbook

    return _extract_xlsx_content(load_workbook(BytesIO(xlsx_bytes), data_only=True))


def _extract_pptx_content(prs) -> str:
    """Extract text content from a python-pptx Presentation object.

    Args:
        prs: python-pptx Presentation instance

    Returns:
        Extracted text content with slide markers
    """
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            # Handle tables in slides
            if shape.has_table:
                table_rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    table_rows.append(" | ".join(cells))
                if table_rows:
                    slide_text.append("\n".join(table_rows))

        if slide_text:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_text))

    return "\n\n".join(slides)


def _extract_pptx_text(path) -> str:
    """Extract text content from a PowerPoint presentation.

    Args:
        path: Path to the .pptx file

    Returns:
        Extracted text content with slide markers
    """
    from pptx import Presentation

    return _extract_pptx_content(Presentation(path))


def _extract_pptx_text_from_bytes(pptx_bytes: bytes) -> str:
    """Extract text content from PowerPoint presentation bytes.

    Args:
        pptx_bytes: Raw .pptx file content

    Returns:
        Extracted text content
    """
    from io import BytesIO
    from pptx import Presentation

    return _extract_pptx_content(Presentation(BytesIO(pptx_bytes)))


def _strip_html_chrome(html: str) -> str:
    """Pre-strip navigation chrome tags from raw HTML.

    Removes <nav>, <aside>, <noscript>, <script>, <style> and elements
    with navigation-related CSS classes/IDs. Uses regex on raw HTML
    since html.parser can't handle unclosed tags reliably.
    """
    import re
    # Strip well-known chrome tags — require closing tag (skip unclosed)
    for tag in ("nav", "aside", "noscript", "script", "style"):
        html = re.sub(
            rf"<{tag}[\s>].*?</{tag}>",
            "", html, flags=re.DOTALL | re.IGNORECASE,
        )
    # Strip divs/sections with navigation-related class/id attributes
    _NAV_ATTR_RE = re.compile(
        r"<(div|section|ul|table)[^>]*?"
        r"(?:class|id|role)\s*=\s*[\"'][^\"']*?"
        r"(?:sidebar|navbox|navbar|navigation|toc\b|catlinks"
        r"|mw-panel|mw-head|mw-editsection"
        r"|menu|breadcrumb|noprint"
        r"|portal|sister-?project|interlanguage|authority-control"
        r"|reflist|references|footnotes|mw-references-wrap|citation)"
        r"[^\"']*?[\"'][^>]*>.*?</\1>",
        re.DOTALL | re.IGNORECASE,
    )
    html = _NAV_ATTR_RE.sub("", html)
    # Strip inline citation markers (<sup class="reference">[1]</sup>)
    html = re.sub(r"<sup[^>]*class=[\"'][^\"']*reference[^\"']*[\"'][^>]*>.*?</sup>", "", html, flags=re.DOTALL | re.IGNORECASE)
    return html


def _convert_html_to_markdown(html: str) -> str:
    """Convert HTML to markdown, preserving heading structure.

    Uses stdlib html.parser — no external dependencies.
    Handles: headings, paragraphs, lists (ul/ol/li), <br>, <pre>/<code>,
    bold, italic, links, and tables.
    Strips navigation chrome (nav, aside, sidebar, navbox, etc.) — link-heavy
    regions are the crawler's job, not NER input.
    """
    from html.parser import HTMLParser
    import re

    html = _strip_html_chrome(html)

    class _MarkdownConverter(HTMLParser):
        def __init__(self):
            super().__init__()
            self._output: list[str] = []
            self._tag_stack: list[str] = []
            self._list_stack: list[str] = []  # "ul" or "ol"
            self._ol_counters: list[int] = []
            self._in_pre = False
            self._href: str | None = None
            self._link_text: list[str] = []
            self._in_link = False

        def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]):
            tag = tag.lower()
            self._tag_stack.append(tag)
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                self._output.append("\n\n")
            elif tag == "p":
                self._output.append("\n\n")
            elif tag == "br":
                self._output.append("\n")
            elif tag == "pre":
                self._in_pre = True
                self._output.append("\n\n```\n")
            elif tag == "ul":
                self._list_stack.append("ul")
            elif tag == "ol":
                self._list_stack.append("ol")
                self._ol_counters.append(0)
            elif tag == "li":
                indent = "  " * (len(self._list_stack) - 1)
                if self._list_stack and self._list_stack[-1] == "ol":
                    self._ol_counters[-1] += 1
                    self._output.append(f"\n{indent}{self._ol_counters[-1]}. ")
                else:
                    self._output.append(f"\n{indent}- ")
            elif tag == "strong" or tag == "b":
                self._output.append("**")
            elif tag == "em" or tag == "i":
                self._output.append("*")
            elif tag == "a":
                attr_dict = dict(attrs)
                self._href = attr_dict.get("href")
                self._in_link = True
                self._link_text = []
            elif tag == "tr":
                self._output.append("\n|")
            elif tag in ("td", "th"):
                self._output.append(" ")

        def handle_endtag(self, tag: str):
            tag = tag.lower()
            if self._tag_stack and self._tag_stack[-1] == tag:
                self._tag_stack.pop()
            if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(tag[1])
                prefix = "#" * level + " "
                # Find the heading text — scan back for the last \n\n
                text_parts: list[str] = []
                while self._output and self._output[-1] != "\n\n":
                    text_parts.append(self._output.pop())
                text = "".join(reversed(text_parts)).strip()
                self._output.append(f"{prefix}{text}\n\n")
            elif tag == "p":
                self._output.append("\n")
            elif tag == "pre":
                self._in_pre = False
                self._output.append("\n```\n\n")
            elif tag == "ul":
                if self._list_stack:
                    self._list_stack.pop()
                self._output.append("\n")
            elif tag == "ol":
                if self._list_stack:
                    self._list_stack.pop()
                if self._ol_counters:
                    self._ol_counters.pop()
                self._output.append("\n")
            elif tag == "strong" or tag == "b":
                self._output.append("**")
            elif tag == "em" or tag == "i":
                self._output.append("*")
            elif tag == "a":
                link_text = "".join(self._link_text).strip()
                if self._href and link_text:
                    self._output.append(f"[{link_text}]({self._href})")
                else:
                    self._output.append(link_text)
                self._in_link = False
                self._href = None
                self._link_text = []
            elif tag in ("td", "th"):
                self._output.append(" |")
            elif tag == "thead":
                # Add separator row after header
                self._output.append("\n|---|")

        def handle_data(self, data: str):
            if self._in_link:
                self._link_text.append(data)
                return
            if self._in_pre:
                self._output.append(data)
            else:
                self._output.append(data)

        def get_markdown(self) -> str:
            text = "".join(self._output)
            # Collapse excessive blank lines
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text.strip()

    converter = _MarkdownConverter()
    converter.feed(html)
    return converter.get_markdown()


def _detect_format(suffix: str) -> str:
    """Detect document format from file extension.

    DEPRECATED: Use _mime.detect_type_from_source() instead.
    Kept for backward compatibility with external callers.
    """
    from ._mime import EXTENSION_TO_SHORT
    return EXTENSION_TO_SHORT.get(suffix.lower(), "text")


def _detect_format_from_content_type(content_type: str) -> str:
    """Detect document format from HTTP content-type.

    DEPRECATED: Use _mime.detect_type_from_source() instead.
    Kept for backward compatibility with external callers.
    """
    from ._mime import detect_type_from_source
    return detect_type_from_source(None, content_type)
